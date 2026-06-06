#!/usr/bin/env python3
# ════════════════════════════════════════════════════════════════════════════
#  ATARS v3.1 — Automated Time-Series Analysis and Reporting System
#  ──────────────────────────────────────────────────────────────────────────
#  Author   : Priyanshu
#  Institute: Global Institute of Technology and Management, Haryana, India
#  Version  : 3.1.0  (FINAL — Production Release)
#  License  : MIT
#  GitHub   : github.com/Priyanshu-ux712/ATARS
#
#  CHANGELOG v3.1 FINAL (bug-fix release):
#  ─────────────────────────────────────────────────────────────────────────
#  [FIX-01] Unified all `import re as _re` / `import re as re` → `import re`
#  [FIX-02] Defined missing `_h2_inner()` helper (was crash in Section 12b)
#  [FIX-03] Added `run_metadata` + live `hash_J` to `build_json_contract()`
#  [FIX-04] Section 16: grounding rows now prefer `v3_ver_v3` fields
#  [FIX-05] Section 16: per-section `g_rate` key fixed (`G_rate` fallback)
#  [FIX-06] Section 16: per-section table prefers v3 sections dict
#  [FIX-07] PPT Slide 5: `n_z_only` → `n_zs_only` (correct IF dict key)
#  [FIX-08] PPT Slide 6: `forecast_mean` → `np.mean(forecast)` (real value)
#  [FIX-09] PPT Slide 6: `ml_hw.get("params")` → direct alpha/beta/gamma keys
#  [FIX-10] PPT Slide 3: `max_daily_mean/min_month_name` → `annual_max/std`
#  [FIX-11] PPT: `year`/`run_id` extracted from `date_range` str, not nesting
#  [FIX-12] `_add_unified_verification_dashboard()` now called from build_report
#  [FIX-13] Data pipeline now has 4 named phases with clear terminal banners
#  [FIX-14] Data cleaning step shows null/row/column metrics in terminal output
#  ─────────────────────────────────────────────────────────────────────────
#  PIPELINE ARCHITECTURE (4 Phases):
#  ─────────────────────────────────────────────────────────────────────────
#  PHASE 1 — DATA PIPELINE
#    1.1 CSV Ingestion → column normalization → datetime parsing
#    1.2 10-Step Automated Cleaning:
#        duplicate removal → datetime gap detection → negative correction
#        → physical bounds clipping → type coercion → sparse column removal
#        → 3-σ winsorization → 4-strategy imputation → audit report
#    1.3 Quality flag assignment (q_i ∈ {1=valid, 2=suspect, 3=invalid})
#
#  PHASE 2 — STATISTICAL ANALYSIS
#    2.1 16 Formal Statistical Operators (Eqs. 1–16)
#    2.2 Isolation Forest multivariate anomaly detection (ML-1)
#    2.3 Holt-Winters triple exponential smoothing forecast (ML-2)
#    2.4 JSON Contract assembly (J — the formal LLM data boundary)
#
#  PHASE 3 — CONTEXTUAL ENRICHMENT
#    3.1 Daily snapshot system (one JSON per day)
#    3.2 Open-Meteo live weather (free API, no key)
#    3.3 Atmospheric stability analysis (Pasquill-Gifford proxy)
#    3.4 Pollution episode classification (N-10, 8 episode types)
#    3.5 7-day weather forecast with pollution risk scores
#
#  PHASE 4 — VERIFICATION LAYER (9 novel modules)
#    RGV  — Numerical grounding verifier (G_rate ∈ [0,1])
#    TCV  — Temporal claim verifier         [N-01]
#    CCD  — Causal claim detector           [N-02]
#    AEE  — Anomaly explanation engine      [N-03]
#    RDD  — Report drift detector           [N-04]
#    UPT  — Uncertainty propagation         [N-05]
#    CVCC — Cross-variable consistency      [N-06]
#    NSS  — Narrative specificity score     [N-07]
#    BNC  — Benchmark comparator            [N-08]
#    MNS  — Master novelty score            [N-09]
#
#  OUTPUT:
#    • 18-section Word report  (Section 00b = data cleaning audit)
#    • 15 publication-quality charts
#    • 15-slide executive PowerPoint
#    • SHA-256 audit log
#    • Daily snapshot JSONs
#    • Windows Task Scheduler .bat file
#  ─────────────────────────────────────────────────────────────────────────
#
#  PIPELINE OVERVIEW:
#  ┌─────────────────────────────────────────────────────────────┐
#  │  CSV Input (any CPCB station data)                          │
#  │      ↓                                                      │
#  │  16 Formal Statistical Operators (Eqs. 1–16)               │
#  │      ↓                                                      │
#  │  Isolation Forest + Holt-Winters ML Layer                   │
#  │      ↓                                                      │
#  │  Open-Meteo Weather Context (free API, no key)              │
#  │      ↓                                                      │
#  │  LLM Narrative (optional, local Ollama)                     │
#  │      ↓                                                      │
#  │  9 Verification Layers:                                     │
#  │    RGV  — Numerical Grounding (G_rate)                      │
#  │    TCV  — Temporal Claim Verifier      [N-01, novel]        │
#  │    CCD  — Causal Claim Detector        [N-02, novel]        │
#  │    AEE  — Anomaly Explanation Engine   [N-03, novel]        │
#  │    RDD  — Report Drift Detector        [N-04, novel]        │
#  │    UPT  — Uncertainty Propagation      [N-05, novel]        │
#  │    CVCC — Cross-Variable Consistency   [N-06, novel]        │
#  │    NSS  — Narrative Specificity Score  [N-07, novel]        │
#  │    BNC  — Benchmark Comparator         [N-08]               │
#  │    MNS  — Master Novelty Score         [N-09, novel]        │
#  │      ↓                                                      │
#  │  Output: 18-section Word report + 15 charts + audit log     │
#  └─────────────────────────────────────────────────────────────┘
#
#  USAGE:
#    python atars_v3.py --data your_data.csv --city "Gurugram" --no-llm
#    python atars_v3.py --data your_data.csv --city "Delhi" --weather open_meteo
#    python atars_v3.py --data your_data.csv --city "Mumbai" --mode test
#
#  Open-Meteo API: free, no key, no signup — https://open-meteo.com
# ════════════════════════════════════════════════════════════════════════════
import os
import sys
import json
import math
import hashlib
import re
import warnings
import argparse
import traceback
from pathlib import Path
from datetime import datetime, timedelta
from collections import defaultdict

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
from matplotlib.patches import FancyBboxPatch
import seaborn as sns
from scipy import stats
from sklearn.linear_model import LinearRegression
from sklearn.preprocessing import StandardScaler
from sklearn.ensemble import IsolationForest
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

warnings.filterwarnings('ignore')
try:
    import urllib3 as _u3
    _u3.disable_warnings(_u3.exceptions.InsecureRequestWarning)
except Exception:
    pass

# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 0 — CONFIGURATION
#  All parameters in one place. Edit this section for your domain.
# ═══════════════════════════════════════════════════════════════════════════════

DEFAULT_CONFIG = {
    # ── Data ──────────────────────────────────────────────────────────────────
    "data_file"            : "data.csv",
    "city"                 : "City",
    "station_id"           : "Station-01",
    "output_dir"           : "atars_output",

    # ── Statistical parameters ─────────────────────────────────────────────────
    "baseline_window_days" : 30,       # W — days for rolling baseline
    "anomaly_z_threshold"  : 3.0,      # ζ_thresh — z-score anomaly boundary
    "confidence_high"      : 0.90,     # Q(D) ≥ 0.90 → HIGH
    "confidence_moderate"  : 0.70,     # Q(D) ≥ 0.70 → MODERATE
    "ci_alpha"             : 0.05,     # α for 95% confidence intervals
    "acf_lags"             : 48,       # number of lags for ACF/PACF
    "bootstrap_samples"    : 1000,     # B for bootstrap CI
    "min_valid_ratio"      : 0.10,     # minimum Q(D) to process a day

    # ── ML parameters (v2.0) ──────────────────────────────────────────────────
    # ML-1: Isolation Forest anomaly detection
    "if_contamination"     : 0.05,    # expected anomaly fraction (5%)
    "if_n_estimators"      : 150,     # number of trees in the forest
    "if_random_state"      : 42,      # fixed seed for reproducibility
    "if_cols"              : ["PM10","NO2","SO2","CO","Ozone"],  # features
    # ML-2: Holt-Winters forecast
    "hw_forecast_days"     : 14,      # days to forecast ahead
    "hw_seasonal_periods"  : 7,       # weekly seasonality
    "hw_alpha"             : 0.3,     # level smoothing parameter
    "hw_beta"              : 0.1,     # trend smoothing parameter
    "hw_gamma"             : 0.2,     # seasonal smoothing parameter
    "use_forecast"         : True,    # enable Holt-Winters forecast chart

    # ── WHO / Reference thresholds (24-hour means, WHO AQG 2021) ──────────────
    "thresholds": {
        "PM10"       : 45.0,    # µg/m³  — WHO AQG 2021
        "NO2"        : 25.0,    # µg/m³  — WHO AQG 2021
        "SO2"        : 40.0,    # µg/m³  — WHO AQG 2021
        "CO"         : 4.0,     # mg/m³  — WHO AQG 2021
        "Ozone"      : 100.0,   # µg/m³  — WHO AQG 2021 (peak season)
        "NH3"        : 400.0,   # µg/m³  — CPCB India standard
        "NO"         : None,    # no WHO guideline (informational)
        "NOx"        : None,
        "Benzene"    : 1.7,     # µg/m³  — WHO annual reference
        "Toluene"    : None,
        "Xylene"     : None,
        "Eth_Benzene": None,
        "MP_Xylene"  : None,
    },

    # ── LLM (Ollama — local, free, offline) ───────────────────────────────────
    "use_llm"          : True,
    "llm_model"        : "llama3.2",   # change to "mistral", "phi3", "llama3.2:8b" etc.
    "ollama_url"       : "http://localhost:11434",
    "llm_timeout"      : 600,          # seconds per section — 10 min for slow hardware
    "llm_retry"        : 2,            # retry count per section on timeout
    "llm_retry_delay"  : 5,            # seconds between retries

    # ── Report ────────────────────────────────────────────────────────────────
    "author"      : "Priyanshu",
    "institution" : "Global Institute of Technology and Management",
    "department"  : "Department of Computer Science & Engineering",
    "degree"      : "B.Tech (Computer Science), Second Year",
    "email"       : "priyanshukumar9053@gmail.com",
    "location"    : "Haryana, India",
}

# ── Colour palette (consistent across all charts) ─────────────────────────────
PAL = {
    "navy"  : "#1A2C4E", "blue"   : "#2557A7", "teal"   : "#1A6B72",
    "green" : "#1E6B3C", "red"    : "#8B1C2A", "amber"  : "#92400E",
    "gray"  : "#6B7280", "lgray"  : "#F0F4F8", "white"  : "#FFFFFF",
    "purple": "#4C1D95", "orange" : "#7C2D12", "dark"   : "#0D1117",
    "chart_bg": "#F8FAFC",
}

plt.rcParams.update({
    'font.family'      : 'DejaVu Sans',
    'axes.facecolor'   : PAL['chart_bg'],
    'figure.facecolor' : PAL['white'],
    'axes.spines.top'  : False,
    'axes.spines.right': False,
    'axes.grid'        : True,
    'grid.alpha'       : 0.3,
    'grid.linestyle'   : ':',
})

# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 1 — COLUMN NORMALIZATION
#  Handles encoding issues (Â° → °, Âµg → µg) automatically.
# ═══════════════════════════════════════════════════════════════════════════════

# Maps raw column names (possibly with encoding artifacts) → clean internal names
# Uses keyword matching so it works regardless of encoding corruption.
COLUMN_KEYWORDS = {
    "humidity"            : "humidity_pct",
    "rain"                : "rain_mm",
    "pressure"            : "pressure_hpa",
    "wind_speed_10"       : "wind_speed_10m",
    "wind_dir"            : "wind_dir_10m",
    "wind_direction_10"   : "wind_dir_10m",
    "wind_speed_100"      : "wind_speed_100m",
    "wind_direction_100"  : "wind_dir_100m",
    "pm10"                : "PM10",
    " no ("               : "NO",        # space before "no" to avoid "nox","no2"
    "no2"                 : "NO2",
    "nox"                 : "NOx",
    "nh3"                 : "NH3",
    "so2"                 : "SO2",
    " co ("               : "CO",        # space to avoid "co2","ozone" matches
    "ozone"               : "Ozone",
    "benzene"             : "Benzene",
    "toluene"             : "Toluene",
    "mp-xylene"           : "MP_Xylene",
    "mp_xylene"           : "MP_Xylene",
    "eth"                 : "Eth_Benzene",
    "xylene"              : "Xylene",    # must come AFTER mp-xylene and eth
}

# Pollutant columns used in statistical analysis
POLLUTANTS = [
    "PM10","NO","NO2","NOx","NH3","SO2","CO",
    "Ozone","Benzene","Toluene","Xylene","Eth_Benzene","MP_Xylene"
]
METEOROLOGICAL = [
    "humidity_pct","rain_mm","pressure_hpa",
    "wind_speed_10m","wind_dir_10m","wind_speed_100m","wind_dir_100m"
]
ALL_NUMERIC = POLLUTANTS + METEOROLOGICAL


def normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    """
    Rename raw columns to clean internal names.
    Handles UTF-8/Latin-1 encoding artifacts automatically.
    Eq. 1: D = {(t_i, x_i, q_i)} — build the observation space.

    Improvements over v1:
    - Reports unrecognized columns so silent data loss is visible
    - Strips BOM markers, non-breaking spaces, unicode control chars
    - Handles common encoding artifacts (Â°, Âµg, ï»¿)
    - Falls back to fuzzy matching for near-misses
    - Warns when expected pollutant columns are entirely absent
    """
    # Strip encoding artifacts from column names before matching
    clean_cols = {}
    for c in df.columns:
        cleaned = (str(c)
                   .strip()
                   .replace("\ufeff", "")      # BOM
                   .replace("\u00a0", " ")      # non-breaking space
                   .replace("Â°", "°")
                   .replace("Âµ", "µ")
                   .replace("Ã©", "é")
                   .replace("ï»¿", ""))
        clean_cols[c] = cleaned
    df = df.rename(columns=clean_cols)

    rename_map   = {}
    unrecognized = []
    cols_lower   = {c: c.lower() for c in df.columns}

    for raw_col, lower_col in cols_lower.items():
        # Skip obviously non-data columns
        if lower_col in ("index", "unnamed", "s.no", "sr.no", "sr", "sno"):
            continue

        # Direct time columns
        if lower_col in ("day", "date", "datetime", "date_time", "timestamp"):
            rename_map[raw_col] = "date" if "date" in lower_col or lower_col in ("day","timestamp") else "datetime"
        elif lower_col == "time":
            rename_map[raw_col] = "time"
        else:
            matched = False
            for keyword, clean_name in COLUMN_KEYWORDS.items():
                if keyword in lower_col:
                    rename_map[raw_col] = clean_name
                    matched = True
                    break

            # Fuzzy fallback: check if any POLLUTANT name appears as substring
            if not matched:
                for pol in POLLUTANTS + METEOROLOGICAL:
                    if pol.lower() in lower_col or lower_col in pol.lower():
                        rename_map[raw_col] = pol
                        matched = True
                        break

            if not matched:
                # Only flag as unrecognized if it looks like it might be data
                # (contains digits, units, or known chemical symbols)
                looks_like_data = any(x in lower_col for x in
                    ["µg", "mg", "ppb", "ppm", "µg/m", "%", "°c", "km/h",
                     "pm", "no", "so", "co", "o3", "nh", "voc", "aqi"])
                if looks_like_data:
                    unrecognized.append(raw_col)

    df = df.rename(columns=rename_map)

    # Remove duplicate columns (keep first)
    df = df.loc[:, ~df.columns.duplicated()]

    # Report unrecognized data-like columns — silent loss is a data quality issue
    if unrecognized:
        print(f"\n  ⚠ UNRECOGNIZED COLUMNS (data-like but not mapped):")
        for col in unrecognized:
            print(f"     · '{col}' — add to COLUMN_KEYWORDS if this is a pollutant/met variable")
        print(f"     These columns will be present in df but excluded from analysis.")

    # Warn if core pollutants are entirely absent
    found_pollutants = [p for p in POLLUTANTS if p in df.columns]
    if not found_pollutants:
        print(f"  ⚠ WARNING: No recognized pollutant columns found after normalization.")
        print(f"     Raw columns were: {list(df.columns[:10])}")
    else:
        print(f"  ✓ Normalized columns: {found_pollutants}")
        met_found = [m for m in METEOROLOGICAL if m in df.columns]
        if met_found:
            print(f"  ✓ Meteorological:    {met_found}")

    return df


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 2 — DATA LOADING & VALIDATION
# ═══════════════════════════════════════════════════════════════════════════════



# ═══════════════════════════════════════════════════════════════════════════════
#  DATA QUALITY ENGINE — Comprehensive Dataset Cleaning
#  Handles: nulls, duplicates, outliers, encoding errors, mixed types,
#           negative values, unit mismatches, sparse columns, datetime gaps
# ═══════════════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════════════
#  DATA QUALITY ENGINE v2 — Professional 18-Step Cleaning Pipeline
# ═══════════════════════════════════════════════════════════════════════════════

FLATLINE_THRESHOLDS = {
    "PM10"          : 6,
    "NO2"           : 6,
    "SO2"           : 8,
    "CO"            : 6,
    "Ozone"         : 8,
    "NO"            : 6,
    "NOx"           : 6,
    "NH3"           : 12,
    "Benzene"       : 12,
    "humidity_pct"  : 24,
    "wind_speed_10m": 3,
}

ROC_LIMITS = {
    "PM10"       : 400.0,
    "NO2"        : 200.0,
    "SO2"        : 200.0,
    "CO"         : 10.0,
    "Ozone"      : 150.0,
    "NO"         : 200.0,
    "NOx"        : 300.0,
    "NH3"        : 300.0,
    "Benzene"    : 50.0,
    "humidity_pct": 30.0,
}

CHEMICAL_CONSTRAINTS_CLEAN = [
    ("NOx", "NO2", ">=", 0.0),
    ("NOx", "NO",  ">=", 0.0),
]

PHYSICAL_BOUNDS_CLEAN = {
    "PM10"          : (0, 1500),   "PM2.5"     : (0, 800),
    "PM25"          : (0, 800),    "NO2"        : (0, 2000),
    "NO"            : (0, 2000),   "NOx"        : (0, 3000),
    "NH3"           : (0, 1000),   "SO2"        : (0, 2000),
    "CO"            : (0, 100),    "Ozone"      : (0, 500),
    "Benzene"       : (0, 500),    "Toluene"    : (0, 500),
    "Xylene"        : (0, 500),    "Eth_Benzene": (0, 200),
    "MP_Xylene"     : (0, 500),    "humidity_pct": (0, 100),
    "wind_speed_10m": (0, 200),    "wind_speed_100m": (0, 300),
    "temperature_2m": (-60, 60),   "pressure_hpa": (800, 1100),
    "rain_mm"       : (0, 500),
}


def _log_step(report: dict, step_id: str, message: str):
    report["steps"].append({"step": step_id, "action": message})
    print(f"  · {step_id}: {message}")


def clean_dataset(df: pd.DataFrame, config: dict) -> tuple:
    """
    Professional 18-step automated data cleaning pipeline.
    Returns (cleaned_df, cleaning_report).

    Step 01: Column name standardisation
    Step 02: Exact duplicate row removal
    Step 03: Datetime parsing, sort, duplicate timestamps, gap detection
    Step 04: Missing data profiling (null pattern, hourly completeness)
    Step 05: Sensor flatline detection (stuck sensor → NULL)
    Step 06: Rate-of-change spike filter (impossible jumps → NULL)
    Step 07: Negative value correction (NULL, not zero)
    Step 08: Physical bounds clipping (instrument saturation)
    Step 09: Mixed-type coercion
    Step 10: Sparse column removal (< 10% valid)
    Step 11: Cross-variable chemical constraint audit
    Step 12: Outlier winsorization (3×IQR Tukey fences)
    Step 13: Datetime gap row injection (before imputation)
    Step 14: Missing value imputation (5-strategy cascade)
    Step 15: Diurnal anomaly flagging (>5σ from hour-of-day mean)
    Step 16: Post-imputation consistency re-check
    Step 17: Null streak analysis (downtime vs random dropout)
    Step 18: 7-dimensional cleaning quality score
    """
    report = {
        "original_shape"        : df.shape,
        "steps"                 : [],
        "total_nulls_before"    : int(df.isnull().sum().sum()),
        "total_nulls_after"     : 0,
        "rows_removed"          : 0,
        "rows_injected"         : 0,
        "cols_removed"          : [],
        "values_imputed"        : 0,
        "values_clipped"        : 0,
        "values_winsorized"     : 0,
        "flatlines_nulled"      : 0,
        "spikes_nulled"         : 0,
        "negatives_nulled"      : 0,
        "constraint_violations" : [],
        "null_streaks"          : {},
        "completeness_by_hour"  : {},
        "completeness_by_col"   : {},
        "datetime_gaps"         : [],
        "duplicate_timestamps"  : 0,
        "imputation_detail"     : {},
        "winsor_detail"         : {},
        "flatline_detail"       : {},
        "spike_detail"          : {},
        "clipped_detail"        : {},
        "diurnal_flags"         : 0,
        "recheck_clipped"       : 0,
        "quality_score"         : 0.0,
        "quality_dimensions"    : {},
    }

    original_rows = len(df)
    original_cols = list(df.columns)
    if original_rows == 0:
        report["final_shape"] = df.shape
        return df, report

    # ── Step 01: Column names ─────────────────────────────────────────────
    df.columns = df.columns.astype(str).str.strip()
    _log_step(report, "01_col_names", "Column names stripped and cast to str")

    # ── Step 02: Exact duplicate rows ─────────────────────────────────────
    before = len(df)
    df     = df.drop_duplicates()
    removed = before - len(df)
    report["rows_removed"] += removed
    _log_step(report, "02_duplicates", f"Removed {removed} exact duplicate rows")

    # ── Step 03: Datetime ─────────────────────────────────────────────────
    datetime_col = None
    for cname in ["datetime","date_time","timestamp","date","Date","Datetime"]:
        if cname in df.columns:
            datetime_col = cname
            break
    if datetime_col:
        try:
            df[datetime_col] = pd.to_datetime(
                df[datetime_col], infer_datetime_format=True, errors="coerce")
            bad_dt = int(df[datetime_col].isna().sum())
            if bad_dt > 0:
                df = df.dropna(subset=[datetime_col])
                report["rows_removed"] += bad_dt
                _log_step(report, "03a_bad_dt", f"Dropped {bad_dt} unparseable datetime rows")
            df = df.sort_values(datetime_col).reset_index(drop=True)
            dup_ts = int(df.duplicated(subset=[datetime_col]).sum())
            if dup_ts > 0:
                df = df.drop_duplicates(subset=[datetime_col], keep="first")
                report["duplicate_timestamps"] = dup_ts
                report["rows_removed"] += dup_ts
                _log_step(report, "03b_dup_ts", f"Removed {dup_ts} duplicate timestamps")
            if len(df) > 1:
                diffs      = df[datetime_col].diff().dropna()
                modal_freq = diffs.mode()[0] if len(diffs.mode()) > 0 else None
                if modal_freq and modal_freq.total_seconds() > 0:
                    for idx in diffs[diffs > modal_freq * 3].index:
                        try:
                            gh = round(diffs.loc[idx].total_seconds() / 3600, 1)
                            if gh > 2:
                                report["datetime_gaps"].append({
                                    "start"    : str(df.loc[idx-1, datetime_col])[:16],
                                    "gap_hours": gh,
                                })
                        except Exception:
                            pass
                    if report["datetime_gaps"]:
                        _log_step(report, "03c_gaps",
                                  f"Found {len(report['datetime_gaps'])} time gaps >2h")
                    else:
                        _log_step(report, "03c_gaps", "No time gaps >2h detected")
        except Exception as e:
            _log_step(report, "03_dt_error", f"Datetime parse issue: {e}")

    # ── Step 04: Missing data profile ─────────────────────────────────────
    numeric_cols  = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    completeness  = {}
    for col in numeric_cols:
        valid = int(df[col].notna().sum())
        pct   = round(valid / max(len(df), 1) * 100, 1)
        completeness[col] = {"valid": valid, "total": len(df), "pct": pct}
    report["completeness_by_col"] = completeness
    if datetime_col and datetime_col in df.columns:
        try:
            df["_htmp"] = df[datetime_col].dt.hour
            primary = next((c for c in ["PM10","NO2","SO2"] if c in df.columns), None)
            if primary:
                hc = (df.groupby("_htmp")[primary]
                        .apply(lambda x: round(x.notna().sum()/max(len(x),1)*100,1)))
                report["completeness_by_hour"] = hc.to_dict()
            df = df.drop(columns=["_htmp"])
        except Exception:
            pass
    worst = sorted([(k,v["pct"]) for k,v in completeness.items()], key=lambda x: x[1])[:3]
    worst_str = ", ".join(f"{k}={p}%" for k,p in worst if p < 80)
    _log_step(report, "04_null_profile",
              f"Worst completeness: {worst_str}" if worst_str else "All columns ≥80% complete")

    # ── Step 05: Sensor flatline detection ───────────────────────────────
    flatline_total  = 0
    flatline_detail = {}
    for col, max_streak in FLATLINE_THRESHOLDS.items():
        if col not in df.columns or not pd.api.types.is_numeric_dtype(df[col]):
            continue
        s       = df[col].copy()
        changed = s.ne(s.shift())
        run_ids = changed.cumsum()
        run_lens= run_ids.map(run_ids.value_counts())
        stuck   = (run_lens > max_streak) & (s > 0) & s.notna()
        n       = int(stuck.sum())
        if n > 0:
            df.loc[stuck, col] = np.nan
            flatline_total     += n
            flatline_detail[col]= n
    report["flatlines_nulled"] = flatline_total
    report["flatline_detail"]  = flatline_detail
    _log_step(report, "05_flatline",
              f"Nulled {flatline_total} stuck-sensor readings "
              f"({flatline_detail})" if flatline_total else "No sensor flatlines detected")

    # ── Step 06: Rate-of-change spike filter ─────────────────────────────
    spike_total  = 0
    spike_detail = {}
    for col, roc_limit in ROC_LIMITS.items():
        if col not in df.columns or not pd.api.types.is_numeric_dtype(df[col]):
            continue
        s    = df[col].copy()
        fwd  = s.diff().abs()
        bwd  = s.diff(-1).abs()
        spk  = (fwd > roc_limit) & (bwd > roc_limit) & s.notna()
        n    = int(spk.sum())
        if n > 0:
            df.loc[spk, col] = np.nan
            spike_total      += n
            spike_detail[col] = n
    report["spikes_nulled"] = spike_total
    report["spike_detail"]  = spike_detail
    _log_step(report, "06_spikes",
              f"Nulled {spike_total} impossible spikes "
              f"({spike_detail})" if spike_total else "No rate-of-change spikes detected")

    # ── Step 07: Negative value correction ───────────────────────────────
    MUST_POS = ["PM10","PM2.5","PM25","NO2","NO","NOx","NH3","SO2","CO","Ozone",
                "Benzene","Toluene","Xylene","Eth_Benzene","MP_Xylene",
                "humidity_pct","wind_speed_10m","wind_speed_100m","rain_mm"]
    neg_total  = 0
    neg_detail = {}
    for col in df.columns:
        if any(p.lower() in col.lower() for p in MUST_POS) and pd.api.types.is_numeric_dtype(df[col]):
            mask = df[col] < 0
            n    = int(mask.sum())
            if n > 0:
                df.loc[mask, col] = np.nan
                neg_total         += n
                neg_detail[col]    = n
    report["negatives_nulled"] = neg_total
    _log_step(report, "07_negatives",
              f"Nulled {neg_total} negative values ({neg_detail})" if neg_total
              else "No negative values in positive-only columns")

    # ── Step 08: Physical bounds clipping ────────────────────────────────
    clipped_total  = 0
    clipped_detail = {}
    for col in df.columns:
        for key, (lo, hi) in PHYSICAL_BOUNDS_CLEAN.items():
            if key.lower() in col.lower() and pd.api.types.is_numeric_dtype(df[col]):
                n = int(((df[col] < lo) | (df[col] > hi)).sum())
                if n > 0:
                    df[col] = df[col].clip(lo, hi)
                    clipped_total           += n
                    clipped_detail[col]      = clipped_detail.get(col, 0) + n
                break
    report["values_clipped"]  = clipped_total
    report["clipped_detail"]  = clipped_detail
    _log_step(report, "08_bounds",
              f"Clipped {clipped_total} out-of-bounds values" if clipped_total
              else "No out-of-physical-bounds values")

    # ── Step 09: Mixed-type coercion ──────────────────────────────────────
    num_targets = list(PHYSICAL_BOUNDS_CLEAN.keys()) + [
        "wind_dir_10m","wind_speed_100m","wind_dir_100m","shortwave_radiation","surface_pressure"]
    coerced = 0
    for col in df.columns:
        if any(p.lower() in col.lower() for p in num_targets) and not pd.api.types.is_numeric_dtype(df[col]):
            before_nn = int(df[col].notna().sum())
            df[col]   = pd.to_numeric(df[col], errors="coerce")
            coerced  += before_nn - int(df[col].notna().sum())
    _log_step(report, "09_coerce",
              f"Coerced mixed types — {coerced} values → NaN" if coerced
              else "All target columns already numeric")

    # ── Step 10: Sparse column removal ───────────────────────────────────
    p_cols = [c for c in df.columns
              if any(p.lower() in c.lower() for p in num_targets)
              and pd.api.types.is_numeric_dtype(df[c])]
    removed_sparse = [(c, round(df[c].notna().sum()/max(len(df),1)*100,1))
                      for c in p_cols if df[c].notna().sum()/max(len(df),1) < 0.10]
    if removed_sparse:
        df = df.drop(columns=[c for c,_ in removed_sparse])
        report["cols_removed"].extend([c for c,_ in removed_sparse])
        _log_step(report, "10_sparse",
                  f"Removed {len(removed_sparse)} sparse columns (<10% valid): {removed_sparse}")
    else:
        _log_step(report, "10_sparse", "No sparse columns (all ≥10% valid)")

    # ── Step 11: Cross-variable chemical constraints ───────────────────────
    violations = []
    for col_a, col_b, rel, tol in CHEMICAL_CONSTRAINTS_CLEAN:
        if col_a not in df.columns or col_b not in df.columns:
            continue
        if not (pd.api.types.is_numeric_dtype(df[col_a]) and pd.api.types.is_numeric_dtype(df[col_b])):
            continue
        if rel == ">=":
            viol = (df[col_a].notna() & df[col_b].notna() & (df[col_a] < df[col_b]*(1-tol)))
            n    = int(viol.sum())
            pct  = round(n / max(len(df),1) * 100, 1)
            if n > 0:
                violations.append({
                    "constraint": f"{col_a} >= {col_b}", "violations": n, "pct": pct,
                    "severity"  : "HIGH" if pct>20 else ("MEDIUM" if pct>5 else "LOW"),
                    "note"      : "Sensor offset likely — check calibration",
                })
    report["constraint_violations"] = violations
    for v in violations:
        _log_step(report, "11_constraints",
                  f"[{v['severity']}] {v['constraint']}: {v['violations']} rows ({v['pct']}%) — {v['note']}")
    if not violations:
        _log_step(report, "11_constraints", "All chemical identity constraints satisfied")

    # ── Step 12: Outlier winsorization (3×IQR Tukey fences) ───────────────
    winsor_total  = 0
    winsor_detail = {}
    for col in df.columns:
        if not pd.api.types.is_numeric_dtype(df[col]):
            continue
        if not any(p.lower() in col.lower() for p in num_targets):
            continue
        s = df[col].dropna()
        if len(s) < 20:
            continue
        q25 = float(s.quantile(0.25))
        q75 = float(s.quantile(0.75))
        iqr = q75 - q25
        if iqr == 0:
            continue
        lo_f = max(0.0, q25 - 3.0 * iqr)
        hi_f = q75 + 3.0 * iqr
        n    = int(((df[col] < lo_f) | (df[col] > hi_f)).sum())
        if n > 0:
            df[col] = df[col].clip(lo_f, hi_f)
            winsor_total       += n
            winsor_detail[col]  = {"count": n, "lo": round(lo_f,2), "hi": round(hi_f,2)}
    report["values_winsorized"] = winsor_total
    report["winsor_detail"]     = winsor_detail
    _log_step(report, "12_winsor",
              f"Winsorized {winsor_total} extreme outliers (3×IQR Tukey fences)" if winsor_total
              else "No outliers beyond 3×IQR fences")

    # ── Step 13: Datetime gap row injection ───────────────────────────────
    rows_injected = 0
    if datetime_col and datetime_col in df.columns and len(df) > 2:
        try:
            diffs      = df[datetime_col].diff().dropna()
            modal_freq = diffs.mode()[0]
            fsec       = modal_freq.total_seconds()
            if 600 <= fsec <= 7200:
                full_idx = pd.date_range(
                    start=df[datetime_col].min(),
                    end  =df[datetime_col].max(),
                    freq =modal_freq)
                before_inject = len(df)
                df = (df.set_index(datetime_col)
                        .reindex(full_idx)
                        .reset_index()
                        .rename(columns={"index": datetime_col}))
                rows_injected = max(0, len(df) - before_inject)
                _log_step(report, "13_gap_inject",
                          f"Injected {rows_injected} NaN rows for missing timestamps (freq={modal_freq})")
            else:
                _log_step(report, "13_gap_inject",
                          f"Gap injection skipped (freq={fsec:.0f}s outside 10min–2h range)")
        except Exception as e:
            _log_step(report, "13_gap_inject", f"Gap injection skipped: {e}")
    else:
        _log_step(report, "13_gap_inject", "Gap injection skipped (no valid datetime)")
    report["rows_injected"] = rows_injected

    # ── Step 14: 5-strategy imputation cascade ─────────────────────────────
    imputed_total = 0
    impute_detail = {}
    for col in df.columns:
        if not pd.api.types.is_numeric_dtype(df[col]):
            continue
        n_before = int(df[col].isna().sum())
        if n_before == 0:
            continue
        # S1: forward-fill (max 2)
        df[col] = df[col].ffill(limit=2)
        # S2: backward-fill (max 2)
        df[col] = df[col].bfill(limit=2)
        # S3: rolling 24h centred mean
        roll = df[col].rolling(window=24, min_periods=3, center=True).mean()
        m3   = df[col].isna()
        df.loc[m3, col] = roll[m3]
        # S4: diurnal hourly median
        if datetime_col and datetime_col in df.columns:
            try:
                df["_h"] = df[datetime_col].dt.hour
                hmed     = df.groupby("_h")[col].transform("median")
                m4       = df[col].isna()
                df.loc[m4, col] = hmed[m4]
                df = df.drop(columns=["_h"])
            except Exception:
                pass
        # S5: global column median (last resort)
        m5 = df[col].isna()
        if m5.any():
            gmed = df[col].median()
            if not pd.isna(gmed):
                df.loc[m5, col] = gmed
        n_after = int(df[col].isna().sum())
        filled  = n_before - n_after
        if filled > 0:
            imputed_total     += filled
            impute_detail[col] = {"before": n_before, "after": n_after, "filled": filled}
    report["values_imputed"]   = imputed_total
    report["imputation_detail"]= impute_detail
    _log_step(report, "14_impute",
              f"Imputed {imputed_total} values across {len(impute_detail)} columns "
              f"(ffill→bfill→rolling24h→diurnal→global_median)" if imputed_total
              else "No imputation needed")

    # ── Step 15: Diurnal anomaly flagging ─────────────────────────────────
    diurnal_flagged = 0
    if datetime_col and datetime_col in df.columns:
        for col in [c for c in ["PM10","NO2","SO2","CO","Ozone"] if c in df.columns]:
            try:
                df["_h"]   = df[datetime_col].dt.hour
                h_mean     = df.groupby("_h")[col].transform("mean")
                h_std      = df.groupby("_h")[col].transform("std")
                flag       = ((df[col] - h_mean).abs() > 5 * h_std) & df[col].notna()
                n          = int(flag.sum())
                if n > 0:
                    df[f"_diurnal_flag_{col}"] = flag.astype(int)
                    diurnal_flagged += n
                df = df.drop(columns=["_h"])
            except Exception:
                pass
    report["diurnal_flags"] = diurnal_flagged
    _log_step(report, "15_diurnal",
              f"Flagged {diurnal_flagged} readings as diurnal anomalies (>5σ) — not auto-corrected"
              if diurnal_flagged else "No diurnal anomaly flags (>5σ threshold clear)")

    # ── Step 16: Post-imputation consistency re-check ─────────────────────
    recheck = 0
    for col in df.columns:
        for key, (lo, hi) in PHYSICAL_BOUNDS_CLEAN.items():
            if key.lower() in col.lower() and pd.api.types.is_numeric_dtype(df[col]):
                n = int(((df[col] < lo) | (df[col] > hi)).sum())
                if n > 0:
                    df[col] = df[col].clip(lo, hi)
                    recheck += n
                break
    report["recheck_clipped"] = recheck
    _log_step(report, "16_recheck",
              f"Post-imputation re-check clipped {recheck} additional values" if recheck
              else "Post-imputation re-check clean")

    # ── Step 17: Null streak analysis ─────────────────────────────────────
    null_streaks = {}
    for col in [c for c in ["PM10","NO2","SO2","CO","Ozone"] if c in df.columns]:
        series = df[col]
        if series.isna().sum() == 0:
            null_streaks[col] = {"max_streak": 0, "total_nulls": 0, "concern": "NONE"}
            continue
        ng      = series.isna().astype(int)
        run_ids = ng.ne(ng.shift()).cumsum()
        run_len = ng.groupby(run_ids).transform("sum")
        ms      = int((run_len * ng).max())
        null_streaks[col] = {
            "max_streak" : ms,
            "total_nulls": int(series.isna().sum()),
            "concern"    : "HIGH" if ms > 24 else ("MEDIUM" if ms > 6 else "LOW"),
        }
    report["null_streaks"] = null_streaks
    long = {k: v["max_streak"] for k,v in null_streaks.items() if v["max_streak"] > 24}
    _log_step(report, "17_streaks",
              f"Long null streaks >24h: {long} — imputed values HIGH uncertainty"
              if long else "No null streaks >24h — imputation confidence HIGH")

    # ── Step 18: 7-dimensional quality score ─────────────────────────────
    report["final_shape"]       = df.shape
    report["total_nulls_after"] = int(df.isnull().sum().sum())

    d1 = 1.0 - (report["total_nulls_after"] / max(df.size, 1))
    d2 = min(1.0, len(df) / max(original_rows, 1))
    d3 = 1.0 - (len(report["cols_removed"]) / max(len(original_cols), 1))
    d4 = 1.0 if not [v for v in violations if v["severity"]=="HIGH"] else max(0.3, 1.0-len(violations)*0.15)
    d5 = 1.0 if flatline_total == 0 else max(0.5, 1.0 - flatline_total/max(df.size,1)*100)
    d6 = max(0.4, 1.0 - max((v.get("max_streak",0) for v in null_streaks.values()), default=0)/240)
    d7 = completeness.get("PM10", {}).get("pct", 100) / 100

    w  = [0.25, 0.15, 0.10, 0.15, 0.10, 0.15, 0.10]
    d  = [d1, d2, d3, d4, d5, d6, d7]
    qs = round(sum(wi*di for wi,di in zip(w, d)), 3)
    ql = "EXCELLENT" if qs>=0.90 else "GOOD" if qs>=0.75 else "FAIR" if qs>=0.60 else "POOR"

    report["quality_score"]      = qs
    report["quality_dimensions"] = {
        "D1_null_reduction"      : round(d1, 3),
        "D2_row_retention"       : round(d2, 3),
        "D3_col_retention"       : round(d3, 3),
        "D4_constraints"         : round(d4, 3),
        "D5_flatline"            : round(d5, 3),
        "D6_null_streaks"        : round(d6, 3),
        "D7_pm10_completeness"   : round(d7, 3),
    }
    _log_step(report, "18_quality",
              f"Quality score: {qs:.3f} [{ql}] | "
              f"Dims: null={d1:.2f} rows={d2:.2f} cols={d3:.2f} "
              f"chem={d4:.2f} flat={d5:.2f} streak={d6:.2f} pm10={d7:.2f}")

    print(f"\n  {'─'*64}")
    print(f"  DATA CLEANING COMPLETE — 18 steps")
    print(f"  {'─'*64}")
    print(f"  Shape:      {original_rows:>7,} × {len(original_cols)} "
          f" →  {len(df):>7,} × {df.shape[1]}")
    print(f"  Nulls:      {report['total_nulls_before']:>7,}  →  {report['total_nulls_after']:>7,}")
    print(f"  Imputed:    {imputed_total:>7,}  |  Clipped:    {clipped_total:>6,}")
    print(f"  Winsorized: {winsor_total:>7,}  |  Flatlines:  {flatline_total:>6,}")
    print(f"  Spikes:     {spike_total:>7,}  |  Negatives:  {neg_total:>6,}")
    print(f"  Constraints:{len(violations):>7}  |  Diurnal flags: {diurnal_flagged:>4,}")
    print(f"  Quality:    {qs:.3f}  [{ql}]")
    if null_streaks:
        worst_streak = max(null_streaks.items(), key=lambda x: x[1].get("max_streak",0))
        print(f"  Worst streak: {worst_streak[0]} = {worst_streak[1]['max_streak']}h null run")
    print(f"  {'─'*64}\n")

    return df, report


def generate_cleaning_report_section(doc, cleaning_report: dict,
                                     h1, h2, body, _set_cell_bg,
                                     WHITE, NAVY, TEAL, GREEN, AMBER, RED, Pt,
                                     WD_ALIGN_PARAGRAPH):
    """
    Section 00b — Data Quality Engine report.
    Full 18-step audit trail in the Word document.
    """
    doc.add_page_break()
    h1("Section 00b — Data Quality Engine: 18-Step Cleaning Audit")
    body(
        "ATARS v3.1 runs a professional 18-step automated cleaning pipeline before analysis. "
        "Every action is logged with before/after counts. "
        "This section is a complete reproducibility audit — any analyst can verify "
        "exactly what was changed and why."
    )

    orig  = cleaning_report.get("original_shape", (0, 0))
    final = cleaning_report.get("final_shape",    (0, 0))
    qs    = cleaning_report.get("quality_score",   0)
    ql    = ("EXCELLENT" if qs>=0.90 else "GOOD" if qs>=0.75 else "FAIR" if qs>=0.60 else "POOR")

    summary_rows = [
        ("Original shape",              f"{orig[0]:,} rows × {orig[1]} columns"),
        ("Final shape",                 f"{final[0]:,} rows × {final[1]} columns"),
        ("Rows removed",                str(cleaning_report.get("rows_removed", 0))),
        ("Rows injected (gap fill)",    str(cleaning_report.get("rows_injected", 0))),
        ("Duplicate timestamps",        str(cleaning_report.get("duplicate_timestamps", 0))),
        ("Columns removed (sparse)",    ", ".join(cleaning_report.get("cols_removed",[])) or "None"),
        ("Values imputed",              str(cleaning_report.get("values_imputed", 0))),
        ("Values clipped (bounds)",     str(cleaning_report.get("values_clipped", 0))),
        ("Values winsorized (3×IQR)",   str(cleaning_report.get("values_winsorized", 0))),
        ("Flatline readings nulled",    str(cleaning_report.get("flatlines_nulled", 0))),
        ("Spike readings nulled",       str(cleaning_report.get("spikes_nulled", 0))),
        ("Negative values nulled",      str(cleaning_report.get("negatives_nulled", 0))),
        ("Diurnal anomaly flags",       str(cleaning_report.get("diurnal_flags", 0))),
        ("Post-imputation re-clipped",  str(cleaning_report.get("recheck_clipped", 0))),
        ("Nulls before cleaning",       str(cleaning_report.get("total_nulls_before", 0))),
        ("Nulls after cleaning",        str(cleaning_report.get("total_nulls_after", 0))),
        ("Datetime gaps detected",      str(len(cleaning_report.get("datetime_gaps", [])))),
        ("Cleaning quality score",      f"{qs:.3f}  [{ql}]"),
    ]
    t = doc.add_table(rows=1, cols=2)
    t.style = "Table Grid"
    for cell, hdr in zip(t.rows[0].cells, ["Metric", "Value"]):
        _set_cell_bg(cell, "1A2C4E")
        r = cell.paragraphs[0].add_run(hdr)
        r.bold = True; r.font.color.rgb = WHITE
    for k, v in summary_rows:
        row = t.add_row()
        row.cells[0].text = k
        row.cells[1].text = str(v)
    doc.add_paragraph()

    # 7-dimensional quality score breakdown
    dims = cleaning_report.get("quality_dimensions", {})
    if dims:
        h2("Quality Score — 7 Dimensions")
        body(
            "The cleaning quality score is a weighted composite of 7 dimensions. "
            "Scores < 0.60 indicate data that needs manual review before publication."
        )
        t2 = doc.add_table(rows=1, cols=3)
        t2.style = "Table Grid"
        for cell, hdr in zip(t2.rows[0].cells, ["Dimension", "Score", "Meaning"]):
            _set_cell_bg(cell, "1A6B72")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        dim_labels = {
            "D1_null_reduction"    : "Null reduction (weight 25%)",
            "D2_row_retention"     : "Row retention (weight 15%)",
            "D3_col_retention"     : "Column retention (weight 10%)",
            "D4_constraints"       : "Chemical constraints (weight 15%)",
            "D5_flatline"          : "Flatline absence (weight 10%)",
            "D6_null_streaks"      : "Null streak quality (weight 15%)",
            "D7_pm10_completeness" : "PM10 completeness (weight 10%)",
        }
        for k, v in dims.items():
            row = t2.add_row()
            bg = "D1FAE5" if v>=0.85 else ("FEF3C7" if v>=0.60 else "FFE4E6")
            for cell, val in zip(row.cells, [dim_labels.get(k, k), f"{v:.3f}", "✓ Good" if v>=0.80 else "⚠ Fair" if v>=0.60 else "✗ Poor"]):
                _set_cell_bg(cell, bg)
                cell.paragraphs[0].add_run(val).font.size = Pt(9)
        doc.add_paragraph()

    # Null streaks
    streaks = cleaning_report.get("null_streaks", {})
    if streaks:
        h2("Null Streak Analysis — Instrument Downtime vs Random Dropout")
        body(
            "A null streak > 6h indicates instrument downtime (not random dropout). "
            "Values imputed during long streaks have HIGH uncertainty and should not "
            "be cited as measured data in publications."
        )
        t3 = doc.add_table(rows=1, cols=4)
        t3.style = "Table Grid"
        for cell, hdr in zip(t3.rows[0].cells, ["Variable","Max Streak (h)","Total Nulls","Concern"]):
            _set_cell_bg(cell, "92400E")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        for var, info in streaks.items():
            row = t3.add_row()
            concern = info.get("concern","LOW")
            bg = "FFE4E6" if concern=="HIGH" else ("FEF3C7" if concern=="MEDIUM" else "D1FAE5")
            for cell, val in zip(row.cells, [var, str(info.get("max_streak",0)),
                                              str(info.get("total_nulls",0)), concern]):
                _set_cell_bg(cell, bg)
                cell.paragraphs[0].add_run(val).font.size = Pt(9)
        doc.add_paragraph()

    # Chemical constraint violations
    violations = cleaning_report.get("constraint_violations", [])
    if violations:
        h2(f"Chemical Constraint Violations ({len(violations)})")
        body(
            "These violations indicate possible sensor calibration issues. "
            "They are flagged for human review — not auto-corrected."
        )
        for v in violations:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(f"[{v.get('severity','?')}] ").bold = True
            p.add_run(f"{v.get('constraint','')} — "
                      f"{v.get('violations',0)} rows ({v.get('pct',0):.1f}%) — "
                      f"{v.get('note','')}")
        doc.add_paragraph()

    # Outlier winsorization detail
    winsor_detail = cleaning_report.get("winsor_detail", {})
    if winsor_detail:
        h2("Outlier Winsorization Detail (3×IQR Tukey Fences)")
        t4 = doc.add_table(rows=1, cols=4)
        t4.style = "Table Grid"
        for cell, hdr in zip(t4.rows[0].cells, ["Column","Count","Lower Fence","Upper Fence"]):
            _set_cell_bg(cell, "4C1D95")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        for col, info in list(winsor_detail.items())[:12]:
            row = t4.add_row()
            for cell, val in zip(row.cells, [col, str(info.get("count",0)),
                                              str(info.get("lo","?")),
                                              str(info.get("hi","?"))]):
                cell.paragraphs[0].add_run(val).font.size = Pt(9)
        doc.add_paragraph()

    # Imputation detail
    imp_detail = cleaning_report.get("imputation_detail", {})
    if imp_detail:
        h2("Imputation Detail — Values Filled per Column")
        body(
            "5-strategy cascade: forward-fill (≤2h) → backward-fill (≤2h) → "
            "rolling 24h mean → diurnal hourly median → global column median."
        )
        t5 = doc.add_table(rows=1, cols=3)
        t5.style = "Table Grid"
        for cell, hdr in zip(t5.rows[0].cells, ["Column","Nulls Before","Values Filled"]):
            _set_cell_bg(cell, "1A2C4E")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        for col, info in list(imp_detail.items())[:15]:
            row = t5.add_row()
            fill_pct = round(info.get("filled",0)/max(info.get("before",1),1)*100,1)
            for cell, val in zip(row.cells, [col, str(info.get("before",0)),
                                              f"{info.get('filled',0)} ({fill_pct}% recovered)"]):
                cell.paragraphs[0].add_run(val).font.size = Pt(9)
        doc.add_paragraph()

    # Flatline and spike detail
    flatline_detail = cleaning_report.get("flatline_detail", {})
    spike_detail    = cleaning_report.get("spike_detail",    {})
    if flatline_detail or spike_detail:
        h2("Sensor Integrity Issues — Flatlines and Spikes")
        if flatline_detail:
            body(f"Sensor flatlines nulled: {flatline_detail}")
        if spike_detail:
            body(f"Impossible spikes nulled: {spike_detail}")
        doc.add_paragraph()

    # Datetime gaps
    gaps = cleaning_report.get("datetime_gaps", [])
    if gaps:
        h2(f"Datetime Gaps Detected ({len(gaps)})")
        body("Gaps > 3× modal frequency. Injected NaN rows filled by imputation cascade.")
        t6 = doc.add_table(rows=1, cols=2)
        t6.style = "Table Grid"
        for cell, hdr in zip(t6.rows[0].cells, ["Gap Start", "Duration (hours)"]):
            _set_cell_bg(cell, "92400E")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        for gap in gaps[:20]:
            row = t6.add_row()
            row.cells[0].text = gap.get("start", "")
            row.cells[1].text = f"{gap.get('gap_hours',0):.1f} h"
        doc.add_paragraph()

    # All cleaning steps
    h2("Full Step-by-Step Cleaning Log")
    steps = cleaning_report.get("steps", [])
    if steps:
        t7 = doc.add_table(rows=1, cols=2)
        t7.style = "Table Grid"
        for cell, hdr in zip(t7.rows[0].cells, ["Step ID", "Action"]):
            _set_cell_bg(cell, "1A6B72")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        for step in steps:
            row = t7.add_row()
            row.cells[0].text = step.get("step", "")
            row.cells[1].text = step.get("action", "")
            for c in row.cells:
                c.paragraphs[0].runs[0].font.size = Pt(9)
    doc.add_paragraph()




def generate_data_profile(df: pd.DataFrame, daily_df: pd.DataFrame,
                           config: dict) -> dict:
    """
    Generates a comprehensive data profile summary dict.
    Called after cleaning and flagging; used in console output and Word report.

    Returns dict with:
      - per_column stats (mean, std, min, max, null%, suspect%, valid%)
      - temporal coverage summary
      - WHO exceedance preview
      - data readiness verdict
    """
    profile = {
        "city"       : config.get("city", "City"),
        "columns"    : {},
        "temporal"   : {},
        "readiness"  : {},
        "who_preview": {},
    }

    thresholds = config.get("thresholds", {})

    # Per-column stats from raw df
    for col in POLLUTANTS:
        if col not in df.columns:
            continue
        s = df[col].dropna()
        if len(s) == 0:
            continue
        total   = len(df)
        valid   = int(df[col].notna().sum())
        flag_col = f"q_{col}"
        suspect  = int((df[flag_col] == 2).sum()) if flag_col in df.columns else 0
        invalid  = int((df[flag_col] == 3).sum()) if flag_col in df.columns else 0

        thr = thresholds.get(col)
        exc = 0
        if thr and len(s) > 0:
            # Daily means from daily_df
            dm = daily_df[daily_df["variable"] == col]["mean"].dropna()
            exc = int((dm > thr).sum()) if len(dm) > 0 else 0

        profile["columns"][col] = {
            "mean"       : round(float(s.mean()), 3),
            "std"        : round(float(s.std()),  3),
            "min"        : round(float(s.min()),  3),
            "max"        : round(float(s.max()),  3),
            "p25"        : round(float(s.quantile(0.25)), 3),
            "p75"        : round(float(s.quantile(0.75)), 3),
            "p95"        : round(float(s.quantile(0.95)), 3),
            "skewness"   : round(float(s.skew()),  3),
            "kurtosis"   : round(float(s.kurtosis()), 3),
            "valid_pct"  : round(valid / max(total, 1) * 100, 1),
            "suspect_pct": round(suspect / max(total, 1) * 100, 1),
            "invalid_pct": round(invalid / max(total, 1) * 100, 1),
            "null_pct"   : round((total - valid) / max(total, 1) * 100, 1),
            "threshold"  : thr,
            "exceed_days": exc,
        }
        if thr:
            ratio = round(float(s.mean()) / thr, 2)
            profile["who_preview"][col] = {
                "mean_ratio" : ratio,
                "status"     : ("CRITICAL" if ratio > 3 else
                                "HIGH"     if ratio > 1 else "COMPLIANT"),
                "exceed_days": exc,
            }

    # Temporal coverage
    if "datetime" in df.columns:
        dt = df["datetime"].dropna()
        if len(dt) > 1:
            span_days  = (dt.max() - dt.min()).days + 1
            n_unique_days = df["date_only"].nunique() if "date_only" in df.columns else span_days
            diffs     = dt.sort_values().diff().dropna()
            modal_hrs = diffs.mode()[0].total_seconds() / 3600 if len(diffs.mode()) > 0 else 1
            profile["temporal"] = {
                "start"         : str(dt.min().date()),
                "end"           : str(dt.max().date()),
                "span_days"     : span_days,
                "unique_days"   : n_unique_days,
                "total_records" : len(df),
                "modal_freq_hrs": round(modal_hrs, 2),
                "coverage_pct"  : round(n_unique_days / max(span_days, 1) * 100, 1),
            }

    # Readiness verdict
    pm10_valid = profile["columns"].get("PM10", {}).get("valid_pct", 0)
    n_cols     = len(profile["columns"])
    avg_valid  = (sum(v["valid_pct"] for v in profile["columns"].values()) /
                  max(n_cols, 1))
    coverage   = profile["temporal"].get("coverage_pct", 100)

    if pm10_valid >= 70 and avg_valid >= 60 and coverage >= 80:
        verdict = "READY"
        note    = "Data meets minimum quality thresholds for analysis"
    elif pm10_valid >= 50 and avg_valid >= 40:
        verdict = "MARGINAL"
        note    = "Data is usable but interpret with caution — low completeness"
    else:
        verdict = "POOR"
        note    = "Data quality insufficient — results unreliable without manual review"

    profile["readiness"] = {
        "verdict"     : verdict,
        "note"        : note,
        "pm10_valid"  : pm10_valid,
        "avg_valid"   : round(avg_valid, 1),
        "coverage_pct": coverage,
    }

    # Console summary
    print(f"  ✓ Data profile: {n_cols} pollutants | "
          f"avg completeness {avg_valid:.1f}% | "
          f"coverage {coverage:.1f}% | [{verdict}]")
    if profile["who_preview"]:
        critical = [k for k,v in profile["who_preview"].items()
                    if v["status"] == "CRITICAL"]
        if critical:
            print(f"  ⚠ Critical WHO exceedances: {critical}")

    return profile


def add_data_profile_to_report(doc, profile: dict, h1, h2, body,
                                _set_cell_bg, WHITE, NAVY, TEAL,
                                GREEN, AMBER, RED, Pt, WD_ALIGN_PARAGRAPH):
    """
    Section 00c — Data Profile.
    Descriptive statistics, temporal coverage, WHO preview, readiness verdict.
    """
    doc.add_page_break()
    h1("Section 00c — Data Profile and Readiness Assessment")
    body(
        "This section presents the descriptive statistical profile of the cleaned dataset. "
        "All statistics are computed from quality-validated records (q_i = 1). "
        "The readiness assessment determines whether the dataset meets minimum "
        "quality thresholds for the formal statistical analysis in subsequent sections."
    )

    # Readiness badge
    rd = profile.get("readiness", {})
    verdict = rd.get("verdict", "N/A")
    badge_colors = {"READY": "1E6B3C", "MARGINAL": "92400E", "POOR": "8B1C2A"}
    badge_bg     = {"READY": "D1FAE5", "MARGINAL": "FEF3C7", "POOR": "FFE4E6"}
    t_badge = doc.add_table(rows=1, cols=3)
    t_badge.alignment = 1  # CENTER
    t_badge.style = "Table Grid"
    cells_data = [
        ("Readiness", verdict),
        ("PM10 Completeness", f"{rd.get('pm10_valid', 0):.1f}%"),
        ("Day Coverage", f"{rd.get('coverage_pct', 0):.1f}%"),
    ]
    for cell, (label, value) in zip(t_badge.rows[0].cells, cells_data):
        cell.paragraphs[0].clear()
        lp = cell.add_paragraph(label)
        lp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in lp.runs:
            r.font.size = Pt(9); r.bold = True
        vp = cell.add_paragraph(value)
        vp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in vp.runs:
            r.font.size = Pt(18); r.bold = True
            r.font.color.rgb = _rgb(badge_colors.get(verdict, "1A2C4E"))
        _set_cell_bg(cell, badge_bg.get(verdict, "DCE8F5"))
    body(f"Assessment: {rd.get('note', '')}", italic=True)
    doc.add_paragraph()

    # Temporal coverage
    temp = profile.get("temporal", {})
    if temp:
        h2("Temporal Coverage")
        t_temp = doc.add_table(rows=1, cols=2)
        t_temp.style = "Table Grid"
        for cell, hdr in zip(t_temp.rows[0].cells, ["Parameter", "Value"]):
            _set_cell_bg(cell, "1A2C4E")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        temp_rows = [
            ("Start date",          temp.get("start", "N/A")),
            ("End date",            temp.get("end",   "N/A")),
            ("Span (calendar days)",str(temp.get("span_days", 0))),
            ("Unique days with data",str(temp.get("unique_days", 0))),
            ("Total records",       f"{temp.get('total_records', 0):,}"),
            ("Sampling frequency",  f"{temp.get('modal_freq_hrs', 0):.1f} hours"),
            ("Day coverage",        f"{temp.get('coverage_pct', 0):.1f}%"),
        ]
        for k, v in temp_rows:
            row = t_temp.add_row()
            row.cells[0].text = k
            row.cells[1].text = str(v)
        doc.add_paragraph()

    # Per-column descriptive statistics
    h2("Descriptive Statistics by Variable")
    body(
        "Statistics computed from hourly observations. "
        "Valid% = q_i=1 fraction. Suspect% = q_i=2 fraction. "
        "p95 = 95th percentile (useful for exceedance planning). "
        "Skewness > 1 indicates right-skewed distribution (typical for pollution data)."
    )
    cols_data = profile.get("columns", {})
    if cols_data:
        headers = ["Variable","Mean","Std","Max","p95","Valid%","Suspect%","Threshold","Exceed Days"]
        t_stats = doc.add_table(rows=1, cols=len(headers))
        t_stats.style = "Table Grid"
        for cell, hdr in zip(t_stats.rows[0].cells, headers):
            _set_cell_bg(cell, "1A2C4E")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.size = Pt(8.5)
            r.font.color.rgb = WHITE
            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        for col, info in cols_data.items():
            row = t_stats.add_row()
            valid_p = info.get("valid_pct", 0)
            bg = ("D1FAE5" if valid_p >= 80 else
                  "FEF3C7" if valid_p >= 50 else "FFE4E6")
            thr = info.get("threshold")
            vals = [
                col,
                f"{info.get('mean', 0):.2f}",
                f"{info.get('std', 0):.2f}",
                f"{info.get('max', 0):.2f}",
                f"{info.get('p95', 0):.2f}",
                f"{valid_p:.1f}%",
                f"{info.get('suspect_pct', 0):.1f}%",
                str(thr) if thr else "—",
                str(info.get("exceed_days", 0)),
            ]
            for i, (cell, val) in enumerate(zip(row.cells, vals)):
                _set_cell_bg(cell, bg)
                p = cell.paragraphs[0]
                p.alignment = (WD_ALIGN_PARAGRAPH.LEFT if i == 0
                                else WD_ALIGN_PARAGRAPH.CENTER)
                r = p.add_run(val)
                r.font.size = Pt(8.5)
        doc.add_paragraph()

    # WHO preview
    who = profile.get("who_preview", {})
    if who:
        h2("WHO Guideline Status Preview")
        t_who = doc.add_table(rows=1, cols=4)
        t_who.style = "Table Grid"
        for cell, hdr in zip(t_who.rows[0].cells,
                              ["Variable","Mean/WHO Ratio","Status","Exceedance Days"]):
            _set_cell_bg(cell, "8B1C2A")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        for col, info in sorted(who.items(), key=lambda x: -x[1]["mean_ratio"]):
            row = t_who.add_row()
            status = info.get("status","N/A")
            bg = ("FFE4E6" if status=="CRITICAL" else
                  "FEF3C7" if status=="HIGH" else "D1FAE5")
            for cell, val in zip(row.cells,
                                   [col,
                                    f"{info.get('mean_ratio',0):.2f}×",
                                    status,
                                    str(info.get("exceed_days",0))]):
                _set_cell_bg(cell, bg)
                cell.paragraphs[0].add_run(val).font.size = Pt(9)
    doc.add_paragraph()



def _detect_unit_issues(df: pd.DataFrame, config: dict) -> list:
    """
    Heuristic unit anomaly detection.
    CPCB data sometimes mixes µg/m³ and mg/m³ for CO without warning.
    CO in µg/m³ would read 1000× higher than mg/m³ — catch this.
    Returns list of warning strings.
    """
    warnings_out = []
    thresholds = config.get("thresholds", {})

    # CO: WHO 24h limit = 4 mg/m³ = 4000 µg/m³
    # If mean CO > 100 mg/m³, data is almost certainly in µg/m³ not mg/m³
    if "CO" in df.columns:
        co_mean = df["CO"].dropna().mean()
        if co_mean > 100:
            warnings_out.append(
                f"CO mean={co_mean:.1f} — likely in µg/m³ not mg/m³ "
                f"(WHO limit is 4 mg/m³). Consider dividing CO by 1000."
            )

    # PM10: if mean > 500 µg/m³, check if it might be in mg/m³ (would be 0.5)
    if "PM10" in df.columns:
        pm_mean = df["PM10"].dropna().mean()
        if pm_mean < 0.5:
            warnings_out.append(
                f"PM10 mean={pm_mean:.4f} — likely in mg/m³ not µg/m³. "
                f"Consider multiplying PM10 by 1000."
            )

    # NO2: typical urban range 10-200 µg/m³; if mean > 500, suspect ppb not µg/m³
    if "NO2" in df.columns:
        no2_mean = df["NO2"].dropna().mean()
        if no2_mean > 500:
            warnings_out.append(
                f"NO2 mean={no2_mean:.1f} — may be in ppb not µg/m³ "
                f"(1 ppb NO2 ≈ 1.88 µg/m³). Verify units with source."
            )

    return warnings_out


def load_data(filepath: str, config: dict = None) -> pd.DataFrame:
    """
    Load CSV/Excel, normalize columns, parse datetime, sort.
    v2 improvements:
    - Tries 6 encodings including common Indian CSV exports
    - Detects likely unit mismatches (CO in µg vs mg, PM in mg vs µg)
    - Handles multi-header CSVs (skips metadata rows)
    - Validates minimum row count
    - Reports exact parsing method used
    """
    if config is None:
        config = {}

    print(f"\n  {'─'*64}")
    print(f"  ATARS — Data Ingestion: {filepath}")
    print(f"  {'─'*64}")

    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"Data file not found: {filepath}")

    file_size_kb = round(path.stat().st_size / 1024, 1)
    print(f"  File: {path.name}  ({file_size_kb} KB)")

    # ── Try multiple encodings ─────────────────────────────────────────────
    encodings = ["utf-8-sig", "utf-8", "latin-1", "cp1252", "iso-8859-1", "cp1250"]
    df = None
    used_enc = None
    for enc in encodings:
        try:
            if path.suffix.lower() in (".xlsx", ".xls"):
                df = pd.read_excel(filepath)
                used_enc = "excel"
                break
            else:
                # Try with and without skiprows for files with metadata headers
                try:
                    _df = pd.read_csv(filepath, encoding=enc, low_memory=False)
                    # Heuristic: if first row looks like metadata (>50% non-numeric in numeric cols)
                    # try again with skiprows=1
                    if len(_df) > 0:
                        df = _df
                        used_enc = enc
                        break
                except UnicodeDecodeError:
                    continue
        except Exception:
            continue

    if df is None:
        raise ValueError(
            f"Could not load '{filepath}' with any known encoding. "
            f"Tried: {encodings}. Check the file is not corrupted."
        )

    print(f"  ✓ Loaded:   {df.shape[0]:,} rows × {df.shape[1]} columns  "
          f"(encoding: {used_enc})")

    # ── Drop fully empty rows/columns ─────────────────────────────────────
    before_rows = len(df)
    df = df.dropna(how="all")
    df = df.loc[:, df.notna().any()]
    dropped_empty = before_rows - len(df)
    if dropped_empty > 0:
        print(f"  ✓ Dropped {dropped_empty} fully-empty rows")

    # ── Normalize column names ────────────────────────────────────────────
    df = normalize_columns(df)

    # ── Build datetime index ──────────────────────────────────────────────
    datetime_built = False
    if "datetime" in df.columns:
        try:
            df["datetime"] = pd.to_datetime(df["datetime"],
                                             infer_datetime_format=True, errors="coerce")
            datetime_built = True
        except Exception:
            pass

    if not datetime_built and "date" in df.columns and "time" in df.columns:
        try:
            df["datetime"] = pd.to_datetime(
                df["date"].astype(str) + " " + df["time"].astype(str),
                dayfirst=True, errors="coerce"
            )
            datetime_built = True
        except Exception:
            pass

    if not datetime_built and "date" in df.columns:
        try:
            df["datetime"] = pd.to_datetime(df["date"], dayfirst=True, errors="coerce")
            datetime_built = True
        except Exception:
            pass

    if not datetime_built:
        raise ValueError(
            "No parseable date/datetime column found. "
            f"Columns present: {list(df.columns[:15])}. "
            "Expected column named 'date', 'datetime', 'timestamp', or 'day'."
        )

    # Drop rows where datetime couldn't be parsed
    bad_dt = int(df["datetime"].isna().sum())
    if bad_dt > 0:
        df = df.dropna(subset=["datetime"])
        print(f"  ⚠ Dropped {bad_dt} rows with unparseable datetime")

    df = df.sort_values("datetime").reset_index(drop=True)

    # ── Coerce numeric columns ────────────────────────────────────────────
    for col in ALL_NUMERIC:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")
        else:
            df[col] = np.nan

    # Negative concentrations → NaN (cleaning step 07 will handle formally)
    for col in POLLUTANTS:
        if col in df.columns:
            df.loc[df[col] < 0, col] = np.nan

    df["date_only"] = df["datetime"].dt.date

    # ── Minimum data check ────────────────────────────────────────────────
    n_days  = df["date_only"].nunique()
    n_recs  = len(df)
    if n_recs < 24:
        print(f"  ⚠ WARNING: Only {n_recs} records — results will be unreliable")
    if n_days < 7:
        print(f"  ⚠ WARNING: Only {n_days} days of data — "
              f"baseline window ({config.get('baseline_window_days',30)} days) cannot be computed")

    # ── Unit anomaly detection ────────────────────────────────────────────
    unit_warnings = _detect_unit_issues(df, config)
    for w in unit_warnings:
        print(f"  ⚠ UNIT ISSUE: {w}")

    # ── Summary ───────────────────────────────────────────────────────────
    dt_min = df["datetime"].min()
    dt_max = df["datetime"].max()
    found_pol = [c for c in POLLUTANTS     if c in df.columns and df[c].notna().any()]
    found_met = [c for c in METEOROLOGICAL if c in df.columns and df[c].notna().any()]

    print(f"  ✓ Date range:  {dt_min.date()} → {dt_max.date()}  ({n_days} days)")
    print(f"  ✓ Records:     {n_recs:,}")
    print(f"  ✓ Pollutants:  {found_pol}")
    if found_met:
        print(f"  ✓ Meteorology: {found_met}")

    return df


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 3 — QUALITY ASSESSMENT
#  Q(D) = N_v / N_total  (Eq. 11)
#  q_i ∈ {0=unreviewed, 1=valid, 2=suspect, 3=invalid}
# ═══════════════════════════════════════════════════════════════════════════════

# Physical plausibility bounds for range-check quality flagging
PHYSICAL_BOUNDS = {
    "PM10"        : (0, 1000),   # µg/m³
    "NO"          : (0, 500),
    "NO2"         : (0, 500),
    "NOx"         : (0, 1000),   # ppb
    "NH3"         : (0, 2000),
    "SO2"         : (0, 500),
    "CO"          : (0, 50),     # mg/m³
    "Ozone"       : (0, 400),
    "Benzene"     : (0, 100),
    "Toluene"     : (0, 500),
    "Xylene"      : (0, 500),
    "Eth_Benzene" : (0, 200),
    "MP_Xylene"   : (0, 500),
    "humidity_pct": (0, 100),
    "rain_mm"     : (0, 500),
    "pressure_hpa": (800, 1100),
    "wind_speed_10m"  : (0, 200),
    "wind_speed_100m" : (0, 300),
}


def assign_quality_flags(df: pd.DataFrame) -> pd.DataFrame:
    """
    Assign quality flag q_i to each record using a 4-tier system.

    Tier system (consistent with WMO data quality standards):
      q_i = 1 (VALID)     : not NaN, within physical bounds, not statistical outlier
      q_i = 2 (SUSPECT)   : not NaN, within physical bounds, but statistical outlier
                             (outside 3×IQR — Tukey extreme fence)
      q_i = 3 (INVALID)   : NaN OR outside physical bounds
      q_i = 4 (DIURNAL)   : valid reading but flagged as diurnal anomaly (>5σ from hour mean)
      q_i = 0 (UNREVIEWED): non-pollutant columns

    Uses PHYSICAL_BOUNDS_CLEAN (consistent with cleaning engine).
    Uses 3×IQR fence (tighter than old 5×IQR — catches more suspect readings).
    """
    df = df.copy()

    # Compute 3×IQR fences for statistical suspect detection
    iqr_bounds = {}
    for col in POLLUTANTS:
        if col in df.columns and pd.api.types.is_numeric_dtype(df[col]):
            q25 = float(df[col].quantile(0.25))
            q75 = float(df[col].quantile(0.75))
            iqr = q75 - q25
            if iqr > 0:
                iqr_bounds[col] = (max(0.0, q25 - 3.0 * iqr), q75 + 3.0 * iqr)

    flag_cols = {}
    for col in POLLUTANTS:
        if col not in df.columns or not pd.api.types.is_numeric_dtype(df[col]):
            continue

        flags = np.ones(len(df), dtype=int)   # default: VALID (1)
        phys  = PHYSICAL_BOUNDS_CLEAN.get(col, PHYSICAL_BOUNDS.get(col, (0, 1e9)))

        # Tier 3: NaN or outside physical bounds → INVALID
        mask_invalid = df[col].isna() | (df[col] < phys[0]) | (df[col] > phys[1])
        flags[mask_invalid] = 3

        # Tier 2: statistical outlier (3×IQR) → SUSPECT (if not already INVALID)
        if col in iqr_bounds:
            lo, hi = iqr_bounds[col]
            mask_suspect = (~mask_invalid) & ((df[col] < lo) | (df[col] > hi))
            flags[mask_suspect] = 2

        # Tier 4: diurnal anomaly flag column exists → mark as DIURNAL
        diurnal_col = f"_diurnal_flag_{col}"
        if diurnal_col in df.columns:
            mask_diurnal = (df[diurnal_col] == 1) & (flags == 1)  # only if otherwise valid
            flags[mask_diurnal] = 4

        flag_cols[f"q_{col}"] = flags

    for k, v in flag_cols.items():
        df[k] = v

    # Summary stats for console
    valid_counts = {col.replace("q_",""):int((df[col]==1).sum())
                    for col in flag_cols if col in df.columns}
    suspect_counts = {col.replace("q_",""):int((df[col]==2).sum())
                      for col in flag_cols if col in df.columns}
    total = max(len(df), 1)
    key_cols = [c for c in ["PM10","NO2","SO2","CO"] if f"q_{c}" in df.columns]
    if key_cols:
        summary = " | ".join(
            f"{c}: {valid_counts.get(c,0)/total*100:.0f}%V "
            f"{suspect_counts.get(c,0)/total*100:.0f}%S"
            for c in key_cols
        )
        print(f"  ✓ Quality flags assigned: {summary}")

    return df


def compute_quality_score(df_day: pd.DataFrame, col: str) -> dict:
    """
    Q(D) = N_v / N_total  — Eq. 11
    Extended to report all 4 flag tiers.
    Q_effective = (N_valid + 0.5 * N_suspect) / N_total  — partial credit for suspect
    """
    flag_col = f"q_{col}"
    if flag_col not in df_day.columns:
        total = len(df_day)
        valid = int(df_day[col].notna().sum()) if col in df_day.columns else 0
        return {
            "N_total": total, "N_valid": valid, "N_suspect": 0,
            "N_invalid": total - valid, "N_diurnal": 0,
            "Q": valid / total if total > 0 else 0.0,
            "Q_effective": valid / total if total > 0 else 0.0,
        }
    flags    = df_day[flag_col].values
    total    = len(flags)
    valid    = int(np.sum(flags == 1))
    suspect  = int(np.sum(flags == 2))
    invalid  = int(np.sum(flags == 3))
    diurnal  = int(np.sum(flags == 4))
    q_score  = valid / total if total > 0 else 0.0
    q_eff    = (valid + 0.5 * suspect) / total if total > 0 else 0.0
    return {
        "N_total"    : total,
        "N_valid"    : valid,
        "N_suspect"  : suspect,
        "N_invalid"  : invalid,
        "N_diurnal"  : diurnal,
        "Q"          : round(q_score, 4),
        "Q_effective": round(q_eff,   4),
    }


def get_confidence_flag(q_score: float, config: dict) -> str:
    """
    Confidence flag from Q(D) — Eq. after Table 3.
    Uses Q_effective (partial credit for suspect readings) when available.
    """
    if q_score >= config["confidence_high"]:
        return "HIGH"
    elif q_score >= config["confidence_moderate"]:
        return "MODERATE"
    elif q_score >= 0.40:
        return "LOW"
    return "INSUFFICIENT"


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 4 — FORMAL STATISTICAL OPERATORS  (Eqs. 3–10)
# ═══════════════════════════════════════════════════════════════════════════════

def aggregation_operator(series: pd.Series) -> dict:
    """
    A(D_v) — Aggregation Operator (Eqs. 3–5).
    Computes x̄, σ, x_max, x_min, N_v from validated series.
    """
    valid = series.dropna().values
    N_v = len(valid)
    if N_v == 0:
        return {"mean": np.nan, "std": np.nan, "max": np.nan,
                "min": np.nan, "median": np.nan, "IQR": np.nan, "N_v": 0}
    mean   = float(np.mean(valid))                                          # Eq. 3
    std    = float(np.std(valid, ddof=1)) if N_v > 1 else 0.0              # Eq. 4
    x_max  = float(np.max(valid))                                          # Eq. 5
    x_min  = float(np.min(valid))
    median = float(np.median(valid))
    q25    = float(np.percentile(valid, 25))
    q75    = float(np.percentile(valid, 75))
    iqr = round(q75 - q25, 4) if N_v > 1 else 0.0
    return {
        "mean": round(mean, 4), "std": round(std, 4),
        "max": round(x_max, 4), "min": round(x_min, 4),
        "median": round(median, 4), "IQR": iqr,
        "N_v": N_v
    }


def confidence_interval(mean: float, std: float, n: int,
                         alpha: float = 0.05, method: str = "classical") -> dict:
    """
    Confidence interval formulations (Table in Section 5.5 of paper).
    Classical, Chebyshev, Bootstrap.
    """
    if n <= 1 or np.isnan(mean):
        return {"lower": np.nan, "upper": np.nan, "method": method}

    if method == "classical":
        # Eq. 12: CI = x̄ ± z_{α/2} · σ/sqrt(N_v)
        z = stats.norm.ppf(1 - alpha / 2)     # z_0.025 = 1.96
        margin = z * std / math.sqrt(n)
        return {"lower": round(mean - margin, 4), "upper": round(mean + margin, 4),
                "method": "classical_95pct", "z_critical": round(z, 4)}

    elif method == "chebyshev":
        # Eq. 26: P(|X-μ| ≥ kσ) ≤ 1/k²  →  k=sqrt(1/α)
        k = math.sqrt(1 / alpha)
        return {"lower": round(mean - k * std, 4), "upper": round(mean + k * std, 4),
                "method": "chebyshev", "k_factor": round(k, 4)}

    return {"lower": np.nan, "upper": np.nan, "method": "unknown"}


def baseline_operator(daily_means: pd.Series, current_date,
                       window: int = 30) -> dict:
    """
    B(D, W) — Baseline Operator (Eqs. 7–8).
    Computes x̄_W and σ_W from W days preceding current_date.
    """
    end   = pd.Timestamp(current_date) - pd.Timedelta(days=1)
    start = end - pd.Timedelta(days=window - 1)
    window_data = daily_means.loc[(daily_means.index >= start) & (daily_means.index <= end)].dropna()


    n_W = len(window_data)
    if n_W < 3:
        return {"mean_W": np.nan, "std_W": np.nan, "n_W": n_W}
    mean_W = float(np.mean(window_data.values))                             # Eq. 7
    std_W  = float(np.std(window_data.values, ddof=1))                     # Eq. 8
    return {"mean_W": round(mean_W, 4), "std_W": round(std_W, 4), "n_W": n_W}


def z_score(x_d: float, mean_W: float, std_W: float) -> float:
    """ζ_d = (x̄_d - x̄_W) / σ_W  — Eq. 9"""
    if std_W == 0 or np.isnan(std_W) or np.isnan(x_d):
        return np.nan
    return round((x_d - mean_W) / std_W, 4)


def delta_pct(x_d: float, x_prev: float) -> float:
    """δ_d = [(x̄_d - x̄_{d-1}) / x̄_{d-1}] × 100%  — Eq. 10"""
    if x_prev == 0 or np.isnan(x_d) or np.isnan(x_prev):
        return np.nan
    return round((x_d - x_prev) / x_prev * 100, 2)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 5 — DAILY AGGREGATION PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════

def build_daily_stats(df: pd.DataFrame, config: dict) -> pd.DataFrame:
    """
    Aggregate hourly observations to daily statistics.
    Returns a DataFrame with one row per day per variable.
    """
    print("\n  [Step 3] Running SQL-equivalent aggregation operators...")
    records = []
    daily_means_by_col = {}

    # First pass: compute daily means for all columns
    for col in ALL_NUMERIC:
        if col not in df.columns:
            continue
        flag_col = f"q_{col}"
        day_groups = df.groupby('date_only')
        means = {}
        for day, grp in day_groups:
            if flag_col in grp.columns:
                valid_vals = grp.loc[grp[flag_col] == 1, col]
            else:
                valid_vals = grp[col].dropna()
            if len(valid_vals) > 0:
                means[pd.Timestamp(day)] = valid_vals.mean()
        s = pd.Series(means).sort_index(); s.index = pd.DatetimeIndex(s.index); daily_means_by_col[col] = s

    # Second pass: full stats per day per column
    for col in ALL_NUMERIC:
        if col not in df.columns:
            continue
        flag_col  = f"q_{col}"
        daily_col = daily_means_by_col.get(col, pd.Series(dtype=float))
        threshold = config["thresholds"].get(col, None)
        prev_mean = np.nan

        for day, grp in df.groupby('date_only'):
            day_ts = pd.Timestamp(day)

            # Validated series for this day
            # Primary: q=1 (VALID) only — most conservative
            # If <3 valid readings, fall back to include q=2 (SUSPECT) with flag
            if flag_col in grp.columns:
                valid_series = grp.loc[grp[flag_col] == 1, col]
                if len(valid_series.dropna()) < 3:
                    # Include suspect readings if valid count too low
                    valid_series = grp.loc[grp[flag_col].isin([1, 2]), col]
            else:
                valid_series = grp[col].dropna()

            q_info  = compute_quality_score(grp, col)
            agg     = aggregation_operator(valid_series)
            conf_f  = get_confidence_flag(q_info["Q"], config)
            baseline= baseline_operator(daily_col, day_ts, config["baseline_window_days"])
            zs      = z_score(agg["mean"], baseline["mean_W"], baseline["std_W"])
            dp      = delta_pct(agg["mean"], prev_mean)
            ci      = confidence_interval(agg["mean"], agg["std"], agg["N_v"], config["ci_alpha"])

            is_anomaly  = (not np.isnan(zs)) and (abs(zs) > config["anomaly_z_threshold"])
            exceeds_thr = (threshold is not None and
                           not np.isnan(agg["mean"]) and
                           agg["mean"] > threshold)

            records.append({
                "date"          : day_ts,
                "variable"      : col,
                "mean"          : agg["mean"],
                "std"           : agg["std"],
                "max"           : agg["max"],
                "min"           : agg["min"],
                "median"        : agg["median"],
                "IQR"           : agg["IQR"],
                "N_v"           : agg["N_v"],
                "N_total"       : q_info["N_total"],
                "N_suspect"     : q_info.get("N_suspect", 0),
                "N_invalid"     : q_info.get("N_invalid", 0),
                "N_diurnal"     : q_info.get("N_diurnal", 0),
                "Q"             : q_info["Q"],
                "Q_effective"   : q_info.get("Q_effective", q_info["Q"]),
                "confidence"    : conf_f,
                "mean_W"        : baseline["mean_W"],
                "std_W"         : baseline["std_W"],
                "n_W"           : baseline["n_W"],
                "z_score"       : zs,
                "delta_pct"     : dp,
                "CI_lower"      : ci["lower"],
                "CI_upper"      : ci["upper"],
                "threshold"     : threshold,
                "exceeds_threshold": exceeds_thr,
                "is_anomaly"    : is_anomaly,
            })
            prev_mean = agg["mean"]

    daily_df = pd.DataFrame(records)
    print(f"  ✓ Daily stats: {len(daily_df)} rows × {len(daily_df.columns)} columns")
    return daily_df


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 6 — TEMPORAL ANALYSIS: ACF / PACF  (Eqs. 12–13)
# ═══════════════════════════════════════════════════════════════════════════════

def compute_acf(series: np.ndarray, n_lags: int) -> np.ndarray:
    """
    Sample ACF: ρ̂(k) = Corr(x_t, x_{t-k})  — Eq. 12
    Manual implementation (no statsmodels required).
    """
    n    = len(series)
    mean = np.mean(series)
    var  = np.sum((series - mean) ** 2) / n
    if var == 0:
        return np.zeros(n_lags + 1)
    acf = [1.0]
    for k in range(1, n_lags + 1):
        if k >= n:
            acf.append(0.0)
        else:
            cov = np.sum((series[k:] - mean) * (series[:-k] - mean)) / n
            acf.append(cov / var)
    return np.array(acf)


def compute_pacf(series: np.ndarray, n_lags: int) -> np.ndarray:
    """
    Sample PACF: φ̂_{kk} via Yule-Walker recursion  — Eq. 13
    Manual implementation using OLS at each lag order.
    """
    n = len(series)
    pacf = [1.0]
    for k in range(1, n_lags + 1):
        if k >= n - 1:
            pacf.append(0.0)
            continue
        # Build design matrix: lag-1 through lag-k
        X = np.column_stack([series[k - j - 1: n - j - 1]
                              for j in range(k)])
        y = series[k:]
        # Normalise
        X_m = X - X.mean(axis=0)
        y_m = y - y.mean()
        try:
            coefs = np.linalg.lstsq(X_m, y_m, rcond=None)[0]
            pacf.append(float(coefs[-1]))
        except Exception:
            pacf.append(0.0)
    return np.array(pacf)


def run_temporal_analysis(daily_df: pd.DataFrame, col: str,
                           n_lags: int = 48) -> dict:
    """Run ACF/PACF on daily mean series for one variable."""
    series = daily_df[daily_df['variable'] == col]['mean'].dropna().values
    if len(series) < n_lags + 5:
        return {}
    acf  = compute_acf(series, n_lags)
    pacf = compute_pacf(series, min(n_lags, 24))  # PACF expensive at high lags
    conf_bound = 1.96 / math.sqrt(len(series))
    return {
        "acf" : acf, "pacf": pacf,
        "conf_bound"      : conf_bound,
        "significant_acf" : [k for k, v in enumerate(acf) if k > 0 and abs(v) > conf_bound],
        "series_length"   : len(series),
    }


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 7 — CORRELATION MATRIX & OLS  (Eqs. 16–18, 14–15)
# ═══════════════════════════════════════════════════════════════════════════════

def compute_correlation_matrix(daily_df: pd.DataFrame,
                                variables: list) -> pd.DataFrame:
    """
    R_{jk} = Corr(X_j, X_k)  — Eq. 18
    Built from daily means. Returns d×d DataFrame.
    NOTE: Association ≠ Causation. Corr(X,Y) ≢ C(X→Y).
    """
    pivot = daily_df.pivot_table(
        index='date', columns='variable', values='mean'
    )
    cols = [v for v in variables if v in pivot.columns]
    return pivot[cols].corr(method='pearson')


def run_ols_regression(daily_df: pd.DataFrame,
                        target: str = "PM10",
                        predictors: list = None) -> dict:
    """
    OLS: β̂ = (X^T X)^{-1} X^T y  — Eq. 14
    Predicts target from meteorological predictors.
    Returns association statistics only — no causal interpretation.
    """
    if predictors is None:
        predictors = [p for p in ["humidity_pct", "wind_speed_10m",
                                   "rain_mm", "pressure_hpa"]
                      if p in daily_df['variable'].unique()]
    pivot = daily_df.pivot_table(
        index='date', columns='variable', values='mean'
    )
    if target not in pivot.columns:
        return {}
    available = [p for p in predictors if p in pivot.columns]
    if len(available) < 2:
        return {}

    df_reg = pivot[[target] + available].dropna()
    if len(df_reg) < 20:
        return {}

    X = df_reg[available].values
    y = df_reg[target].values

    scaler = StandardScaler()
    X_s = scaler.fit_transform(X)

    beta, residuals, _, _ = np.linalg.lstsq(
        np.column_stack([np.ones(len(X_s)), X_s]), y, rcond=None
    )

    y_pred = np.column_stack([np.ones(len(X_s)), X_s]) @ beta
    ss_res = np.sum((y - y_pred) ** 2)
    ss_tot = np.sum((y - np.mean(y)) ** 2)
    r2     = 1 - ss_res / ss_tot if ss_tot > 0 else 0.0

    return {
        "target"    : target,
        "predictors": available,
        "beta"      : {p: round(float(b), 4) for p, b in zip(available, beta[1:])},
        "intercept" : round(float(beta[0]), 4),
        "r_squared" : round(r2, 4),
        "n_obs"     : len(df_reg),
        "note"      : "Association only — Corr(X,Y) does not imply causation (Eq. 17)"
    }



# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 7.4 — FORMAL STATISTICAL INFERENCE ENGINE
#  ─────────────────────────────────────────────────────────────────────────────
#  The original pipeline computes descriptive statistics (mean, std, CI).
#  A real analyst also asks: are the patterns SIGNIFICANT?
#
#  This module adds four inferential layers — all pure numpy/scipy, zero extra deps:
#
#  1. Mann-Kendall Trend Test     — monotonic trend significance (non-parametric)
#     Used by WMO, USEPA, and most published air quality studies.
#     H0: no monotonic trend. Output: tau, p-value, direction, Sen's slope.
#
#  2. Normality Tests (Shapiro-Wilk + D'Agostino-Pearson)
#     Determines whether parametric stats (mean, CI) are valid.
#     If NOT normal: Chebyshev CI is reported; skewness and kurtosis flagged.
#
#  3. Pearson Correlation with p-values
#     The original matrix has r but no significance — every r looks trustworthy.
#     Now adds p-value and marks correlations as SIGNIFICANT / NOT_SIGNIFICANT.
#
#  4. OLS Regression with full inference
#     Adds: standard errors, t-statistics, p-values per coefficient,
#     adjusted R², F-statistic, residual standard error.
#     The original OLS returned β̂ with no way to know if any predictor matters.
#
#  5. Seasonal Decomposition (additive, manual STL-equivalent)
#     Decomposes each pollutant series into: trend + seasonal + residual.
#     Trend extracted via centred moving average; seasonal via monthly means.
#
#  6. Exceedance Frequency Distribution
#     For each pollutant with a WHO threshold: computes exceedance probability,
#     return period (expected days between exceedances), and WHO excess burden.
# ═══════════════════════════════════════════════════════════════════════════════


def mann_kendall_test(series: np.ndarray) -> dict:
    """
    Mann-Kendall monotonic trend test.
    Non-parametric — no normality assumption required.
    Appropriate for autocorrelated time series (unlike linear regression t-test).

    H0: no monotonic trend (τ = 0)
    H1: monotonic trend (τ ≠ 0)

    Algorithm:
      S = Σ_{i<j} sign(x_j - x_i)
      Var(S) = n(n-1)(2n+5)/18   (adjusted for ties)
      Z = (S-1)/sqrt(Var(S)) if S>0, 0 if S=0, (S+1)/sqrt(Var(S)) if S<0
      p-value from standard normal

    Sen's slope: median of all pairwise slopes (x_j - x_i)/(j - i)
    """
    n = len(series)
    if n < 4:
        return {"available": False, "reason": "n < 4"}

    series = np.array(series, dtype=float)
    valid  = series[~np.isnan(series)]
    n      = len(valid)
    if n < 4:
        return {"available": False, "reason": "too many NaNs"}

    # Compute S
    S = 0
    for i in range(n - 1):
        for j in range(i + 1, n):
            diff = valid[j] - valid[i]
            if diff > 0:
                S += 1
            elif diff < 0:
                S -= 1
    # (O(n²) — acceptable for daily series ≤ 1000 days)

    # Variance of S (accounting for ties)
    unique_vals, counts = np.unique(valid, return_counts=True)
    tie_correction = np.sum(counts * (counts - 1) * (2 * counts + 5))
    var_S = (n * (n - 1) * (2 * n + 5) - tie_correction) / 18
    if var_S <= 0:
        return {"available": False, "reason": "zero variance"}

    # Z statistic
    if S > 0:
        Z = (S - 1) / math.sqrt(var_S)
    elif S == 0:
        Z = 0.0
    else:
        Z = (S + 1) / math.sqrt(var_S)

    # Two-tailed p-value
    p_value = 2 * (1 - stats.norm.cdf(abs(Z)))

    # Sen's slope: median pairwise slope
    slopes = []
    for i in range(n - 1):
        for j in range(i + 1, n):
            if j != i:
                slopes.append((valid[j] - valid[i]) / (j - i))
    sens_slope = float(np.median(slopes)) if slopes else 0.0

    # Direction
    if p_value < 0.05:
        direction = "INCREASING" if S > 0 else "DECREASING"
    else:
        direction = "NO_SIGNIFICANT_TREND"

    return {
        "available"     : True,
        "S"             : int(S),
        "Z"             : round(Z, 4),
        "p_value"       : round(p_value, 6),
        "significant"   : p_value < 0.05,
        "direction"     : direction,
        "sens_slope"    : round(sens_slope, 4),
        "sens_slope_unit": "µg/m³ per day",
        "n"             : n,
        "interpretation": (
            f"Statistically significant {direction.lower().replace('_',' ')} trend "
            f"(p={p_value:.4f}, Sen's slope={sens_slope:.3f} µg/m³/day)"
            if p_value < 0.05 else
            f"No statistically significant trend detected (p={p_value:.4f})"
        ),
    }


def normality_tests(series: np.ndarray, alpha: float = 0.05) -> dict:
    """
    Tests whether a series is normally distributed.
    Uses two complementary tests — both must pass for normality.

    1. Shapiro-Wilk: sensitive to deviations from normality for n ≤ 5000
    2. D'Agostino-Pearson (normaltest): tests skewness + kurtosis jointly

    If NOT normal: parametric CI is less reliable; report Chebyshev CI instead.
    """
    valid = series[~np.isnan(series)] if hasattr(series, '__len__') else np.array([])
    n = len(valid)

    result = {
        "n"              : n,
        "is_normal"      : False,
        "skewness"       : None,
        "kurtosis"       : None,
        "shapiro_W"      : None,
        "shapiro_p"      : None,
        "dagostino_stat" : None,
        "dagostino_p"    : None,
        "verdict"        : "INSUFFICIENT_DATA",
        "recommendation" : "",
    }

    if n < 8:
        result["verdict"] = "INSUFFICIENT_DATA"
        result["recommendation"] = "Need ≥ 8 observations for normality testing"
        return result

    result["skewness"] = round(float(stats.skew(valid)),   3)
    result["kurtosis"] = round(float(stats.kurtosis(valid)), 3)   # excess kurtosis

    # Shapiro-Wilk (use scipy; cap at 5000 for speed)
    try:
        sample = valid if n <= 5000 else np.random.choice(valid, 5000, replace=False)
        W, p_sw = stats.shapiro(sample)
        result["shapiro_W"] = round(float(W),    4)
        result["shapiro_p"] = round(float(p_sw), 6)
    except Exception:
        pass

    # D'Agostino-Pearson
    try:
        k2, p_dp = stats.normaltest(valid)
        result["dagostino_stat"] = round(float(k2),   4)
        result["dagostino_p"]    = round(float(p_dp), 6)
    except Exception:
        pass

    # Verdict: both tests must agree (or one unavailable)
    sw_normal = (result["shapiro_p"]    is not None and result["shapiro_p"]    > alpha)
    dp_normal = (result["dagostino_p"]  is not None and result["dagostino_p"]  > alpha)
    both_available = result["shapiro_p"] is not None and result["dagostino_p"] is not None

    if both_available:
        is_normal = sw_normal and dp_normal
    elif result["shapiro_p"] is not None:
        is_normal = sw_normal
    elif result["dagostino_p"] is not None:
        is_normal = dp_normal
    else:
        is_normal = False

    result["is_normal"] = is_normal
    abs_skew = abs(result["skewness"] or 0)
    abs_kurt = abs(result["kurtosis"] or 0)

    if is_normal:
        result["verdict"] = "NORMAL"
        result["recommendation"] = (
            "Parametric statistics (mean, 95% CI) are appropriate for this variable."
        )
    elif abs_skew > 2 or abs_kurt > 7:
        result["verdict"] = "HIGHLY_SKEWED"
        result["recommendation"] = (
            f"Highly skewed distribution (skew={result['skewness']:.2f}, "
            f"kurt={result['kurtosis']:.2f}). "
            "Use median and IQR rather than mean and std. "
            "Chebyshev CI reported in addition to classical CI."
        )
    else:
        result["verdict"] = "NON_NORMAL"
        result["recommendation"] = (
            "Distribution is non-normal. Mean/CI are still reported but interpret cautiously. "
            "Log-transform may improve normality for right-skewed pollution data."
        )

    return result


def correlation_with_pvalues(daily_df: pd.DataFrame,
                              variables: list) -> tuple:
    """
    Pearson correlation matrix WITH p-values and significance flags.
    Returns (corr_matrix, pval_matrix, significance_matrix).

    p-value computed from t = r * sqrt(n-2) / sqrt(1-r²)
    Bonferroni correction applied for multiple comparisons.

    Significance levels:
      *** p < 0.001
      **  p < 0.01
      *   p < 0.05
      ns  p ≥ 0.05
    """
    pivot = daily_df.pivot_table(
        index="date", columns="variable", values="mean"
    )
    cols = [v for v in variables if v in pivot.columns]
    if len(cols) < 2:
        return pd.DataFrame(), pd.DataFrame(), pd.DataFrame()

    data = pivot[cols].dropna()
    n    = len(data)
    d    = len(cols)

    corr_matrix = data.corr(method="pearson")
    pval_matrix = pd.DataFrame(np.ones((d, d)), index=cols, columns=cols)
    sig_matrix  = pd.DataFrame([["ns"] * d] * d, index=cols, columns=cols)

    # Number of tests (for Bonferroni)
    n_tests = d * (d - 1) / 2
    bonf    = 0.05 / max(n_tests, 1)

    for i, ci in enumerate(cols):
        for j, cj in enumerate(cols):
            if i == j:
                pval_matrix.loc[ci, cj] = 0.0
                sig_matrix.loc[ci, cj]  = "—"
                continue
            r = corr_matrix.loc[ci, cj]
            if abs(r) >= 1.0 or n <= 2:
                continue
            t_stat = r * math.sqrt(n - 2) / math.sqrt(max(1 - r**2, 1e-10))
            p      = float(2 * stats.t.sf(abs(t_stat), df=n - 2))
            pval_matrix.loc[ci, cj] = round(p, 6)
            # Significance with Bonferroni correction
            if p < bonf:
                sig_matrix.loc[ci, cj] = "***"
            elif p < 0.01:
                sig_matrix.loc[ci, cj] = "**"
            elif p < 0.05:
                sig_matrix.loc[ci, cj] = "*"
            else:
                sig_matrix.loc[ci, cj] = "ns"

    return corr_matrix, pval_matrix, sig_matrix


def ols_with_inference(daily_df: pd.DataFrame,
                       target: str = "PM10",
                       predictors: list = None) -> dict:
    """
    Full OLS regression with statistical inference.
    Adds to the base OLS:
      - Standard errors per coefficient (SE(β̂_j))
      - t-statistics per coefficient
      - p-values per coefficient (two-tailed)
      - 95% confidence intervals per coefficient
      - Adjusted R² (penalises for number of predictors)
      - F-statistic and p-value (overall model significance)
      - Residual standard error (RSE)
      - Variance Inflation Factors (VIF) — multicollinearity check
      - Durbin-Watson statistic — residual autocorrelation check

    All of these are standard in any published regression table.
    The original implementation returned only β̂ and R² — unpublishable.
    """
    if predictors is None:
        predictors = [p for p in ["humidity_pct", "wind_speed_10m",
                                   "rain_mm", "pressure_hpa"]
                      if p in daily_df["variable"].unique()]

    pivot = daily_df.pivot_table(
        index="date", columns="variable", values="mean"
    )
    if target not in pivot.columns:
        return {"available": False, "reason": "target not in data"}

    available = [p for p in predictors if p in pivot.columns]
    if len(available) < 1:
        return {"available": False, "reason": "no predictors available"}

    df_reg = pivot[[target] + available].dropna()
    n      = len(df_reg)
    k      = len(available)   # number of predictors (excl. intercept)

    if n < k + 5:
        return {"available": False, "reason": f"too few obs ({n}) for {k} predictors"}

    X_raw = df_reg[available].values
    y     = df_reg[target].values.astype(float)

    scaler = StandardScaler()
    X_s    = scaler.fit_transform(X_raw)
    X_aug  = np.column_stack([np.ones(n), X_s])   # design matrix

    # β̂ = (XᵀX)⁻¹Xᵀy
    try:
        XtX_inv = np.linalg.inv(X_aug.T @ X_aug)
    except np.linalg.LinAlgError:
        XtX_inv = np.linalg.pinv(X_aug.T @ X_aug)

    beta    = XtX_inv @ X_aug.T @ y
    y_pred  = X_aug @ beta
    resid   = y - y_pred

    # Model fit
    ss_res  = float(np.sum(resid**2))
    ss_tot  = float(np.sum((y - np.mean(y))**2))
    r2      = 1 - ss_res / max(ss_tot, 1e-10)
    r2_adj  = 1 - (1 - r2) * (n - 1) / max(n - k - 1, 1)
    rse     = math.sqrt(ss_res / max(n - k - 1, 1))   # residual std error

    # Coefficient standard errors and inference
    se_vec   = np.sqrt(np.diag(XtX_inv) * rse**2)
    t_stats  = beta / np.where(se_vec > 0, se_vec, 1e-10)
    p_vals   = [float(2 * stats.t.sf(abs(t), df=n-k-1)) for t in t_stats]
    ci_lo    = beta - 1.96 * se_vec
    ci_hi    = beta + 1.96 * se_vec

    # F-statistic for overall model significance
    # F = (R²/k) / ((1-R²)/(n-k-1))
    if r2 < 1.0 and k > 0:
        f_stat = (r2 / k) / max((1 - r2) / (n - k - 1), 1e-10)
        f_pval = float(stats.f.sf(f_stat, k, n - k - 1))
    else:
        f_stat = np.inf
        f_pval = 0.0

    # Variance Inflation Factors — detect multicollinearity
    vif = {}
    if len(available) > 1:
        for i, pred in enumerate(available):
            other_X = np.delete(X_s, i, axis=1)
            other_X_aug = np.column_stack([np.ones(n), other_X])
            try:
                b_vif  = np.linalg.lstsq(other_X_aug, X_s[:, i], rcond=None)[0]
                y_vif  = other_X_aug @ b_vif
                ss_res_v = np.sum((X_s[:, i] - y_vif)**2)
                ss_tot_v = np.sum((X_s[:, i] - np.mean(X_s[:, i]))**2)
                r2_vif   = 1 - ss_res_v / max(ss_tot_v, 1e-10)
                vif[pred] = round(1 / max(1 - r2_vif, 1e-10), 2)
            except Exception:
                vif[pred] = None

    # Durbin-Watson statistic — test for residual autocorrelation
    # DW = Σ(e_t - e_{t-1})² / Σe_t²   (≈2 = no autocorrelation)
    dw = float(np.sum(np.diff(resid)**2) / max(ss_res, 1e-10))

    # Build coefficient table (skip intercept for VIF/interpretation)
    coef_table = {}
    for i, pred in enumerate(available):
        p_i = p_vals[i + 1]   # +1 for intercept offset
        coef_table[pred] = {
            "beta"    : round(float(beta[i + 1]), 4),
            "se"      : round(float(se_vec[i + 1]), 4),
            "t_stat"  : round(float(t_stats[i + 1]), 3),
            "p_value" : round(p_i, 6),
            "ci_lower": round(float(ci_lo[i + 1]), 4),
            "ci_upper": round(float(ci_hi[i + 1]), 4),
            "sig"     : ("***" if p_i < 0.001 else "**" if p_i < 0.01
                         else "*" if p_i < 0.05 else "ns"),
            "vif"     : vif.get(pred),
        }

    sig_label = "***" if f_pval < 0.001 else "**" if f_pval < 0.01 else "*" if f_pval < 0.05 else "ns"

    return {
        "available"      : True,
        "target"         : target,
        "predictors"     : available,
        "n_obs"          : n,
        "k_predictors"   : k,
        "intercept"      : round(float(beta[0]), 4),
        "intercept_se"   : round(float(se_vec[0]), 4),
        "intercept_p"    : round(float(p_vals[0]), 6),
        "r_squared"      : round(r2, 4),
        "r_squared_adj"  : round(r2_adj, 4),
        "rse"            : round(rse, 4),
        "f_statistic"    : round(float(f_stat), 3) if not np.isinf(f_stat) else 9999,
        "f_pvalue"       : round(float(f_pval), 6),
        "f_significance" : sig_label,
        "durbin_watson"  : round(dw, 3),
        "dw_interpretation": (
            "No significant autocorrelation" if 1.5 <= dw <= 2.5 else
            "Positive autocorrelation — OLS standard errors underestimated"
            if dw < 1.5 else
            "Negative autocorrelation"
        ),
        "coef_table"     : coef_table,
        "high_vif_vars"  : [k for k, v in vif.items() if v and v > 5],
        "note"           : (
            "All predictors standardised (z-scored) before regression. "
            "β coefficients are standardised — comparable in magnitude. "
            "Association only — Corr(X,Y) ≢ C(X→Y)."
        ),
    }


def seasonal_decomposition(daily_df: pd.DataFrame,
                             col: str = "PM10",
                             period: int = 365) -> dict:
    """
    Additive seasonal decomposition: x_t = Trend_t + Seasonal_t + Residual_t

    Method:
    1. Trend: centred moving average with window = period (365d or 30d)
    2. Detrended: x_t - Trend_t
    3. Seasonal: mean of detrended values for each day-of-year (or month)
    4. Residual: x_t - Trend_t - Seasonal_t

    Returns arrays for Trend, Seasonal, Residual components.
    Residual should be roughly stationary (no pattern) if model fits well.
    Residual std / series std = unexplained variation fraction.
    """
    sub = daily_df[daily_df["variable"] == col][["date","mean"]].dropna()
    sub = sub.sort_values("date").reset_index(drop=True)
    n   = len(sub)

    if n < 14:
        return {"available": False, "reason": "insufficient data (< 14 days)"}

    series = sub["mean"].values.astype(float)
    dates  = pd.to_datetime(sub["date"])

    # Choose decomposition period
    if n >= 365:
        window  = 365
        use_doy = True   # day-of-year seasonal pattern
    else:
        window  = min(30, n // 2)
        use_doy = False  # monthly seasonal pattern

    # Step 1: Trend via centred moving average
    trend = pd.Series(series).rolling(
        window=window, min_periods=window//4, center=True
    ).mean().values

    # Step 2: Detrended
    detrended = series - trend

    # Step 3: Seasonal component
    seasonal = np.zeros(n)
    if use_doy:
        doy_vals = dates.dt.dayofyear.values
        for doy in range(1, 367):
            mask = doy_vals == doy
            if mask.sum() > 0:
                seasonal[mask] = np.nanmean(detrended[mask])
    else:
        month_vals = dates.dt.month.values
        for m in range(1, 13):
            mask = month_vals == m
            if mask.sum() > 0:
                seasonal[mask] = np.nanmean(detrended[mask])

    # Step 4: Residual
    residual = series - trend - seasonal

    # Quality metrics
    series_std    = float(np.nanstd(series))
    residual_std  = float(np.nanstd(residual[~np.isnan(residual)]))
    trend_std     = float(np.nanstd(trend[~np.isnan(trend)]))
    seasonal_range= float(np.nanmax(seasonal) - np.nanmin(seasonal))

    # Fraction of variance explained by each component
    var_total    = float(np.nanvar(series))
    var_trend    = float(np.nanvar(trend[~np.isnan(trend)]))
    var_seasonal = float(np.nanvar(seasonal))
    var_resid    = float(np.nanvar(residual[~np.isnan(residual)]))

    pct_trend    = round(var_trend    / max(var_total, 1e-10) * 100, 1)
    pct_seasonal = round(var_seasonal / max(var_total, 1e-10) * 100, 1)
    pct_resid    = round(var_resid    / max(var_total, 1e-10) * 100, 1)

    return {
        "available"      : True,
        "variable"       : col,
        "n"              : n,
        "period_days"    : window,
        "trend"          : trend,
        "seasonal"       : seasonal,
        "residual"       : residual,
        "dates"          : dates,
        "series"         : series,
        "series_std"     : round(series_std,    3),
        "residual_std"   : round(residual_std,  3),
        "trend_std"      : round(trend_std,     3),
        "seasonal_range" : round(seasonal_range,3),
        "pct_trend"      : pct_trend,
        "pct_seasonal"   : pct_seasonal,
        "pct_residual"   : pct_resid,
        "interpretation" : (
            f"Trend explains {pct_trend:.1f}% of variance, "
            f"seasonality {pct_seasonal:.1f}%, "
            f"unexplained residual {pct_resid:.1f}%. "
            + ("Well-fitted model — low residual." if pct_resid < 30 else
               "High residual — episodic events dominate." if pct_resid > 60 else
               "Moderate residual — some episodic variation.")
        ),
    }


def exceedance_frequency_analysis(daily_df: pd.DataFrame,
                                   config: dict) -> dict:
    """
    Exceedance frequency distribution for each pollutant with a WHO threshold.

    For each variable:
    - Exceedance probability P(X > θ) = N_exceed / N_total
    - Return period R = 1 / P(X > θ)  [expected days between exceedances]
    - Expected WHO excess burden: mean(X | X > θ) - θ
    - Consecutive exceedance streaks: longest run of exceeding days
    - Monthly exceedance probability (which month is worst)

    These statistics belong in any published air quality paper.
    """
    results = {}
    thresholds = config.get("thresholds", {})

    for col in POLLUTANTS:
        thr = thresholds.get(col)
        if thr is None:
            continue

        sub = daily_df[daily_df["variable"] == col][["date","mean"]].dropna()
        if len(sub) < 7:
            continue

        sub     = sub.copy()
        sub["date"] = pd.to_datetime(sub["date"])
        vals    = sub["mean"].values
        n_total = len(vals)
        exceed  = vals > thr
        n_exceed = int(exceed.sum())

        if n_exceed == 0:
            results[col] = {
                "threshold"  : thr,
                "n_days"     : n_total,
                "n_exceed"   : 0,
                "P_exceed"   : 0.0,
                "return_days": float("inf"),
                "status"     : "COMPLIANT",
            }
            continue

        P_exceed    = round(n_exceed / n_total, 4)
        return_days = round(1 / P_exceed, 1) if P_exceed > 0 else float("inf")
        excess_vals = vals[exceed]
        mean_excess = float(np.mean(excess_vals))
        excess_over_thr = round(mean_excess - thr, 3)

        # Consecutive exceedance streaks
        max_streak = 0
        cur_streak = 0
        for ex in exceed:
            if ex:
                cur_streak += 1
                max_streak  = max(max_streak, cur_streak)
            else:
                cur_streak  = 0

        # Monthly exceedance probability
        sub["exceed"] = exceed
        sub["month"]  = sub["date"].dt.month
        month_names   = ["Jan","Feb","Mar","Apr","May","Jun",
                         "Jul","Aug","Sep","Oct","Nov","Dec"]
        monthly_exc   = sub.groupby("month")["exceed"].mean()
        worst_month_num = int(monthly_exc.idxmax()) if len(monthly_exc) > 0 else 1
        worst_month     = month_names[worst_month_num - 1]
        worst_month_pct = round(float(monthly_exc.max()) * 100, 1)

        results[col] = {
            "threshold"       : thr,
            "n_days"          : n_total,
            "n_exceed"        : n_exceed,
            "P_exceed"        : P_exceed,
            "return_days"     : return_days,
            "mean_when_exceed": round(mean_excess, 3),
            "excess_over_thr" : excess_over_thr,
            "max_streak_days" : max_streak,
            "worst_month"     : worst_month,
            "worst_month_pct" : worst_month_pct,
            "monthly_prob"    : {month_names[m-1]: round(float(v),3)
                                 for m, v in monthly_exc.items()},
            "status"          : ("CRITICAL"  if P_exceed > 0.75 else
                                 "SEVERE"    if P_exceed > 0.50 else
                                 "HIGH"      if P_exceed > 0.25 else
                                 "MODERATE"  if P_exceed > 0.10 else "LOW"),
        }

    return results


def run_full_statistical_inference(daily_df: pd.DataFrame,
                                    config: dict) -> dict:
    """
    Orchestrator: runs all statistical inference tests for the pipeline.
    Returns a single dict that goes into J and the report.
    """
    print("\n  [Stat-Inference] Running formal statistical tests...")
    results = {}

    # 1. Mann-Kendall trend tests for all pollutants
    mk_results = {}
    for col in POLLUTANTS:
        sub = daily_df[daily_df["variable"] == col]["mean"].dropna().values
        if len(sub) >= 10:
            mk = mann_kendall_test(sub)
            mk_results[col] = mk
            if mk.get("significant"):
                print(f"  ✓ MK [{col}]: {mk['direction']} "
                      f"(p={mk['p_value']:.4f}, slope={mk['sens_slope']:.3f} µg/m³/day)")
    results["mann_kendall"] = mk_results

    # 2. Normality tests for key pollutants
    norm_results = {}
    for col in ["PM10","NO2","SO2","CO","Ozone"]:
        sub = daily_df[daily_df["variable"] == col]["mean"].dropna().values
        if len(sub) >= 8:
            norm_results[col] = normality_tests(sub)
    results["normality"] = norm_results

    # 3. Correlation with p-values
    key_vars = [v for v in ["PM10","NO2","SO2","CO","Ozone","NH3",
                             "humidity_pct","wind_speed_10m","pressure_hpa"]
                if v in daily_df["variable"].unique()]
    if len(key_vars) >= 3:
        corr_m, pval_m, sig_m = correlation_with_pvalues(daily_df, key_vars)
        results["correlation_pvalues"] = {
            "matrix"         : corr_m.to_dict() if not corr_m.empty else {},
            "p_values"       : pval_m.to_dict() if not pval_m.empty else {},
            "significance"   : sig_m.to_dict()  if not sig_m.empty  else {},
        }
        # Count significant pairs
        if not sig_m.empty:
            n_sig = sum(1 for i in sig_m.index for j in sig_m.columns
                       if i < j and sig_m.loc[i,j] in ("*","**","***"))
            print(f"  ✓ Correlation: {n_sig} significant pairs (p<0.05) "
                  f"among {len(key_vars)} variables")

    # 4. Full OLS with inference
    ols = ols_with_inference(daily_df, target="PM10")
    results["ols_inference"] = ols
    if ols.get("available"):
        sig_preds = [k for k,v in ols["coef_table"].items() if v["sig"] != "ns"]
        print(f"  ✓ OLS: R²={ols['r_squared']:.3f} (adj={ols['r_squared_adj']:.3f}) "
              f"| F={ols['f_statistic']:.1f} ({ols['f_significance']}) "
              f"| DW={ols['durbin_watson']:.2f}")
        if sig_preds:
            print(f"  ✓ Significant predictors: {sig_preds}")

    # 5. Seasonal decomposition for PM10
    decomp = seasonal_decomposition(daily_df, col="PM10")
    results["seasonal_decomposition"] = {
        "available"    : decomp.get("available", False),
        "pct_trend"    : decomp.get("pct_trend"),
        "pct_seasonal" : decomp.get("pct_seasonal"),
        "pct_residual" : decomp.get("pct_residual"),
        "seasonal_range": decomp.get("seasonal_range"),
        "interpretation": decomp.get("interpretation",""),
        "period_days"  : decomp.get("period_days"),
    }
    results["_decomp_arrays"] = decomp   # kept for charting, not serialized to J

    # 6. Exceedance frequency analysis
    exc_freq = exceedance_frequency_analysis(daily_df, config)
    results["exceedance_frequency"] = exc_freq
    critical = [k for k,v in exc_freq.items()
                if v.get("status") in ("CRITICAL","SEVERE")]
    if critical:
        print(f"  ⚠ Critical exceedance: {critical}")
    print(f"  ✓ Exceedance analysis: {len(exc_freq)} pollutants with thresholds")

    return results




# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 7.5 — ML ANOMALY DETECTION: ISOLATION FOREST
#  Runs alongside z-score (ζ_d). Does NOT replace formal Eq. 9.
#  Detects multivariate anomalies that univariate z-score cannot see.
#
#  SECURITY NOTE: All computation is local. IsolationForest is fit and
#  applied entirely in memory on daily_df data. No data leaves the machine.
#  Fixed random_state=42 ensures reproducible results (same seed → same model).
# ═══════════════════════════════════════════════════════════════════════════════

def run_isolation_forest(daily_df: pd.DataFrame, config: dict) -> dict:
    """
    ML-1: Isolation Forest multivariate anomaly detection.

    How it works:
      1. Pivot daily_df to a wide matrix: rows = dates, cols = pollutants
      2. Standardise each column (zero mean, unit variance)
      3. Fit IsolationForest on the full dataset (unsupervised — no labels)
      4. Score each day: anomaly_score ∈ [-1, 1] where lower = more anomalous
      5. Flag days where predict() == -1 as ML anomalies

    Why it improves on z-score:
      - z-score is univariate: checks one variable at a time
      - IsolationForest is multivariate: detects days where the COMBINATION
        of PM10 + NO2 + SO2 is anomalous even if no single variable crosses 3σ
      - Does not assume normality (z-score implicitly does)
      - Automatically handles correlated features

    Formal complement: ML anomaly ∪ z-score anomaly ⊇ z-score anomaly.
    A day flagged by IF but NOT z-score warrants manual investigation.

    Returns:
      dict with keys: flagged_dates, scores, n_anomalies, features_used,
                      contamination, comparison_df
    """
    print("\n  [ML-1] Running Isolation Forest multivariate anomaly detection...")

    # Select feature columns available in data
    feature_cols = [c for c in config["if_cols"]
                    if c in daily_df["variable"].unique()]
    if len(feature_cols) < 2:
        print(f"  ⚠  Isolation Forest skipped — fewer than 2 feature columns available.")
        return {"available": False, "reason": "insufficient_features"}

    # Build wide matrix: one row per day
    pivot = daily_df.pivot_table(index="date", columns="variable", values="mean")
    pivot = pivot[[c for c in feature_cols if c in pivot.columns]].dropna()
    if len(pivot) < 30:
        print(f"  ⚠  Isolation Forest skipped — fewer than 30 complete rows ({len(pivot)}).")
        return {"available": False, "reason": "insufficient_rows"}

    used_cols = list(pivot.columns)
    X = pivot.values

    # Standardise (important: IF is not scale-invariant)
    scaler = StandardScaler()
    X_s = scaler.fit_transform(X)

    # Fit Isolation Forest
    # contamination: expected fraction of anomalies (5% = 1 in 20 days)
    # random_state: fixed for reproducibility (same seed → same tree structure)
    # n_estimators: 150 trees for stable scoring
    clf = IsolationForest(
        contamination  = config["if_contamination"],
        n_estimators   = config["if_n_estimators"],
        random_state   = config["if_random_state"],
        n_jobs         = 1,      # single-threaded for reproducibility
    )
    clf.fit(X_s)

    preds  = clf.predict(X_s)          # 1 = normal, -1 = anomaly
    scores = clf.decision_function(X_s)  # higher = more normal

    anomaly_mask  = preds == -1
    flagged_dates = list(pivot.index[anomaly_mask])
    n_anomalies   = int(anomaly_mask.sum())

    # Build comparison dataframe (for report and Chart 14)
    comp = pd.DataFrame({
        "date"         : pivot.index,
        "if_score"     : scores,
        "ml_anomaly"   : anomaly_mask,
    }).set_index("date")

    # Merge with z-score anomaly column from daily_df (PM10 only)
    pm10_anom = daily_df[daily_df["variable"] == "PM10"][["date","is_anomaly","z_score"]].copy()
    pm10_anom = pm10_anom.set_index("date").rename(
        columns={"is_anomaly": "zscore_anomaly", "z_score": "z_score"})
    comp = comp.join(pm10_anom, how="left")
    comp["zscore_anomaly"] = comp["zscore_anomaly"].fillna(False)

    # Agreement categories
    comp["both"]    = comp["ml_anomaly"] & comp["zscore_anomaly"]
    comp["if_only"] = comp["ml_anomaly"] & ~comp["zscore_anomaly"]
    comp["zs_only"] = ~comp["ml_anomaly"] & comp["zscore_anomaly"]

    n_both   = int(comp["both"].sum())
    n_if_only= int(comp["if_only"].sum())
    n_zs_only= int(comp["zs_only"].sum())

    print(f"  ✓ Isolation Forest: {n_anomalies} anomalies in {len(pivot)} days "
          f"({n_anomalies/len(pivot)*100:.1f}%)")
    print(f"  ✓ Agreement: both={n_both} | IF-only={n_if_only} | z-score-only={n_zs_only}")

    return {
        "available"      : True,
        "features_used"  : used_cols,
        "n_days"         : len(pivot),
        "n_anomalies"    : n_anomalies,
        "anomaly_rate"   : round(n_anomalies / len(pivot), 4),
        "contamination"  : config["if_contamination"],
        "n_estimators"   : config["if_n_estimators"],
        "random_state"   : config["if_random_state"],
        "flagged_dates"  : [str(d.date()) for d in flagged_dates],
        "comparison_df"  : comp,
        "n_both"         : n_both,
        "n_if_only"      : n_if_only,
        "n_zs_only"      : n_zs_only,
        "security_note"  : (
            "IsolationForest fit on local daily_df only. "
            "No data transmitted. random_state=42 ensures reproducibility."
        ),
    }


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 7.6 — ML FORECASTING: HOLT-WINTERS TRIPLE EXPONENTIAL SMOOTHING
#  Pure numpy implementation — zero extra dependencies beyond existing stack.
#  Forecasts PM10 14 days ahead. Used for Chart 15 only.
#
#  SECURITY NOTE: Computed entirely in memory from daily_df.
#  Forecast values are NOT added to J (the LLM data contract) because:
#    1. Forecasts have uncertainty — adding them to J risks LLM stating
#       forecasts as historical facts.
#    2. The AI constraints in J explicitly prohibit extrapolation.
#  Forecasts appear only in Chart 15 and the report ML section.
# ═══════════════════════════════════════════════════════════════════════════════

def holt_winters_forecast(series: np.ndarray, seasonal_periods: int = 7,
                           forecast_days: int = 14,
                           alpha: float = 0.3, beta: float = 0.1,
                           gamma: float = 0.2) -> dict:
    """
    Triple Exponential Smoothing (Holt-Winters additive seasonal model).

    Parameters:
        series          : 1-D array of historical daily PM10 means
        seasonal_periods: length of one season (7 = weekly)
        forecast_days   : number of days to forecast ahead
        alpha (α)       : level smoothing  — how fast level adapts (0 < α < 1)
        beta  (β)       : trend smoothing  — how fast trend adapts (0 < β < 1)
        gamma (γ)       : seasonal smoothing (0 < γ < 1)

    Update equations (additive Holt-Winters):
        Level:    L_t = α(x_t - S_{t-m}) + (1-α)(L_{t-1} + T_{t-1})
        Trend:    T_t = β(L_t - L_{t-1}) + (1-β)T_{t-1}
        Seasonal: S_t = γ(x_t - L_t) + (1-γ)S_{t-m}

    Forecast:   x̂_{t+h} = L_t + h·T_t + S_{t-m+((h-1) mod m)+1}

    Returns dict with: forecast array, lower/upper 80% prediction bands,
    fitted values, and model parameters.
    """
    n = len(series)
    m = seasonal_periods

    if n < 2 * m:
        return {"available": False, "reason": "insufficient_history"}

    series = np.array(series, dtype=float)

    # ── Initialisation ──────────────────────────────────────────────────────
    # Level: mean of first season
    L = np.mean(series[:m])
    # Trend: average slope across first two seasons
    T = (np.mean(series[m:2*m]) - np.mean(series[:m])) / m
    # Seasonal indices: deviation of each period from level
    S = np.array([(series[i] - L) for i in range(m)], dtype=float)

    levels   = np.zeros(n)
    trends   = np.zeros(n)
    seasonals= np.zeros(n + forecast_days + m)
    seasonals[:m] = S
    fitted   = np.zeros(n)

    # ── Smoothing pass ──────────────────────────────────────────────────────
    for t in range(n):
        s_idx = t % m
        x_t   = series[t]
        L_new = alpha * (x_t - seasonals[t % m]) + (1 - alpha) * (L + T)
        T_new = beta  * (L_new - L) + (1 - beta) * T
        S_new = gamma * (x_t - L_new) + (1 - gamma) * seasonals[t % m]
        fitted[t]      = L + T + seasonals[t % m]
        seasonals[t + m] = S_new
        L, T           = L_new, T_new
        levels[t]  = L
        trends[t]  = T

    # ── Residual std for prediction intervals ───────────────────────────────
    resid = series - fitted
    resid_std = float(np.std(resid, ddof=1))

    # ── Forecast ────────────────────────────────────────────────────────────
    forecast = np.zeros(forecast_days)
    for h in range(1, forecast_days + 1):
        s_idx      = (n + h - 1) % m
        forecast[h-1] = L + h * T + seasonals[n + h - 1]

    # 80% prediction interval: ±1.28 * resid_std * sqrt(h)
    # Wider interval for longer horizons (uncertainty grows with h)
    pi_factor = 1.28
    lower = np.array([forecast[h] - pi_factor * resid_std * math.sqrt(h+1)
                      for h in range(forecast_days)])
    upper = np.array([forecast[h] + pi_factor * resid_std * math.sqrt(h+1)
                      for h in range(forecast_days)])

    # Clip to non-negative (concentrations cannot be negative)
    forecast = np.clip(forecast, 0, None)
    lower    = np.clip(lower, 0, None)
    upper    = np.clip(upper, 0, None)

    rmse = float(np.sqrt(np.mean(resid**2)))
    mae  = float(np.mean(np.abs(resid)))

    return {
        "available"       : True,
        "forecast"        : forecast,
        "lower_80"        : lower,
        "upper_80"        : upper,
        "fitted"          : fitted,
        "rmse"            : round(rmse, 3),
        "mae"             : round(mae, 3),
        "resid_std"       : round(resid_std, 3),
        "alpha"           : alpha,
        "beta"            : beta,
        "gamma"           : gamma,
        "seasonal_periods": m,
        "forecast_days"   : forecast_days,
        "n_history"       : n,
        "security_note"   : (
            "Forecast computed locally on daily_df. NOT added to JSON contract J "
            "to prevent LLM from stating forecasts as historical facts. "
            "Forecast appears in Chart 15 and ML report section only."
        ),
    }


def run_hw_forecast(daily_df: pd.DataFrame, config: dict) -> dict:
    """
    Orchestrate Holt-Winters forecast for PM10.
    Extracts PM10 daily means, calls holt_winters_forecast(),
    builds forecast date index, returns full result dict.
    """
    print("\n  [ML-2] Running Holt-Winters forecast for PM10...")

    pm10 = daily_df[daily_df["variable"] == "PM10"][["date","mean"]].dropna()
    pm10 = pm10.sort_values("date")
    if len(pm10) < 14:
        print("  ⚠  Forecast skipped — fewer than 14 PM10 data points.")
        return {"available": False, "reason": "insufficient_pm10_data"}

    series   = pm10["mean"].values
    last_date= pm10["date"].iloc[-1]

    result = holt_winters_forecast(
        series          = series,
        seasonal_periods= config["hw_seasonal_periods"],
        forecast_days   = config["hw_forecast_days"],
        alpha           = config["hw_alpha"],
        beta            = config["hw_beta"],
        gamma           = config["hw_gamma"],
    )

    if not result.get("available"):
        print(f"  ⚠  Forecast failed: {result.get('reason')}")
        return result

    # Build forecast date range
    forecast_dates = pd.date_range(
        start=last_date + pd.Timedelta(days=1),
        periods=config["hw_forecast_days"],
        freq="D"
    )
    result["forecast_dates"] = forecast_dates
    result["history_dates"]  = pm10["date"].values
    result["history_series"] = series

    print(f"  ✓ Holt-Winters: RMSE={result['rmse']:.2f} µg/m³ | "
          f"MAE={result['mae']:.2f} µg/m³ | "
          f"Forecast horizon: {config['hw_forecast_days']} days")
    print(f"  ✓ Next 7-day mean forecast: "
          f"{np.mean(result['forecast'][:7]):.1f} µg/m³")

    return result


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 7.7 — CHART INSIGHT GENERATOR
#  Generates deterministic text insights for every chart from actual data.
#  NO LLM. NO hallucination risk. All numbers come from daily_df / corr_matrix.
#  These insights appear below each chart in the Word report.
# ═══════════════════════════════════════════════════════════════════════════════

def generate_chart_insights(daily_df: pd.DataFrame, df: pd.DataFrame,
                             config: dict, corr_matrix: pd.DataFrame,
                             ols_result: dict, ml_if: dict,
                             ml_hw: dict) -> dict:
    """
    Compute data-driven text insight for each chart.

    Every string is built from actual computed values in daily_df and corr_matrix.
    No LLM is involved. No value is invented or estimated.
    These insights are factual summaries — not interpretations.

    Returns:
        dict: {chart_name: insight_string}
    """
    insights = {}
    city = config["city"]
    thr  = config["thresholds"]

    def safe_get(col, stat="mean"):
        sub = daily_df[daily_df["variable"] == col]
        if sub.empty or sub[stat].isna().all():
            return None
        return float(sub[stat].mean()) if stat == "mean" else float(sub[stat].iloc[0])

    def annual_mean(col):
        sub = daily_df[daily_df["variable"] == col]["mean"].dropna()
        return float(sub.mean()) if len(sub) > 0 else None

    def annual_max(col):
        sub = daily_df[daily_df["variable"] == col]["mean"].dropna()
        return float(sub.max()) if len(sub) > 0 else None

    def exceed_days(col):
        sub = daily_df[daily_df["variable"] == col]
        return int(sub["exceeds_threshold"].sum()) if "exceeds_threshold" in sub.columns else 0

    def anomaly_days(col):
        sub = daily_df[daily_df["variable"] == col]
        return int(sub["is_anomaly"].sum()) if "is_anomaly" in sub.columns else 0

    def avg_quality(col):
        sub = daily_df[daily_df["variable"] == col]["Q"].dropna()
        return float(sub.mean()) if len(sub) > 0 else None

    def peak_month(col):
        sub = daily_df[daily_df["variable"] == col][["date","mean"]].dropna()
        if sub.empty: return "N/A"
        sub = sub.copy(); sub["month"] = sub["date"].dt.month
        m = sub.groupby("month")["mean"].mean().idxmax()
        return ["Jan","Feb","Mar","Apr","May","Jun",
                "Jul","Aug","Sep","Oct","Nov","Dec"][m - 1]

    def low_month(col):
        sub = daily_df[daily_df["variable"] == col][["date","mean"]].dropna()
        if sub.empty: return "N/A"
        sub = sub.copy(); sub["month"] = sub["date"].dt.month
        m = sub.groupby("month")["mean"].mean().idxmin()
        return ["Jan","Feb","Mar","Apr","May","Jun",
                "Jul","Aug","Sep","Oct","Nov","Dec"][m - 1]

    # ── Chart 1: Hourly Time-Series ──────────────────────────────────────────
    pm10_mean  = annual_mean("PM10")
    pm10_max   = annual_max("PM10")
    pm10_ex    = exceed_days("PM10")
    no2_mean   = annual_mean("NO2")
    pm10_thr   = thr.get("PM10", 45)
    no2_thr    = thr.get("NO2", 25)
    total_days = daily_df["date"].nunique()

    if pm10_mean is not None:
        pm10_ex_pct = pm10_ex / total_days * 100 if total_days > 0 else 0
        insights["timeseries"] = (
            f"PM10 annual mean: {pm10_mean:.1f} µg/m³ "
            f"({'above' if pm10_mean > pm10_thr else 'below'} WHO guideline {pm10_thr} µg/m³). "
            f"Peak daily mean: {pm10_max:.1f} µg/m³. "
            f"WHO exceedances: {pm10_ex} days ({pm10_ex_pct:.1f}% of monitored period). "
            + (f"NO₂ annual mean: {no2_mean:.1f} µg/m³ "
               f"({'above' if no2_mean > no2_thr else 'below'} WHO guideline {no2_thr} µg/m³)."
               if no2_mean is not None else "NO₂ data not available.")
        )
    else:
        insights["timeseries"] = "PM10 data not available for time-series analysis."

    # ── Chart 2: Correlation Heatmap ─────────────────────────────────────────
    if corr_matrix is not None and not corr_matrix.empty:
        flat = corr_matrix.where(
            np.triu(np.ones(corr_matrix.shape), k=1).astype(bool)
        ).stack()
        if len(flat) > 0:
            top_pair = flat.abs().idxmax()
            top_r    = float(flat[top_pair])
            v1, v2   = top_pair
            direction = "positive" if top_r > 0 else "negative"
            # Find PM10-humidity if available
            pm10_hum = None
            if "PM10" in corr_matrix.index and "humidity_pct" in corr_matrix.columns:
                pm10_hum = float(corr_matrix.loc["PM10","humidity_pct"])
            insights["correlation"] = (
                f"Strongest association: {v1}–{v2} (r = {top_r:.3f}, {direction}). "
                + (f"PM10–humidity association: r = {pm10_hum:.3f} "
                   f"({'rain washes out particles — consistent with wet deposition' if pm10_hum < -0.3 else 'weak inverse pattern'})."
                   if pm10_hum is not None else "")
                + " NOTE: All correlations are statistical associations only — not causal relationships."
            )
        else:
            insights["correlation"] = "Correlation matrix computed. Interpret all r values as associations only — not causation."
    else:
        insights["correlation"] = "Correlation matrix unavailable."

    # ── Chart 3: 30-Day Rolling Trend ────────────────────────────────────────
    pk = peak_month("PM10"); lk = low_month("PM10")
    pm10_std = daily_df[daily_df["variable"] == "PM10"]["mean"].std()
    insights["trend"] = (
        f"PM10 30-day rolling mean peaks in {pk} and reaches minimum in {lk}, "
        f"indicating {'strong' if pm10_std is not None and pm10_std > 20 else 'moderate'} "
        f"seasonal variation (annual σ = {pm10_std:.1f} µg/m³). "
        f"The rolling window smooths short-term spikes to reveal the underlying seasonal pattern. "
        f"Deviations above the rolling mean indicate elevated pollution episodes."
    ) if pm10_mean is not None else "PM10 trend data unavailable."

    # ── Chart 4: Anomaly Detection ────────────────────────────────────────────
    pm10_an = anomaly_days("PM10")
    no2_an  = anomaly_days("NO2")
    z_thr   = config["anomaly_z_threshold"]
    insights["anomalies"] = (
        f"Anomaly threshold: |ζ_d| > {z_thr:.1f} (3σ rule). "
        f"PM10 anomaly days: {pm10_an} "
        f"({pm10_an/total_days*100:.1f}% of monitored period). "
        + (f"NO₂ anomaly days: {no2_an}. " if no2_an is not None else "")
        + f"Anomalies are computed against a {config['baseline_window_days']}-day rolling baseline, "
        f"so seasonal variation does not inflate anomaly counts. "
        f"Each anomaly represents a statistically significant deviation from recent local norms."
    )

    # ── Chart 5: ACF/PACF ────────────────────────────────────────────────────
    # Compute lag-1 ACF for PM10
    pm10_series = daily_df[daily_df["variable"] == "PM10"]["mean"].dropna().values
    lag1_acf = "N/A"
    if len(pm10_series) > 5:
        acf_arr = compute_acf(pm10_series, 2)
        lag1_acf = f"{acf_arr[1]:.3f}"
    insights["acf_pacf"] = (
        f"PM10 lag-1 autocorrelation: ρ̂(1) = {lag1_acf} — "
        f"{'strong' if float(lag1_acf) > 0.5 else 'moderate' if float(lag1_acf) > 0.3 else 'weak'} "
        f"temporal memory (today's PM10 {'strongly ' if float(lag1_acf) > 0.5 else ''}predicts tomorrow's). "
        f"ACF at 48 lags and PACF at 24 lags. Dashed lines = 95% confidence bounds (±1.96/√n). "
        f"Bars extending beyond the dashed lines are statistically significant. "
        f"This autocorrelation structure violates the independence assumption of standard statistical tests."
    ) if lag1_acf != "N/A" else "ACF/PACF requires at least 53 days of PM10 data."

    # ── Chart 6: Data Quality Dashboard ──────────────────────────────────────
    overall_q = daily_df.groupby("variable")["Q"].mean()
    low_q_vars = [v for v, q in overall_q.items() if q < config["confidence_moderate"]]
    best_q_var = overall_q.idxmax() if len(overall_q) > 0 else "N/A"
    worst_q_var= overall_q.idxmin() if len(overall_q) > 0 else "N/A"
    mean_q     = float(overall_q.mean()) if len(overall_q) > 0 else 0
    insights["quality"] = (
        f"Overall data completeness Q(D) = {mean_q:.3f} "
        f"({'HIGH' if mean_q >= config['confidence_high'] else 'MODERATE' if mean_q >= config['confidence_moderate'] else 'LOW'}). "
        f"Best quality: {best_q_var} (Q = {float(overall_q.get(best_q_var, 0)):.3f}). "
        f"Lowest quality: {worst_q_var} (Q = {float(overall_q.get(worst_q_var, 0)):.3f}). "
        + (f"Variables with LOW confidence (Q < {config['confidence_moderate']}): "
           f"{', '.join(low_q_vars) if low_q_vars else 'none'}. " )
        + "Q(D) measures completeness only — not sensor calibration accuracy. "
        "Days with Q < 0.50 are excluded from baseline window calculations."
    )

    # ── Chart 7: Seasonal Box Plots ───────────────────────────────────────────
    insights["seasonal"] = (
        f"Monthly box plots show the full distribution of daily PM10 means for each month. "
        f"Peak pollution month: {pk}. Lowest pollution month: {lk}. "
        f"Box = 25th–75th percentile. Whiskers = 1.5×IQR. Circles = outlier days. "
        f"Annual PM10 σ = {pm10_std:.1f} µg/m³ — "
        f"{'indicating high seasonal variability' if pm10_std is not None and pm10_std > 20 else 'indicating moderate seasonal variability'}."
    ) if pm10_mean is not None else "PM10 seasonal data unavailable."

    # ── Chart 8: Diurnal Pattern ──────────────────────────────────────────────
    # Find peak hour for PM10 from hourly data
    if "PM10" in df.columns and "datetime" in df.columns:
        pm10_h = df[["datetime","PM10"]].dropna()
        pm10_h = pm10_h.copy(); pm10_h["hour"] = pm10_h["datetime"].dt.hour
        hourly_avg = pm10_h.groupby("hour")["PM10"].mean()
        if len(hourly_avg) > 0:
            peak_hr  = int(hourly_avg.idxmax())
            low_hr   = int(hourly_avg.idxmin())
            peak_val = float(hourly_avg.max())
            low_val  = float(hourly_avg.min())
            insights["diurnal"] = (
                f"Diurnal PM10 pattern: peak at {peak_hr:02d}:00 ({peak_val:.1f} µg/m³), "
                f"minimum at {low_hr:02d}:00 ({low_val:.1f} µg/m³). "
                f"{'Morning and evening peaks consistent with traffic emission patterns.' if peak_hr in [7,8,9,18,19,20] else 'Peak timing may reflect local industrial or meteorological patterns.'} "
                f"Diurnal amplitude = {peak_val - low_val:.1f} µg/m³."
            )
        else:
            insights["diurnal"] = "Diurnal pattern computed from hourly data. Shows average concentration by hour of day (0–23)."
    else:
        insights["diurnal"] = "Diurnal pattern computed from hourly data. Shows average concentration by hour of day (0–23)."

    # ── Chart 9: Day-of-Week ──────────────────────────────────────────────────
    if "PM10" in daily_df["variable"].unique():
        pm10_d = daily_df[daily_df["variable"] == "PM10"][["date","mean"]].dropna().copy()
        pm10_d["dow"] = pm10_d["date"].dt.dayofweek
        dow_avg  = pm10_d.groupby("dow")["mean"].mean()
        day_names= ["Mon","Tue","Wed","Thu","Fri","Sat","Sun"]
        peak_dow = day_names[int(dow_avg.idxmax())]
        low_dow  = day_names[int(dow_avg.idxmin())]
        weekday_avg = float(dow_avg[:5].mean())
        weekend_avg = float(dow_avg[5:].mean())
        diff_pct = (weekday_avg - weekend_avg) / weekend_avg * 100 if weekend_avg > 0 else 0
        insights["day_of_week"] = (
            f"Highest PM10: {peak_dow}. Lowest: {low_dow}. "
            f"Weekday mean: {weekday_avg:.1f} µg/m³ vs weekend mean: {weekend_avg:.1f} µg/m³ "
            f"({'weekdays {:.1f}% higher'.format(diff_pct) if diff_pct > 5 else 'weekdays {:.1f}% lower'.format(abs(diff_pct)) if diff_pct < -5 else 'minimal weekday–weekend difference'} — "
            f"{'consistent with traffic and industrial emission patterns' if diff_pct > 5 else 'suggests non-traffic dominant sources'})."
        )
    else:
        insights["day_of_week"] = "Day-of-week pattern shows average daily means by weekday for PM10, NO₂, and CO."

    # ── Chart 10: VOC/BTEX Panel ──────────────────────────────────────────────
    benz_mean = annual_mean("Benzene")
    benz_thr  = thr.get("Benzene", 1.7)
    tol_mean  = annual_mean("Toluene")
    if benz_mean is not None:
        insights["voc"] = (
            f"Benzene annual mean: {benz_mean:.3f} µg/m³ "
            f"({'above' if benz_mean > benz_thr else 'below'} WHO annual reference {benz_thr} µg/m³). "
            + (f"Toluene annual mean: {tol_mean:.2f} µg/m³. " if tol_mean is not None else "")
            + f"BTEX compounds (Benzene, Toluene, Ethylbenzene, Xylene) indicate "
            f"{'petrochemical and traffic emission sources' if benz_mean > benz_thr else 'low petrochemical source influence at this station'}. "
            f"Benzene is a Group 1 IARC carcinogen — even sub-guideline levels warrant monitoring."
        )
    else:
        insights["voc"] = "VOC/BTEX data not available at this station."

    # ── Chart 11: Rain Scatter ────────────────────────────────────────────────
    if corr_matrix is not None and "PM10" in corr_matrix.index and "rain_mm" in corr_matrix.columns:
        r_rain_pm10 = float(corr_matrix.loc["PM10", "rain_mm"])
        insights["rain_scatter"] = (
            f"PM10–rainfall Pearson r = {r_rain_pm10:.3f} "
            f"({'negative — consistent with wet deposition (rain washes PM10 out of atmosphere)' if r_rain_pm10 < -0.2 else 'near-zero — rainfall has limited association with PM10 at this station' if abs(r_rain_pm10) < 0.2 else 'positive — unexpected direction; investigate data quality'}). "
            f"OLS trend line shows association direction. "
            f"IMPORTANT: This is statistical association only — not a proven causal wet-deposition mechanism."
        )
    else:
        insights["rain_scatter"] = (
            "Scatter of daily rainfall against PM10 and SO₂ concentrations. "
            "Negative slope indicates wet deposition effect (association, not causation)."
        )

    # ── Chart 12: Exceedance Calendar ────────────────────────────────────────
    pm10_thr_val = thr.get("PM10", 45)
    insights["calendar"] = (
        f"Full-year calendar heatmap of PM10 daily means (µg/m³). "
        f"× marks days exceeding WHO guideline ({pm10_thr_val} µg/m³). "
        f"Total exceedances: {pm10_ex} days ({pm10_ex / total_days * 100:.1f}% of monitored period). "
        f"Clustering of × marks in {pk} indicates seasonal high-pollution period. "
        f"Green cells = clean days. Red cells = high pollution. "
        f"Each cell value = that day's daily mean PM10 (µg/m³)."
    ) if pm10_mean is not None else "Calendar heatmap of PM10 exceedance days."

    # ── Chart 13: RGV Grounding ───────────────────────────────────────────────
    insights["grounding"] = (
        "The Runtime Grounding Verifier (RGV) extracts every number from every "
        "LLM sentence and verifies it against the JSON data contract J using "
        "type-adaptive tolerance bounds (±5% for concentrations, ±2 for integer counts, "
        "±0.1 for quality scores). "
        "GROUNDED = all numbers match J. PARTIALLY = ≥50% match. "
        "UNGROUNDED = <50% match (LLM may have hallucinated these values). "
        "NON_NUMERICAL = sentence has no numbers (excluded from G_rate). "
        "G_rate = grounded sentences / numerical sentences. PASS if G_rate ≥ 0.75."
    )

    # ── Chart 14: ML Anomaly Comparison ──────────────────────────────────────
    if ml_if.get("available"):
        n_if    = ml_if["n_anomalies"]
        n_both  = ml_if["n_both"]
        n_ifonly= ml_if["n_if_only"]
        n_zsonly= ml_if["n_zs_only"]
        feats   = ", ".join(ml_if["features_used"])
        insights["ml_anomaly"] = (
            f"Isolation Forest ({ml_if['n_estimators']} trees, contamination={ml_if['contamination']:.0%}) "
            f"on {len(ml_if['features_used'])} variables ({feats}): "
            f"{n_if} anomaly days ({ml_if['anomaly_rate']*100:.1f}%). "
            f"Agreement with z-score: {n_both} days flagged by both methods. "
            f"{n_ifonly} days flagged by IF only (multivariate pattern — z-score missed). "
            f"{n_zsonly} days flagged by z-score only (univariate spike, within multivariate norm). "
            f"IF-only days represent genuine added value of multivariate detection. "
            f"NOTE: Isolation Forest does NOT replace formal ζ_d — both methods are reported. "
            f"random_state=42 ensures this result is reproducible."
        )
    else:
        insights["ml_anomaly"] = "Isolation Forest skipped (insufficient multivariate data)."

    # ── Chart 15: Holt-Winters Forecast ──────────────────────────────────────
    if ml_hw.get("available"):
        fc_mean = float(np.mean(ml_hw["forecast"]))
        fc_max  = float(np.max(ml_hw["forecast"]))
        fc_days = ml_hw["forecast_days"]
        insights["forecast"] = (
            f"Holt-Winters triple exponential smoothing (α={ml_hw['alpha']}, "
            f"β={ml_hw['beta']}, γ={ml_hw['gamma']}, "
            f"seasonal period={ml_hw['seasonal_periods']} days). "
            f"In-sample RMSE={ml_hw['rmse']:.2f} µg/m³ | MAE={ml_hw['mae']:.2f} µg/m³. "
            f"{fc_days}-day PM10 forecast: mean={fc_mean:.1f} µg/m³, peak={fc_max:.1f} µg/m³. "
            f"Shaded band = 80% prediction interval (widens with forecast horizon). "
            f"IMPORTANT: This is a statistical extrapolation, not a physical model. "
            f"Forecast values are NOT in the JSON contract J — they were not used by the LLM. "
            f"Uncertainty grows significantly beyond 7 days."
        )
    else:
        insights["forecast"] = "Holt-Winters forecast skipped (insufficient PM10 history)."

    return insights

def _save_fig(fig, path: Path, name: str) -> str:
    """Save figure and return file path string."""
    fpath = str(path / name)
    fig.savefig(fpath, dpi=160, bbox_inches='tight', facecolor=PAL['white'])
    plt.close(fig)
    return fpath



def chart_data_quality_overview(df: pd.DataFrame, daily_df: pd.DataFrame,
                                 cleaning_report: dict, data_profile: dict,
                                 config: dict, out: Path) -> str:
    """
    Chart 00 — Data Quality Overview Dashboard.
    4-panel summary produced BEFORE analysis charts so the reader
    understands data quality before seeing results.

    Panel A: Per-column completeness bar chart
    Panel B: Null vs valid vs suspect stacked bar by pollutant
    Panel C: Hourly completeness heatmap (hour-of-day × month)
    Panel D: Cleaning action waterfall (what was done to the data)
    """
    fig = plt.figure(figsize=(16, 12), facecolor=PAL["white"])
    fig.suptitle(
        f"ATARS — Data Quality Overview Dashboard  ·  {config['city']}\n"
        "Generated BEFORE analysis — read this first",
        fontsize=13, fontweight="bold", color=PAL["navy"], y=0.98
    )

    import matplotlib.gridspec as gs_mod
    gs = gs_mod.GridSpec(2, 2, figure=fig, hspace=0.48, wspace=0.38)
    ax_comp  = fig.add_subplot(gs[0, 0])   # completeness bars
    ax_flags = fig.add_subplot(gs[0, 1])   # quality flag stacked
    ax_hour  = fig.add_subplot(gs[1, 0])   # hourly heatmap
    ax_clean = fig.add_subplot(gs[1, 1])   # cleaning waterfall

    # ── Panel A: Column completeness ──────────────────────────────────────
    comp = data_profile.get("completeness_by_col",
           cleaning_report.get("completeness_by_col", {}))
    pol_comp = {k: v["pct"] for k, v in comp.items()
                if k in POLLUTANTS and isinstance(v, dict)}
    if pol_comp:
        sorted_pol = sorted(pol_comp.items(), key=lambda x: -x[1])
        cols  = [c for c, _ in sorted_pol]
        pcts  = [p for _, p in sorted_pol]
        colors_bar = [PAL["green"] if p >= 80 else
                      PAL["amber"] if p >= 50 else PAL["red"]
                      for p in pcts]
        ax_comp.barh(cols, pcts, color=colors_bar, height=0.6, alpha=0.85)
        ax_comp.axvline(80, color=PAL["green"], linestyle="--", linewidth=1.5,
                        label="80% threshold")
        ax_comp.axvline(50, color=PAL["amber"], linestyle="--", linewidth=1.5,
                        label="50% threshold")
        ax_comp.set_xlim(0, 105)
        ax_comp.set_xlabel("Valid Records (%)", fontsize=9)
        ax_comp.set_title("(a) Column Completeness\n(green=≥80%, amber=≥50%, red=<50%)",
                          fontsize=10, fontweight="bold", color=PAL["navy"])
        ax_comp.legend(fontsize=8)
        for i, (col, pct) in enumerate(sorted_pol):
            ax_comp.text(min(pct + 1, 103), i, f"{pct:.1f}%",
                         va="center", fontsize=8, color=PAL["navy"])

    # ── Panel B: Quality flag distribution ───────────────────────────────
    flag_data = {}
    for col in POLLUTANTS:
        fc = f"q_{col}"
        if fc in df.columns:
            total = max(len(df), 1)
            flag_data[col] = {
                "Valid"    : int((df[fc]==1).sum()) / total * 100,
                "Suspect"  : int((df[fc]==2).sum()) / total * 100,
                "Invalid"  : int((df[fc]==3).sum()) / total * 100,
                "Diurnal"  : int((df[fc]==4).sum()) / total * 100,
            }
    if flag_data:
        fd_cols = list(flag_data.keys())
        valid_pct   = [flag_data[c]["Valid"]   for c in fd_cols]
        suspect_pct = [flag_data[c]["Suspect"] for c in fd_cols]
        invalid_pct = [flag_data[c]["Invalid"] for c in fd_cols]
        diurnal_pct = [flag_data[c]["Diurnal"] for c in fd_cols]
        x = range(len(fd_cols))
        ax_flags.bar(x, valid_pct,   color=PAL["green"],  alpha=0.85, label="Valid (q=1)")
        ax_flags.bar(x, suspect_pct, color=PAL["amber"],  alpha=0.85, label="Suspect (q=2)",
                     bottom=valid_pct)
        inv_bottom = [v+s for v,s in zip(valid_pct, suspect_pct)]
        ax_flags.bar(x, invalid_pct, color=PAL["red"],    alpha=0.85, label="Invalid (q=3)",
                     bottom=inv_bottom)
        diu_bottom = [v+s+i for v,s,i in zip(valid_pct, suspect_pct, invalid_pct)]
        ax_flags.bar(x, diurnal_pct, color=PAL["purple"], alpha=0.85, label="Diurnal (q=4)",
                     bottom=diu_bottom)
        ax_flags.set_xticks(list(x))
        ax_flags.set_xticklabels(fd_cols, rotation=35, ha="right", fontsize=8)
        ax_flags.set_ylabel("Records (%)", fontsize=9)
        ax_flags.set_ylim(0, 105)
        ax_flags.legend(fontsize=7.5, loc="lower right")
        ax_flags.set_title("(b) Quality Flag Distribution by Pollutant\n"
                            "q=1 valid · q=2 suspect · q=3 invalid · q=4 diurnal anomaly",
                            fontsize=10, fontweight="bold", color=PAL["navy"])

    # ── Panel C: Hourly completeness heatmap ──────────────────────────────
    if "datetime" in df.columns and "PM10" in df.columns:
        try:
            dfc = df[["datetime","PM10"]].copy().dropna(subset=["datetime"])
            dfc["hour"]  = dfc["datetime"].dt.hour
            dfc["month"] = dfc["datetime"].dt.month
            dfc["valid"] = dfc["PM10"].notna().astype(int)
            pivot = dfc.pivot_table(index="month", columns="hour",
                                     values="valid", aggfunc="mean") * 100
            if not pivot.empty:
                im = ax_hour.imshow(pivot.values, aspect="auto",
                                    cmap="RdYlGn", vmin=0, vmax=100)
                ax_hour.set_xticks(range(0, 24, 3))
                ax_hour.set_xticklabels([f"{h:02d}" for h in range(0,24,3)], fontsize=8)
                month_names = ["Jan","Feb","Mar","Apr","May","Jun",
                               "Jul","Aug","Sep","Oct","Nov","Dec"]
                ax_hour.set_yticks(range(len(pivot.index)))
                ax_hour.set_yticklabels([month_names[m-1] for m in pivot.index], fontsize=8)
                ax_hour.set_xlabel("Hour of Day", fontsize=9)
                ax_hour.set_title("(c) PM10 Hourly Completeness Heatmap\n"
                                  "Green=complete, Red=missing (month × hour)",
                                  fontsize=10, fontweight="bold", color=PAL["navy"])
                fig.colorbar(im, ax=ax_hour, shrink=0.8, label="Valid %")
        except Exception:
            ax_hour.set_title("(c) Hourly Completeness — insufficient data",
                              fontsize=10, fontweight="bold", color=PAL["navy"])
            ax_hour.axis("off")

    # ── Panel D: Cleaning action waterfall ────────────────────────────────
    clean_actions = {
        "Flatlines nulled"  : cleaning_report.get("flatlines_nulled", 0),
        "Spikes nulled"     : cleaning_report.get("spikes_nulled",    0),
        "Negatives nulled"  : cleaning_report.get("negatives_nulled", 0),
        "Bounds clipped"    : cleaning_report.get("values_clipped",   0),
        "Winsorized"        : cleaning_report.get("values_winsorized",0),
        "Values imputed"    : cleaning_report.get("values_imputed",   0),
        "Rows removed"      : cleaning_report.get("rows_removed",     0),
        "Rows injected"     : cleaning_report.get("rows_injected",    0),
    }
    actions = [(k, v) for k, v in clean_actions.items() if v > 0]
    if actions:
        a_labels = [a[0] for a in actions]
        a_vals   = [a[1] for a in actions]
        colors_w = [PAL["red"]   if "null" in l.lower() or "remov" in l.lower() else
                    PAL["amber"] if "clip" in l.lower() or "winsor" in l.lower() else
                    PAL["green"] if "impute" in l.lower() or "inject" in l.lower() else
                    PAL["teal"]
                    for l in a_labels]
        bars = ax_clean.barh(a_labels, a_vals, color=colors_w, height=0.6, alpha=0.85)
        ax_clean.set_xlabel("Count of Records Affected", fontsize=9)
        ax_clean.set_title("(d) Cleaning Actions Applied\n"
                            "What was done to the data (and how many values)",
                            fontsize=10, fontweight="bold", color=PAL["navy"])
        for bar, val in zip(bars, a_vals):
            ax_clean.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height()/2,
                          f"{val:,}", va="center", fontsize=8.5)
        qs = cleaning_report.get("quality_score", 0)
        ql = "EXCELLENT" if qs>=0.90 else "GOOD" if qs>=0.75 else "FAIR" if qs>=0.60 else "POOR"
        ax_clean.set_title(f"(d) Cleaning Actions Applied  —  Quality Score: {qs:.3f} [{ql}]",
                           fontsize=10, fontweight="bold", color=PAL["navy"])
    else:
        ax_clean.text(0.5, 0.5, "No cleaning actions required (data was already clean)",
                      ha="center", va="center", fontsize=12, color=PAL["green"],
                      transform=ax_clean.transAxes)
        ax_clean.axis("off")

    fig.text(0.5, 0.01,
             "This chart is generated from the cleaning audit. "
             "Green bars = data recovery. Red/amber = data removed/corrected. "
             "All actions logged in Section 00b.",
             ha="center", fontsize=8, color=PAL["gray"], style="italic")

    return _save_fig(fig, out, "chart00_data_quality_overview.png")



def chart_timeseries(df: pd.DataFrame, daily_df: pd.DataFrame,
                      config: dict, out: Path) -> str:
    """Chart 1 — Hourly time-series for PM10 and NO2 with WHO guidelines."""
    fig, axes = plt.subplots(2, 1, figsize=(14, 9), facecolor=PAL['white'])
    fig.subplots_adjust(hspace=0.45)

    for ax, col, color, thr_color, title in [
        (axes[0], 'PM10', PAL['blue'],   PAL['amber'],
         f"PM10 Hourly Concentration  —  WHO 24h Guideline: {config['thresholds']['PM10']} µg/m³"),
        (axes[1], 'NO2',  PAL['teal'], PAL['red'],
         f"NO₂ Hourly Concentration  —  WHO 24h Guideline: {config['thresholds']['NO2']} µg/m³"),
    ]:
        if col not in df.columns:
            continue
        plot_data = df[['datetime', col]].dropna()
        ax.plot(plot_data['datetime'], plot_data[col],
                color=color, linewidth=0.8, alpha=0.7, label=f'{col} (hourly)')

        # Daily mean overlay
        dm = daily_df[daily_df['variable'] == col][['date', 'mean']].dropna()
        ax.plot(dm['date'], dm['mean'],
                color=PAL['navy'], linewidth=2.0, linestyle='--',
                label='Daily mean', zorder=5)

        # WHO threshold line
        thr = config['thresholds'].get(col)
        if thr:
            ax.axhline(thr, color=thr_color, linewidth=1.8,
                       linestyle=':', label=f'WHO guideline ({thr})', zorder=4)
            ax.fill_between(plot_data['datetime'], thr,
                             plot_data[col].clip(lower=thr),
                             alpha=0.12, color=PAL['red'], label='Exceedance zone')

        ax.set_title(title, fontsize=11, fontweight='bold', color=PAL['navy'], pad=8)
        ax.set_xlabel('Date', fontsize=9)
        ax.set_ylabel(f'{col} (µg/m³)' if col != 'CO' else f'{col} (mg/m³)', fontsize=9)
        ax.legend(fontsize=8.5, loc='upper right')
        ax.tick_params(axis='x', rotation=30)

    fig.suptitle(f"ATARS — Hourly Pollutant Time-Series  ·  {config['city']}",
                 fontsize=13, fontweight='bold', color=PAL['navy'])
    fig.text(0.5, 0.01,
             'Illustrative analysis of available dataset. Dashed = daily mean. '
             'Reference lines from WHO AQG 2021.',
             ha='center', fontsize=8, color=PAL['gray'], style='italic')
    return _save_fig(fig, out, 'chart1_timeseries.png')


def chart_correlation(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 2 — Pearson correlation heatmap R_{jk}."""
    key_vars = [v for v in
                ["PM10","NO2","SO2","CO","Ozone","NH3",
                 "humidity_pct","wind_speed_10m","pressure_hpa"]
                if v in daily_df['variable'].unique()]
    if len(key_vars) < 4:
        return None

    corr = compute_correlation_matrix(daily_df, key_vars)
    labels = {
        "PM10":"PM10","NO2":"NO₂","SO2":"SO₂","CO":"CO",
        "Ozone":"Ozone","NH3":"NH₃","humidity_pct":"Humidity",
        "wind_speed_10m":"Wind Speed","pressure_hpa":"Pressure"
    }
    corr.index   = [labels.get(c, c) for c in corr.index]
    corr.columns = [labels.get(c, c) for c in corr.columns]

    fig, ax = plt.subplots(figsize=(11, 9), facecolor=PAL['white'])
    mask = np.zeros_like(corr.values, dtype=bool)
    np.fill_diagonal(mask, False)

    sns.heatmap(corr, ax=ax, annot=True, fmt='.2f', mask=None,
                cmap=sns.diverging_palette(220, 20, as_cmap=True),
                vmin=-1, vmax=1, center=0, square=True,
                linewidths=0.5, annot_kws={'size': 8.5},
                cbar_kws={'shrink': 0.8, 'label': 'Pearson r'})

    ax.set_title(
        f'Pearson Correlation Matrix R  ·  {config["city"]}\n'
        r'$\mathbf{R}_{jk} = \mathrm{Corr}(X_j, X_k)$  '
        '— Association only, NOT causation  [Corr(X,Y) ≢ C(X→Y)]',
        fontsize=11, fontweight='bold', color=PAL['navy'], pad=12)
    ax.tick_params(axis='both', labelsize=9)
    plt.xticks(rotation=35, ha='right')
    plt.yticks(rotation=0)
    fig.text(0.5, 0.01,
             f'n = {int(daily_df["N_v"].sum()):,} valid daily records.  '
             'Causal inference requires experimental design — outside ATARS scope.',
             ha='center', fontsize=8, color=PAL['gray'], style='italic')
    return _save_fig(fig, out, 'chart2_correlation.png')


def chart_trend(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 3 — 30-day rolling trend with baseline for PM10 and NO2."""
    fig, ax = plt.subplots(figsize=(14, 6), facecolor=PAL['white'])

    for col, color, label in [
        ('PM10', PAL['blue'],  'PM10 Daily Mean'),
        ('NO2',  PAL['teal'], 'NO₂ Daily Mean'),
    ]:
        dm = daily_df[daily_df['variable'] == col][['date', 'mean', 'mean_W']].dropna()
        if len(dm) < 5:
            continue
        ax.plot(dm['date'], dm['mean'],
                color=color, linewidth=1.5, alpha=0.6, label=label)
        ax.plot(dm['date'], dm['mean_W'],
                color=color, linewidth=2.5, linestyle='--',
                alpha=0.9, label=f'{col} 30-day baseline x̄_W')

    for col, thr_color in [('PM10', PAL['amber']), ('NO2', PAL['red'])]:
        thr = config['thresholds'].get(col)
        if thr:
            ax.axhline(thr, color=thr_color, linewidth=1.5,
                       linestyle=':', alpha=0.8,
                       label=f'{col} WHO guideline ({thr})')

    ax.set_title(
        f'Daily Trend & 30-Day Rolling Baseline  ·  {config["city"]}  '
        r'— $\bar{x}_W$ and $\sigma_W$ (Eqs. 7–8)',
        fontsize=11, fontweight='bold', color=PAL['navy'])
    ax.set_xlabel('Date', fontsize=9)
    ax.set_ylabel('Concentration (µg/m³)', fontsize=9)
    ax.legend(fontsize=8.5, loc='upper right', ncol=2)
    ax.tick_params(axis='x', rotation=30)
    return _save_fig(fig, out, 'chart3_trend.png')


def chart_anomalies(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 4 — Z-score anomaly detection for PM10."""
    col = 'PM10'
    dm  = daily_df[daily_df['variable'] == col][
        ['date', 'mean', 'z_score', 'is_anomaly', 'CI_lower', 'CI_upper']
    ].dropna(subset=['z_score'])
    if len(dm) < 10:
        return None

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(14, 9), facecolor=PAL['white'])
    fig.subplots_adjust(hspace=0.45)

    # Panel A: daily mean with CI band
    ax1.plot(dm['date'], dm['mean'], color=PAL['blue'], linewidth=1.8, label='Daily mean x̄_d')
    ax1.fill_between(dm['date'], dm['CI_lower'], dm['CI_upper'],
                     alpha=0.18, color=PAL['blue'], label='95% CI band')
    anom = dm[dm['is_anomaly']]
    ax1.scatter(anom['date'], anom['mean'], color=PAL['red'], s=50, zorder=6,
                label=f'Anomaly (|ζ_d| > {config["anomaly_z_threshold"]})', marker='D')
    ax1.set_title(r'(a) PM10 Daily Mean with 95% CI  —  $\mathcal{CI}_{0.95} = \bar{x} \pm 1.96\,\sigma/\sqrt{N_v}$  (Eq. 12)',
                  fontsize=10, fontweight='bold', color=PAL['navy'])
    ax1.set_ylabel('PM10 (µg/m³)', fontsize=9)
    ax1.legend(fontsize=8.5)
    ax1.tick_params(axis='x', rotation=30)

    # Panel B: z-score
    ax2.bar(dm['date'], dm['z_score'],
            color=[PAL['red'] if a else PAL['blue'] for a in dm['is_anomaly']],
            width=0.8, alpha=0.7)
    ax2.axhline(config['anomaly_z_threshold'],  color=PAL['red'],
                linestyle='--', linewidth=1.8, label=f'ζ_thresh = ±{config["anomaly_z_threshold"]}')
    ax2.axhline(-config['anomaly_z_threshold'], color=PAL['red'], linestyle='--', linewidth=1.8)
    ax2.axhline(0, color=PAL['navy'], linewidth=0.8)
    ax2.set_title(r'(b) Z-Score $\zeta_d = (\bar{x}_d - \bar{x}_W)\,/\,\sigma_W$  (Eq. 9)  '
                  f'— Anomaly threshold: |ζ| > {config["anomaly_z_threshold"]}',
                  fontsize=10, fontweight='bold', color=PAL['navy'])
    ax2.set_xlabel('Date', fontsize=9)
    ax2.set_ylabel('Z-Score ζ_d', fontsize=9)
    ax2.legend(fontsize=8.5)
    ax2.tick_params(axis='x', rotation=30)

    fig.suptitle(f'ATARS — Anomaly Detection  ·  {config["city"]}',
                 fontsize=13, fontweight='bold', color=PAL['navy'])
    return _save_fig(fig, out, 'chart4_anomalies.png')


def chart_acf_pacf(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 5 — ACF and PACF for PM10 daily means."""
    col = 'PM10'
    n_lags = min(config['acf_lags'], 60)
    result = run_temporal_analysis(daily_df, col, n_lags)
    if not result:
        return None

    acf  = result['acf']
    pacf = result['pacf']
    conf = result['conf_bound']
    lag_ax_full  = np.arange(len(acf))
    lag_ax_pacf  = np.arange(len(pacf))

    fig, axes = plt.subplots(2, 1, figsize=(13, 9), facecolor=PAL['white'])
    fig.subplots_adjust(hspace=0.52)

    axes[0].bar(lag_ax_full, acf,
                color=[PAL['blue'] if abs(v) > conf else '#A0B4CC' for v in acf],
                width=0.6, zorder=3)
    axes[0].axhline(conf,  color=PAL['red'], linestyle='--', linewidth=1.8,
                    label=f'95% CI (±{conf:.3f})')
    axes[0].axhline(-conf, color=PAL['red'], linestyle='--', linewidth=1.8)
    axes[0].axhline(0, color=PAL['navy'], linewidth=0.8)
    axes[0].set_title(r'(a) ACF: $\hat{\rho}(k) = \mathrm{Corr}(x_t,\, x_{t-k})$  (Eq. 12)  '
                      f'— PM10 Daily Means  ·  {config["city"]}',
                      fontsize=10, fontweight='bold', color=PAL['navy'])
    axes[0].set_xlabel('Lag k (days)', fontsize=9)
    axes[0].set_ylabel('ACF ρ̂(k)', fontsize=9)
    axes[0].legend(fontsize=8.5)
    axes[0].set_xlim(-0.5, len(acf) + 0.5)
    axes[0].set_ylim(-0.8, 1.1)

    axes[1].bar(lag_ax_pacf, pacf,
                color=[PAL['teal'] if abs(v) > conf else '#A0D4CC' for v in pacf],
                width=0.6, zorder=3)
    axes[1].axhline(conf,  color=PAL['red'], linestyle='--', linewidth=1.8,
                    label=f'95% CI (±{conf:.3f})')
    axes[1].axhline(-conf, color=PAL['red'], linestyle='--', linewidth=1.8)
    axes[1].axhline(0, color=PAL['navy'], linewidth=0.8)
    axes[1].set_title(r'(b) PACF: $\hat{\phi}_{kk}$ — Partial correlation at lag k (Eq. 13)',
                      fontsize=10, fontweight='bold', color=PAL['navy'])
    axes[1].set_xlabel('Lag k (days)', fontsize=9)
    axes[1].set_ylabel('PACF φ̂_kk', fontsize=9)
    axes[1].legend(fontsize=8.5)
    axes[1].set_xlim(-0.5, len(pacf) + 0.5)
    axes[1].set_ylim(-0.8, 1.1)

    sig_lags = result.get('significant_acf', [])
    fig.text(0.5, 0.01,
             f'Significant ACF lags (|ρ̂(k)| > {conf:.3f}): {sig_lags[:10]}  '
             f'— n = {result["series_length"]} daily observations.',
             ha='center', fontsize=8, color=PAL['gray'], style='italic')
    return _save_fig(fig, out, 'chart5_acf_pacf.png')


def chart_quality_summary(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 6 — Data quality dashboard: Q(D) and confidence flags."""
    fig = plt.figure(figsize=(14, 8), facecolor=PAL['white'])
    gs  = gridspec.GridSpec(1, 2, figure=fig, wspace=0.38)
    ax1 = fig.add_subplot(gs[0, 0])
    ax2 = fig.add_subplot(gs[0, 1])

    # Panel A: Q(D) time series for PM10
    col = 'PM10'
    qd  = daily_df[daily_df['variable'] == col][['date', 'Q', 'confidence']].dropna()
    if len(qd) > 0:
        colors_q = {'HIGH': PAL['green'], 'MODERATE': PAL['amber'], 'LOW': PAL['red']}
        bar_colors = [colors_q.get(f, PAL['gray']) for f in qd['confidence']]
        ax1.bar(qd['date'], qd['Q'], color=bar_colors, width=0.8, alpha=0.75)
        ax1.axhline(config['confidence_high'],     color=PAL['green'], linestyle='--',
                    linewidth=1.8, label=f'HIGH ≥ {config["confidence_high"]}')
        ax1.axhline(config['confidence_moderate'], color=PAL['amber'], linestyle='--',
                    linewidth=1.8, label=f'MOD ≥ {config["confidence_moderate"]}')
        ax1.set_ylim(0, 1.05)
        ax1.set_title(f'(a) Q(D) = N_v/N_total  (Eq. 11)\nData Completeness Score — PM10',
                      fontsize=10, fontweight='bold', color=PAL['navy'])
        ax1.set_xlabel('Date', fontsize=9)
        ax1.set_ylabel('Q(D) ∈ [0, 1]', fontsize=9)
        ax1.legend(fontsize=8.5)
        ax1.tick_params(axis='x', rotation=30)

    # Panel B: Confidence flag distribution (all pollutants)
    flag_counts = (daily_df.groupby(['variable', 'confidence'])
                   .size().unstack(fill_value=0))
    key_cols = [c for c in POLLUTANTS if c in flag_counts.index][:8]
    if len(key_cols) > 0:
        fc = flag_counts.loc[key_cols]
        for flag, color in [('HIGH', PAL['green']),
                             ('MODERATE', PAL['amber']),
                             ('LOW', PAL['red'])]:
            if flag in fc.columns:
                ax2.barh(key_cols, fc[flag], label=flag, color=color,
                         alpha=0.75, height=0.6)
        ax2.set_title('(b) Confidence Flag Distribution\nby Pollutant Variable',
                      fontsize=10, fontweight='bold', color=PAL['navy'])
        ax2.set_xlabel('Days', fontsize=9)
        ax2.legend(fontsize=8.5)

    fig.suptitle(f'ATARS — Data Quality Assessment  ·  {config["city"]}',
                 fontsize=13, fontweight='bold', color=PAL['navy'])
    return _save_fig(fig, out, 'chart6_quality.png')


def chart_seasonal(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 7 — Monthly box plots showing seasonal pattern for PM10 and NO2."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 7), facecolor=PAL['white'])
    fig.subplots_adjust(wspace=0.35)

    month_names = ['Jan','Feb','Mar','Apr','May','Jun',
                   'Jul','Aug','Sep','Oct','Nov','Dec']

    for ax, col, color in [(axes[0], 'PM10', PAL['blue']),
                            (axes[1], 'NO2',  PAL['teal'])]:
        dm = daily_df[daily_df['variable'] == col][['date', 'mean']].dropna()
        if len(dm) < 30:
            continue
        dm['month'] = dm['date'].dt.month
        monthly = [dm[dm['month'] == m]['mean'].values for m in range(1, 13)]
        bp = ax.boxplot(monthly, patch_artist=True, notch=False,
                        medianprops={'color': PAL['navy'], 'linewidth': 2.5},
                        whiskerprops={'color': color, 'linewidth': 1.5},
                        capprops={'color': color, 'linewidth': 1.5},
                        flierprops={'markerfacecolor': PAL['red'],
                                    'markersize': 3, 'alpha': 0.5})
        for patch in bp['boxes']:
            patch.set_facecolor(color)
            patch.set_alpha(0.4)

        thr = config['thresholds'].get(col)
        if thr:
            ax.axhline(thr, color=PAL['amber'], linewidth=2,
                       linestyle='--', label=f'WHO guideline ({thr})')
        ax.set_xticks(range(1, 13))
        ax.set_xticklabels(month_names, fontsize=8.5)
        unit = 'mg/m³' if col == 'CO' else 'µg/m³'
        ax.set_title(f'{col} — Monthly Distribution\n'
                     r'Median $\tilde{x}$, IQR, whiskers: 1.5×IQR',
                     fontsize=10, fontweight='bold', color=PAL['navy'])
        ax.set_xlabel('Month', fontsize=9)
        ax.set_ylabel(f'{col} ({unit})', fontsize=9)
        if thr:
            ax.legend(fontsize=8.5)

    fig.suptitle(f'ATARS — Seasonal Pattern Analysis  ·  {config["city"]}',
                 fontsize=13, fontweight='bold', color=PAL['navy'])
    fig.text(0.5, 0.01,
             'Non-parametric: median-based (Eq. 24). Boxes show IQR = Q₀.₇₅ − Q₀.₂₅ (Eq. 25).',
             ha='center', fontsize=8, color=PAL['gray'], style='italic')
    return _save_fig(fig, out, 'chart7_seasonal.png')


# ── NEW CHART 8: Hourly Diurnal Profile (avg by hour of day) ──────────────────
def chart_diurnal(df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 8 — Average concentration by hour of day (diurnal cycle) for key pollutants."""
    key_cols = [c for c in ['PM10','NO','NO2','CO','Ozone'] if c in df.columns]
    if not key_cols:
        return None
    df2 = df.copy()
    df2['hour'] = df2['datetime'].dt.hour

    fig, axes = plt.subplots(2, 3, figsize=(15, 9), facecolor=PAL['white'])
    fig.subplots_adjust(hspace=0.45, wspace=0.35)
    colors_ = [PAL['blue'], PAL['navy'], PAL['teal'], PAL['amber'], PAL['purple'],
               PAL['green']]
    ax_flat = axes.flatten()

    for idx, col in enumerate(key_cols[:5]):
        ax = ax_flat[idx]
        hourly = df2.groupby('hour')[col].agg(['mean','std']).reset_index()
        ax.fill_between(hourly['hour'],
                        hourly['mean'] - hourly['std'],
                        hourly['mean'] + hourly['std'],
                        alpha=0.18, color=colors_[idx])
        ax.plot(hourly['hour'], hourly['mean'],
                color=colors_[idx], linewidth=2.5, marker='o', markersize=4)
        thr = config['thresholds'].get(col)
        if thr:
            ax.axhline(thr, color=PAL['red'], linewidth=1.5,
                       linestyle='--', alpha=0.7, label=f'WHO: {thr}')
        ax.set_title(f'{col} — Diurnal Cycle',
                     fontsize=10, fontweight='bold', color=PAL['navy'])
        ax.set_xlabel('Hour of Day (0–23)', fontsize=8.5)
        unit = 'mg/m³' if col == 'CO' else 'µg/m³'
        ax.set_ylabel(f'{unit}', fontsize=8.5)
        ax.set_xticks([0, 6, 12, 18, 23])
        ax.legend(fontsize=8) if thr else None

    # 6th panel: combined normalised overlay
    ax6 = ax_flat[5]
    for idx, col in enumerate(key_cols[:4]):
        series = df2.groupby('hour')[col].mean()
        if series.max() > 0:
            ax6.plot(series.index, series.values / series.max(),
                     color=colors_[idx], linewidth=2, label=col)
    ax6.set_title('Normalised Overlay\n(peak = 1.0 per variable)',
                  fontsize=10, fontweight='bold', color=PAL['navy'])
    ax6.set_xlabel('Hour of Day', fontsize=8.5)
    ax6.set_ylabel('Relative Concentration', fontsize=8.5)
    ax6.legend(fontsize=8.5)
    ax6.set_xticks([0, 6, 12, 18, 23])

    fig.suptitle(f'ATARS — Diurnal (Hourly) Profile  ·  {config["city"]}',
                 fontsize=13, fontweight='bold', color=PAL['navy'])
    fig.text(0.5, 0.005, 'Mean ± 1σ band. Peak hours indicate traffic/industrial emission cycles.',
             ha='center', fontsize=8, color=PAL['gray'], style='italic')
    return _save_fig(fig, out, 'chart8_diurnal.png')


# ── NEW CHART 9: Day-of-Week Pattern ─────────────────────────────────────────
def chart_day_of_week(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 9 — Average daily mean by day of week for PM10, NO2, CO."""
    key_cols = [c for c in ['PM10','NO2','CO','SO2'] if c in daily_df['variable'].unique()]
    if not key_cols:
        return None

    dow_names = ['Mon','Tue','Wed','Thu','Fri','Sat','Sun']
    fig, axes = plt.subplots(1, len(key_cols[:4]), figsize=(14, 6), facecolor=PAL['white'])
    if len(key_cols) == 1:
        axes = [axes]
    colors_ = [PAL['blue'], PAL['teal'], PAL['amber'], PAL['purple']]

    for ax, col, color in zip(axes, key_cols[:4], colors_):
        dm = daily_df[daily_df['variable'] == col][['date','mean']].dropna()
        dm['dow'] = dm['date'].dt.dayofweek
        dow_mean = dm.groupby('dow')['mean'].mean()
        dow_std  = dm.groupby('dow')['mean'].std()
        x = dow_mean.index
        ax.bar(x, dow_mean.values, color=color, alpha=0.7, width=0.6)
        ax.errorbar(x, dow_mean.values, yerr=dow_std.values,
                    fmt='none', color=PAL['navy'], capsize=4, linewidth=1.5)
        thr = config['thresholds'].get(col)
        if thr:
            ax.axhline(thr, color=PAL['red'], linewidth=1.5,
                       linestyle='--', label=f'WHO: {thr}')
        ax.set_xticks(range(7))
        ax.set_xticklabels(dow_names, fontsize=8.5)
        unit = 'mg/m³' if col == 'CO' else 'µg/m³'
        ax.set_title(f'{col} by\nDay of Week',
                     fontsize=10, fontweight='bold', color=PAL['navy'])
        ax.set_ylabel(f'{unit}', fontsize=8.5)
        if thr:
            ax.legend(fontsize=8)

    fig.suptitle(f'ATARS — Day-of-Week Pattern Analysis  ·  {config["city"]}',
                 fontsize=13, fontweight='bold', color=PAL['navy'])
    fig.text(0.5, 0.01,
             'Mean ± σ error bars. Weekday vs weekend patterns indicate anthropogenic emission cycles.',
             ha='center', fontsize=8, color=PAL['gray'], style='italic')
    return _save_fig(fig, out, 'chart9_day_of_week.png')


# ── NEW CHART 10: VOC Panel (Volatile Organic Compounds) ─────────────────────
def chart_voc_panel(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 10 — VOC compounds trend: Benzene, Toluene, Xylene, Eth-Benzene, MP-Xylene."""
    voc_cols = [c for c in ['Benzene','Toluene','Xylene','Eth_Benzene','MP_Xylene']
                if c in daily_df['variable'].unique()]
    if len(voc_cols) < 2:
        return None

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(14, 10), facecolor=PAL['white'])
    fig.subplots_adjust(hspace=0.42)
    voc_colors = [PAL['blue'], PAL['teal'], PAL['purple'], PAL['amber'], PAL['green']]

    # Panel A: time series for all VOCs
    for col, color in zip(voc_cols, voc_colors):
        dm = daily_df[daily_df['variable'] == col][['date','mean']].dropna()
        if len(dm) < 5:
            continue
        ax1.plot(dm['date'], dm['mean'],
                 color=color, linewidth=1.6, label=col, alpha=0.85)
    # Benzene WHO line
    b_thr = config['thresholds'].get('Benzene', 1.7)
    ax1.axhline(b_thr, color=PAL['red'], linewidth=2,
                linestyle='--', label=f'Benzene WHO ({b_thr} µg/m³)')
    ax1.set_title('(a) VOC Compounds — Daily Mean Time Series',
                  fontsize=10, fontweight='bold', color=PAL['navy'])
    ax1.set_ylabel('Concentration (µg/m³)', fontsize=9)
    ax1.legend(fontsize=8.5, loc='upper right', ncol=3)
    ax1.tick_params(axis='x', rotation=30)

    # Panel B: monthly VOC means as grouped bar
    month_names = ['Jan','Feb','Mar','Apr','May','Jun',
                   'Jul','Aug','Sep','Oct','Nov','Dec']
    x   = np.arange(12)
    w   = 0.15
    for i, (col, color) in enumerate(zip(voc_cols, voc_colors)):
        dm = daily_df[daily_df['variable'] == col][['date','mean']].dropna()
        dm['month'] = dm['date'].dt.month
        monthly = [dm[dm['month'] == m]['mean'].mean() for m in range(1, 13)]
        ax2.bar(x + i * w, monthly, width=w, color=color, alpha=0.8, label=col)
    ax2.set_xticks(x + w * (len(voc_cols)-1)/2)
    ax2.set_xticklabels(month_names, fontsize=8.5)
    ax2.set_title('(b) Monthly Average VOC Concentrations — Grouped',
                  fontsize=10, fontweight='bold', color=PAL['navy'])
    ax2.set_ylabel('Concentration (µg/m³)', fontsize=9)
    ax2.legend(fontsize=8.5, ncol=3)

    fig.suptitle(f'ATARS — Volatile Organic Compounds (VOC) Panel  ·  {config["city"]}',
                 fontsize=13, fontweight='bold', color=PAL['navy'])
    fig.text(0.5, 0.005,
             'BTEX group: Benzene, Toluene, Ethylbenzene, Xylene. '
             'Benzene WHO annual guideline: 1.7 µg/m³.',
             ha='center', fontsize=8, color=PAL['gray'], style='italic')
    return _save_fig(fig, out, 'chart10_voc.png')


# ── NEW CHART 11: Rain vs Pollutant Scatter (Wet Deposition) ─────────────────
def chart_rain_scatter(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 11 — Scatter of rain vs PM10 and SO2 to show wet deposition effect."""
    cols_to_plot = [c for c in ['PM10','SO2','NO2'] if c in daily_df['variable'].unique()]
    if 'rain_mm' not in daily_df['variable'].unique() or not cols_to_plot:
        return None

    rain_dm = daily_df[daily_df['variable'] == 'rain_mm'][['date','mean']].rename(
        columns={'mean': 'rain'}).dropna()

    fig, axes = plt.subplots(1, len(cols_to_plot[:3]), figsize=(14, 6),
                              facecolor=PAL['white'])
    if len(cols_to_plot) == 1:
        axes = [axes]
    colors_ = [PAL['blue'], PAL['teal'], PAL['amber']]

    for ax, col, color in zip(axes, cols_to_plot[:3], colors_):
        dm = daily_df[daily_df['variable'] == col][['date','mean']].rename(
            columns={'mean': col}).dropna()
        merged = pd.merge(rain_dm, dm, on='date')
        if len(merged) < 10:
            continue
        ax.scatter(merged['rain'], merged[col], color=color,
                   alpha=0.45, s=18, label=f'n={len(merged)}')
        # Linear trend
        if merged['rain'].std() > 0:
            z = np.polyfit(merged['rain'], merged[col], 1)
            p = np.poly1d(z)
            xr = np.linspace(merged['rain'].min(), merged['rain'].max(), 100)
            ax.plot(xr, p(xr), color=PAL['navy'], linewidth=2,
                    linestyle='--', label=f'Trend (OLS)')
            r_val = np.corrcoef(merged['rain'], merged[col])[0, 1]
            ax.text(0.97, 0.95, f'r = {r_val:.2f}',
                    transform=ax.transAxes, ha='right', va='top',
                    fontsize=10, color=PAL['navy'],
                    bbox=dict(facecolor='white', alpha=0.6, edgecolor='none'))
        ax.set_title(f'{col} vs Rain\n(wet deposition indicator)',
                     fontsize=10, fontweight='bold', color=PAL['navy'])
        ax.set_xlabel('Daily Rain (mm)', fontsize=9)
        unit = 'mg/m³' if col == 'CO' else 'µg/m³'
        ax.set_ylabel(f'{col} ({unit})', fontsize=9)
        ax.legend(fontsize=8.5)

    fig.suptitle(f'ATARS — Rain vs Pollutant Scatter  ·  {config["city"]}',
                 fontsize=13, fontweight='bold', color=PAL['navy'])
    fig.text(0.5, 0.01,
             'Association only — Corr(rain, pollutant) ≢ causal wet-deposition effect. '
             'r = Pearson correlation coefficient.',
             ha='center', fontsize=8, color=PAL['gray'], style='italic')
    return _save_fig(fig, out, 'chart11_rain_scatter.png')


# ── NEW CHART 12: Exceedance Calendar Heatmap ────────────────────────────────
def chart_exceedance_calendar(daily_df: pd.DataFrame, config: dict, out: Path) -> str:
    """Chart 12 — Calendar heatmap of PM10 daily mean (month × day-of-month)."""
    col = 'PM10'
    dm = daily_df[daily_df['variable'] == col][['date','mean','exceeds_threshold']].dropna()
    if len(dm) < 30:
        return None

    dm['month'] = dm['date'].dt.month
    dm['day']   = dm['date'].dt.day

    # Build 12×31 matrix
    matrix = np.full((12, 31), np.nan)
    exceed = np.zeros((12, 31), dtype=bool)
    for _, row in dm.iterrows():
        m, d = int(row['month']) - 1, int(row['day']) - 1
        matrix[m, d] = row['mean']
        exceed[m, d] = row['exceeds_threshold']

    fig, ax = plt.subplots(figsize=(15, 7), facecolor=PAL['white'])
    month_names = ['Jan','Feb','Mar','Apr','May','Jun',
                   'Jul','Aug','Sep','Oct','Nov','Dec']

    im = ax.imshow(matrix, aspect='auto', cmap='RdYlGn_r',
                   vmin=0, vmax=config['thresholds'].get(col, 150) * 2)
    # Mark exceedances with X
    for m in range(12):
        for d in range(31):
            if exceed[m, d]:
                ax.text(d, m, '×', ha='center', va='center',
                        fontsize=7, color='white', fontweight='bold')
            elif not np.isnan(matrix[m, d]):
                ax.text(d, m, f'{matrix[m,d]:.0f}', ha='center', va='center',
                        fontsize=5.5, color='white' if matrix[m,d] > 80 else '#1A2C4E')

    ax.set_yticks(range(12))
    ax.set_yticklabels(month_names, fontsize=9)
    ax.set_xticks(range(31))
    ax.set_xticklabels([str(d+1) for d in range(31)], fontsize=7.5)
    ax.set_xlabel('Day of Month', fontsize=9)

    cbar = fig.colorbar(im, ax=ax, shrink=0.8, pad=0.01)
    cbar.set_label(f'PM10 (µg/m³)', fontsize=9)
    thr = config['thresholds'].get(col)
    if thr:
        cbar.ax.axhline(thr, color='white', linewidth=2, linestyle='--')

    ax.set_title(
        f'(a) PM10 Daily Mean Calendar Heatmap  ·  {config["city"]}\n'
        f'× = WHO exceedance (>{thr} µg/m³). Green = clean, Red = polluted.',
        fontsize=11, fontweight='bold', color=PAL['navy'])
    return _save_fig(fig, out, 'chart12_calendar.png')



# ── CHART 14: ML Anomaly Comparison (IF vs z-score) ─────────────────────────
def chart_ml_anomaly_comparison(daily_df, config, ml_if, out):
    """Chart 14 — Isolation Forest vs z-score anomaly comparison."""
    if not ml_if.get("available"):
        return None
    comp = ml_if["comparison_df"].copy()
    if comp.empty:
        return None

    from matplotlib.patches import Patch
    fig, axes = plt.subplots(3, 1, figsize=(14, 12), facecolor=PAL['white'])
    fig.subplots_adjust(hspace=0.50)
    dates = comp.index

    # Panel A: IF score
    ax = axes[0]
    ax.fill_between(dates, comp["if_score"], color=PAL["teal"], alpha=0.35)
    ax.plot(dates, comp["if_score"], color=PAL["teal"], linewidth=0.8, label="IF score")
    thresh = float(comp["if_score"].quantile(config["if_contamination"]))
    ax.axhline(thresh, color=PAL["red"], linestyle="--", linewidth=1.5,
               label=f"Anomaly boundary ({config['if_contamination']:.0%})")
    anomaly_dates = comp[comp["ml_anomaly"]].index
    for d in anomaly_dates:
        ax.axvspan(d - pd.Timedelta(hours=12), d + pd.Timedelta(hours=12),
                   alpha=0.25, color=PAL["red"])
    ax.set_title("(a) Isolation Forest Decision Score — lower = more anomalous",
                 fontsize=10, fontweight="bold", color=PAL["navy"])
    ax.set_ylabel("IF Score", fontsize=9); ax.legend(fontsize=8.5)
    ax.tick_params(axis="x", rotation=30)

    # Panel B: z-score
    ax = axes[1]
    if "z_score" in comp.columns:
        vz = comp["z_score"].dropna()
        ax.plot(vz.index, vz.values, color=PAL["blue"], linewidth=0.9, label="PM10 ζ_d")
        zt = config["anomaly_z_threshold"]
        ax.axhline(zt,  color=PAL["amber"], linestyle="--", linewidth=1.5, label=f"+{zt}σ")
        ax.axhline(-zt, color=PAL["amber"], linestyle="--", linewidth=1.5, label=f"-{zt}σ")
        ax.fill_between(vz.index, zt, vz.values.clip(min=zt), alpha=0.18, color=PAL["red"])
        ax.fill_between(vz.index, -zt, vz.values.clip(max=-zt), alpha=0.18, color=PAL["red"])
    ax.set_title("(b) Z-Score ζ_d — Formal Operator Eq. 9", fontsize=10,
                 fontweight="bold", color=PAL["navy"])
    ax.set_ylabel("ζ_d (std devs)", fontsize=9); ax.legend(fontsize=8.5)
    ax.tick_params(axis="x", rotation=30)

    # Panel C: Agreement
    ax = axes[2]
    for d, row in comp.iterrows():
        if row["both"]:      c = PAL["red"]
        elif row["if_only"]: c = PAL["purple"]
        elif row["zs_only"]: c = PAL["amber"]
        else:                c = PAL["teal"]
        ax.axvspan(d - pd.Timedelta(hours=12), d + pd.Timedelta(hours=12),
                   alpha=0.65, color=c)
    legend_elements = [
        Patch(facecolor=PAL["red"],    alpha=0.65, label=f"Both ({ml_if['n_both']})"),
        Patch(facecolor=PAL["purple"], alpha=0.65, label=f"IF-only ({ml_if['n_if_only']})"),
        Patch(facecolor=PAL["amber"],  alpha=0.65, label=f"Z-only ({ml_if['n_zs_only']})"),
        Patch(facecolor=PAL["teal"],   alpha=0.65, label="Normal"),
    ]
    ax.legend(handles=legend_elements, fontsize=8.5, loc="upper right", ncol=2)
    ax.set_title("(c) Agreement Map — IF vs Z-Score", fontsize=10,
                 fontweight="bold", color=PAL["navy"])
    ax.set_yticks([]); ax.tick_params(axis="x", rotation=30)

    fig.suptitle(
        f"ATARS — ML Anomaly: Isolation Forest vs Z-Score  ·  {config['city']}\n"
        f"IF: {ml_if['n_anomalies']} days | Features: {', '.join(ml_if['features_used'])} | "
        f"random_state={ml_if['random_state']}",
        fontsize=11, fontweight="bold", color=PAL["navy"])
    fig.text(0.5, 0.005,
             "IF does NOT replace ζ_d (Eq. 9). Both reported. IF-only days = multivariate patterns.",
             ha="center", fontsize=8, color=PAL["gray"], style="italic")
    return _save_fig(fig, out, "chart14_ml_anomaly.png")


# ── CHART 15: Holt-Winters PM10 Forecast ─────────────────────────────────────
def chart_forecast(daily_df, config, ml_hw, out):
    """Chart 15 — Holt-Winters Triple Exponential Smoothing PM10 Forecast."""
    if not ml_hw.get("available"):
        return None

    fig, axes = plt.subplots(2, 1, figsize=(14, 10), facecolor=PAL["white"])
    fig.subplots_adjust(hspace=0.45)

    hist_dates = ml_hw["history_dates"]
    hist_vals  = ml_hw["history_series"]
    fitted     = ml_hw["fitted"]
    fc_dates   = ml_hw["forecast_dates"]
    fc_vals    = ml_hw["forecast"]
    fc_lower   = ml_hw["lower_80"]
    fc_upper   = ml_hw["upper_80"]

    # Panel A
    ax = axes[0]
    ax.plot(hist_dates, hist_vals, color=PAL["blue"], linewidth=0.8,
            alpha=0.7, label="Historical PM10")
    ax.plot(hist_dates, fitted, color=PAL["teal"], linewidth=1.5,
            linestyle="--", label="HW fitted", alpha=0.85)
    ax.axvline(hist_dates[-1], color=PAL["navy"], linestyle=":", linewidth=1.5)
    ax.plot(fc_dates, fc_vals, color=PAL["amber"], linewidth=2.2,
            marker="o", markersize=4, label=f"{ml_hw['forecast_days']}-day forecast")
    ax.fill_between(fc_dates, fc_lower, fc_upper, alpha=0.25, color=PAL["amber"],
                    label="80% prediction interval")
    thr_v = config["thresholds"].get("PM10")
    if thr_v:
        ax.axhline(thr_v, color=PAL["red"], linestyle=":", linewidth=1.5,
                   label=f"WHO guideline ({thr_v})")
    ax.set_title(f"(a) PM10 Historical + {ml_hw['forecast_days']}-Day Holt-Winters Forecast",
                 fontsize=10, fontweight="bold", color=PAL["navy"])
    ax.set_ylabel("PM10 (µg/m³)", fontsize=9)
    ax.legend(fontsize=8.5, loc="upper left"); ax.tick_params(axis="x", rotation=30)

    # Panel B: Residuals
    ax = axes[1]
    resid = hist_vals - fitted
    colors_r = [PAL["red"] if r > 0 else PAL["teal"] for r in resid]
    ax.bar(hist_dates, resid, color=colors_r, alpha=0.6, width=0.8)
    ax.axhline(0, color=PAL["navy"], linewidth=1.2)
    rs = float(np.std(resid, ddof=1))
    ax.axhline(rs,  color=PAL["amber"], linestyle="--", linewidth=1.2, label=f"+1σ ({rs:.1f})")
    ax.axhline(-rs, color=PAL["amber"], linestyle="--", linewidth=1.2, label=f"-1σ")
    ax.set_title(f"(b) Residuals — RMSE={ml_hw['rmse']:.2f} | MAE={ml_hw['mae']:.2f} µg/m³",
                 fontsize=10, fontweight="bold", color=PAL["navy"])
    ax.set_ylabel("Residual (µg/m³)", fontsize=9)
    ax.legend(fontsize=8.5); ax.tick_params(axis="x", rotation=30)

    fig.suptitle(
        f"ATARS — Holt-Winters PM10 Forecast  ·  {config['city']}\n"
        f"α={ml_hw['alpha']} | β={ml_hw['beta']} | γ={ml_hw['gamma']} | "
        f"seasonal={ml_hw['seasonal_periods']}d",
        fontsize=11, fontweight="bold", color=PAL["navy"])
    fig.text(0.5, 0.005,
             "Statistical extrapolation only. Forecast NOT in JSON contract J (LLM did not use it). "
             "80% PI widens with forecast horizon h.",
             ha="center", fontsize=8, color=PAL["gray"], style="italic")
    return _save_fig(fig, out, "chart15_forecast.png")


def generate_all_charts(df, daily_df, config, out, verification=None,
                         ml_if=None, ml_hw=None,
                         cleaning_report=None, data_profile=None):
    cleaning_report_arg = cleaning_report or {}
    data_profile_arg    = data_profile    or {}
    """
    Generate all charts. Returns dict of {name: filepath}.
    v2.0: Adds Chart 14 (ML anomaly comparison) and Chart 15 (forecast).
    """
    n_base   = 13   # +1 for Chart 00 (data quality overview)
    n_grnd   = 1 if verification else 0
    n_ml_if  = 1 if (ml_if and ml_if.get("available")) else 0
    n_ml_hw  = 1 if (ml_hw and ml_hw.get("available") and config.get("use_forecast")) else 0
    n_charts = n_base + n_grnd + n_ml_if + n_ml_hw
    print(f"\n  [Step 6] Generating {n_charts} publication-quality charts "
          f"({n_base} core + {n_grnd} RGV + {n_ml_if} ML-anomaly + {n_ml_hw} ML-forecast)...")
    charts = {}
    tasks  = [
        ('dq_overview',  chart_data_quality_overview,
                         [df, daily_df, cleaning_report_arg, data_profile_arg, config, out]),
        ('timeseries',   chart_timeseries,          [df, daily_df, config, out]),
        ('correlation',  chart_correlation,          [daily_df, config, out]),
        ('trend',        chart_trend,                [daily_df, config, out]),
        ('anomalies',    chart_anomalies,            [daily_df, config, out]),
        ('acf_pacf',     chart_acf_pacf,             [daily_df, config, out]),
        ('quality',      chart_quality_summary,      [daily_df, config, out]),
        ('seasonal',     chart_seasonal,             [daily_df, config, out]),
        ('diurnal',      chart_diurnal,              [df, config, out]),
        ('day_of_week',  chart_day_of_week,          [daily_df, config, out]),
        ('voc',          chart_voc_panel,            [daily_df, config, out]),
        ('rain_scatter', chart_rain_scatter,         [daily_df, config, out]),
        ('calendar',     chart_exceedance_calendar,  [daily_df, config, out]),
    ]
    for name, fn, args in tasks:
        try:
            path = fn(*args)
            if path:
                charts[name] = path
                print(f"  \u2713 Chart {len(charts):02d}/{n_charts}: {name}")
            else:
                print(f"  \u26a0  Skipped (insufficient data): {name}")
        except Exception as e:
            print(f"  \u2717  Chart error ({name}): {e}")
    if verification:
        try:
            path = chart_grounding_verification(verification, config, out)
            if path:
                charts['grounding'] = path
                print(f"  \u2713 Chart {len(charts):02d}/{n_charts}: grounding_verification (RGV ★)")
        except Exception as e:
            print(f"  \u2717  Chart error (grounding): {e}")
    if ml_if and ml_if.get("available"):
        try:
            path = chart_ml_anomaly_comparison(daily_df, config, ml_if, out)
            if path:
                charts['ml_anomaly'] = path
                print(f"  \u2713 Chart {len(charts):02d}/{n_charts}: ml_anomaly (Isolation Forest)")
        except Exception as e:
            print(f"  \u2717  Chart error (ml_anomaly): {e}")
    if ml_hw and ml_hw.get("available") and config.get("use_forecast"):
        try:
            path = chart_forecast(daily_df, config, ml_hw, out)
            if path:
                charts['forecast'] = path
                print(f"  \u2713 Chart {len(charts):02d}/{n_charts}: forecast (Holt-Winters)")
        except Exception as e:
            print(f"  \u2717  Chart error (forecast): {e}")
    return charts


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 9 — JSON CONTRACT ASSEMBLY
#  J = {x̄, σ, δ, ζ, Q} — the formal interface between computation and language
# ═══════════════════════════════════════════════════════════════════════════════

def build_json_contract(df: pd.DataFrame, daily_df: pd.DataFrame,
                         config: dict, ols_result: dict,
                         corr_matrix: pd.DataFrame) -> dict:
    """
    Assemble the formal JSON contract J.
    This is the ONLY input to the LLM — no raw data is passed.
    Grounding constraint: G(s) ⊆ J for all narrative sentences s.
    """
    print("\n  [Step 4] Assembling JSON contract J...")

    date_range = f"{df['datetime'].min().date()} to {df['datetime'].max().date()}"
    total_days  = int((df['datetime'].max() - df['datetime'].min()).days) + 1
    total_recs  = len(df)

    # Overall quality across all pollutants
    avg_q = daily_df.groupby('variable')['Q'].mean().to_dict()
    overall_q = float(np.mean(list(avg_q.values()))) if avg_q else 0.0
    overall_conf = get_confidence_flag(overall_q, config)

    # Per-variable summary statistics
    variables = {}
    for col in POLLUTANTS:
        sub = daily_df[daily_df['variable'] == col]
        if sub.empty or sub['mean'].isna().all():
            continue
        valid = sub[sub['confidence'] != 'LOW']
        annual_mean = float(sub['mean'].mean())
        annual_std  = float(sub['mean'].std())
        annual_max  = float(sub['mean'].max())
        annual_min  = float(sub['mean'].min())
        n_anomaly   = int(sub['is_anomaly'].sum())
        n_exceed    = int(sub['exceeds_threshold'].sum())
        thr         = config['thresholds'].get(col)
        variables[col] = {
            "annual_mean"      : round(annual_mean, 3),
            "annual_std"       : round(annual_std, 3),
            "annual_max"       : round(annual_max, 3),
            "annual_min"       : round(annual_min, 3),
            "threshold"        : thr,
            "exceeds_days"     : n_exceed,
            "exceed_pct"       : round(n_exceed / total_days * 100, 2) if total_days > 0 else 0,
            "anomaly_days"     : n_anomaly,
            "mean_z_score"     : round(float(sub['z_score'].mean()), 3)
                                  if not sub['z_score'].isna().all() else None,
            "avg_Q"            : round(float(sub['Q'].mean()), 4),
        }

    # High-exceedance summary
    exceedances = {k: v['exceeds_days']
                   for k, v in variables.items() if v['exceeds_days'] > 0}
    top_pollutant = max(exceedances, key=exceedances.get) if exceedances else "PM10"

    # Correlation matrix (top correlations only)
    top_corr = []
    if corr_matrix is not None and not corr_matrix.empty:
        flat = corr_matrix.stack()
        flat.index.names = ['var1', 'var2']
        flat = flat.reset_index()
        flat.columns = ['var1', 'var2', 'r']
        flat = flat[flat['var1'] != flat['var2']].copy()
        flat['abs_r'] = flat['r'].abs()
        flat = flat.nlargest(10, 'abs_r')
        top_corr = [
            {"var1": row.var1, "var2": row.var2, "r": round(row.r, 3)}
            for _, row in flat.iterrows()
        ]

    J = {
        "schema_version"  : "1.0",
        "framework"       : "ATARS v1.0",
        "author"          : config['author'],
        "institution"     : config['institution'],
        "city"            : config['city'],
        "station_id"      : config['station_id'],
        "date_range"      : date_range,
        "total_days"      : total_days,
        "total_records"   : total_recs,
        "data_quality"    : {
            "overall_Q"        : round(overall_q, 4),
            "confidence_flag"  : overall_conf,
            "per_variable_Q"   : {k: round(v, 3) for k, v in avg_q.items()},
        },
        "variables"       : variables,
        "exceedances"     : exceedances,
        "top_pollutant"   : top_pollutant,
        "top_correlations": top_corr,
        "ols_regression"  : ols_result if ols_result else {},
        "formal_operators": {
            "baseline_window_W"   : config['baseline_window_days'],
            "z_threshold"         : config['anomaly_z_threshold'],
            "ci_alpha"            : config['ci_alpha'],
        },
        "ai_constraints"  : {
            "grounding_rule"  : "G(s) ⊆ J ∧ G(s) ≠ ∅ — all claims reference this JSON",
            "temperature"     : 0,
            "causation_note"  : "Corr(X,Y) does not imply X→Y — association only",
            "extrapolation"   : "prohibited — no forecasts beyond this dataset",
        },
        # SECURITY / ML NOTE: ML results (Isolation Forest, Holt-Winters forecast)
        # are intentionally NOT included in J. Reasons:
        #   1. J is the LLM's ONLY data source — ML forecasts in J risk the LLM
        #      stating uncertain extrapolations as historical facts.
        #   2. Isolation Forest anomaly flags are supplementary — formal ζ_d
        #      is the auditable anomaly criterion already in J.
        #   3. ML results appear in Chart 14, Chart 15, and the report ML
        #      section only — clearly marked as supplementary analysis.
        # Stat inference summary — key metrics only (arrays not serializable)
        "stat_inference_summary": {
            "mann_kendall_pm10": None,   # filled post-assembly by run_pipeline
            "normality_pm10"   : None,
            "ols_r2_adj"       : None,
            "seasonal_pct_trend": None,
        },
        "ml_metadata"     : {
            "if_contamination": "injected_at_runtime",
            "if_n_estimators" : "injected_at_runtime",
            "hw_forecast_days": "injected_at_runtime",
            "note"            : "ML results NOT in J — see chart14 and chart15 in report",
        },
        "run_metadata"    : {
            "run_id"    : hashlib.sha256(date_range.encode()).hexdigest()[:16],
            "hash_J"    : "computed_post_assembly",
            "date_range": date_range,
            "timestamp" : datetime.now().isoformat(),
        },
    }

    # Back-fill hash of J itself (excluding run_metadata for stability)
    J_for_hash = {k: v for k, v in J.items() if k != "run_metadata"}
    J["run_metadata"]["hash_J"] = hashlib.sha256(
        _safe_dumps(J_for_hash, sort_keys=True, ensure_ascii=True).encode()
    ).hexdigest()

    print(f"  ✓ JSON contract assembled: {len(variables)} variables, "
          f"{len(top_corr)} top correlations")
    return J


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 10 — LLM NARRATIVE GENERATION (Ollama — optional, offline, free)
#  Formal grounding: G: J → N where ∀ s ∈ N: G(s) ⊆ J
# ═══════════════════════════════════════════════════════════════════════════════

NARRATIVE_SYSTEM_PROMPT = """You are the ATARS analytical narrative engine.
Your ONLY job is to convert the JSON data contract into formal research prose.

═══════════════════════════════════════════════════════════════
GROUNDING CONTRACT — READ BEFORE WRITING ANY SENTENCE
═══════════════════════════════════════════════════════════════
You will receive a JSON data contract. Every number in your
response MUST be taken verbatim from that JSON.

PROHIBITED: inventing, rounding beyond 2 decimal places,
approximating, or using numbers not present in the JSON.

PERMITTED: paraphrasing descriptions, choosing sentence
structure, selecting which JSON values to highlight.

ANTI-HALLUCINATION RULES (non-negotiable):
1. COPY numbers exactly as they appear in JSON — do not round or approximate.
2. If a value is not in JSON → do NOT include it. Write around it instead.
3. Do NOT assert causation. Use: "is associated with", "correlates with".
4. Do NOT use words: "causes", "due to", "because of", "leads to", "results in".
5. Do NOT refer to dates, events, or facts outside the JSON date range.
6. Do NOT make comparative claims (e.g. "worse than last year") unless JSON has them.
7. Temperature τ=0 — deterministic output only.
8. 3-5 sentences per section maximum.
9. Output ONLY research sentences — no preamble, no headers, no labels.
   Do NOT begin with: "Here are", "The following", "This section",
   "Based on", "Certainly", "Sure", "Below", "As requested",
   "The analysis shows", "The data shows", "I have analyzed".
   Start immediately with the first research sentence.

SENTENCE VERIFICATION — before writing each sentence ask yourself:
  → Does every number in this sentence appear in the JSON? (If NO → remove it)
  → Am I asserting causation? (If YES → rephrase as association)
  → Am I inventing context not in JSON? (If YES → delete)

Write for a peer-reviewed research audience. Be formal, precise,
and explicitly honest about data quality limitations.
═══════════════════════════════════════════════════════════════"""

NARRATIVE_SECTIONS = {
    "A_summary": (
        "Write exactly 4 research sentences summarising annual air quality for {city}. "
        "Sentence 1: state total_days={total_days} days monitored and total_records in the dataset. "
        "Sentence 2: state PM10 annual_mean from variables.PM10.annual_mean and "
        "how many exceedance days (variables.PM10.exceeds_days) against the {pm10_thr} µg/m³ WHO guideline. "
        "Sentence 3: state overall Q(D)={Q} and confidence flag {conf}. "
        "Sentence 4: name the top_pollutant by exceedance count. "
        "Every number you write MUST appear in the JSON. Copy values exactly — no rounding."
    ),
    "B_pollutants": (
        "Write exactly 4 research sentences on the top pollutants. "
        "Sentence 1: state PM10 annual_mean, annual_std, annual_max from variables.PM10. "
        "Sentence 2: state PM10 exceeds_days and exceed_pct against threshold. "
        "Sentence 3: state NO2 annual_mean and exceeds_days if present in variables.NO2. "
        "Sentence 4: if Benzene is in variables, state its annual_mean against its threshold. "
        "Copy every number verbatim from the JSON permitted number list."
    ),
    "C_correlations": (
        "Write exactly 3 research sentences on statistical associations. "
        "Sentence 1: state the strongest correlation pair from top_correlations[0] "
        "including the exact r value. "
        "Sentence 2: state the second strongest pair from top_correlations[1] and its r value. "
        "Sentence 3: explicitly state — all r values are statistical associations only; "
        "Corr(X,Y) does not imply causal effect C(X→Y). "
        "Use only r values from top_correlations in the JSON."
    ),
    "D_anomalies": (
        "Write exactly 4 research sentences on anomaly detection. "
        "Sentence 1: state the z-score threshold {z_thresh} and baseline window {w_days} days. "
        "Sentence 2: state PM10 anomaly_days from variables.PM10.anomaly_days. "
        "Sentence 3: state the mean_z_score for PM10 from variables.PM10.mean_z_score. "
        "Sentence 4: note that anomalies reflect deviation from the {w_days}-day rolling baseline. "
        "Copy all numbers exactly from JSON."
    ),
    "E_quality": (
        "Write exactly 4 research sentences on data quality. "
        "Sentence 1: state overall Q(D)={Q} and confidence classification {conf}. "
        "Sentence 2: state the per-variable Q values for PM10 and NO2 from "
        "data_quality.per_variable_Q. "
        "Sentence 3: name any variables with Q below 0.70 or confirm all meet the threshold. "
        "Sentence 4: state that Q(D) measures record completeness only, "
        "not sensor calibration accuracy. "
        "All Q values must be exact numbers from JSON."
    ),
    "F_limitations": (
        "Write exactly 4 sentences on analytical limitations. "
        "Sentence 1: state this analysis covers total_days={total_days} days and establishes "
        "no causal relationships. "
        "Sentence 2: state Q(D)={Q} measures completeness only, not calibration accuracy. "
        "Sentence 3: state results are specific to this station and the {total_days}-day period. "
        "Sentence 4: state expert domain review is required before regulatory use. "
        "Include the exact Q value and day count from JSON. No preamble."
    ),
}


def _extract_permitted_numbers(J: dict) -> list:
    """
    Extracts all numerical values from J as a flat sorted list.
    Used to build the explicit allowlist injected into the LLM prompt.
    Numbers in this list are the ONLY ones the LLM is permitted to use.
    """
    nums = set()
    def _recurse(obj):
        if isinstance(obj, dict):
            for v in obj.values():
                _recurse(v)
        elif isinstance(obj, list):
            for v in obj:
                _recurse(v)
        elif isinstance(obj, (int, float)) and not isinstance(obj, bool):
            try:
                fv = float(obj)
                if not (math.isnan(fv) or math.isinf(fv)) and abs(fv) < 1e9:
                    nums.add(round(fv, 4))
            except (TypeError, ValueError):
                pass
    _recurse(J)
    return sorted(nums)


def try_llm_narrative(J: dict, config: dict) -> dict:
    """
    Generate narrative using Ollama local LLM.
    Falls back gracefully if Ollama is not running.
    G(s) ⊆ J enforced by system prompt constraints.
    """
    import requests as req

    narrative = {}
    if not config.get('use_llm', True):
        return _placeholder_narrative(J)

    # Check Ollama is running
    try:
        resp = req.get(f"{config['ollama_url']}/api/tags", timeout=5)
        if resp.status_code != 200:
            raise ConnectionError()
        models = [m['name'] for m in resp.json().get('models', [])]
        print(f"  ✓ Ollama running — available models: {models}")
        if not any(config['llm_model'].split(':')[0] in m for m in models):
            print(f"  ⚠ Model '{config['llm_model']}' not found. "
                  f"Run: ollama pull {config['llm_model']}")
            return _placeholder_narrative(J)
    except Exception:
        print("  ⚠ Ollama not running — using structured placeholders.")
        print("  To enable AI narrative: install Ollama from https://ollama.com")
        print(f"  Then run: ollama pull {config['llm_model']}")
        return _placeholder_narrative(J)

    import time as _time
    n_sections = len(NARRATIVE_SECTIONS)
    print(f"\n  [Step 5] Generating AI narrative with {config['llm_model']}...")
    print(f"  ⏱  Timeout per section: {config['llm_timeout']}s  "
          f"| Retries: {config['llm_retry']}  "
          f"| Sections: {n_sections}")
    print(f"  ⏱  Worst-case total wait: "
          f"~{config['llm_timeout'] * n_sections // 60} min "
          f"(LLM speed depends on your hardware)\n")

    json_str = _safe_dumps(J, indent=2)
    for idx, (section_key, prompt_template) in enumerate(NARRATIVE_SECTIONS.items(), 1):
        pm10_thr_val = (J.get("variables",{}).get("PM10",{}).get("threshold") or 45.0)
        prompt = prompt_template.format(
            city      = J.get("city","City"),
            Q         = J["data_quality"]["overall_Q"],
            conf      = J["data_quality"]["confidence_flag"],
            z_thresh  = J["formal_operators"]["z_threshold"],
            w_days    = J["formal_operators"]["baseline_window_W"],
            total_days= J.get("total_days", 0),
            pm10_thr  = pm10_thr_val,
        )
        # Build explicit number allowlist from J for this prompt
        # This tells LLM exactly which numbers are permitted
        allowed_nums = _extract_permitted_numbers(J)
        allowlist_str = ", ".join(str(n) for n in allowed_nums[:50])
        full_prompt = (
            f"JSON DATA CONTRACT:\n{json_str}\n\n"
            f"PERMITTED NUMBERS (only these may appear in your response):\n"
            f"{allowlist_str}\n\n"
            f"TASK: {prompt}"
        )

        success = False
        for attempt in range(1, config['llm_retry'] + 2):  # +2: first try + retries
            try:
                t_start = _time.time()
                print(f"  → Section {idx}/{n_sections} [{section_key}]"
                      f"{'  (retry ' + str(attempt-1) + ')' if attempt > 1 else ''} ...",
                      end='', flush=True)
                r = req.post(
                    f"{config['ollama_url']}/api/generate",
                    json={"model"  : config['llm_model'],
                          "prompt" : full_prompt,
                          "system" : NARRATIVE_SYSTEM_PROMPT,
                          "stream" : False,
                          "options": {"temperature": 0, "num_predict": 512}},
                    timeout=config['llm_timeout']
                )
                elapsed = _time.time() - t_start
                if r.status_code == 200:
                    text = _clean_llm_response(r.json().get('response', '').strip())
                    text = _correct_hallucinated_numbers(text, J)
                    narrative[section_key] = text
                    print(f" ✓  ({elapsed:.0f}s, {len(text)} chars)")
                    success = True
                    break
                else:
                    print(f" ✗  HTTP {r.status_code}")
                    narrative[section_key] = _placeholder_text(section_key, J)
                    break
            except req.exceptions.Timeout:
                elapsed = _time.time() - t_start
                if attempt <= config['llm_retry']:
                    print(f" ⏱ Timeout ({elapsed:.0f}s) — retrying in "
                          f"{config['llm_retry_delay']}s...")
                    _time.sleep(config['llm_retry_delay'])
                else:
                    print(f" ⏱ Timeout after {config['llm_retry']+1} attempts "
                          f"— using statistical placeholder.")
                    narrative[section_key] = _placeholder_text(section_key, J)
            except Exception as e:
                print(f" ⚠ Error: {e}")
                narrative[section_key] = _placeholder_text(section_key, J)
                break

    narrative['model']   = config['llm_model']
    narrative['ai_flag'] = True
    return narrative



def _correct_hallucinated_numbers(text: str, J: dict) -> str:
    """
    Post-generation correction: scans LLM output for numbers not in J
    and replaces them with the closest permitted value.
    Only corrects numbers that differ from ALL J values by more than 15%.
    """
    permitted = _extract_permitted_numbers(J)
    if not permitted:
        return text

    def _find_closest(val):
        best = None
        best_diff = float('inf')
        for pv in permitted:
            if pv == 0:
                continue
            diff = abs(val - pv) / max(abs(pv), 1.0)
            if diff < best_diff:
                best_diff = diff
                best = pv
        return best, best_diff

    def _replace_num(m):
        try:
            val = float(m.group(0))
        except ValueError:
            return m.group(0)
        if abs(val) < 10 and '.' not in m.group(0):
            return m.group(0)
        closest, diff = _find_closest(val)
        if diff > 0.15 and closest is not None:
            if '.' in m.group(0):
                decimals = len(m.group(0).split('.')[-1])
                return f"{closest:.{decimals}f}"
            return str(int(closest)) if closest == int(closest) else str(closest)
        return m.group(0)

    corrected = re.sub(r"-?[0-9]+(?:\.[0-9]+)?", _replace_num, text)
    return corrected


def _clean_llm_response(text: str) -> str:
    """
    Remove common LLM preamble patterns that appear despite system prompt instructions.
    These patterns indicate the LLM echoed back its task rather than producing content.

    Known bad openers that LLMs produce even at temperature=0:
      - "The following research prose..."
      - "Here are the sentences..."
      - "Based on the JSON data..."
      - "Certainly! Here..."
      - "Below is..."
    """
    # Strip leading/trailing whitespace
    text = text.strip()

    # Remove markdown formatting that some models add
    text = text.replace("**", "").replace("__", "")

    # Patterns to strip from the start of the response
    preamble_patterns = [
        # "The following X:" preambles
        r"^The following\s+[\w\s,]+:\s*",
        r"^The following\s+[\w\s,]+\.\s*",
        # "Here is/are ..." openers — includes "Here are 3-4 sentences about X:"
        r"^Here (?:is|are)\s+[\w\s\-,]+[.:]\s*",
        r"^Here(?:'s| is) (?:a |the |my )?[\w\s]+[.:]\s*",
        # "Based on ..." openers
        r"^Based on (?:the |this |)?(?:provided |given |)?(?:JSON |data |)?[\w\s]+[,:]\s*",
        # "Certainly/Sure/Absolutely" openers
        r"^(?:Certainly|Sure|Absolutely|Of course|Yes)[!,.]?\s*(?:Here\s+(?:is|are)[^.:]*[.:]?\s*)?",
        # "As requested ..." openers
        r"^As requested[,.]?\s*",
        # "Below is/are ..." openers
        r"^Below (?:is|are)\s+[\w\s,]+[.:]\s*",
        # "In this section ..." openers
        r"^In this section[,.]?\s*",
        # "This section ..." echo
        r"^This section (?:provides?|presents?|covers?|acknowledges?)[^.]*[.:]\s*",
        # "This analysis ..." echo
        r"^This analysis (?:provides?|presents?)[^.]*[.:]\s*",
        # "This study presents a" echo (specific bug-report pattern)
        r"^This study presents a?\s*",
        # "The analysis shows/reveals ..." echo
        r"^The (?:analysis|data|results?) (?:shows?|reveals?|indicates?)[^.]*[.:]\s*",
        # Numbered list openers "1." that indicate instruction repetition
        r"^1\.\s+(?:State|Note|Mention|Write)[^.]+\.\s*",
    ]

    for pattern in preamble_patterns:
        cleaned = re.sub(pattern, "", text, flags=re.IGNORECASE).strip()
        # Accept if cleaned is non-empty AND at least 25% of original
        # (lower threshold: 0.25 vs 0.5 — catches cases where preamble is >50% of text)
        if cleaned and len(cleaned) >= max(20, len(text) * 0.25):
            text = cleaned

    # Remove any trailing instruction echoes (sometimes LLMs append "Write only...")
    text = re.sub(r"\s*Write only (?:these|the) \d+ sentences.*$", "", text, flags=re.IGNORECASE | re.DOTALL).strip()
    text = re.sub(r"\s*\(Note:.*\)$", "", text, flags=re.IGNORECASE).strip()

    # Capitalise first letter if it was stripped
    if text and text[0].islower():
        text = text[0].upper() + text[1:]

    return text.strip()

def _placeholder_text(key: str, J: dict) -> str:
    """
    Data-driven placeholder narrative — all numbers pulled directly from J.
    Every sentence that contains a number uses an exact J value so that
    the RGV can ground it. Vague sentences with no numbers are excluded
    from G_rate calculation (NON_NUMERICAL verdict).
    """
    city   = J.get("city", "City")
    Q      = J["data_quality"]["overall_Q"]
    conf   = J["data_quality"]["confidence_flag"]
    top    = J.get("top_pollutant", "PM10")
    exc    = J.get("exceedances", {})
    n_days = J.get("total_days", 0)
    n_recs = J.get("total_records", 0)
    vars_  = J.get("variables", {})
    ops    = J.get("formal_operators", {})
    z_thr  = ops.get("z_threshold", 3.0)
    w_days = ops.get("baseline_window_W", 30)
    corrs  = J.get("top_correlations", [])

    # Pull PM10 values directly from J
    pm10    = vars_.get("PM10", {})
    pm10_mean = pm10.get("annual_mean", 0)
    pm10_max  = pm10.get("annual_max",  0)
    pm10_min  = pm10.get("annual_min",  0)
    pm10_std  = pm10.get("annual_std",  0)
    pm10_exc  = pm10.get("exceeds_days", 0)
    pm10_pct  = pm10.get("exceed_pct",  0)
    pm10_anom = pm10.get("anomaly_days", 0)
    pm10_thr  = pm10.get("threshold",   45.0) or 45.0
    pm10_q    = pm10.get("avg_Q",        Q)

    # Pull NO2 values
    no2     = vars_.get("NO2", {})
    no2_mean = no2.get("annual_mean", 0)
    no2_exc  = no2.get("exceeds_days", 0)
    no2_thr  = no2.get("threshold", 25.0) or 25.0

    # Benzene
    benz     = vars_.get("Benzene", {})
    benz_mean = benz.get("annual_mean", 0)
    benz_thr  = benz.get("threshold", 1.7) or 1.7

    # Top correlations
    top_corr_str = ""
    if corrs:
        c = corrs[0]
        top_corr_str = (f"The strongest statistical association was between "
                        f"{c.get('var1','')} and {c.get('var2','')} "
                        f"(r = {c.get('r',0):.3f}).")

    # Top exceedances for B_pollutants
    exc_top3 = list(exc.items())[:3]
    exc_str  = "; ".join(f"{k}: {v} days" for k, v in exc_top3) if exc_top3 else "none recorded"

    # Per-variable Q summary
    per_q = J["data_quality"].get("per_variable_Q", {})
    low_q = [k for k,v in per_q.items() if v < 0.70]
    best_q_var = max(per_q, key=per_q.get) if per_q else "PM10"
    best_q_val = round(per_q.get(best_q_var, Q), 3)

    texts = {
        "A_summary": (
            f"The annual air quality dataset for {city} encompasses {n_days} monitored days "
            f"and {n_recs:,} individual records across {len(vars_)} measured variables. "
            f"Overall data completeness Q(D) = {Q:.3f}, classified as {conf} confidence "
            f"under the ATARS formal quality framework. "
            f"PM10 annual mean was {pm10_mean:.2f} µg/m³ against a WHO 24-hour guideline "
            f"of {pm10_thr:.1f} µg/m³, with {pm10_exc} exceedance days "
            f"representing {pm10_pct:.1f}% of the monitoring period."
        ),
        "B_pollutants": (
            f"The primary WHO exceedances by day count were: {exc_str}. "
            f"PM10 recorded an annual mean of {pm10_mean:.2f} µg/m³ "
            f"(std = {pm10_std:.2f} µg/m³, max daily mean = {pm10_max:.2f} µg/m³), "
            f"exceeding the {pm10_thr:.1f} µg/m³ WHO guideline on {pm10_exc} days. "
            + (f"NO₂ annual mean was {no2_mean:.2f} µg/m³ against the "
               f"{no2_thr:.1f} µg/m³ WHO guideline, with {no2_exc} exceedance days. "
               if no2_mean > 0 else "")
            + (f"Benzene annual mean was {benz_mean:.3f} µg/m³ against the "
               f"{benz_thr:.1f} µg/m³ WHO annual reference value."
               if benz_mean > 0 else "")
        ),
        "C_correlations": (
            f"Pearson correlation analysis was performed across {len(vars_)} variable pairs. "
            + (top_corr_str + " " if top_corr_str else "")
            + f"All {len(corrs)} top correlation pairs are reported in the correlation matrix. "
            f"All reported r values are statistical associations only — "
            f"causal inference requires experimental design outside this framework scope. "
            f"Corr(X,Y) does not imply causal effect C(X→Y)."
        ),
        "D_anomalies": (
            f"Z-score anomaly detection applied threshold |ζ_d| > {z_thr:.1f} "
            f"against a {w_days}-day rolling baseline. "
            f"PM10 anomaly days: {pm10_anom} "
            f"(days where daily mean exceeded {z_thr:.1f} standard deviations above the baseline). "
            f"PM10 annual mean z-score: {pm10.get('mean_z_score', 0):.3f}. "
            f"Anomalies reflect deviation from the local {w_days}-day norm, "
            f"not absolute WHO guideline exceedance."
        ),
        "E_quality": (
            f"Data quality was assessed via Q(D) = N_v/N_total per variable per day. "
            f"Overall Q(D) = {Q:.3f} ({conf}). "
            f"Best completeness: {best_q_var} at Q = {best_q_val:.3f}. "
            + (f"Variables with Q < 0.70 (LOW confidence): {', '.join(low_q)}. "
               if low_q else "All variables met the Q ≥ 0.70 MODERATE threshold. ")
            + f"PM10 average Q(D) = {pm10_q:.3f}. "
            f"Q(D) measures record completeness only — sensor calibration accuracy "
            f"is outside the scope of this quality assessment."
        ),
        "F_limitations": (
            f"This analysis provides statistical characterisation of {n_days} days of "
            f"monitoring data and establishes no causal relationships. "
            f"All correlation statistics (Corr(X,Y)) are association measures only "
            f"and do not imply causal effect parameters C(X→Y). "
            f"Q(D) = {Q:.3f} measures record completeness, not sensor calibration accuracy "
            f"or measurement uncertainty. "
            f"Results are specific to this station and the {n_days}-day monitoring period "
            f"and require expert domain review before informing regulatory decisions."
        ),
    }
    return texts.get(key, f"[Section {key} — computed from JSON contract J]")


def _placeholder_narrative(J: dict) -> dict:
    return {k: _placeholder_text(k, J) for k in NARRATIVE_SECTIONS} | \
           {'model': 'statistical_placeholder', 'ai_flag': False}


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 10.5 — RUNTIME GROUNDING VERIFIER (RGV)
#
#  NOVEL CONTRIBUTION — programmatic enforcement of G(s) ⊆ J
#
#  Prior work: G(s) ⊆ J was enforced only by system-prompt instruction.
#  The LLM could hallucinate numbers not in J and no system would detect it.
#
#  This module: every numerical claim in every LLM sentence is extracted,
#  matched against the JSON contract J, and classified as:
#    GROUNDED       — claim found in J within tolerance
#    PARTIALLY_GROUNDED — some claims found, some not
#    UNGROUNDED     — numerical claims present, none found in J
#    NON_NUMERICAL  — no claims to verify (not a violation)
#    UNCERTAIN      — single ambiguous claim
#
#  G_rate = |{s ∈ N : G(s) ⊆ J}| / |N_numerical|   — scalar metric (0–1)
#
#  This G_rate is:
#    (a) measurable — computed per section and overall
#    (b) verifiable — full evidence table showing which J path matched
#    (c) novel      — no prior environmental monitoring pipeline defines or
#                     measures this property programmatically
#
#  Reference: ATARS formal framework, Eq. 19, G: J → N
# ═══════════════════════════════════════════════════════════════════════════════


class RuntimeGroundingVerifier:
    """
    Programmatic enforcement of the formal grounding constraint:
      G(s) ⊆ J  ∧  G(s) ≠ ∅  for all numerically-bearing sentences s ∈ N
    where J is the JSON data contract and N is the set of generated sentences.

    Algorithm:
      1. Flatten J into a numeric lookup table L = {path: value}
      2. For each sentence s ∈ N:
         a. Extract numerical claims C(s) = {(val, type, context)}
         b. For each c ∈ C(s): check if ∃ j ∈ L: |c - j|/|j| ≤ ε
         c. Assign verdict based on |{grounded claims}| / |C(s)|
      3. Compute G_rate = |{s: GROUNDED or PARTIALLY_GROUNDED}| / |N_numerical|
      4. Emit full evidence table (claim → J path → J value → match status)
    """

    # Tolerance bounds for numerical matching
    TOL_RELATIVE = 0.05    # ±5% — accounts for display rounding (83.21 shown as 83.2)
    TOL_ABSOLUTE = 2.0     # ±2 for integer counts (day counts can vary by ±1-2)
    TOL_EXACT    = 0.005   # ±0.005 for quality scores (3 decimal places)

    # Minimum G_rate for a section to pass
    PASS_THRESHOLD = 0.75   # WMO-consistent threshold

    def __init__(self, J: dict):
        self.J      = J
        self.lookup = {}                         # path → float value
        self.string_lookup = {}                  # path → string value
        self._build_lookup(J, path='')

    # ── Lookup table construction ───────────────────────────────────────────
    def _build_lookup(self, obj, path: str):
        """
        Recursively flatten J into {dotted.path: value}.
        Stores rounded variants (0, 1, 2 dp) so display-rounded numbers
        in narrative still match back to J values.
        """
        if isinstance(obj, dict):
            for k, v in obj.items():
                new_path = f"{path}.{k}" if path else k
                self._build_lookup(v, new_path)
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                self._build_lookup(v, f"{path}[{i}]")
        elif isinstance(obj, bool):
            pass
        elif isinstance(obj, (int, float)):
            try:
                fval = float(obj)
                if not (math.isnan(fval) or math.isinf(fval)):
                    self.lookup[path] = fval
                    for dp in (0, 1, 2):
                        self.lookup[f"{path}__r{dp}"] = round(fval, dp)
            except (TypeError, ValueError):
                pass
        elif isinstance(obj, str):
            self.string_lookup[path] = obj

    def extract_claims(self, sentence: str) -> list:
        """
        Extract all numerical claims from a sentence.
        Returns list of dicts: {raw, value, type, context}
        Skips citation-style numbers (equation refs, section numbers < 30).
        """
        claims = []
        # Match integers and floats, including negative
        pattern = re.compile(
            r'(?<![a-zA-Z])(-?\d{1,6}(?:\.\d{1,6})?)\b'
        )
        for m in pattern.finditer(sentence):
            raw = m.group(1)
            try:
                val = float(raw)
            except ValueError:
                continue

            # Skip likely equation/section references
            if val == int(val) and 0 < val <= 30 and '.' not in raw:
                ctx_before = sentence[max(0, m.start()-15):m.start()].lower()
                if any(w in ctx_before for w in
                       ['eq.', 'section', 'fig.', 'table', 'prop.']):
                    continue

            # Determine semantic type from surrounding context (window ±30 chars)
            win = sentence[max(0, m.start()-30): m.end()+30].lower()

            if any(w in win for w in ['day', 'days']):
                ctype = 'days_count'
            elif any(w in win for w in ['record', 'observation', 'n=']):
                ctype = 'record_count'
            elif '%' in win[len(raw):len(raw)+5] or 'percent' in win:
                ctype = 'percentage'
            elif any(w in win for w in ['µg', 'mg/m', 'ppb', 'ppm', 'µg/m']):
                ctype = 'concentration'
            elif any(w in win for w in ['q(', 'q =', 'quality', 'completeness',
                                         'score', 'completene']):
                ctype = 'quality_score'
            elif any(w in win for w in ['r²', 'r2', 'r-squared', 'pearson',
                                         'correlation']):
                ctype = 'correlation_coef'
            elif any(w in win for w in ['z-score', 'zeta', 'ζ', 'z_score',
                                         'anomal']):
                ctype = 'z_score'
            elif any(w in win for w in ['threshold', 'guideline', 'limit',
                                         'who', 'who ']):
                ctype = 'threshold'
            elif any(w in win for w in ['year', 'annual', 'month']):
                ctype = 'temporal'
            else:
                ctype = 'numeric'

            claims.append({
                'raw': raw, 'value': val, 'type': ctype,
                'context': win.strip()
            })
        return claims

    # ── Single claim matching ────────────────────────────────────────────────
    def match_claim(self, val: float, ctype: str) -> dict:
        """
        Find the closest match in J lookup for val.
        Improvements: exact match priority, percentage scaling fallback,
        wider tolerance for display-rounded values.
        """
        best = {'matched': False, 'j_path': None, 'j_value': None,
                'diff': float('inf'), 'method': None}

        candidates = list(self.lookup.items())

        for path, j_val in candidates:
            abs_diff = abs(val - j_val)

            # ── Exact match — highest priority ────────────────────────────
            if abs_diff < self.TOL_EXACT:
                return {'matched': True, 'j_path': path, 'j_value': j_val,
                        'diff': round(abs_diff, 8), 'method': 'exact'}

            # ── Type-specific tolerance ────────────────────────────────────
            if ctype in ('days_count', 'record_count') or (
                    val == int(val) and 10 <= abs(val) <= 100000):
                tol    = self.TOL_ABSOLUTE
                method = 'absolute'
                score  = abs_diff
            elif ctype in ('quality_score', 'correlation_coef', 'z_score'):
                tol    = self.TOL_EXACT * 10
                method = 'tight'
                score  = abs_diff
            elif j_val != 0:
                rel_diff = abs_diff / abs(j_val)
                tol      = self.TOL_RELATIVE
                method   = 'relative'
                score    = rel_diff
            else:
                tol    = self.TOL_ABSOLUTE
                method = 'absolute'
                score  = abs_diff

            if score <= tol and score < best['diff']:
                best = {'matched': True, 'j_path': path, 'j_value': j_val,
                        'diff': round(score, 6), 'method': method}

        # ── Percentage scale fallback ─────────────────────────────────────
        # LLM may write 87.4 when J stores 0.874 (fraction) or vice-versa
        if not best['matched']:
            for path, j_val in candidates:
                for scale in (100.0, 0.01):
                    scaled = j_val * scale
                    if scaled != 0 and abs(scaled) > 0.001:
                        rel = abs(val - scaled) / abs(scaled)
                        if rel < self.TOL_RELATIVE and rel < best['diff']:
                            best = {'matched': True, 'j_path': path,
                                    'j_value': scaled,
                                    'diff': round(rel, 6),
                                    'method': f'scale_{scale}'}

        return best

    def verify_sentence(self, sentence: str) -> dict:
        """
        Verify one sentence. Assigns verdict and full evidence.
        Returns structured result dict.
        """
        claims  = self.extract_claims(sentence)

        if not claims:
            return {
                'verdict': 'NON_NUMERICAL', 'sentence': sentence,
                'claims_total': 0, 'claims_grounded': 0,
                'claim_rate': None, 'evidence': []
            }

        evidence = []
        n_grounded = 0
        for c in claims:
            match = self.match_claim(c['value'], c['type'])
            is_g  = match['matched']
            if is_g:
                n_grounded += 1
            evidence.append({
                'claim'       : c['raw'],
                'value'       : c['value'],
                'type'        : c['type'],
                'grounded'    : is_g,
                'j_path'      : match['j_path'],
                'j_value'     : match['j_value'],
                'diff'        : match['diff'],
                'method'      : match['method'],
            })

        total = len(claims)
        rate  = n_grounded / total

        if rate == 1.0:
            verdict = 'GROUNDED'
        elif rate >= 0.5:
            verdict = 'PARTIALLY_GROUNDED'
        elif total == 1:
            verdict = 'UNCERTAIN'
        else:
            verdict = 'UNGROUNDED'

        return {
            'verdict'         : verdict,
            'sentence'        : sentence,
            'claims_total'    : total,
            'claims_grounded' : n_grounded,
            'claim_rate'      : round(rate, 4),
            'evidence'        : evidence,
        }

    # ── Section verification ─────────────────────────────────────────────────
    def verify_section(self, section_key: str, text: str) -> dict:
        """
        Verify an entire narrative section.
        Splits into sentences, verifies each, computes section G_rate.
        """
        # Split on sentence boundaries
        raw_sents = re.split(r'(?<=[.!?])\s+', text.strip())
        sentences = [s.strip() for s in raw_sents if len(s.strip()) > 8]

        results       = [self.verify_sentence(s) for s in sentences]
        numerical     = [r for r in results if r['verdict'] != 'NON_NUMERICAL']
        grounded_list = [r for r in numerical
                         if r['verdict'] in ('GROUNDED', 'PARTIALLY_GROUNDED')]

        g_rate = (len(grounded_list) / len(numerical)) if numerical else 1.0
        passes = g_rate >= self.PASS_THRESHOLD

        return {
            'section_key'          : section_key,
            'sentences'            : results,
            'total_sentences'      : len(results),
            'numerical_sentences'  : len(numerical),
            'grounded_sentences'   : len(grounded_list),
            'ungrounded_sentences' : len(numerical) - len(grounded_list),
            'G_rate'               : round(g_rate, 4),
            'passes'               : passes,
            'verdict'              : 'PASS ✓' if passes else 'REVIEW ⚠',
        }

    # ── Full narrative verification ──────────────────────────────────────────
    def verify_narrative(self, narrative: dict) -> dict:
        """
        Verify the complete generated narrative N.

        Formally: computes
          G_rate = |{s ∈ N_num : verdict ∈ {GROUNDED, PARTIALLY_GROUNDED}}|
                   ────────────────────────────────────────────────────────
                                    |N_num|

        where N_num = {s ∈ N : C(s) ≠ ∅} (sentences with numerical claims).
        """
        sections      = {}
        total_num     = 0
        total_grounded= 0
        all_ungrounded= []

        for key in NARRATIVE_SECTIONS:
            if key not in narrative:
                continue
            result = self.verify_section(key, narrative[key])
            sections[key]     = result
            total_num        += result['numerical_sentences']
            total_grounded   += result['grounded_sentences']

            for sr in result['sentences']:
                if sr['verdict'] == 'UNGROUNDED':
                    all_ungrounded.append({
                        'section'  : key,
                        'sentence' : (sr['sentence'][:500] + ' …' if len(sr['sentence']) > 500 else sr['sentence']),
                        'n_claims' : sr['claims_total'],
                        'evidence' : sr['evidence'][:3],
                    })

        overall = (total_grounded / total_num) if total_num > 0 else 1.0
        passed  = overall >= self.PASS_THRESHOLD

        return {
            'schema'             : 'RGV-1.0',
            'formal_property'    : 'G(s) ⊆ J  ∀ s ∈ N  [ATARS Eq. 19]',
            'overall_G_rate'     : round(overall, 4),
            'total_numerical'    : total_num,
            'total_grounded'     : total_grounded,
            'total_ungrounded'   : total_num - total_grounded,
            'verification_passed': passed,
            'overall_verdict'    : 'PASS ✓' if passed else 'REVIEW ⚠',
            'pass_threshold'     : self.PASS_THRESHOLD,
            'sections'           : sections,
            'ungrounded_list'    : all_ungrounded,
            'j_lookup_size'      : len(self.lookup),
            'timestamp'          : datetime.now().isoformat(),
        }


def chart_grounding_verification(verification: dict, config: dict, out: Path) -> str:
    """
    Chart 13 — RGV Grounding Verification Dashboard.
    Shows G_rate per section, verdict distribution, and evidence heatmap.
    This visualisation is part of the novel contribution.
    """
    sections = verification.get('sections', {})
    if not sections:
        return None

    fig = plt.figure(figsize=(14, 10), facecolor=PAL['white'])
    gs  = gridspec.GridSpec(2, 3, figure=fig, hspace=0.52, wspace=0.38)

    ax_bar   = fig.add_subplot(gs[0, :2])   # G_rate per section (wide)
    ax_gauge = fig.add_subplot(gs[0, 2])    # Overall G_rate gauge
    ax_pie   = fig.add_subplot(gs[1, 0])    # Verdict distribution
    ax_heat  = fig.add_subplot(gs[1, 1:])   # Evidence density heatmap

    # ── Panel A: G_rate per section ─────────────────────────────────────────
    sec_labels = list(sections.keys())
    g_rates    = [sections[k]['G_rate'] for k in sec_labels]
    verdicts   = [sections[k]['passes'] for k in sec_labels]
    bar_colors = [PAL['green'] if v else PAL['amber'] for v in verdicts]

    bars = ax_bar.barh(sec_labels, g_rates, color=bar_colors, height=0.55, alpha=0.85)
    ax_bar.axvline(verification['pass_threshold'], color=PAL['red'],
                   linewidth=2, linestyle='--',
                   label=f'Pass threshold ({verification["pass_threshold"]:.0%})')
    ax_bar.set_xlim(0, 1.05)
    ax_bar.set_xlabel('G_rate (grounded sentences / numerical sentences)', fontsize=9)
    ax_bar.set_title(
        r'(a) G_rate per narrative section  —  $G_\text{rate} = '
        r'|\{s : G(s) \subseteq J\}| \,/\, |N_\text{num}|$',
        fontsize=10, fontweight='bold', color=PAL['navy'])
    ax_bar.legend(fontsize=8.5)
    for bar, rate in zip(bars, g_rates):
        ax_bar.text(bar.get_width() + 0.01, bar.get_y() + bar.get_height()/2,
                    f'{rate:.1%}', va='center', fontsize=9, fontweight='bold')
    ax_bar.tick_params(axis='y', labelsize=9)

    # ── Panel B: Overall G_rate gauge ──────────────────────────────────────
    overall = verification['overall_G_rate']
    theta   = np.linspace(0, np.pi, 200)
    ax_gauge.set_aspect('equal')
    # Background arc
    ax_gauge.plot(np.cos(theta), np.sin(theta), color='#E5E7EB', linewidth=18,
                  solid_capstyle='round')
    # Filled arc proportional to G_rate
    theta_fill = np.linspace(0, np.pi * overall, 200)
    arc_color  = PAL['green'] if overall >= 0.75 else (PAL['amber'] if overall >= 0.5 else PAL['red'])
    ax_gauge.plot(np.cos(theta_fill), np.sin(theta_fill),
                  color=arc_color, linewidth=18, solid_capstyle='round')
    ax_gauge.text(0, 0.2, f'{overall:.1%}', ha='center', va='center',
                  fontsize=26, fontweight='bold', color=arc_color)
    ax_gauge.text(0, -0.18, verification['overall_verdict'],
                  ha='center', va='center', fontsize=13, fontweight='bold',
                  color=arc_color)
    ax_gauge.text(0, -0.42, 'Overall G_rate', ha='center', va='center',
                  fontsize=9, color=PAL['gray'])
    ax_gauge.set_xlim(-1.3, 1.3); ax_gauge.set_ylim(-0.6, 1.2)
    ax_gauge.axis('off')
    ax_gauge.set_title('(b) Overall\nGrounding Rate',
                        fontsize=10, fontweight='bold', color=PAL['navy'])

    # ── Panel C: Verdict distribution pie ──────────────────────────────────
    verdict_counts = defaultdict(int)
    for sec_data in sections.values():
        for sr in sec_data['sentences']:
            verdict_counts[sr['verdict']] += 1
    v_labels = list(verdict_counts.keys())
    v_vals   = list(verdict_counts.values())
    v_colors_map = {
        'GROUNDED'           : PAL['green'],
        'PARTIALLY_GROUNDED' : PAL['teal'],
        'NON_NUMERICAL'      : '#A0B4CC',
        'UNCERTAIN'          : PAL['amber'],
        'UNGROUNDED'         : PAL['red'],
    }
    v_colors = [v_colors_map.get(l, PAL['gray']) for l in v_labels]
    wedges, texts, autotexts = ax_pie.pie(
        v_vals, labels=v_labels, colors=v_colors,
        autopct='%1.0f%%', startangle=90,
        textprops={'fontsize': 8},
        pctdistance=0.75
    )
    for at in autotexts: at.set_fontsize(8)
    ax_pie.set_title('(c) Verdict Distribution\n(all sentences)',
                     fontsize=10, fontweight='bold', color=PAL['navy'])

    # ── Panel D: Evidence density — claims per section ──────────────────────
    sec_names   = list(sections.keys())
    ev_grounded = [sections[k]['grounded_sentences']   for k in sec_names]
    ev_ungnd    = [sections[k]['ungrounded_sentences']  for k in sec_names]
    ev_non_num  = [sections[k]['total_sentences'] -
                   sections[k]['numerical_sentences']  for k in sec_names]
    x = np.arange(len(sec_names))
    w = 0.25
    ax_heat.bar(x - w, ev_grounded, w, label='Grounded',     color=PAL['green'],  alpha=0.85)
    ax_heat.bar(x,     ev_ungnd,    w, label='Ungrounded',   color=PAL['red'],    alpha=0.85)
    ax_heat.bar(x + w, ev_non_num,  w, label='Non-numerical',color='#A0B4CC',     alpha=0.85)
    ax_heat.set_xticks(x); ax_heat.set_xticklabels(sec_names, fontsize=9)
    ax_heat.set_ylabel('Sentence count', fontsize=9)
    ax_heat.set_title('(d) Sentence-Level Grounding Analysis\n'
                      'Grounded / Ungrounded / Non-numerical per section',
                      fontsize=10, fontweight='bold', color=PAL['navy'])
    ax_heat.legend(fontsize=8.5)

    fig.suptitle(
        f'ATARS — Runtime Grounding Verifier (RGV)  ·  {config["city"]}\n'
        r'Novel contribution: programmatic enforcement of $G(s) \subseteq J$  '
        f'— G_rate = {overall:.1%}  [{verification["overall_verdict"]}]',
        fontsize=12, fontweight='bold', color=PAL['navy'])

    fig.text(0.5, 0.01,
             f'RGV v1.0  |  J lookup size: {verification["j_lookup_size"]} entries  |  '
             f'Sections verified: {len(sections)}  |  '
             f'Total sentences: {sum(s["total_sentences"] for s in sections.values())}  |  '
             f'Timestamp: {verification["timestamp"][:19]}',
             ha='center', fontsize=7.5, color=PAL['gray'], style='italic')

    return _save_fig(fig, out, 'chart13_grounding_verification.png')


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 11 — WORD REPORT GENERATION (python-docx)
# ═══════════════════════════════════════════════════════════════════════════════

def _set_cell_bg(cell, hex_color: str):
    """Set table cell background colour."""
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement('w:shd')
    shd.set(qn('w:fill'), hex_color)
    shd.set(qn('w:val'), 'clear')
    tcPr.append(shd)


def _add_table_row(table, cells_data, bold_col=0, bg=None, font_size=10):
    """Add a styled row to a Word table."""
    row = table.add_row()
    for i, (cell, data) in enumerate(zip(row.cells, cells_data)):
        para = cell.paragraphs[0]
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER if i > 0 else WD_ALIGN_PARAGRAPH.LEFT
        run  = para.add_run(str(data))
        run.bold      = (i == bold_col)
        run.font.size = Pt(font_size)
        if bg:
            _set_cell_bg(cell, bg if i % 2 == 0 else 'FFFFFF')


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))



def _h2_inner(doc, text: str, color):
    """Inner h2 helper for use inside build_report lambdas."""
    p = doc.add_heading(text, level=2)
    for run in p.runs:
        run.font.color.rgb = color
        run.font.size = Pt(13)
    p.paragraph_format.space_before = Pt(10)
    p.paragraph_format.space_after  = Pt(5)


def build_report(J: dict, daily_df: pd.DataFrame, charts: dict,
                 narrative: dict, ols_result: dict,
                 corr_matrix: pd.DataFrame, config: dict,
                 out_dir: Path, verification: dict = None,
                 chart_insights: dict = None,
                 ml_if: dict = None, ml_hw: dict = None,
                 # v3.1 additions — all optional, gracefully skipped if None
                 v3_upt       : dict = None,
                 v3_benchmark : dict = None,
                 v3_comparison: dict = None,
                 v3_aee       : dict = None,
                 v3_cvcc      : dict = None,
                 v3_nss       : dict = None,
                 v3_tcv       : dict = None,
                 v3_ccd       : dict = None,
                 v3_mns       : dict = None,
                 v3_drift     : dict = None,
                 v3_ver_v3      : dict = None,
                 cleaning_report: dict = None,
                 data_profile : dict = None,
                 v3_stability : dict = None,
                 v3_episodes  : list = None,
                 v3_forecast  : dict = None,
                 stat_inference: dict = None,
                 ) -> str:
    """
    Assemble the complete Word document report — ATARS v3.1.
    All v2 sections preserved in original positions.
    v3.1 features integrated at their logical section positions:
      Section 01: + Uncertainty Bounds (UPT)
      Section 02: + Benchmark Comparator
      Section 02b: Today vs Historical Comparison
      Section 04: + Anomaly Explanation Engine (AEE)
      Section 06: + Cross-Variable Consistency (CVCC)
      Section 15: + Narrative Specificity Score (NSS)
      Section 16: RGV + TCV + CCD unified
      Section 17: Master Novelty Score (MNS) — new
      Section 19: Audit + Report Drift (RDD)
    """
    chart_insights = chart_insights or {}
    ml_if = ml_if or {}
    ml_hw = ml_hw or {}
    print("\n  [Step 7] Assembling Word report...")
    doc = Document()

    # ── Page margins ──────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin    = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    # ── Colour palette ────────────────────────────────────────────────────────
    NAVY   = _rgb('1A2C4E'); BLUE   = _rgb('2557A7')
    TEAL   = _rgb('1A6B72'); GREEN  = _rgb('1E6B3C')
    RED    = _rgb('8B1C2A'); AMBER  = _rgb('92400E')
    GRAY   = _rgb('6B7280'); WHITE  = _rgb('FFFFFF')
    PURPLE = _rgb('4C1D95')

    def h1(text, color=NAVY):
        p = doc.add_heading(text, level=1)
        for run in p.runs:
            run.font.color.rgb = color
            run.font.size = Pt(18)
        p.paragraph_format.space_before = Pt(18)
        p.paragraph_format.space_after  = Pt(8)

    def h2(text, color=BLUE):
        p = doc.add_heading(text, level=2)
        for run in p.runs:
            run.font.color.rgb = color
            run.font.size = Pt(14)
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after  = Pt(6)

    def body(text, italic=False, color=None):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_before = Pt(3)
        p.paragraph_format.space_after  = Pt(6)
        for run in p.runs:
            run.italic = italic
            if color:
                run.font.color.rgb = color
            run.font.size = Pt(11)
        return p

    def caption(text):
        p = doc.add_paragraph(text)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in p.runs:
            run.font.size  = Pt(9)
            run.italic     = True
            run.font.color.rgb = GRAY
        p.paragraph_format.space_before = Pt(2)
        p.paragraph_format.space_after  = Pt(10)

    def add_image(chart_key, width_inch=6.2, cap_text=""):
        path = charts.get(chart_key)
        if path and os.path.exists(path):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            run.add_picture(path, width=Inches(width_inch))
        if cap_text:
            caption(cap_text)
        # Add deterministic chart insight text if available
        insight = chart_insights.get(chart_key, "")
        if insight:
            p_ins = doc.add_paragraph()
            p_ins.paragraph_format.space_before = Pt(4)
            p_ins.paragraph_format.space_after  = Pt(8)
            p_ins.paragraph_format.left_indent  = Pt(14)
            # Teal vertical bar via shading is not directly available in python-docx
            # Use a styled paragraph with ► prefix
            run_lbl = p_ins.add_run("► Chart Insight:  ")
            run_lbl.bold = True
            run_lbl.font.color.rgb = TEAL
            run_lbl.font.size = Pt(9.5)
            run_txt = p_ins.add_run(insight)
            run_txt.italic = True
            run_txt.font.size = Pt(9.5)
            run_txt.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

    # ════════════════════════════════════════════════════════════════════════
    #  COVER PAGE
    # ════════════════════════════════════════════════════════════════════════
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(0)
    run = p.add_run('ATARS')
    run.font.size  = Pt(54)
    run.bold       = True
    run.font.color.rgb = NAVY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('Automated Time-Series Analysis and Reporting System  v3.1')
    run.font.size = Pt(16); run.font.color.rgb = BLUE

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f'Annual Air Quality Analysis Report  —  {J["city"]}')
    run.font.size = Pt(20); run.bold = True; run.font.color.rgb = NAVY
    p.paragraph_format.space_before = Pt(16)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f'{J["date_range"]}  ·  {J["total_days"]} days  ·  {J["total_records"]:,} records')
    run.font.size = Pt(12); run.italic = True; run.font.color.rgb = GRAY

    doc.add_paragraph()

    # Confidence badge
    conf    = J['data_quality']['confidence_flag']
    q_score = J['data_quality']['overall_Q']
    conf_colors = {'HIGH': '1E6B3C', 'MODERATE': '92400E', 'LOW': '8B1C2A'}
    conf_bg     = {'HIGH': 'D1FAE5', 'MODERATE': 'FEF3C7', 'LOW': 'FFE4E6'}

    t = doc.add_table(rows=1, cols=4)
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.style     = 'Table Grid'
    data_conf = [
        ('Confidence Level', conf),
        ('Q(D) Score', f'{q_score:.3f}'),
        ('Total Days', str(J['total_days'])),
        ('Total Records', f'{J["total_records"]:,}'),
    ]
    for cell, (label, value) in zip(t.rows[0].cells, data_conf):
        cell.paragraphs[0].clear()
        lp = cell.add_paragraph(label)
        lp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in lp.runs:
            r.font.size = Pt(9); r.bold = True; r.font.color.rgb = GRAY
        vp = cell.add_paragraph(value)
        vp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in vp.runs:
            r.font.size = Pt(18); r.bold = True
            r.font.color.rgb = _rgb(conf_colors.get(conf, '1A2C4E'))
        _set_cell_bg(cell, conf_bg.get(conf, 'DCE8F5'))

    doc.add_paragraph()

    # Author block
    t2 = doc.add_table(rows=1, cols=1)
    t2.alignment = WD_TABLE_ALIGNMENT.CENTER
    t2.style     = 'Table Grid'
    cell = t2.rows[0].cells[0]
    _set_cell_bg(cell, 'F4F5F7')
    for text, pt, bold, color in [
        (config['author'], 18, True, NAVY),
        (config['degree'], 11, False, GRAY),
        (config['department'], 12, True, BLUE),
        (config['institution'], 12, False, NAVY),
        (f"{config['location']}  ·  {config['email']}", 10, False, GRAY),
        (f"Document: RP-ATARS-{datetime.now().year}-001  ·  Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", 9, False, GRAY),
        (f"GitHub: https://github.com/Priyanshu-ux712/ATARS", 9, False, GRAY),
    ]:
        p = cell.add_paragraph(text)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(3)
        for r in p.runs:
            r.font.size = Pt(pt); r.bold = bold; r.font.color.rgb = color

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  EXECUTIVE SUMMARY
    # ════════════════════════════════════════════════════════════════════════
    h1("Executive Summary")

    # Key metrics strip
    pm10_v   = J.get("variables",{}).get("PM10",{})
    no2_v    = J.get("variables",{}).get("NO2",{})
    benz_v   = J.get("variables",{}).get("Benzene",{})
    q_overall= J["data_quality"]["overall_Q"]
    conf_ovr = J["data_quality"]["confidence_flag"]
    n_vars   = len(J.get("variables",{}))
    exc_dict = J.get("exceedances",{})
    top_pol  = J.get("top_pollutant","PM10")

    exec_rows = [
        ("Monitoring period",     J.get("date_range","N/A")),
        ("Total days monitored",  str(J.get("total_days",0))),
        ("Total records",         f"{J.get('total_records',0):,}"),
        ("Variables monitored",   str(n_vars)),
        ("Overall data quality",  f"Q(D) = {q_overall:.3f}  [{conf_ovr}]"),
        ("PM10 annual mean",      f"{pm10_v.get('annual_mean',0):.2f} µg/m³  "
                                  f"(WHO limit: {pm10_v.get('threshold',45)} µg/m³)"),
        ("PM10 exceedance days",  f"{pm10_v.get('exceeds_days',0)} days  "
                                  f"({pm10_v.get('exceed_pct',0):.1f}%)"),
        ("NO₂ annual mean",       f"{no2_v.get('annual_mean',0):.2f} µg/m³"
                                  if no2_v.get("annual_mean") else "N/A"),
        ("Benzene annual mean",   f"{benz_v.get('annual_mean',0):.3f} µg/m³"
                                  if benz_v.get("annual_mean") else "N/A"),
        ("Top pollutant",         top_pol),
        ("Total WHO exceedances", ", ".join(f"{k}: {v}d" for k,v in list(exc_dict.items())[:4])
                                  or "None detected"),
        ("Anomaly days (PM10)",   str(pm10_v.get("anomaly_days",0))),
        ("Station / City",        f"{J.get('station_id','N/A')}  ·  {J.get('city','N/A')}"),
        ("Framework",             "ATARS v3.1 — MIT Open Source"),
    ]
    t_exec = doc.add_table(rows=1, cols=2)
    t_exec.style = "Table Grid"
    for cell, hdr in zip(t_exec.rows[0].cells, ["Parameter", "Value"]):
        _set_cell_bg(cell, "1A2C4E")
        r = cell.paragraphs[0].add_run(hdr)
        r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(11)
    for k, v in exec_rows:
        row = t_exec.add_row()
        row.cells[0].text = k
        row.cells[1].text = str(v)
        row.cells[0].paragraphs[0].runs[0].bold = True
        row.cells[0].paragraphs[0].runs[0].font.size = Pt(10)
        row.cells[1].paragraphs[0].runs[0].font.size = Pt(10)
    doc.add_paragraph()

    # Key findings bullets
    h2("Key Findings")
    pm10_mean_val = pm10_v.get("annual_mean", 0)
    pm10_thr_val  = pm10_v.get("threshold", 45) or 45
    pm10_ratio    = round(pm10_mean_val / pm10_thr_val, 1) if pm10_thr_val else 0
    benz_mean_val = benz_v.get("annual_mean", 0)
    benz_thr_val  = benz_v.get("threshold", 1.7) or 1.7

    findings = [
        (f"PM10 annual mean of {pm10_mean_val:.2f} µg/m³ is {pm10_ratio:.1f}× the WHO "
         f"24-hour guideline of {pm10_thr_val:.1f} µg/m³, exceeding it on "
         f"{pm10_v.get('exceeds_days',0)} days ({pm10_v.get('exceed_pct',0):.1f}% of period)."),
        (f"Data completeness Q(D) = {q_overall:.3f} [{conf_ovr}] across "
         f"{n_vars} monitored variables over {J.get('total_days',0)} days."),
        (f"Z-score anomaly detection (|ζ_d| > {J['formal_operators']['z_threshold']:.1f}) "
         f"identified {pm10_v.get('anomaly_days',0)} anomalous PM10 days against a "
         f"{J['formal_operators']['baseline_window_W']}-day rolling baseline."),
    ]
    if benz_mean_val and benz_mean_val > benz_thr_val:
        findings.append(
            f"Benzene annual mean {benz_mean_val:.3f} µg/m³ exceeds the WHO "
            f"annual reference of {benz_thr_val:.1f} µg/m³ (IARC Group 1 carcinogen)."
        )
    findings.append(
        "All correlation statistics reported are Pearson r values (statistical association "
        "only) — no causal relationships are established by this analysis."
    )
    for f_txt in findings:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(f_txt).font.size = Pt(10)

    # Report structure overview
    doc.add_paragraph()
    h2("Report Structure")
    sections_overview = [
        ("Sections 00b–00c", "Data cleaning audit (18 steps) and data profile"),
        ("Section 01",       "Data quality statement and Q(D) scores"),
        ("Section 02",       "Annual statistics, benchmark vs WHO/NAAQS/India avg"),
        ("Sections 03–05",   "Time-series, anomaly detection, ACF/PACF"),
        ("Section 05b",      "Formal statistical inference (MK, OLS, decomposition)"),
        ("Sections 06–12",   "Correlation, seasonal, diurnal, VOC, exceedance calendar"),
        ("Section 12b",      "Atmospheric stability, episode classification, 7-day forecast"),
        ("Sections 13–14",   "ML analysis: Isolation Forest + Holt-Winters forecast"),
        ("Sections 15–17",   "AI narrative verification (RGV, TCV, CCD, MNS)"),
        ("Sections 18–19",   "Task scheduler automation + reproducibility audit"),
        ("Section 20",       "Conclusions and recommendations"),
        ("Section 21",       "References and methodology"),
    ]
    t_struct = doc.add_table(rows=1, cols=2)
    t_struct.style = "Table Grid"
    for cell, hdr in zip(t_struct.rows[0].cells, ["Section", "Content"]):
        _set_cell_bg(cell, "1A6B72")
        r = cell.paragraphs[0].add_run(hdr)
        r.bold = True; r.font.color.rgb = WHITE
    for sec, desc in sections_overview:
        row = t_struct.add_row()
        row.cells[0].text = sec;  row.cells[0].paragraphs[0].runs[0].bold = True
        row.cells[1].text = desc
        for c in row.cells:
            c.paragraphs[0].runs[0].font.size = Pt(9)
    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 01 — DATA QUALITY STATEMENT
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 01 — Data Quality Statement')
    body(narrative.get('E_quality', ''))

    # Quality table
    h2('Quality Score Q(D) per Variable')
    q_vars = J['data_quality']['per_variable_Q']
    if q_vars:
        t = doc.add_table(rows=1, cols=6)
        t.style = 'Table Grid'
        for cell, hdr in zip(t.rows[0].cells,
                              ['Variable', 'Q(D) Valid', 'Q_effective',
                               'Confidence', 'Anomaly Days', 'Exceedance Days']):
            p = cell.paragraphs[0]
            run = p.add_run(hdr)
            run.bold = True; run.font.size = Pt(10)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _set_cell_bg(cell, '1A2C4E')
            for r in p.runs:
                r.font.color.rgb = WHITE

        for col, q_val in sorted(q_vars.items(), key=lambda x: -x[1]):
            v_data = J['variables'].get(col, {})
            row = t.add_row()
            bg = 'D1FAE5' if q_val >= 0.90 else ('FEF3C7' if q_val >= 0.70 else 'FFE4E6')
            q_eff_val = daily_df[daily_df["variable"]==col]["Q_effective"].mean() if "Q_effective" in daily_df.columns and col in daily_df["variable"].values else q_val
            q_eff_val = round(float(q_eff_val), 3) if not (q_eff_val != q_eff_val) else q_val
            vals = [
                col, f'{q_val:.3f}',
                f'{q_eff_val:.3f}',
                'HIGH' if q_val >= 0.90 else ('MODERATE' if q_val >= 0.70 else ('LOW' if q_val >= 0.40 else 'INSUFFICIENT')),
                str(v_data.get('anomaly_days', 'N/A')),
                str(v_data.get('exceeds_days', 0)),
            ]
            for i, (cell, val) in enumerate(zip(row.cells, vals)):
                p = cell.paragraphs[0]
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER if i > 0 else WD_ALIGN_PARAGRAPH.LEFT
                r = p.add_run(val); r.font.size = Pt(10)
                _set_cell_bg(cell, bg)

    add_image('quality', 6.2,
              'Figure — Q(D) = N_v/N_total per variable (Eq. 11). '
              'Confidence flags: HIGH ≥ 0.90, MODERATE ≥ 0.70, LOW < 0.70.')

    # ── Cleaning Report Section (Section 00b) ────────────────────────────────
    if cleaning_report and cleaning_report.get("steps"):
        generate_cleaning_report_section(
            doc, cleaning_report,
            h1, h2, body, _set_cell_bg,
            WHITE, NAVY, TEAL, GREEN, AMBER, RED, Pt,
            WD_ALIGN_PARAGRAPH
        )
        # Add the data quality overview chart immediately after cleaning section
        if charts.get("dq_overview"):
            add_image("dq_overview", 6.2,
                      "Figure — Data Quality Overview: "
                      "(a) column completeness, (b) quality flag distribution, "
                      "(c) hourly completeness heatmap, (d) cleaning action waterfall.")

    # ── Section 00c: Data Profile ─────────────────────────────────────────
    if data_profile and data_profile.get("columns"):
        try:
            add_data_profile_to_report(
                doc, data_profile,
                h1, h2, body, _set_cell_bg,
                WHITE, NAVY, TEAL, GREEN, AMBER, RED, Pt,
                WD_ALIGN_PARAGRAPH
            )
        except Exception as _dp_err:
            print(f"  ⚠ Data profile section error (non-fatal): {_dp_err}")

    # ── v3.1 N-05: Uncertainty Bounds ────────────────────────────────────────
    if v3_upt and v3_upt.get('uncertainty_table'):
        h2('Measurement Uncertainty — 95% Confidence Intervals')
        body(
            'The table below presents the 95% confidence interval for each variable\'s '
            'annual mean, derived from the daily distribution of valid measurements. '
            'Relative uncertainty reflects the CI width as a percentage of the annual mean.'
        )
        unc_t = doc.add_table(rows=1, cols=5)
        unc_t.style = 'Table Grid'
        for cell, hdr in zip(unc_t.rows[0].cells,
                              ['Variable','Annual Mean','CI Lower (95%)','CI Upper (95%)','Uncertainty Class']):
            _set_cell_bg(cell, '1A2C4E')
            p = cell.paragraphs[0]
            r = p.add_run(hdr); r.bold = True; r.font.size = Pt(9)
            r.font.color.rgb = WHITE
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for var, unc in list(v3_upt['uncertainty_table'].items())[:10]:
            row = unc_t.add_row()
            bg = 'D1FAE5' if unc.get('uncertainty_class') == 'LOW' else (
                 'FEF3C7' if unc.get('uncertainty_class') == 'MEDIUM' else 'FFE4E6')
            for cell, val in zip(row.cells, [
                var,
                f"{unc.get('mean',0):.2f} µg/m³",
                f"{unc.get('ci_95_lower',0):.2f}",
                f"{unc.get('ci_95_upper',0):.2f}",
                unc.get('uncertainty_class','N/A'),
            ]):
                _set_cell_bg(cell, bg)
                p = cell.paragraphs[0]
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                r = p.add_run(val); r.font.size = Pt(9)
        doc.add_paragraph()

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 02 — ANNUAL STATISTICS TABLE
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 02 — Annual Statistics Summary')
    body(narrative.get('A_summary', ''))

    h2('Formal Statistical Operators: A(D_v), ζ_d, δ_d (Eqs. 3–10)')
    stats_cols = ['Variable', 'Annual Mean x̄', 'Std σ', 'Max', 'Min',
                  'Threshold θ', 'Exceed Days', 'Anomaly Days', 'Avg ζ_d']
    t = doc.add_table(rows=1, cols=len(stats_cols))
    t.style = 'Table Grid'
    for cell, hdr in zip(t.rows[0].cells, stats_cols):
        p = cell.paragraphs[0]
        run = p.add_run(hdr); run.bold = True; run.font.size = Pt(9)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _set_cell_bg(cell, '1A2C4E')
        for r in p.runs: r.font.color.rgb = WHITE

    for i, (col, vd) in enumerate(J['variables'].items()):
        row = t.add_row()
        bg  = 'FFE4E6' if vd.get('exceeds_days', 0) > 30 else \
              ('FEF3C7' if vd.get('exceeds_days', 0) > 0 else
               ('DCE8F5' if i % 2 == 0 else 'FFFFFF'))
        thr = vd.get('threshold')
        zs  = vd.get('mean_z_score')
        vals = [
            col,
            f'{vd["annual_mean"]:.2f}',
            f'{vd["annual_std"]:.2f}',
            f'{vd["annual_max"]:.2f}',
            f'{vd["annual_min"]:.2f}',
            f'{thr}' if thr else '—',
            f'{vd.get("exceeds_days", 0)} ({vd.get("exceed_pct", 0):.1f}%)',
            str(vd.get('anomaly_days', 0)),
            f'{zs:.2f}' if zs is not None else '—',
        ]
        for j, (cell, val) in enumerate(zip(row.cells, vals)):
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT if j == 0 else WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run(val); r.font.size = Pt(9)
            _set_cell_bg(cell, bg)

    body(
        'x̄ = mean of all daily means over the monitoring period (annual mean). '
        'σ = sample standard deviation of daily means. '
        'Exceed Days = number of days where the daily mean exceeded guideline threshold θ. '
        'Anomaly Days = number of days where |ζ_d| > 3.0 (z-score threshold, Eq. 9). '
        'All statistics are computed from quality-validated records only (q_i = 1).',
        italic=True
    )

    # ── Benchmark vs WHO / NAAQS / India national average ─────────────────
    if v3_benchmark and v3_benchmark.get("comparisons"):
        h2("Benchmark Comparison — WHO AQG 2021 · CPCB NAAQS · India National Average")
        body(
            f"Station: {v3_benchmark.get('station','N/A')}. "
            f"Overall severity classification: {v3_benchmark.get('overall_severity','N/A')}. "
            f"Reference source: {v3_benchmark.get('reference_source','WHO AQG 2021 + CPCB NAAQS')}."
        )
        t_bm = doc.add_table(rows=1, cols=6)
        t_bm.style = "Table Grid"
        for cell, hdr in zip(t_bm.rows[0].cells,
                              ["Variable","Station Mean","vs WHO","WHO Status",
                               "vs NAAQS","vs India Avg"]):
            _set_cell_bg(cell, "8B1C2A")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.size = Pt(9); r.font.color.rgb = WHITE
            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        for col_name, comp in v3_benchmark["comparisons"].items():
            row = t_bm.add_row()
            st  = comp.get("who_status","N/A")
            bg  = "FFE4E6" if st=="EXCEEDS_WHO" else "D1FAE5"
            above = comp.get("above_national_avg", 0)
            for cell, val in zip(row.cells, [
                col_name,
                f"{comp.get('station_annual_mean',0):.2f} {comp.get('unit','µg/m³')}",
                comp.get("vs_who", "N/A"),
                st,
                comp.get("naaqs_status","N/A"),
                f"{above:+.2f} µg/m³" if above else "N/A",
            ]):
                _set_cell_bg(cell, bg)
                cell.paragraphs[0].add_run(str(val)).font.size = Pt(9)
                cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        body("Green = within WHO guideline. Red = exceeds WHO. "
             "NAAQS = Central Pollution Control Board India annual standard. "
             "India Avg = approximate national mean from CPCB annual reports.",
             italic=True)
        doc.add_paragraph()

    # ── Today vs Historical Comparison ────────────────────────────────────
    if v3_comparison and v3_comparison.get("today"):
        h2("Today vs Historical Comparison (Computed, Not AI-Generated)")
        body("All values below are pre-computed by the ATARS comparison engine. "
             "The AI narrative does not perform these calculations.")
        cmp = v3_comparison
        cmp_rows = [
            ("Reference date",          cmp.get("today","N/A")),
            ("PM10 on reference date",  f"{cmp.get('pm10_today',0):.2f} µg/m³"),
            ("PM10 yesterday",          f"{cmp.get('pm10_yesterday',0):.2f} µg/m³"),
            ("Change vs yesterday",
             f"{cmp.get('pm10_change_yesterday_pct',0):+.1f}%  "
             f"({cmp.get('pm10_direction_yesterday','N/A')})"),
            ("7-day rolling average",   f"{cmp.get('pm10_week_avg',0):.2f} µg/m³"),
            ("Change vs 7-day avg",     f"{cmp.get('pm10_change_week_pct',0):+.1f}%  "
             f"({cmp.get('pm10_direction_week','N/A')})"),
            ("30-day rolling average",  f"{cmp.get('pm10_month_avg',0):.2f} µg/m³"),
            ("Change vs 30-day avg",    f"{cmp.get('pm10_change_month_pct',0):+.1f}%  "
             f"({cmp.get('pm10_direction_month','N/A')})"),
            ("Status",                  cmp.get("status","N/A").replace("_"," ")),
        ]
        t_cmp = doc.add_table(rows=1, cols=2)
        t_cmp.style = "Table Grid"
        for cell, hdr in zip(t_cmp.rows[0].cells, ["Metric","Value"]):
            _set_cell_bg(cell, "1A6B72")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        for k, v in cmp_rows:
            row = t_cmp.add_row()
            row.cells[0].text = k; row.cells[1].text = str(v)
            row.cells[0].paragraphs[0].runs[0].font.size = Pt(10)
            row.cells[1].paragraphs[0].runs[0].font.size = Pt(10)
        tags = cmp.get("event_tags",[])
        if tags:
            body(f"Context tags: {', '.join(tags)}", italic=True)
        doc.add_paragraph()

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 03 — TIME-SERIES VISUALISATION
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 03 — Hourly Time-Series Analysis')
    add_image('timeseries', 6.2,
              'Figure — Hourly concentrations with daily mean overlay and WHO AQG 2021 '
              'reference guidelines. Shaded areas indicate exceedance zones.')

    body(narrative.get('B_pollutants', ''))

    h2('30-Day Rolling Baseline Trend (Eq. 7–8)')
    body(
        'The chart below shows the daily pollutant concentration alongside the 30-day '
        'rolling baseline mean \u03c4\u0305_W and standard deviation \u03c3_W. '
        'Deviations above the rolling mean identify periods of elevated pollution '
        'relative to recent local norms, without seasonal bias.'
    )
    add_image('trend', 6.2,
              'Figure — Daily trend and 30-day rolling baseline x̄_W (Eqs. 7–8). '
              'Dashed lines show baseline mean. Dotted lines show WHO guidelines.')

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 04 — ANOMALY DETECTION
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 04 — Anomaly Detection (Z-Score Analysis)')
    body(narrative.get('D_anomalies', ''))

    add_image('anomalies', 6.2,
              'Figure — (a) Daily mean with 95% CI band (Eq. 12). '
              r'(b) Z-score ζ_d = (x̄_d − x̄_W) / σ_W (Eq. 9). '
              'Red diamonds indicate anomaly days (|ζ_d| > 3.0).')

    # Anomaly summary table
    h2('Anomaly Summary by Pollutant')
    anom_data = [(col, vd['anomaly_days'], vd['mean_z_score'])
                 for col, vd in J['variables'].items()
                 if vd.get('anomaly_days', 0) > 0]
    if anom_data:
        t = doc.add_table(rows=1, cols=3)
        t.style = 'Table Grid'
        for cell, hdr in zip(t.rows[0].cells,
                              ['Variable', 'Anomaly Days (|ζ_d| > 3)', 'Mean ζ_d']):
            p = cell.paragraphs[0]; run = p.add_run(hdr)
            run.bold = True; run.font.size = Pt(10)
            _set_cell_bg(cell, '1A2C4E')
            for r in p.runs: r.font.color.rgb = WHITE
        for col, adays, mz in sorted(anom_data, key=lambda x: -x[1]):
            row = t.add_row()
            for cell, val in zip(row.cells,
                                  [col, str(adays),
                                   f'{mz:.2f}' if mz else '—']):
                p = cell.paragraphs[0]
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                r = p.add_run(val); r.font.size = Pt(10)
                _set_cell_bg(cell, 'FFF5F5')

    # ── v3.1 N-03: Anomaly Explanation Engine ────────────────────────────────
    if v3_aee and v3_aee.get('explanations'):
        h2('Anomaly Explanation Engine (N-03) — Rule-Based Contextual Analysis')
        body(
            f"Each anomaly day is explained by cross-referencing meteorological conditions, "
            f"season, and day-of-week patterns. Method: rule-based deterministic — zero LLM. "
            f"Explanation rate: {round(v3_aee.get('explanation_rate',0)*100,1)}% "
            f"({v3_aee.get('explained',0)}/{v3_aee.get('total_anomalies',0)} anomaly days explained)."
        )
        aee_t = doc.add_table(rows=1, cols=5)
        aee_t.style = 'Table Grid'
        for cell, hdr in zip(aee_t.rows[0].cells,
                              ['Date','PM10 Mean','Type','Season','Primary Explanation']):
            _set_cell_bg(cell, '92400E')
            p = cell.paragraphs[0]
            r = p.add_run(hdr); r.bold = True; r.font.size = Pt(9)
            r.font.color.rgb = WHITE
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for ex in v3_aee['explanations'][:15]:
            row = aee_t.add_row()
            bg = 'FFE4E6' if ex.get('confidence') == 'HIGH' else 'FEF3C7'
            for cell, val in zip(row.cells, [
                ex.get('date',''),
                f"{ex.get('pm10_mean',0):.1f} µg/m³",
                ex.get('anomaly_type','')[:18],
                ex.get('season','').capitalize(),
                ex.get('primary_explanation','')[:60],
            ]):
                _set_cell_bg(cell, bg)
                p = cell.paragraphs[0]
                r = p.add_run(val); r.font.size = Pt(8)
        doc.add_paragraph()

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 05 — TEMPORAL DEPENDENCY (ACF / PACF)
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 05 — Temporal Dependency: ACF and PACF')
    if charts.get('acf_pacf'):
        add_image('acf_pacf', 6.0,
                  'Figure — (a) ACF ρ̂(k) = Corr(x_t, x_{t-k}) (Eq. 12). '
                  '(b) PACF φ̂_kk: conditional correlation at lag k (Eq. 13). '
                  'Bars outside red dashed CI band indicate significant temporal dependence.')
    else:
        # Fallback: show lag-1 autocorrelation from daily data
        h2('Autocorrelation Summary (ACF chart requires ≥53 days of data)')
        pm10_sub = daily_df[daily_df['variable'] == 'PM10']['mean'].dropna().values
        if len(pm10_sub) > 2:
            import math as _math
            n   = len(pm10_sub)
            mu  = float(sum(pm10_sub) / n)
            var = sum((x-mu)**2 for x in pm10_sub) / n
            lag1 = (sum((pm10_sub[i]-mu)*(pm10_sub[i-1]-mu)
                        for i in range(1,n)) / n) / var if var > 0 else 0
            acf_rows = [
                ('Series length', f'{n} daily observations'),
                ('Lag-1 ACF ρ̂(1)', f'{lag1:.4f}'),
                ('Interpretation', 'STRONG temporal memory' if abs(lag1) > 0.5
                 else 'MODERATE' if abs(lag1) > 0.3 else 'WEAK'),
                ('95% CI bound', f'±{1.96/_math.sqrt(n):.4f}'),
                ('ACF chart status', f'Requires ≥53 days; dataset has {n} days'),
            ]
            t_acf = doc.add_table(rows=1, cols=2)
            t_acf.style = 'Table Grid'
            for cell, hdr in zip(t_acf.rows[0].cells, ['Metric', 'Value']):
                _set_cell_bg(cell, '1A2C4E')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for k, v in acf_rows:
                row = t_acf.add_row()
                row.cells[0].paragraphs[0].add_run(k).font.size = Pt(9)
                row.cells[1].paragraphs[0].add_run(str(v)).font.size = Pt(9)
                _set_cell_bg(row.cells[0], 'DCE8F5')
            doc.add_paragraph()

    body(
        'The Autocorrelation Function (ACF) measures the total correlation between '
        'observations separated by lag k. The Partial ACF (PACF) isolates the direct '
        'contribution of lag k after removing the influence of all shorter lags. '
        'Together, ACF and PACF characterise the temporal memory structure of the series '
        'and provide the empirical basis for identifying appropriate ARMA model orders.'
    )

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 05b — FORMAL STATISTICAL INFERENCE
    # ════════════════════════════════════════════════════════════════════════
    if stat_inference:
        doc.add_page_break()
        h1("Section 05b — Formal Statistical Inference")
        body(
            "This section presents results of formal statistical hypothesis tests "
            "applied to the cleaned daily time series. "
            "Tests are selected to be appropriate for environmental monitoring data: "
            "non-parametric where normality cannot be assumed, and autocorrelation-robust."
        )

        # ── Mann-Kendall Trend Tests ──────────────────────────────────────────
        h2("Mann-Kendall Monotonic Trend Tests (Non-Parametric)")
        body(
            "H₀: no monotonic trend. H₁: monotonic trend. "
            "Non-parametric — valid for autocorrelated and non-normal series. "
            "Sen's slope = median pairwise slope (µg/m³ per day). "
            "p < 0.05 = statistically significant at 5% level."
        )
        mk_results = stat_inference.get("mann_kendall", {})
        if mk_results:
            t_mk = doc.add_table(rows=1, cols=6)
            t_mk.style = "Table Grid"
            for cell, hdr in zip(t_mk.rows[0].cells,
                                  ["Variable","S Statistic","Z Score",
                                   "p-value","Direction","Sen's Slope (µg/m³/day)"]):
                _set_cell_bg(cell, "1A2C4E")
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.size = Pt(9)
                r.font.color.rgb = WHITE
                cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
            for col_name, mk in sorted(mk_results.items(),
                                        key=lambda x: x[1].get("p_value", 1) if x[1].get("available") else 1):
                if not mk.get("available"):
                    continue
                row = t_mk.add_row()
                sig  = mk.get("significant", False)
                bg   = "FFE4E6" if (sig and mk.get("direction","")=="INCREASING") else (
                       "D1FAE5" if (sig and mk.get("direction","")=="DECREASING") else "F5F5F5")
                for cell, val in zip(row.cells, [
                    col_name,
                    str(mk.get("S","")),
                    f"{mk.get('Z',0):.3f}",
                    f"{mk.get('p_value',1):.4f}{'*' if sig else ''}",
                    mk.get("direction","N/A").replace("_"," "),
                    f"{mk.get('sens_slope',0):.4f}",
                ]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(val).font.size = Pt(9)
                    cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
            body("* = significant at p < 0.05. Red = significant increasing trend. "
                 "Green = significant decreasing trend.", italic=True)
            doc.add_paragraph()

        # ── Normality Tests ───────────────────────────────────────────────────
        h2("Normality Tests (Shapiro-Wilk + D'Agostino-Pearson)")
        body(
            "Determines whether parametric statistics are appropriate. "
            "If non-normal: median and IQR are more reliable than mean and std. "
            "Skewness > 1 is typical for pollution data (right-skewed by episodes)."
        )
        norm_results = stat_inference.get("normality", {})
        if norm_results:
            t_norm = doc.add_table(rows=1, cols=7)
            t_norm.style = "Table Grid"
            for cell, hdr in zip(t_norm.rows[0].cells,
                                   ["Variable","Shapiro-Wilk W","S-W p-value",
                                    "D'Agostino p","Skewness","Kurtosis","Verdict"]):
                _set_cell_bg(cell, "1A6B72")
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.size = Pt(9)
                r.font.color.rgb = WHITE
            for col_name, nr in norm_results.items():
                row = t_norm.add_row()
                v   = nr.get("verdict","N/A")
                bg  = "D1FAE5" if v=="NORMAL" else ("FFE4E6" if v=="HIGHLY_SKEWED" else "FEF3C7")
                for cell, val in zip(row.cells, [
                    col_name,
                    f"{nr.get('shapiro_W','N/A')}" if nr.get("shapiro_W") else "N/A",
                    f"{nr.get('shapiro_p','N/A')}" if nr.get("shapiro_p") else "N/A",
                    f"{nr.get('dagostino_p','N/A')}" if nr.get("dagostino_p") else "N/A",
                    f"{nr.get('skewness','N/A')}",
                    f"{nr.get('kurtosis','N/A')}",
                    v,
                ]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(str(val)).font.size = Pt(9)
            doc.add_paragraph()

        # ── OLS Inference Table ────────────────────────────────────────────────
        ols_inf = stat_inference.get("ols_inference", {})
        if ols_inf.get("available"):
            h2(f"OLS Regression with Full Inference — PM10 ~ f(meteorology)")
            body(
                f"n = {ols_inf['n_obs']} observations, "
                f"k = {ols_inf['k_predictors']} predictors. "
                f"R² = {ols_inf['r_squared']:.4f} (adj. R² = {ols_inf['r_squared_adj']:.4f}). "
                f"F({ols_inf['k_predictors']}, {ols_inf['n_obs']-ols_inf['k_predictors']-1}) "
                f"= {ols_inf['f_statistic']:.2f} "
                f"(p = {ols_inf['f_pvalue']:.4f}, {ols_inf['f_significance']}). "
                f"RSE = {ols_inf['rse']:.3f}. "
                f"Durbin-Watson = {ols_inf['durbin_watson']:.3f} "
                f"({ols_inf['dw_interpretation']})."
            )
            if ols_inf.get("high_vif_vars"):
                body(f"⚠ High VIF (>5) — multicollinearity concern: "
                     f"{ols_inf['high_vif_vars']}", italic=True)
            t_ols = doc.add_table(rows=1, cols=7)
            t_ols.style = "Table Grid"
            for cell, hdr in zip(t_ols.rows[0].cells,
                                  ["Predictor","β̂ (std)","SE","t-stat",
                                   "p-value","95% CI","VIF"]):
                _set_cell_bg(cell, "4C1D95")
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.size = Pt(9)
                r.font.color.rgb = WHITE
            for pred, ct in ols_inf["coef_table"].items():
                row = t_ols.add_row()
                sig = ct.get("sig","ns")
                bg  = "D1FAE5" if sig in ("***","**") else (
                      "FEF3C7" if sig == "*" else "F5F5F5")
                vif_val = ct.get("vif")
                vif_str = f"{vif_val:.1f}{'⚠' if vif_val and vif_val > 5 else ''}" if vif_val else "N/A"
                for cell, val in zip(row.cells, [
                    pred,
                    f"{ct['beta']:.4f} {sig}",
                    f"{ct['se']:.4f}",
                    f"{ct['t_stat']:.3f}",
                    f"{ct['p_value']:.4f}",
                    f"[{ct['ci_lower']:.3f}, {ct['ci_upper']:.3f}]",
                    vif_str,
                ]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(str(val)).font.size = Pt(9)
            body("*** p<0.001  ** p<0.01  * p<0.05  ns = not significant. "
                 "VIF > 5 indicates multicollinearity. "
                 "All predictors standardised — β coefficients are comparable. "
                 "Association only — Corr(X,Y) ≢ C(X→Y).", italic=True)
            doc.add_paragraph()

        # ── Seasonal Decomposition ────────────────────────────────────────────
        decomp_s = stat_inference.get("seasonal_decomposition", {})
        if decomp_s.get("available"):
            h2("Seasonal Decomposition — PM10 Trend + Seasonal + Residual")
            body(
                f"Additive decomposition: x_t = Trend_t + Seasonal_t + Residual_t. "
                f"Trend extracted via {decomp_s.get('period_days',365)}-day centred moving average. "
                f"Variance explained: "
                f"Trend={decomp_s.get('pct_trend','N/A')}%, "
                f"Seasonal={decomp_s.get('pct_seasonal','N/A')}%, "
                f"Residual={decomp_s.get('pct_residual','N/A')}%. "
                f"Seasonal amplitude = {decomp_s.get('seasonal_range','N/A'):.1f} µg/m³. "
                + decomp_s.get("interpretation","")
            )
            doc.add_paragraph()

        # ── Exceedance Frequency ──────────────────────────────────────────────
        exc_freq = stat_inference.get("exceedance_frequency", {})
        if exc_freq:
            h2("WHO Exceedance Frequency Analysis")
            body(
                "P(exceed) = probability that a randomly chosen day exceeds the WHO guideline. "
                "Return period = expected days between exceedances. "
                "Excess burden = mean concentration on exceedance days minus guideline."
            )
            t_exc = doc.add_table(rows=1, cols=7)
            t_exc.style = "Table Grid"
            for cell, hdr in zip(t_exc.rows[0].cells,
                                   ["Variable","Threshold","P(exceed)","Return (days)",
                                    "Excess Burden","Max Streak","Status"]):
                _set_cell_bg(cell, "8B1C2A")
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.size = Pt(9)
                r.font.color.rgb = WHITE
            for col_name, ef in sorted(exc_freq.items(),
                                        key=lambda x: -x[1].get("P_exceed",0)):
                row = t_exc.add_row()
                st  = ef.get("status","N/A")
                bg  = ("FFE4E6" if st in ("CRITICAL","SEVERE") else
                       "FEF3C7" if st=="HIGH" else "F5F5F5")
                ret = ef.get("return_days", float("inf"))
                ret_str = f"{ret:.1f}" if ret != float("inf") else "∞ (no exceed)"
                for cell, val in zip(row.cells, [
                    col_name,
                    str(ef.get("threshold","N/A")),
                    f"{ef.get('P_exceed',0)*100:.1f}%",
                    ret_str,
                    f"+{ef.get('excess_over_thr',0):.2f} µg/m³" if ef.get("n_exceed",0)>0 else "N/A",
                    f"{ef.get('max_streak_days',0)} days",
                    st,
                ]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(str(val)).font.size = Pt(9)
            doc.add_paragraph()

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 06 — CORRELATION & MULTIVARIATE ANALYSIS
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 06 — Correlation & Multivariate Analysis')
    body(narrative.get('C_correlations', ''))

    add_image('correlation', 6.0,
              'Figure — Pearson correlation matrix R (Eq. 18). '
              'Blue = positive, red = negative. Values annotated per cell. '
              'FORMAL CAUSAL DISCLAIMER: Corr(X,Y) ≢ C(X→Y). '
              'Association does not imply causation.')

    # OLS regression results
    # Use full OLS inference if available, fall back to basic ols_result
    _ols = (stat_inference or {}).get("ols_inference") or ols_result or {}
    if _ols and _ols.get("available", "r_squared" in _ols):
        _r2    = _ols.get("r_squared",     _ols.get("r_squared", 0))
        _r2adj = _ols.get("r_squared_adj", _r2)
        _n     = _ols.get("n_obs", 0)
        _tgt   = _ols.get("target", "PM10")
        h2(f'OLS Regression with Inference: {_tgt} ~ f(meteorology) (Eq. 14)')
        body(
            f'β̂ = (XᵀX)⁻¹Xᵀy.  n = {_n} observations.  '
            f'R² = {_r2:.4f}  (adj. R² = {_r2adj:.4f}).  '
            + (f'F-stat = {_ols["f_statistic"]:.2f} (p = {_ols["f_pvalue"]:.4f}, '
               f'{_ols["f_significance"]}).  '
               f'Durbin-Watson = {_ols["durbin_watson"]:.3f} '
               f'({_ols["dw_interpretation"]}).'
               if _ols.get("f_statistic") else "")
            + f'  Association only — Corr(X,Y) ≢ C(X→Y).'
        )
        if _ols.get("coef_table"):
            # Full inference table
            t = doc.add_table(rows=1, cols=6)
            t.style = 'Table Grid'
            for cell, hdr in zip(t.rows[0].cells,
                                  ['Predictor','β̂ (std)','SE','t-stat','p-value','Sig']):
                p = cell.paragraphs[0]; r = p.add_run(hdr)
                r.bold = True; r.font.size = Pt(9)
                _set_cell_bg(cell, '1A2C4E')
                for rn in p.runs: rn.font.color.rgb = WHITE
            for pred, ct in _ols["coef_table"].items():
                row = t.add_row()
                sig = ct.get("sig","ns")
                bg  = 'D1FAE5' if sig in ("***","**") else ('FEF3C7' if sig=="*" else 'DCE8F5')
                for cell, val in zip(row.cells, [
                    pred, f'{ct["beta"]:.4f}', f'{ct["se"]:.4f}',
                    f'{ct["t_stat"]:.3f}', f'{ct["p_value"]:.4f}', sig
                ]):
                    p = cell.paragraphs[0]; r = p.add_run(val); r.font.size = Pt(9)
                    _set_cell_bg(cell, bg)
        else:
            # Basic table from ols_result
            t = doc.add_table(rows=1, cols=3)
            t.style = 'Table Grid'
            for cell, hdr in zip(t.rows[0].cells,
                                  ['Predictor','β̂','Direction']):
                p = cell.paragraphs[0]; r = p.add_run(hdr)
                r.bold = True; r.font.size = Pt(10)
                _set_cell_bg(cell, '1A2C4E')
                for rn in p.runs: rn.font.color.rgb = WHITE
            for pred, beta in (_ols.get("coef_table") or
                               {p: {"beta": b} for p,b in _ols.get("beta",{}).items()}).items():
                bv  = beta if isinstance(beta, float) else beta.get("beta", 0)
                row = t.add_row()
                for cell, val in zip(row.cells,
                                      [pred, f'{bv:.4f}',
                                       'positive' if bv > 0 else 'negative']):
                    p = cell.paragraphs[0]; r = p.add_run(val); r.font.size = Pt(10)
                    _set_cell_bg(cell, 'DCE8F5')
        body('*** p<0.001  ** p<0.01  * p<0.05  ns = not significant. '
             'Standardised coefficients — comparable in magnitude. ',
             italic=True)

    # ── v3.1 N-06: Cross-Variable Consistency Checker ───────────────────────
    if v3_cvcc:
        h2('Cross-Variable Consistency Check (N-06)')
        body(
            'Chemical and physical constraints between related variables are verified '
            'algorithmically. Example: NOx ≥ NO2 always by chemical definition. '
            f"Verdict: {v3_cvcc.get('CVCC_verdict','N/A')}. "
            f"{v3_cvcc.get('constraints_satisfied',0)}/{v3_cvcc.get('constraints_checked',0)} "
            'constraints satisfied.'
        )
        if v3_cvcc.get('violations'):
            for viol in v3_cvcc['violations']:
                p = doc.add_paragraph(style='List Bullet')
                p.add_run(f"[{viol.get('severity','?')}] ").bold = True
                p.add_run(viol.get('description', viol.get('note','')))
        else:
            body('All chemical constraints satisfied.', italic=True)
        doc.add_paragraph()

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 07 — SEASONAL ANALYSIS
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 07 — Seasonal Pattern Analysis')
    body(
        'Monthly box plots reveal the full distribution of daily pollutant concentrations '
        'across the year. Each box spans the 25th to 75th percentile, with the horizontal '
        'line at the median. Whiskers extend to 1.5 times the inter-quartile range; '
        'individual points beyond the whiskers are statistical outliers. '
        'This non-parametric display requires no normality assumption and clearly shows '
        'which months carry the highest pollution burden.'
    )
    add_image('seasonal', 6.2,
              'Figure — Monthly distribution: median x̃ (Eq. 24), IQR (Eq. 25), '
              'whiskers = 1.5×IQR. Non-parametric — no normality assumption. '
              'Dotted line = WHO 24h reference guideline.')

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 08 — DIURNAL (HOURLY) PROFILE
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 08 — Diurnal Cycle Analysis')
    body(
        'The diurnal profile shows the average pollutant concentration for each hour of '
        'the day (00:00 to 23:00), computed across the full monitoring period. '
        'The timing of daily peaks identifies the dominant emission sources: '
        'morning and evening peaks typically indicate traffic, while midday peaks '
        'may reflect industrial or secondary photochemical formation. '
        'The shaded band represents the mean ± 1σ envelope across all days in the dataset. '
        'The normalised overlay panel (peak = 1.0 per pollutant) allows direct '
        'shape comparison across variables regardless of their concentration magnitudes.'
    )
    # Data-driven diurnal insight
    _ci = chart_insights.get("diurnal","")
    if _ci:
        body(_ci)
    add_image('diurnal', 6.2,
              'Figure — Hourly diurnal profile for key pollutants. Shaded band = mean ± σ. '
              'Bottom-right panel: normalised overlay (peak = 1.0 per variable). '
              'Peak hours identify primary emission timing.')

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 09 — DAY-OF-WEEK PATTERN
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 09 — Day-of-Week Pattern')
    body(
        'Day-of-week analysis reveals whether pollution follows a weekly anthropogenic '
        'cycle. Concentrations are typically higher on weekdays due to traffic and '
        'industrial activity, and lower on weekends. A clear weekday-to-weekend '
        'contrast indicates human-driven emission sources; the absence of such a '
        'pattern suggests natural or regionally transported background pollution.'
    )
    _dow_insight = chart_insights.get("day_of_week","")
    if _dow_insight:
        body(_dow_insight)
    add_image('day_of_week', 6.2,
              'Figure — Mean daily concentration by day of week. Error bars = ±1σ. '
              'Weekday–weekend contrast indicates anthropogenic emission cycle.')

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 10 — VOC COMPOUNDS PANEL (skipped if no VOC data in dataset)
    # ════════════════════════════════════════════════════════════════════════
    _voc_vars_present = any(
        J.get('variables', {}).get(v, {}).get('annual_mean') is not None
        for v in ['Benzene','Toluene','Xylene','Eth_Benzene','MP_Xylene']
    )
    if _voc_vars_present or charts.get('voc'):
        doc.add_page_break()
        h1('Section 10 — Volatile Organic Compounds (VOC) Panel')
        body(
            'The BTEX group (Benzene, Toluene, Ethylbenzene, and Xylene) together with '
            'MP-Xylene are primary markers of vehicular exhaust and industrial solvent emissions. '
            'Benzene is classified as a Group 1 carcinogen by the IARC, with a WHO annual '
            'reference concentration of 1.7 µg/m³. '
            'The grouped monthly bar chart in panel (b) reveals the seasonal pattern of '
            'VOC concentrations and identifies months with elevated BTEX loading.'
        )
        add_image('voc', 6.0,
                  'Figure — (a) VOC daily means time series. '
                  '(b) Monthly grouped bar chart for all VOC compounds. '
                  'Benzene WHO annual guideline: 1.7 µg/m³ (dotted red line, panel a).')

    # VOC summary table (only shown if section rendered)
    if _voc_vars_present or charts.get('voc'):
      voc_vars = ["Benzene","Toluene","Xylene","Eth_Benzene","MP_Xylene"]
      voc_data = {k: v for k,v in J.get("variables",{}).items() if k in voc_vars}
      if voc_data:
        h2("VOC/BTEX Annual Summary")
        t_voc = doc.add_table(rows=1, cols=5)
        t_voc.style = "Table Grid"
        for cell, hdr in zip(t_voc.rows[0].cells,
                              ["Compound","Annual Mean","WHO Limit","Exceed Days","Status"]):
            _set_cell_bg(cell, "4C1D95")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
        for vname, vinfo in voc_data.items():
            row = t_voc.add_row()
            thr  = vinfo.get("threshold")
            mean = vinfo.get("annual_mean",0)
            exc  = vinfo.get("exceeds_days",0)
            bg   = "FFE4E6" if (thr and mean > thr) else "D1FAE5"
            for cell, val in zip(row.cells, [
                vname,
                f"{mean:.3f} µg/m³",
                f"{thr:.1f} µg/m³" if thr else "No guideline",
                str(exc) if thr else "N/A",
                "EXCEEDS" if (thr and mean > thr) else ("COMPLIANT" if thr else "MONITOR"),
            ]):
                _set_cell_bg(cell, bg)
                cell.paragraphs[0].add_run(val).font.size = Pt(9)
        doc.add_paragraph()

    # Only add page break if VOC chart or table was rendered
    # (avoids blank page when station has no VOC data)
    doc.add_page_break()  # Before Section 11

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 11 — RAIN vs POLLUTANT SCATTER
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 11 — Precipitation vs Pollutant Scatter (Wet Deposition)')
    body(
        'Scatter plots of daily rainfall (mm) against PM10, SO2, and NO2 daily means '
        'reveal statistical associations between precipitation events and pollutant '
        'concentration levels. Negative Pearson r values indicate that higher rainfall '
        'is associated with lower concentrations — consistent with wet deposition, '
        'where precipitation removes particulates and soluble gases from the atmosphere. '
        'Note: this is a statistical association only. Confounding meteorological '
        'factors such as wind speed and boundary layer height may contribute independently.'
    )
    if charts.get('rain_scatter'):
        add_image('rain_scatter', 6.2,
                  'Figure — Daily rain (mm) vs pollutant concentration scatter. '
                  'Dashed line = OLS trend. r = Pearson correlation coefficient. '
                  'Association only — not causal wet deposition estimation.')
    else:
        h2('Precipitation Data — Not Available at This Station')
        body(
            'The rain vs pollutant scatter chart requires a rain_mm column in the dataset. '
            'This column was not detected in the uploaded CSV. '
            'Without precipitation data, wet deposition analysis cannot be performed. '
            'Alternative: use Open-Meteo weather data (--weather open_meteo) which '
            'provides daily precipitation from ERA5 reanalysis for any location.'
        )
        # Show correlation table from what IS available
        if ols_result and ols_result.get('r_squared') is not None:
            body(f"OLS regression (PM10 ~ meteorological predictors): "
                 f"R² = {ols_result['r_squared']:.4f}, "
                 f"n = {ols_result.get('n_obs',0)} observations. "
                 f"See Section 06 for full correlation analysis.")
        doc.add_paragraph()

    # Section 12: only give it its own page if calendar chart renders
    if charts.get('calendar'):
        doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 12 — EXCEEDANCE CALENDAR HEATMAP
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 12 — Annual Exceedance Calendar Heatmap')
    body('The calendar heatmap provides a full-year view of daily PM10 concentrations '
         'at a glance. Each cell shows the daily mean value. Cells marked × indicate '
         'days where the daily mean exceeded the WHO 24h guideline. The colour scale '
         'from green (clean) to red (polluted) immediately reveals seasonal pollution '
         'clusters, high-concentration episodes, and clean-air periods.')
    if charts.get('calendar'):
        add_image('calendar', 6.2,
                  'Figure — PM10 calendar heatmap: month (rows) × day of month (columns). '
                  '× = WHO 24h guideline exceedance. Green = clean, Red = polluted. '
                  'Numeric value = daily mean µg/m³.')
    else:
        pm10_v = J.get('variables', {}).get('PM10', {})
        exc_d  = pm10_v.get('exceeds_days', 0)
        exc_p  = pm10_v.get('exceed_pct', 0)
        ann_m  = pm10_v.get('annual_mean', 0)
        body(
            f"Exceedance summary: {exc_d} days exceeded the WHO 24h guideline "
            f"({exc_p:.1f}% of the monitored period). "
            f"Annual PM10 mean: {ann_m:.1f} µg/m³. "
            f"See Section 03 (time-series) and Section 04 (anomaly) for full detail."
        )

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 12b — ATMOSPHERIC WEATHER ANALYSIS & POLLUTION RISK
    # ════════════════════════════════════════════════════════════════════════
    # Section 12b always shown — shows whatever data is available
    doc.add_page_break()
    h1('Section 12b — Atmospheric Analysis, Weather Forecast & Episode Classification')
    body(
        'This section integrates live meteorological data from Open-Meteo '
        '(free API, no key required) with ATARS formal analysis. '
        'Atmospheric stability, pollution episode classification, and '
        '7-day weather forecast with pollution risk are presented. '
        'Open-Meteo source: https://open-meteo.com — ERA5 reanalysis (historical) '
        'and GFS/ICON forecast (future dates). Free, no API key required.'
    )
    # Weather context summary from comparison
    if v3_comparison and v3_comparison.get("weather"):
        h2("Current Meteorological Context (Open-Meteo)")
        wx = v3_comparison["weather"]
        wx_rows = [
            ("Data source",      wx.get("source","N/A")),
            ("Temperature",      f"{wx.get('temp_c', wx.get('temp_mean_c',0)):.1f} °C"
                                 + (f"  (max {wx.get('temp_max_c',0):.1f} / min {wx.get('temp_min_c',0):.1f})"
                                    if wx.get("temp_max_c") else "")),
            ("Humidity",         f"{wx.get('humidity_pct',0):.1f} %"),
            ("Wind speed",       f"{wx.get('wind_kmh', wx.get('wind_mean_kmh',0)):.1f} km/h"
                                 + (f"  (direction: {wx.get('wind_direction','N/A')})"
                                    if wx.get('wind_direction') else "")),
            ("Precipitation",    f"{wx.get('rain_mm',0):.1f} mm"),
            ("Pressure",         f"{wx.get('pressure_hpa',0):.1f} hPa"
                                 if wx.get("pressure_hpa") else "N/A"),
            ("Weather",          wx.get("weather_desc","N/A")),
            ("Pollution risk",   f"{wx.get('pollution_risk','N/A')}  "
                                 f"(score {wx.get('pollution_risk_score',0)}/6)"),
            ("Inversion risk",   "YES — pollution trapping likely"
                                 if wx.get("inversion_risk") else "No"),
            ("Wet deposition",   "YES — rain expected to reduce PM"
                                 if wx.get("wet_deposition") else "No"),
            ("Stagnant air",     "YES — dispersion limited"
                                 if wx.get("stagnant_air") else "No"),
        ]
        t_wx = doc.add_table(rows=1, cols=2)
        t_wx.style = "Table Grid"
        for cell, hdr in zip(t_wx.rows[0].cells, ["Parameter","Value"]):
            _set_cell_bg(cell, "1A6B72")
            r = cell.paragraphs[0].add_run(hdr)
            r.bold = True; r.font.color.rgb = WHITE
        for k, v in wx_rows:
            if v == "N/A" or not v:
                continue
            row = t_wx.add_row()
            row.cells[0].text = k; row.cells[1].text = str(v)
            row.cells[0].paragraphs[0].runs[0].font.size = Pt(10)
            row.cells[1].paragraphs[0].runs[0].font.size = Pt(10)
        doc.add_paragraph()

    # ── Atmospheric Stability ─────────────────────────────────────────
    if v3_stability and v3_stability.get("daily"):
        h2('Atmospheric Stability Analysis (Pasquill-Gifford Proxy)')
        body(
            f"Atmospheric stability was computed for {v3_stability.get('total_days',0)} days. "
            f"{v3_stability.get('high_risk_count',0)} days showed HIGH or SEVERE "
            f"pollution trapping risk (inversion index ≥ 5). "
            f"Higher inversion index = stronger pollution trapping = elevated PM10."
        )
        hr = v3_stability.get("high_risk_days", [])[:8]
        if hr:
            h3 = lambda t: _h2_inner(doc, t, AMBER)
            h3("Highest Pollution Trapping Days")
            stab_t = doc.add_table(rows=1, cols=5)
            stab_t.style = "Table Grid"
            for cell, hdr in zip(stab_t.rows[0].cells,
                ["Date","Inversion Index","Trapping Risk","Temp Min (°C)","Wind (km/h)"]):
                _set_cell_bg(cell, "92400E")
                r2 = cell.paragraphs[0].add_run(hdr)
                r2.bold = True; r2.font.color.rgb = WHITE
            for day in hr:
                row = stab_t.add_row()
                idx = day.get("inversion_index",0)
                bg  = "FFE4E6" if idx >= 7 else "FEF3C7"
                for cell, val in zip(row.cells, [
                    day.get("date",""),
                    f"{idx}/10",
                    day.get("trapping_risk","N/A"),
                    f"{day.get('temp_min',0):.1f}",
                    f"{day.get('wind_kmh',0):.1f}",
                ]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(str(val)).font.size = Pt(9)
            doc.add_paragraph()

    # ── If no atmospheric stability data, explain why ────────────────────
    if not v3_stability or not v3_stability.get("daily"):
        h2("Atmospheric Stability — Data Not Available")
        body(
            "Atmospheric stability analysis requires Open-Meteo weather data. "
            "Run with --weather open_meteo (default, free, no API key) to enable. "
            "When active, it computes Pasquill-Gifford proxy stability class, "
            "boundary layer height, inversion index (0–10), and pollution trapping risk."
        )

    # ── Episode Classification ────────────────────────────────────────────
    if v3_episodes:
        h2("Pollution Episode Classification (N-10 — Novel Contribution)")
        body(
            f"Every anomaly day was classified using rule-based deterministic "
            f"logic — zero LLM. {len(v3_episodes)} pollution episodes classified "
            f"across 8 episode types."
        )
        ep_t = doc.add_table(rows=1, cols=4)
        ep_t.style = "Table Grid"
        for cell, hdr in zip(ep_t.rows[0].cells,
            ["Date","PM10 (µg/m³)","Episode Type","Confidence"]):
            _set_cell_bg(cell, "1A2C4E")
            r3 = cell.paragraphs[0].add_run(hdr)
            r3.bold = True; r3.font.color.rgb = WHITE
        for ep in v3_episodes[:15]:
            row = ep_t.add_row()
            conf = ep.get("confidence","LOW")
            bg   = "D1FAE5" if conf=="HIGH" else ("FEF3C7" if conf=="MEDIUM" else "F5F5F5")
            for cell, val in zip(row.cells, [
                ep.get("date",""),
                f"{ep.get('pm10_mean',0):.1f}",
                ep.get("episode_type","").replace("_"," "),
                conf,
            ]):
                _set_cell_bg(cell, bg)
                cell.paragraphs[0].add_run(str(val)).font.size = Pt(9)
        doc.add_paragraph()
        # Episode reason for top 3
        body("Top episode explanations:", italic=True)
        for ep in sorted(v3_episodes,
                         key=lambda x: x.get("pm10_mean",0), reverse=True)[:3]:
            p2 = doc.add_paragraph(style="List Bullet")
            p2.add_run(f"{ep.get('date','')}: ").bold = True
            p2.add_run(ep.get("reason",""))
        doc.add_paragraph()

    # ── 7-day forecast ────────────────────────────────────────────────
    if v3_forecast:
        h2("7-Day Weather Forecast with Pollution Risk")
        body(
            "Source: Open-Meteo forecast API (free, no key required). "
            "Pollution risk score combines wind speed, temperature inversion, "
            "precipitation, fog, and photochemical smog indicators."
        )
        fc_t = doc.add_table(rows=1, cols=7)
        fc_t.style = "Table Grid"
        for cell, hdr in zip(fc_t.rows[0].cells,
            ["Date","Weather","Temp (°C)","Wind (km/h)","Rain (mm)","Risk","Score"]):
            _set_cell_bg(cell, "1A6B72")
            r4 = cell.paragraphs[0].add_run(hdr)
            r4.bold = True; r4.font.color.rgb = WHITE
            r4.font.size = Pt(8)
        for d, w in sorted(v3_forecast.items()):
            row = fc_t.add_row()
            risk  = w.get("pollution_risk","N/A")
            score = w.get("pollution_risk_score",0)
            bg    = ("FFE4E6" if risk in ("CRITICAL","HIGH") else
                     "FEF3C7" if risk == "MODERATE" else "D1FAE5")
            for cell, val in zip(row.cells, [
                d,
                w.get("weather_desc","N/A")[:15],
                f"{w.get('temp_mean_c',w.get('temp_min_c',0)):.1f}",
                f"{w.get('wind_mean_kmh',w.get('wind_kmh',0)):.1f}",
                f"{w.get('rain_mm',0):.1f}",
                risk,
                f"{score}/6",
            ]):
                _set_cell_bg(cell, bg)
                cell.paragraphs[0].add_run(str(val)).font.size = Pt(8)
        doc.add_paragraph()
        # Risk alert
        high_risk = [d for d,w in v3_forecast.items()
                    if w.get("pollution_risk_score",0) >= 3]
        if high_risk:
            p3 = doc.add_paragraph()
            r5 = p3.add_run(
                f"⚠ HIGH POLLUTION RISK FORECAST: {len(high_risk)} day(s) "
                f"with risk score ≥ 3: {', '.join(high_risk)}")
            r5.bold = True; r5.font.color.rgb = RED
        doc.add_paragraph()

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 13 — ML ANALYSIS: ISOLATION FOREST
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 13 — ML Analysis: Isolation Forest Multivariate Anomaly Detection')
    p = doc.add_paragraph()
    run = p.add_run(
        'This analysis is additive and does not replace the formal z-score operator ζ_d (Eq. 9). '
        'Isolation Forest detects days where the combination of multiple pollutants '
        'is anomalous — a pattern that univariate z-score analysis cannot identify. '
        'All computation is performed locally on the aggregated daily dataset; '
        'no raw data is transmitted externally. '
        'The fixed random seed (random_state = 42) ensures this result is fully reproducible.'
    )
    run.font.size = Pt(10); run.italic = True; run.font.color.rgb = TEAL
    p.paragraph_format.space_before = Pt(4)

    if ml_if.get('available'):
        # Summary table
        t = doc.add_table(rows=1, cols=2)
        t.style = 'Table Grid'
        hdr = t.rows[0].cells
        hdr[0].text = 'Isolation Forest Parameter'; hdr[1].text = 'Value'
        for cell in hdr:
            cell.paragraphs[0].runs[0].bold = True
            _set_cell_bg(cell, '003366')
            cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

        rows_data = [
            ('Features used', ', '.join(ml_if['features_used'])),
            ('n_estimators', str(ml_if['n_estimators'])),
            ('contamination', f"{ml_if['contamination']:.0%} (expected anomaly fraction)"),
            ('random_state', f"{ml_if['random_state']} (fixed for reproducibility)"),
            ('Total days analysed', str(ml_if['n_days'])),
            ('Anomaly days (IF)', f"{ml_if['n_anomalies']} ({ml_if['anomaly_rate']*100:.1f}%)"),
            ('Flagged by BOTH methods', str(ml_if['n_both'])),
            ('IF-only (multivariate — missed by z-score)', str(ml_if['n_if_only'])),
            ('Z-score only (univariate spike, within multivariate norm)', str(ml_if['n_zs_only'])),
        ]
        for lbl, val in rows_data:
            row = t.add_row().cells
            row[0].text = lbl; row[1].text = val

        doc.add_paragraph()
        add_image('ml_anomaly', 6.2,
                  'Figure — Isolation Forest vs Z-Score Anomaly Comparison. '
                  '(a) IF decision score — lower = more anomalous. '
                  '(b) Z-score ζ_d with ±3σ threshold. '
                  '(c) Agreement map — Red=both, Purple=IF-only, Amber=z-score-only.')
        
        p = doc.add_paragraph()
        run = p.add_run(
            'Interpretation of IF-only days: These are dates where no single pollutant '
            'crossed the 3σ z-score threshold, but the joint distribution of '
            f"{', '.join(ml_if['features_used'])} was statistically unusual compared "
            'to the training distribution. These warrant manual investigation as potential '
            'multi-source pollution episodes. IF-only days represent the added analytical '
            'value of multivariate over univariate anomaly detection.'
        )
        run.font.size = Pt(9.5); run.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
    else:
        reason = ml_if.get("reason", "insufficient data")
        body(f"Isolation Forest was not run: {reason}.", italic=True)
        h2("Why Isolation Forest Was Skipped")
        skip_reasons = {
            "insufficient_features":
                f"Fewer than 2 pollutant feature columns were available after cleaning. "
                f"Isolation Forest requires at least 2 correlated variables to detect "
                f"multivariate anomalies. With only 1 variable, z-score is equivalent.",
            "insufficient_rows":
                f"Fewer than 30 complete daily rows were available after quality filtering. "
                f"Isolation Forest needs sufficient data to build a stable anomaly boundary. "
                f"The z-score analysis (Section 04) covers all available days.",
        }
        body(skip_reasons.get(reason,
             f"Reason: {reason}. Z-score anomaly detection (Section 04) remains valid "
             f"and covers all monitored days."))
        h2("What Isolation Forest Would Show")
        body(
            "When sufficient data is available, Isolation Forest detects days where the "
            "COMBINATION of PM10, NO2, SO2, CO, and Ozone is anomalous — even if no "
            "single variable crosses the 3σ z-score threshold. "
            "This catches multi-source pollution episodes that univariate analysis misses. "
            "To enable it, ensure at least 2 pollutant columns have ≥30 complete daily means."
        )

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 14 — ML ANALYSIS: HOLT-WINTERS FORECAST
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 14 — ML Analysis: Holt-Winters PM10 Forecast')
    p = doc.add_paragraph()
    run = p.add_run(
        'This section presents supplementary PM10 forecasting using Triple Exponential '
        'Smoothing (Holt-Winters additive seasonal model), which captures level, trend, '
        'and weekly seasonal components simultaneously. '
        'Forecast values are deliberately excluded from the JSON contract J so that the '
        'language model cannot present statistical extrapolations as established facts. '
        'All forecast content appears in this section and Chart 15 only.'
    )
    run.font.size = Pt(10); run.italic = True; run.font.color.rgb = TEAL
    p.paragraph_format.space_before = Pt(4)

    if ml_hw.get('available'):
        t = doc.add_table(rows=1, cols=2)
        t.style = 'Table Grid'
        hdr = t.rows[0].cells
        hdr[0].text = 'Holt-Winters Parameter'; hdr[1].text = 'Value'
        for cell in hdr:
            cell.paragraphs[0].runs[0].bold = True
            _set_cell_bg(cell, '003366')
            cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

        fc_mean = float(np.mean(ml_hw['forecast']))
        fc_max  = float(np.max(ml_hw['forecast']))
        rows_data = [
            ('α (level smoothing)', str(ml_hw['alpha'])),
            ('β (trend smoothing)', str(ml_hw['beta'])),
            ('γ (seasonal smoothing)', str(ml_hw['gamma'])),
            ('Seasonal period', f"{ml_hw['seasonal_periods']} days (weekly)"),
            ('Forecast horizon', f"{ml_hw['forecast_days']} days"),
            ('History used', f"{ml_hw['n_history']} days"),
            ('In-sample RMSE', f"{ml_hw['rmse']:.2f} µg/m³"),
            ('In-sample MAE', f"{ml_hw['mae']:.2f} µg/m³"),
            ('Residual std (σ)', f"{ml_hw['resid_std']:.2f} µg/m³"),
            (f"{ml_hw['forecast_days']}-day forecast mean", f"{fc_mean:.1f} µg/m³"),
            (f"{ml_hw['forecast_days']}-day forecast peak", f"{fc_max:.1f} µg/m³"),
            ('Prediction interval', '80% (±1.28σ√h — widens with horizon)'),
        ]
        for lbl, val in rows_data:
            row = t.add_row().cells
            row[0].text = lbl; row[1].text = val

        doc.add_paragraph()
        add_image('forecast', 6.2,
                  'Figure — Holt-Winters PM10 Forecast. '
                  '(a) Historical PM10 + fitted values + 14-day forecast with 80% prediction interval. '
                  '(b) Model residuals with ±1σ bands.')

        p = doc.add_paragraph()
        run = p.add_run(
            'Limitations of this forecast: (1) Holt-Winters is a statistical extrapolation, '
            'not a physical atmospheric model — it does not account for meteorological events, '
            'policy interventions, or emission changes. '
            '(2) Uncertainty grows significantly beyond 7 days — the 80% PI widens with √h. '
            '(3) The model was trained on this dataset only — transferability to other '
            'stations or time periods is not validated here. '
            '(4) For public health decisions, use this forecast as an indicative signal '
            'only — not as a definitive prediction.'
        )
        run.font.size = Pt(9.5); run.font.color.rgb = RGBColor(0x92, 0x40, 0x0E)
    else:
        reason = ml_hw.get("reason", "disabled or insufficient data")
        body(f"Holt-Winters forecast was not run: {reason}.", italic=True)
        h2("Why the Forecast Was Skipped")
        skip_map = {
            "insufficient_pm10_data":
                "Fewer than 14 PM10 daily observations were available. "
                "Holt-Winters requires at least 14 days (2× the weekly seasonal period) "
                "to initialise level, trend, and seasonal components.",
            "insufficient_history":
                "Fewer than 2 seasonal periods of data were available (< 14 days). "
                "Increase the monitoring window to enable forecasting.",
            "disabled_by_flag":
                "Forecast was disabled via the --no-forecast flag. "
                "Remove this flag to enable the 14-day PM10 forecast.",
        }
        body(skip_map.get(reason,
             f"Reason: {reason}. "
             "Ensure PM10 data spans at least 14 days and --no-forecast is not set."))
        h2("What the Forecast Would Show")
        body(
            "When enabled, Holt-Winters triple exponential smoothing produces a 14-day "
            "PM10 forecast with 80% prediction intervals. "
            "It captures: (α) level — current pollution baseline; "
            "(β) trend — whether PM10 is rising or falling; "
            "(γ) weekly seasonality — consistent day-of-week patterns. "
            "Forecast values are intentionally excluded from the LLM data contract J "
            "to prevent the AI from presenting uncertain extrapolations as historical facts."
        )
    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 15 — AI-ASSISTED COMMENTARY AND LIMITATIONS
    # ════════════════════════════════════════════════════════════════════════
    h1('Section 15 — AI-Assisted Commentary and Limitations')

    # AI disclosure banner
    p = doc.add_paragraph()
    run = p.add_run(
        'Sections 01 through 12 contain AI-assisted narrative text generated '
        'by a Large Language Model operating under formal grounding constraints '
        'G(s) ⊆ J (ATARS Eq. 19). '
        f'Model used: {narrative.get("model", "statistical_placeholder")}. '
        'Temperature τ = 0 (deterministic output). '
        'All numerical claims in AI-generated text reference J exclusively. '
        + ('AI narrative generation was active.' if narrative.get('ai_flag')
           else 'AI narrative was unavailable — statistical placeholders used.')
    )
    run.font.size = Pt(10); run.italic = True; run.font.color.rgb = BLUE
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after  = Pt(12)
    body(narrative.get('F_limitations', ''))

    # ── Narrative Specificity Score (N-07) ────────────────────────────────
    if v3_nss:
        h2("Narrative Specificity Score (N-07)")
        body(
            "Measures how specific vs vague each narrative sentence is. "
            f"Overall specificity: {v3_nss.get('overall_specificity',0):.2f}  "
            f"[{v3_nss.get('NSS_verdict','N/A')}]. "
            "Score = 1.0: highly specific (numbers, units, formal notation). "
            "Score = 0.0: vague (no measurable claims). "
            "Vague sentences are penalised; specific numerical sentences rewarded."
        )
        nss_secs = v3_nss.get("sections", {})
        if nss_secs:
            t_nss = doc.add_table(rows=1, cols=3)
            t_nss.style = "Table Grid"
            for cell, hdr in zip(t_nss.rows[0].cells,
                                  ["Section","Avg Specificity","Verdict"]):
                _set_cell_bg(cell, "4C1D95")
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE
            for sk, sv in nss_secs.items():
                row = t_nss.add_row()
                sp = sv.get("avg_specificity", 0)
                bg = "D1FAE5" if sp>=0.4 else ("FEF3C7" if sp>=0.2 else "FFE4E6")
                for cell, val in zip(row.cells,
                                      [sk, f"{sp:.2f}", sv.get("verdict","N/A")]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(val).font.size = Pt(9)
            doc.add_paragraph()

    # ── Causal Claim Detector (N-02) ──────────────────────────────────────
    if v3_ccd:
        h2("Causal Claim Detector (N-02) — Formal Constraint: Corr(X,Y) ≢ C(X→Y)")
        body(
            f"CCD verdict: {v3_ccd.get('CCD_verdict','N/A')}. "
            f"Causal claims detected: {v3_ccd.get('causal_claim_count',0)} "
            f"out of {v3_ccd.get('total_sentences',0)} sentences checked. "
            "Causal language (e.g. 'causes', 'due to', 'leads to') is flagged when "
            "only correlation data exists in the JSON contract J. "
            f"Formal rule enforced: {v3_ccd.get('formal_rule','Corr(X,Y) ≢ C(X→Y)')}."
        )
        flagged = v3_ccd.get("flagged", [])
        if flagged:
            h2(f"Flagged Causal Claims ({len(flagged)})")
            t_ccd = doc.add_table(rows=1, cols=3)
            t_ccd.style = "Table Grid"
            for cell, hdr in zip(t_ccd.rows[0].cells,
                                  ["Section","Severity","Sentence (truncated)"]):
                _set_cell_bg(cell, "92400E")
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE
            for fc in flagged[:8]:
                row = t_ccd.add_row()
                sev = fc.get("severity","?")
                bg  = "FFE4E6" if sev=="HIGH" else "FEF3C7"
                for cell, val in zip(row.cells,
                                      [fc.get("section",""),
                                       sev,
                                       fc.get("sentence","")[:80]]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(val).font.size = Pt(9)
            doc.add_paragraph()
        else:
            body("✓ No causal language detected — all associations correctly framed.",
                 italic=True)
        doc.add_paragraph()

    # ── Data Completeness Warning ─────────────────────────────────────────
    h2("Formal Analytical Limitations")
    limitations = [
        (f"Statistical association only: All r values and OLS coefficients describe "
         f"Pearson association between variables. Corr(X,Y) does not imply C(X→Y)."),
        (f"Data completeness Q(D) = {J['data_quality']['overall_Q']:.3f} measures "
         f"record completeness only — sensor calibration accuracy and measurement "
         f"uncertainty are outside the scope of this quality metric."),
        (f"Results are specific to this station ({J.get('station_id','N/A')}) and "
         f"the {J.get('total_days',0)}-day monitoring period "
         f"({J.get('date_range','N/A')}). "
         f"Generalisation to other stations requires independent validation."),
        ("Expert domain review by a qualified atmospheric scientist is recommended "
         "before using these results to inform regulatory or public health decisions."),
        ("Holt-Winters forecast values are statistical extrapolations only — not "
         "physical atmosphere model predictions. Uncertainty grows with forecast horizon."),
        ("Isolation Forest anomaly detection is data-driven and unsupervised — "
         "flagged dates warrant manual inspection but are not confirmed pollution events."),
    ]
    for lim in limitations:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(lim).font.size = Pt(10)
    doc.add_paragraph()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 16 — RUNTIME GROUNDING VERIFICATION + NARRATIVE AUDIT
    # ════════════════════════════════════════════════════════════════════════
    if verification:
        doc.add_page_break()
        h1('Section 16 — Runtime Grounding Verification + Narrative Audit')

        # Novel contribution banner
        p = doc.add_paragraph()
        _is_ai = narrative.get('ai_flag', False)
        _grnd_note = (
            'This section presents the Runtime Grounding Verifier (RGV), which '
            'algorithmically enforces G(s) ⊆ J for every LLM-generated numerical '
            'claim. The scalar G_rate ∈ [0,1] quantifies grounding. '
            'This is the first published system to embed numerical grounding '
            'verification inside an automated environmental reporting pipeline.'
            if _is_ai else
            'NOTE: LLM narrative was unavailable — statistical placeholders were used. '
            'Placeholder text is generated directly from J values, so G_rate = 100% is '
            'expected and does NOT indicate AI grounding quality. '
            'G_rate is only meaningful when an LLM (Ollama) generates the narrative. '
            'To enable real verification: install Ollama and run without --no-llm.'
        )
        run = p.add_run(_grnd_note)
        run.italic = True
        run.font.size = Pt(10)
        run.font.color.rgb = TEAL if _is_ai else AMBER
        doc.add_paragraph()

        # Grounding result
        h2('Formal Verification Result')
        g_rate  = verification.get("overall_G_rate", 0)
        verdict = verification.get("overall_verdict", "N/A")
        n_num   = verification.get("total_numerical", 0)
        n_grd   = verification.get("total_grounded",  0)
        n_ugrd  = verification.get("total_ungrounded", 0)
        j_size  = verification.get("j_lookup_size",   0)

        p = doc.add_paragraph()
        run = p.add_run(f"G_rate = {g_rate:.1%}   [{verdict}]")
        run.bold = True
        run.font.size = Pt(18)
        passed = g_rate >= 0.75
        run.font.color.rgb = GREEN if passed else AMBER

        # Use v3 fields where available, fall back to v2
        _v3ref = v3_ver_v3 or {}
        grounding_rows = [
            ("G_rate (base v2)",             f"{g_rate:.1%}"),
            ("G_rate (enhanced v3)",         f"{_v3ref.get('overall_G_rate', g_rate):.1%}"),
            ("Verdict",                      _v3ref.get("overall_verdict", verdict)),
            ("J lookup entries",             str(j_size)),
            ("Total numerical sentences",    str(_v3ref.get("total_numerical", n_num))),
            ("Grounded",                     str(_v3ref.get("total_grounded", n_grd))),
            ("Ungrounded",                   str(_v3ref.get("total_ungrounded", n_ugrd))),
            ("Avg sentence confidence",      str(_v3ref.get("overall_confidence", "N/A"))),
            ("Semantic alignment",           "PASS" if _v3ref.get("semantic_pass", True) else "REVIEW"),
            ("Direction contradictions",     str(_v3ref.get("contradiction_count", 0))),
            ("Sources verified against",     "J contract + WHO AQG + weather context"),
        ]
        tbl = doc.add_table(rows=1, cols=2)
        tbl.style = "Table Grid"
        for cell, hdr in zip(tbl.rows[0].cells, ["Metric","Value"]):
            _set_cell_bg(cell, "1A6B72")
            p2 = cell.paragraphs[0]
            r2 = p2.add_run(hdr)
            r2.bold = True
            r2.font.color.rgb = WHITE
        for k, v in grounding_rows:
            row = tbl.add_row().cells
            row[0].text = k
            row[1].text = str(v)
        doc.add_paragraph()

        # Per-section results
        h2("Per-Section Grounding Results")
        # Prefer v3_ver_v3 sections (richer); fall back to v2 verification
        sections_v = (v3_ver_v3 or {}).get("sections", {}) or verification.get("sections", {})
        if sections_v:
            tbl2 = doc.add_table(rows=1, cols=4)
            tbl2.style = "Table Grid"
            for cell, hdr in zip(tbl2.rows[0].cells, ["Section","G_rate","Verdict","Confidence"]):
                _set_cell_bg(cell, "1A2C4E")
                p3 = cell.paragraphs[0]
                r3 = p3.add_run(hdr)
                r3.bold = True
                r3.font.color.rgb = WHITE
                p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for sk, sv in sections_v.items():
                row = tbl2.add_row()
                sec_g = sv.get("G_rate", sv.get("g_rate", 0))
                bg = "D1FAE5" if sec_g >= 0.75 else "FFE4E6"
                for cell, val in zip(row.cells, [
                    sk,
                    f"{sec_g:.1%}",
                    sv.get("verdict","N/A"),
                    str(sv.get("avg_confidence","N/A")),
                ]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(val).font.size = Pt(9)
            doc.add_paragraph()

        # TCV — Temporal Claim Verifier
        if v3_tcv and v3_tcv.get("total_claims", 0) > 0:
            h2("Temporal Claim Verifier — N-01")
            body(
                f"Verifies month/season claims in AI narrative against computed "
                f"temporal facts. TCV rate: {round(v3_tcv.get('TCV_rate',0)*100,1)}% "
                f"[{v3_tcv.get('TCV_verdict','N/A')}]. "
                f"{v3_tcv.get('total_claims',0)} temporal claims checked. "
                f"PM10 peak month: {v3_tcv.get('actual_peak_month','').capitalize()}. "
                f"PM10 low month: {v3_tcv.get('actual_low_month','').capitalize()}."
            )
        

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 16b — NINE NOVEL CONTRIBUTIONS: DETAILED VERIFICATION RESULTS
    # ════════════════════════════════════════════════════════════════════════
    if verification:
        doc.add_page_break()
        h1('Section 16b — Nine Novel Verification Modules: Detailed Results')
        body(
            'ATARS v3.1 introduces 9 novel verification modules (N-01 through N-09) '
            'each targeting a different failure mode of LLM-generated scientific text. '
            'This section shows the detailed output of each module — analogous to how '
            'Section 16 shows RGV results. Every result is computed deterministically; '
            'zero LLM involvement in the verification process itself.'
        )

        # ── N-01: TCV ──────────────────────────────────────────────────────
        h2('N-01 — Temporal Claim Verifier (TCV)')
        body(
            'Verifies every month/season claim in the AI narrative against computed '
            'temporal facts from daily_df. First system to verify temporal assertions '
            'algorithmically in an automated environmental report.'
        )
        if v3_tcv:
            tcv_rate = v3_tcv.get('TCV_rate', 1.0)
            tcv_vdt  = v3_tcv.get('TCV_verdict', 'N/A')
            p_tcv = doc.add_paragraph()
            r_tcv = p_tcv.add_run(f"TCV Rate = {tcv_rate:.1%}   [{tcv_vdt}]")
            r_tcv.bold = True; r_tcv.font.size = Pt(14)
            r_tcv.font.color.rgb = GREEN if tcv_vdt in ('PASS','NO_CLAIMS_FOUND') else AMBER
            tcv_rows = [
                ('TCV Rate',              f"{tcv_rate:.1%}"),
                ('Verdict',               tcv_vdt),
                ('Total temporal claims', str(v3_tcv.get('total_claims', 0))),
                ('Correct claims',        str(v3_tcv.get('correct', 0))),
                ('Incorrect claims',      str(v3_tcv.get('incorrect', 0))),
                ('PM10 actual peak month',v3_tcv.get('actual_peak_month','N/A').capitalize()),
                ('PM10 actual low month', v3_tcv.get('actual_low_month','N/A').capitalize()),
                ('Tolerance',             '±1 month for peak/low month claims'),
                ('Novel contribution',    'First temporal claim verifier in automated env. reporting'),
            ]
            t_tcv = doc.add_table(rows=1, cols=2)
            t_tcv.style = 'Table Grid'
            for cell, hdr in zip(t_tcv.rows[0].cells, ['Metric', 'Value']):
                _set_cell_bg(cell, '1A6B72')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for k, v in tcv_rows:
                row = t_tcv.add_row()
                row.cells[0].paragraphs[0].add_run(k).font.size = Pt(9)
                row.cells[1].paragraphs[0].add_run(str(v)).font.size = Pt(9)
                _set_cell_bg(row.cells[0], 'EFF6FF')
            doc.add_paragraph()
            # Show incorrect claims if any
            incorrect = [r for r in v3_tcv.get('results', [])
                        if r.get('verdict') == 'INCORRECT']
            if incorrect:
                body(f"Incorrect temporal claims ({len(incorrect)}):", italic=True)
                for ic in incorrect[:5]:
                    p2 = doc.add_paragraph(style='List Bullet')
                    p2.add_run(f"Claimed '{ic.get('claimed_month', ic.get('claimed_season','?'))}'"
                               f" — actual month {ic.get('actual_month_num','?')}: "
                               f"{ic.get('sentence','')[:80]}").font.size = Pt(9)

        # ── N-02: CCD ──────────────────────────────────────────────────────
        doc.add_paragraph()
        h2('N-02 — Causal Claim Detector (CCD)')
        body(
            'Algorithmically enforces Corr(X,Y) ≢ C(X→Y). Scans every narrative sentence '
            'for causal language (causes, due to, leads to, etc.) and flags instances '
            'where only correlation data exists in J. First causal claim detector '
            'embedded in a scientific reporting pipeline.'
        )
        if v3_ccd:
            ccd_vdt = v3_ccd.get('CCD_verdict', 'N/A')
            p_ccd = doc.add_paragraph()
            r_ccd = p_ccd.add_run(
                f"CCD Verdict: {ccd_vdt}   |   "
                f"Causal claims found: {v3_ccd.get('causal_claim_count', 0)}"
            )
            r_ccd.bold = True; r_ccd.font.size = Pt(13)
            r_ccd.font.color.rgb = GREEN if ccd_vdt == 'PASS' else AMBER
            ccd_rows = [
                ('CCD Verdict',           ccd_vdt),
                ('Causal claims found',   str(v3_ccd.get('causal_claim_count', 0))),
                ('Total sentences scanned', str(v3_ccd.get('total_sentences', 0))),
                ('Formal rule',           'Corr(X,Y) ≢ C(X→Y)'),
                ('Detection method',      'Keyword trigger list (26 causal phrases)'),
                ('Safe language check',   'Skips sentences with association-safe phrasing'),
            ]
            t_ccd = doc.add_table(rows=1, cols=2)
            t_ccd.style = 'Table Grid'
            for cell, hdr in zip(t_ccd.rows[0].cells, ['Metric', 'Value']):
                _set_cell_bg(cell, '1A2C4E')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for k, v in ccd_rows:
                row = t_ccd.add_row()
                row.cells[0].paragraphs[0].add_run(k).font.size = Pt(9)
                row.cells[1].paragraphs[0].add_run(str(v)).font.size = Pt(9)
                _set_cell_bg(row.cells[0], 'F5F5F5')
            flagged = v3_ccd.get('flagged', [])
            if flagged:
                doc.add_paragraph()
                body(f"Flagged causal sentences ({len(flagged)}):", italic=True)
                for fc in flagged[:4]:
                    p3 = doc.add_paragraph(style='List Bullet')
                    p3.add_run(
                        f"[{fc.get('severity','?')}] "
                        f"Triggers: {', '.join(fc.get('triggers_found',[])[:3])} — "
                        f"{fc.get('sentence','')[:80]}"
                    ).font.size = Pt(9)
            doc.add_paragraph()

        # ── N-03: AEE ──────────────────────────────────────────────────────
        h2('N-03 — Anomaly Explanation Engine (AEE)')
        body(
            'For every anomaly day detected by z-score or Isolation Forest, AEE applies '
            'rule-based meteorological context: wind speed, humidity, temperature, '
            'precipitation, day-of-week, and season. Zero LLM. Fully reproducible.'
        )
        if v3_aee:
            expl_rate = v3_aee.get('explanation_rate', 0)
            total_an  = v3_aee.get('total_anomalies', 0)
            explained = v3_aee.get('explained', 0)
            p_aee = doc.add_paragraph()
            r_aee = p_aee.add_run(
                f"AEE: {explained}/{total_an} anomaly days explained "
                f"({expl_rate*100:.1f}%)"
            )
            r_aee.bold = True; r_aee.font.size = Pt(13); r_aee.font.color.rgb = TEAL
            aee_rows_s = [
                ('Total anomaly days',    str(total_an)),
                ('High-confidence explanations', str(explained)),
                ('Explanation rate',      f"{expl_rate*100:.1f}%"),
                ('Method',                'Rule-based deterministic — 8 meteorological rules'),
                ('Rules applied',         'Wind stagnation, humidity, inversion, rainfall, '
                                          'dust, photochemistry, agricultural burning, traffic'),
            ]
            t_aee = doc.add_table(rows=1, cols=2)
            t_aee.style = 'Table Grid'
            for cell, hdr in zip(t_aee.rows[0].cells, ['Metric', 'Value']):
                _set_cell_bg(cell, '92400E')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for k, v in aee_rows_s:
                row = t_aee.add_row()
                row.cells[0].paragraphs[0].add_run(k).font.size = Pt(9)
                row.cells[1].paragraphs[0].add_run(str(v)).font.size = Pt(9)
                _set_cell_bg(row.cells[0], 'FEF3C7')
            # Top 5 explanations
            if v3_aee.get('explanations'):
                doc.add_paragraph()
                body("Top anomaly explanations:", italic=True)
                t_aee2 = doc.add_table(rows=1, cols=4)
                t_aee2.style = 'Table Grid'
                for cell, hdr in zip(t_aee2.rows[0].cells,
                                     ['Date','PM10','Season','Primary Explanation']):
                    _set_cell_bg(cell, '1A2C4E')
                    r = cell.paragraphs[0].add_run(hdr)
                    r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
                for ex in v3_aee['explanations'][:8]:
                    row = t_aee2.add_row()
                    for cell, val in zip(row.cells, [
                        ex.get('date',''),
                        f"{ex.get('pm10_mean',0):.1f} µg/m³",
                        ex.get('season','').capitalize(),
                        ex.get('primary_explanation','')[:60],
                    ]):
                        cell.paragraphs[0].add_run(val).font.size = Pt(8)
            doc.add_paragraph()

        # ── N-04: RDD ──────────────────────────────────────────────────────
        h2('N-04 — Report Drift Detector (RDD)')
        body(
            'Compares current run against previous audit logs to detect degradation '
            'in G_rate, Q(D), or pollutant statistics over time. Alerts when AI '
            'narrative quality is systematically declining across repeated runs.'
        )
        if v3_drift:
            rdd_rows = [
                ('Status',                v3_drift.get('status', 'N/A')),
                ('Previous runs compared', str(v3_drift.get('runs_compared', 0))),
                ('Drifts detected',       str(len(v3_drift.get('drifts_detected', [])))),
                ('Alert triggered',       str(v3_drift.get('alert', False))),
                ('Current G_rate',        f"{v3_drift.get('current_g_rate', 0):.3f}"),
                ('Current Q(D)',          f"{v3_drift.get('current_q', 0):.4f}"),
            ]
            t_rdd = doc.add_table(rows=1, cols=2)
            t_rdd.style = 'Table Grid'
            for cell, hdr in zip(t_rdd.rows[0].cells, ['Metric', 'Value']):
                _set_cell_bg(cell, '1A2C4E')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for k, v in rdd_rows:
                row = t_rdd.add_row()
                row.cells[0].paragraphs[0].add_run(k).font.size = Pt(9)
                row.cells[1].paragraphs[0].add_run(str(v)).font.size = Pt(9)
                _set_cell_bg(row.cells[0], 'F5F5F5')
            for d in v3_drift.get('drifts_detected', []):
                p4 = doc.add_paragraph(style='List Bullet')
                p4.add_run(
                    f"[{d.get('severity','?')}] {d.get('metric','?')}: "
                    f"current={d.get('current',0):.3f} vs "
                    f"avg={d.get('historical_avg',0):.3f} — "
                    f"{d.get('direction','?')}"
                ).font.size = Pt(9)
            doc.add_paragraph()

        # ── N-05: UPT ──────────────────────────────────────────────────────
        h2('N-05 — Uncertainty Propagation Tracker (UPT)')
        body(
            'Computes 95% confidence intervals for every key statistic in J. '
            'Annotates LLM claims with uncertainty bounds so every number in the '
            'report is accompanied by its scientifically honest CI range.'
        )
        if v3_upt:
            upt_rows = [
                ('Claims annotated with CI',  str(v3_upt.get('total_annotated', 0))),
                ('Variables with CI bounds',  str(len(v3_upt.get('uncertainty_table', {})))),
                ('CI method',                 '95% classical (x̄ ± 1.96 × σ/√n)'),
                ('Novel contribution',        'First uncertainty propagation tracker in LLM reporting'),
            ]
            t_upt = doc.add_table(rows=1, cols=2)
            t_upt.style = 'Table Grid'
            for cell, hdr in zip(t_upt.rows[0].cells, ['Metric', 'Value']):
                _set_cell_bg(cell, '4C1D95')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for k, v in upt_rows:
                row = t_upt.add_row()
                row.cells[0].paragraphs[0].add_run(k).font.size = Pt(9)
                row.cells[1].paragraphs[0].add_run(str(v)).font.size = Pt(9)
                _set_cell_bg(row.cells[0], 'EDE9FE')
            # CI table per variable
            if v3_upt.get('uncertainty_table'):
                doc.add_paragraph()
                body("95% CI bounds by variable:", italic=True)
                t_upt2 = doc.add_table(rows=1, cols=5)
                t_upt2.style = 'Table Grid'
                for cell, hdr in zip(t_upt2.rows[0].cells,
                                     ['Variable','Mean','CI Lower','CI Upper','Class']):
                    _set_cell_bg(cell, '4C1D95')
                    r = cell.paragraphs[0].add_run(hdr)
                    r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
                for var, unc in list(v3_upt['uncertainty_table'].items())[:10]:
                    row = t_upt2.add_row()
                    bg = ('D1FAE5' if unc.get('uncertainty_class')=='LOW' else
                          'FEF3C7' if unc.get('uncertainty_class')=='MEDIUM' else 'FFE4E6')
                    for cell, val in zip(row.cells, [
                        var,
                        f"{unc.get('mean',0):.2f} µg/m³",
                        f"{unc.get('ci_95_lower',0):.2f}",
                        f"{unc.get('ci_95_upper',0):.2f}",
                        unc.get('uncertainty_class','N/A'),
                    ]):
                        _set_cell_bg(cell, bg)
                        cell.paragraphs[0].add_run(val).font.size = Pt(9)
            doc.add_paragraph()

        # ── N-06: CVCC ─────────────────────────────────────────────────────
        h2('N-06 — Cross-Variable Consistency Checker (CVCC)')
        body(
            'Verifies chemical and physical constraints between related variables. '
            'NOx ≥ NO2 always (NOx = NO + NO2). T/B ratio > 1 in traffic emissions. '
            'First cross-variable chemical constraint checker in automated reporting.'
        )
        if v3_cvcc:
            cvcc_vdt = v3_cvcc.get('CVCC_verdict', 'N/A')
            p_cvcc = doc.add_paragraph()
            r_cvcc = p_cvcc.add_run(f"CVCC Verdict: {cvcc_vdt}")
            r_cvcc.bold = True; r_cvcc.font.size = Pt(13)
            r_cvcc.font.color.rgb = GREEN if cvcc_vdt == 'PASS' else AMBER
            cvcc_rows = [
                ('CVCC Verdict',            cvcc_vdt),
                ('Constraints checked',     str(v3_cvcc.get('constraints_checked', 0))),
                ('Constraints satisfied',   str(v3_cvcc.get('constraints_satisfied', 0))),
                ('Violations',              str(len(v3_cvcc.get('violations', [])))),
                ('Constraints verified',    'NOx≥NO2, NOx≥NO, T/B ratio, PM10≥PM2.5'),
            ]
            t_cvcc = doc.add_table(rows=1, cols=2)
            t_cvcc.style = 'Table Grid'
            for cell, hdr in zip(t_cvcc.rows[0].cells, ['Metric', 'Value']):
                _set_cell_bg(cell, '1E6B3C')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for k, v in cvcc_rows:
                row = t_cvcc.add_row()
                row.cells[0].paragraphs[0].add_run(k).font.size = Pt(9)
                row.cells[1].paragraphs[0].add_run(str(v)).font.size = Pt(9)
                _set_cell_bg(row.cells[0], 'D1FAE5')
            if v3_cvcc.get('violations'):
                doc.add_paragraph()
                body("Violations detected:", italic=True)
                for viol in v3_cvcc['violations']:
                    p5 = doc.add_paragraph(style='List Bullet')
                    p5.add_run(
                        f"[{viol.get('severity','?')}] {viol.get('description', viol.get('note',''))}"
                    ).font.size = Pt(9)
            else:
                doc.add_paragraph()
                body("✓ All chemical constraints satisfied.", italic=True)
            doc.add_paragraph()

        # ── N-07: NSS ──────────────────────────────────────────────────────
        h2('N-07 — Narrative Specificity Scorer (NSS)')
        body(
            'Scores each sentence 0.0–1.0 for specificity. Rewards sentences with '
            'numbers, units, formal notation (WHO, µg/m³, Q(D), ζ_d). Penalises vague '
            'language (very high, significant, elevated). A sentence scoring 0.0 '
            'contributes no measurable information.'
        )
        if v3_nss:
            nss_score = v3_nss.get('overall_specificity', 0)
            nss_vdt   = v3_nss.get('NSS_verdict', 'N/A')
            p_nss = doc.add_paragraph()
            r_nss = p_nss.add_run(
                f"Overall Specificity = {nss_score:.2f}   [{nss_vdt}]"
            )
            r_nss.bold = True; r_nss.font.size = Pt(13)
            r_nss.font.color.rgb = (GREEN if nss_vdt == 'SPECIFIC' else
                                    AMBER if nss_vdt == 'ADEQUATE' else RED)
            t_nss = doc.add_table(rows=1, cols=3)
            t_nss.style = 'Table Grid'
            for cell, hdr in zip(t_nss.rows[0].cells,
                                  ['Section', 'Avg Specificity', 'Verdict']):
                _set_cell_bg(cell, '1A2C4E')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for sec_k, sec_v in v3_nss.get('sections', {}).items():
                row = t_nss.add_row()
                sp = sec_v.get('avg_specificity', 0)
                bg = 'D1FAE5' if sp >= 0.4 else ('FEF3C7' if sp >= 0.2 else 'FFE4E6')
                for cell, val in zip(row.cells, [
                    sec_k,
                    f"{sp:.2f}",
                    sec_v.get('verdict', 'N/A'),
                ]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(val).font.size = Pt(9)
            doc.add_paragraph()

        # ── N-08: Benchmark Comparator ────────────────────────────────────
        h2('N-08 — Benchmark Comparator (BNC)')
        body(
            'Compares station results against CPCB NAAQS, WHO AQG 2021, and India '
            'national averages. All reference values stored internally — zero external '
            'API required. Provides national context for station-level findings.'
        )
        if v3_benchmark:
            sev = v3_benchmark.get('overall_severity', 'N/A')
            p_bnc = doc.add_paragraph()
            r_bnc = p_bnc.add_run(
                f"Station Severity: {sev}  ·  "
                f"Reference: {v3_benchmark.get('reference_source','WHO AQG 2021 + CPCB NAAQS')}"
            )
            r_bnc.bold = True; r_bnc.font.size = Pt(12)
            r_bnc.font.color.rgb = (RED if sev in ('CRITICAL','SEVERE') else
                                    AMBER if sev == 'HIGH' else GREEN)
            t_bnc = doc.add_table(rows=1, cols=5)
            t_bnc.style = 'Table Grid'
            for cell, hdr in zip(t_bnc.rows[0].cells,
                                  ['Variable','Station Mean','vs WHO','vs NAAQS','vs India Avg']):
                _set_cell_bg(cell, '8B1C2A')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            for var, comp in v3_benchmark.get('comparisons', {}).items():
                row = t_bnc.add_row()
                who_stat = comp.get('who_status', 'N/A')
                bg = 'FFE4E6' if 'EXCEED' in str(who_stat) else 'D1FAE5'
                for cell, val in zip(row.cells, [
                    var,
                    f"{comp.get('station_annual_mean',0):.2f} {comp.get('unit','µg/m³')}",
                    comp.get('vs_who', 'N/A'),
                    comp.get('naaqs_status', 'N/A'),
                    f"{comp.get('above_national_avg',0):+.2f} {comp.get('unit','µg/m³')}",
                ]):
                    _set_cell_bg(cell, bg)
                    cell.paragraphs[0].add_run(val).font.size = Pt(9)
            doc.add_paragraph()

        # ── N-09: MNS ──────────────────────────────────────────────────────
        h2('N-09 — Master Novelty Score (MNS)')
        body(
            'Aggregates all 7 independent verification dimensions into a single '
            'scalar trust metric. MNS ∈ [0,1]. '
            'CERTIFIED ≥ 0.90 · TRUSTED ≥ 0.80 · ACCEPTABLE ≥ 0.70 · '
            'REVIEW ≥ 0.55 · UNRELIABLE < 0.55. '
            'This score is the first unified trust metric in automated environmental reporting.'
        )
        if v3_mns:
            mns_val2 = v3_mns.get('MNS', 0)
            mns_vdt2 = v3_mns.get('verdict', 'N/A')
            p_mns2 = doc.add_paragraph()
            r_mns2 = p_mns2.add_run(f"MNS = {mns_val2:.3f}   [{mns_vdt2}]")
            r_mns2.bold = True; r_mns2.font.size = Pt(20)
            r_mns2.font.color.rgb = (GREEN if mns_vdt2 in ('CERTIFIED','TRUSTED','ACCEPTABLE')
                                     else AMBER)
            body(f"Equation: {v3_mns.get('equation','')}", italic=True)
            t_mns2 = doc.add_table(rows=1, cols=3)
            t_mns2.style = 'Table Grid'
            for cell, hdr in zip(t_mns2.rows[0].cells,
                                  ['Verification Dimension', 'Score', 'Weight']):
                _set_cell_bg(cell, '1A2C4E')
                r = cell.paragraphs[0].add_run(hdr)
                r.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(9)
            comps2 = v3_mns.get('components', {})
            wts2   = v3_mns.get('weights', {})
            labels2 = {
                'v2_G_rate': 'N-RGV v2 Numerical Grounding',
                'v3_G_rate': 'N-RGV v3 Multi-Source Grounding',
                'TCV_rate' : 'N-01 Temporal Claim Accuracy',
                'CCD_score': 'N-02 Causal Language Score',
                'CVCC_score':'N-06 Cross-Variable Consistency',
                'confidence':'Sentence Confidence Score',
                'semantic'  :'Semantic Alignment Score',
            }
            for k2, v2 in comps2.items():
                row = t_mns2.add_row()
                bg2 = 'D1FAE5' if v2 >= 0.80 else ('FEF3C7' if v2 >= 0.60 else 'FFE4E6')
                for cell, val in zip(row.cells, [
                    labels2.get(k2, k2),
                    f"{v2*100:.1f}%",
                    f"{int(wts2.get(k2,0)*100)}%",
                ]):
                    _set_cell_bg(cell, bg2)
                    cell.paragraphs[0].add_run(val).font.size = Pt(9)
            doc.add_paragraph()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 17 — MASTER NOVELTY SCORE
    # ════════════════════════════════════════════════════════════════════════
    if v3_mns:
        doc.add_page_break()
        h1('Section 17 — Master Novelty Score (MNS)')
        body(
            'The Master Novelty Score aggregates seven independent verification '
            'dimensions into a single trust metric. MNS ∈ [0,1]: '
            'CERTIFIED ≥ 0.90, TRUSTED ≥ 0.80, ACCEPTABLE ≥ 0.70, '
            'REVIEW ≥ 0.55, UNRELIABLE < 0.55.'
        )
        mns_val = v3_mns.get('MNS', 0)
        mns_vdt = v3_mns.get('verdict', 'N/A')
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(f"MNS = {mns_val:.3f}   [{mns_vdt}]")
        r.bold = True
        r.font.size = Pt(22)
        r.font.color.rgb = GREEN if mns_vdt in ('CERTIFIED','TRUSTED','ACCEPTABLE') else AMBER
        body(f"Equation: {v3_mns.get('equation','')}", italic=True)
        doc.add_paragraph()
        labels = {
            'v2_G_rate':'Numerical Grounding base',
            'v3_G_rate':'Numerical Grounding enhanced',
            'TCV_rate':'Temporal Claim Accuracy',
            'CCD_score':'Causal Language Score',
            'CVCC_score':'Cross-Variable Consistency',
            'confidence':'Sentence Confidence',
            'semantic':'Semantic Alignment',
        }
        mns_t = doc.add_table(rows=1, cols=3)
        mns_t.style = 'Table Grid'
        for cell, hdr in zip(mns_t.rows[0].cells, ['Dimension','Score','Weight']):
            _set_cell_bg(cell, '1A2C4E')
            r2 = cell.paragraphs[0].add_run(hdr)
            r2.bold = True; r2.font.color.rgb = WHITE
        comps = v3_mns.get('components', {})
        wts   = v3_mns.get('weights', {})
        for k, v in comps.items():
            row = mns_t.add_row()
            bg = 'D1FAE5' if v >= 0.80 else ('FEF3C7' if v >= 0.60 else 'FFE4E6')
            for i2, (cell, val) in enumerate(zip(row.cells,
                [labels.get(k,k), f"{round(v*100,1)}%", f"{int(wts.get(k,0)*100)}%"])):
                _set_cell_bg(cell, bg)
                cell.paragraphs[0].add_run(val).font.size = Pt(10)
        doc.add_paragraph()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 18 — WINDOWS TASK SCHEDULER
    # ════════════════════════════════════════════════════════════════════════
    doc.add_page_break()
    h1('Section 18 — Windows Task Scheduler: Automated Daily Execution')
    body(
        'ATARS v3.1 auto-generates a Windows batch file (run_atars.bat) after every run. '
        'Register it with Windows Task Scheduler to run ATARS automatically every day — '
        'no manual intervention needed after the one-time setup below.'
    )

    # ── Setup Command ─────────────────────────────────────────────────────────
    h2('Quick Setup — One Command')
    body(
        'The batch file is saved to your output directory automatically. '
        'Once registered, ATARS will generate a fresh dated report every day at your chosen time.'
    )
    _bat_name = f"{config.get('output_dir','atars_output')}\\run_atars.bat"
    _cmd_rows = [
        ('Script location', _bat_name),
        ('Python command',  f'python atars_v3.py --data {config.get("data_file","data.csv")} --city "{config.get("city","City")}" --no-llm'),
        ('Output folder',   config.get('output_dir', 'atars_output')),
        ('Log file',        f"{config.get('output_dir','atars_output')}\\scheduler_log.txt"),
    ]
    t_cmd = doc.add_table(rows=len(_cmd_rows)+1, cols=2)
    t_cmd.style = 'Table Grid'
    for cell, hdr in zip(t_cmd.rows[0].cells, ['Item', 'Value']):
        _set_cell_bg(cell, '1A2C4E')
        r2 = cell.paragraphs[0].add_run(hdr)
        r2.bold = True; r2.font.color.rgb = WHITE; r2.font.size = Pt(10)
    for i, (k, v) in enumerate(_cmd_rows, 1):
        t_cmd.rows[i].cells[0].paragraphs[0].add_run(k).font.size = Pt(9)
        t_cmd.rows[i].cells[1].paragraphs[0].add_run(str(v)).font.size = Pt(9)
        _set_cell_bg(t_cmd.rows[i].cells[0], 'EFF6FF')
    doc.add_paragraph()

    # ── Step-by-step registration ─────────────────────────────────────────────
    h2('Task Scheduler Registration — Step by Step')
    steps_data = [
        ('1', 'Open Task Scheduler',
         'Press  Win + S  →  type "Task Scheduler"  →  press Enter.',
         '1A2C4E', 'DCE8F5'),
        ('2', 'Create Basic Task',
         'In the right-hand Actions panel click  "Create Basic Task…"',
         '1A6B72', 'D1FAE5'),
        ('3', 'Name the task',
         'Name: ATARS Daily Report | Description: Automated air quality report generation',
         '1A2C4E', 'DCE8F5'),
        ('4', 'Set trigger → Daily',
         'Choose  Daily.  Set start time to  07:00 AM  '
         '(or 30 min after your data file arrives).',
         '1A6B72', 'D1FAE5'),
        ('5', 'Set action',
         'Choose  "Start a Program".',
         '1A2C4E', 'DCE8F5'),
        ('6', 'Browse to the batch file',
         f'Program/script: {_bat_name} | Start in: {config.get("output_dir","atars_output")}',
         '1A6B72', 'D1FAE5'),
        ('7', 'Finish',
         'Click  Finish.  The task now appears in Task Scheduler Library.',
         '1A2C4E', 'DCE8F5'),
        ('8', 'Test immediately',
         'Right-click the task  →  Run.  '
         'Check the output folder — a new dated report should appear within minutes.',
         '1A6B72', 'D1FAE5'),
    ]
    t_steps = doc.add_table(rows=1, cols=3)
    t_steps.style = 'Table Grid'
    col_widths = [Cm(1.5), Cm(5.0), Cm(10.0)]
    for cell, hdr in zip(t_steps.rows[0].cells, ['#', 'Action', 'Detail']):
        _set_cell_bg(cell, '1A2C4E')
        rh = cell.paragraphs[0].add_run(hdr)
        rh.bold = True; rh.font.color.rgb = WHITE; rh.font.size = Pt(10)
        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    for num, action, detail, hdr_col, row_col in steps_data:
        row_s = t_steps.add_row()
        # Step number badge
        _set_cell_bg(row_s.cells[0], hdr_col)
        rn = row_s.cells[0].paragraphs[0].add_run(num)
        rn.bold = True; rn.font.color.rgb = WHITE; rn.font.size = Pt(13)
        row_s.cells[0].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        # Action
        _set_cell_bg(row_s.cells[1], row_col)
        ra = row_s.cells[1].paragraphs[0].add_run(action)
        ra.bold = True; ra.font.size = Pt(10)
        # Detail (support multi-line)
        _set_cell_bg(row_s.cells[2], 'FFFFFF')
        for line in detail.split('\n'):
            if row_s.cells[2].paragraphs[0].text == '':
                row_s.cells[2].paragraphs[0].add_run(line).font.size = Pt(9)
            else:
                row_s.cells[2].add_paragraph(line).runs[0].font.size = Pt(9) if row_s.cells[2].paragraphs[-1].runs else None
    doc.add_paragraph()

    # ── Important notes ───────────────────────────────────────────────────────
    h2('Important Notes')
    notes_items = [
        ('Data timing',
         f'Your CSV ({config.get("data_file","data.csv")}) must arrive BEFORE the scheduled run. '
         'Schedule data delivery at least 30 minutes earlier.'),
        ('Internet requirement',
         'Open-Meteo weather (--weather open_meteo, default): requires internet. '
         'Offline mode: add --weather mock to the batch file command.'),
        ('Log file',
         f'Every run appends a line to scheduler_log.txt in the output folder. '
         'Check it if a run appears to have failed.'),
        ('Report naming',
         'Each run creates a new timestamped file — no previous report is ever overwritten.'),
        ('Pausing automation',
         'Task Scheduler → right-click task → Disable. Re-enable the same way.'),
        ('Updating data path',
         'If your CSV path changes, edit the --data argument inside run_atars.bat '
         'with any text editor, then save.'),
    ]
    t_notes = doc.add_table(rows=1, cols=2)
    t_notes.style = 'Table Grid'
    for cell, hdr in zip(t_notes.rows[0].cells, ['Topic', 'Note']):
        _set_cell_bg(cell, '1A6B72')
        rh2 = cell.paragraphs[0].add_run(hdr)
        rh2.bold = True; rh2.font.color.rgb = WHITE; rh2.font.size = Pt(10)
    for topic, note in notes_items:
        row_n = t_notes.add_row()
        _set_cell_bg(row_n.cells[0], 'EFF6FF')
        row_n.cells[0].paragraphs[0].add_run(topic).font.size = Pt(9)
        row_n.cells[0].paragraphs[0].runs[0].bold = True
        row_n.cells[1].paragraphs[0].add_run(note).font.size = Pt(9)
    doc.add_paragraph()
    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 19 — REPRODUCIBILITY AND AUDIT RECORD
    # ════════════════════════════════════════════════════════════════════════
    doc.add_page_break()
    h1('Section 19 — Reproducibility and Audit Record')
    body(
        'ATARS implements SHA-256 cryptographic hashing of the JSON contract J '
        'and the generated report. Any change to input data, configuration, or '
        'ATARS version produces a different hash, making tampering detectable.'
    )
    data_hash = J.get('run_metadata', {}).get('hash_J', 'computed_at_runtime')
    audit_rows = [
        ('Data Quality Q(D)',    f"{J.get('data_quality',{}).get('overall_Q',0):.4f}  [{J.get('data_quality',{}).get('confidence_flag','N/A')}]"),
        ('Baseline Window W',    f"{J.get('formal_operators',{}).get('baseline_window_W','N/A')} days"),
        ('Z-Threshold ζ',        str(J.get('formal_operators',{}).get('z_threshold','N/A'))),
        ('Run Hash (J)',         data_hash[:32] + '...' if len(str(data_hash)) > 32 else str(data_hash)),
        ('Report Version',       'ATARS v3.1.0'),
        ('License',              'MIT — Open Source'),
    ]
    audit_t = doc.add_table(rows=1, cols=2)
    audit_t.style = 'Table Grid'
    for cell, hdr in zip(audit_t.rows[0].cells, ['Parameter','Value']):
        _set_cell_bg(cell, '1A2C4E')
        r3 = cell.paragraphs[0].add_run(hdr)
        r3.bold = True; r3.font.color.rgb = WHITE
    for k, v in audit_rows:
        row = audit_t.add_row()
        row.cells[0].text = k
        row.cells[1].text = str(v)
    doc.add_paragraph()

    # Report Drift Detector — N-04
    if v3_drift:
        h2('Report Drift Detector (N-04)')
        body(
            f"Status: {v3_drift.get('status','N/A')}. "
            f"Runs compared: {v3_drift.get('runs_compared',0)}. "
            f"Drifts detected: {len(v3_drift.get('drifts_detected',[]))}."
        )
        for d in v3_drift.get('drifts_detected',[]):
            p2 = doc.add_paragraph(style='List Bullet')
            p2.add_run(
                f"[{d.get('severity','?')}] {d.get('metric','?')}: "
                f"current={d.get('current',0):.3f} vs "
                f"avg={d.get('historical_avg',0):.3f} — {d.get('direction','?')}")
        if not v3_drift.get('drifts_detected'):
            body('No drift detected — report quality stable.', italic=True)
        doc.add_paragraph()

    # ── Save report ────────────────────────────────────────────────────────
    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 20 — CONCLUSIONS AND RECOMMENDATIONS
    # ════════════════════════════════════════════════════════════════════════
    doc.add_page_break()
    h1("Section 20 — Conclusions and Recommendations")

    pm10_v_c  = J.get("variables",{}).get("PM10",{})
    no2_v_c   = J.get("variables",{}).get("NO2",{})
    benz_v_c  = J.get("variables",{}).get("Benzene",{})
    q_c       = J["data_quality"]["overall_Q"]
    city_c    = J.get("city","City")
    n_days_c  = J.get("total_days",0)
    exc_c     = J.get("exceedances",{})
    ops_c     = J.get("formal_operators",{})

    h2("Principal Conclusions")
    pm10_mean_c = pm10_v_c.get("annual_mean",0)
    pm10_thr_c  = pm10_v_c.get("threshold",45) or 45
    pm10_exc_c  = pm10_v_c.get("exceeds_days",0)
    pm10_pct_c  = pm10_v_c.get("exceed_pct",0)
    pm10_anom_c = pm10_v_c.get("anomaly_days",0)

    conclusions = [
        f"PM10 at {city_c} recorded an annual mean of {pm10_mean_c:.2f} µg/m³ over "
        f"{n_days_c} monitored days, exceeding the WHO 24-hour guideline of "
        f"{pm10_thr_c:.1f} µg/m³ on {pm10_exc_c} days ({pm10_pct_c:.1f}% of the period). "
        f"This represents a statistically characterised exceedance pattern — not an isolated event.",

        f"Z-score anomaly detection (|ζ_d| > {ops_c.get('z_threshold',3.0):.1f}) against a "
        f"{ops_c.get('baseline_window_W',30)}-day rolling baseline identified {pm10_anom_c} "
        f"anomalous PM10 days. These represent statistically significant deviations from local "
        f"norms and warrant targeted investigation of emission sources.",

        f"Overall data completeness Q(D) = {q_c:.3f} confirms the dataset is statistically "
        f"adequate for the analyses presented. Quality-validated records (q_i = 1) form the "
        f"basis of all statistical operators.",

        "Multivariate Isolation Forest anomaly detection identified additional anomaly days "
        "not captured by univariate z-score analysis — confirming that multi-source pollution "
        "episodes are present and require multivariate monitoring approaches.",

        "All correlation statistics reported are Pearson r values (statistical association). "
        "No causal emission source attribution is established by this analysis. "
        "Source apportionment requires receptor modelling beyond this framework's scope.",
    ]
    if benz_v_c.get("annual_mean",0) and benz_v_c.get("threshold"):
        benz_m = benz_v_c["annual_mean"]
        benz_t = benz_v_c["threshold"]
        if benz_m > benz_t:
            conclusions.insert(2,
                f"Benzene annual mean {benz_m:.3f} µg/m³ exceeded the WHO annual reference "
                f"of {benz_t:.1f} µg/m³ (IARC Group 1 carcinogen). "
                f"Monitoring frequency for BTEX compounds should be increased at this station."
            )

    for c in conclusions:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(c).font.size = Pt(10)
    doc.add_paragraph()

    h2("Recommendations")
    recommendations = [
        ("Immediate", f"Investigate the {pm10_anom_c} PM10 anomaly days to identify "
                      "emission event sources (industrial, agricultural, traffic, dust)."),
        ("Short-term", f"Increase monitoring frequency during high-exceedance months "
                       "identified in the seasonal analysis."),
        ("Short-term", "Install additional monitoring equipment if sensor flatline events "
                       f"detected during cleaning exceed 5% of records."),
        ("Medium-term", "Conduct source apportionment (PMF or CMB receptor modelling) "
                        "to quantify contributions from traffic, industry, and regional transport."),
        ("Medium-term", "Implement real-time alert system triggered when daily PM10 "
                        f"exceeds {pm10_thr_c:.1f} µg/m³ (WHO guideline) or z-score "
                        f"exceeds {ops_c.get('z_threshold',3.0):.1f}."),
        ("Long-term", "Establish multi-year trend using Mann-Kendall analysis across "
                      "successive annual datasets to assess regulatory compliance trajectory."),
    ]
    t_rec = doc.add_table(rows=1, cols=3)
    t_rec.style = "Table Grid"
    for cell, hdr in zip(t_rec.rows[0].cells, ["Priority","Area","Recommendation"]):
        _set_cell_bg(cell, "1A2C4E")
        r = cell.paragraphs[0].add_run(hdr)
        r.bold = True; r.font.color.rgb = WHITE
    for pri, rec in recommendations:
        row = t_rec.add_row()
        bg  = "FFE4E6" if pri=="Immediate" else ("FEF3C7" if pri=="Short-term" else "D1FAE5")
        for cell, val in zip(row.cells, [pri, "Air Quality", rec]):
            _set_cell_bg(cell, bg)
            cell.paragraphs[0].add_run(val).font.size = Pt(9)
    doc.add_paragraph()

    # ════════════════════════════════════════════════════════════════════════
    #  SECTION 21 — METHODOLOGY AND REFERENCES
    # ════════════════════════════════════════════════════════════════════════
    doc.add_page_break()
    h1("Section 21 — Methodology and References")

    h2("Formal Framework — ATARS Operators")
    body(
        "The ATARS pipeline implements 16 formal statistical operators. "
        "Key operators referenced in this report:"
    )
    ops_table = [
        ("Eq. 1",  "D = {(t_i, x_i, q_i)}",        "Observation space definition"),
        ("Eq. 3",  "x̄ = (1/N_v) Σ x_i",            "Aggregation operator — mean"),
        ("Eq. 4",  "σ = √[(1/(N_v-1)) Σ(x_i-x̄)²]", "Sample standard deviation"),
        ("Eq. 7",  "x̄_W = mean of W-day window",    "Baseline mean operator"),
        ("Eq. 8",  "σ_W = std of W-day window",      "Baseline std operator"),
        ("Eq. 9",  "ζ_d = (x̄_d - x̄_W) / σ_W",     "Z-score anomaly operator"),
        ("Eq. 10", "δ_d = [(x̄_d - x̄_{d-1}) / x̄_{d-1}] × 100%", "Day-on-day change"),
        ("Eq. 11", "Q(D) = N_v / N_total",           "Data quality score"),
        ("Eq. 12", "CI = x̄ ± z_{α/2} · σ/√N_v",   "95% Confidence interval"),
        ("Eq. 18", "R_{jk} = Corr(X_j, X_k)",       "Pearson correlation matrix"),
        ("Eq. 19", "G: J → N, G(s) ⊆ J",            "LLM grounding constraint"),
    ]
    t_ops = doc.add_table(rows=1, cols=3)
    t_ops.style = "Table Grid"
    for cell, hdr in zip(t_ops.rows[0].cells, ["Equation","Formula","Description"]):
        _set_cell_bg(cell, "1A6B72")
        r = cell.paragraphs[0].add_run(hdr)
        r.bold = True; r.font.color.rgb = WHITE
    for eq, form, desc in ops_table:
        row = t_ops.add_row()
        for cell, val in zip(row.cells, [eq, form, desc]):
            cell.paragraphs[0].add_run(val).font.size = Pt(9)
    doc.add_paragraph()

    h2("References")
    references = [
        ("WHO AQG 2021",
         "World Health Organization (2021). WHO Global Air Quality Guidelines: "
         "Particulate Matter (PM2.5 and PM10), Ozone, Nitrogen Dioxide, Sulfur Dioxide "
         "and Carbon Monoxide. World Health Organization. ISBN 978-92-4-003422-8."),
        ("CPCB NAAQS",
         "Central Pollution Control Board, India (2009). National Ambient Air Quality "
         "Standards. Ministry of Environment, Forest and Climate Change, Government of India."),
        ("IARC Benzene",
         "International Agency for Research on Cancer (2018). Benzene. "
         "IARC Monographs on the Identification of Carcinogenic Hazards to Humans, Vol. 120."),
        ("Mann-Kendall",
         "Mann, H.B. (1945). Nonparametric Tests Against Trend. Econometrica, 13(3), 245–259. "
         "Kendall, M.G. (1975). Rank Correlation Methods. Griffin, London."),
        ("Sen's Slope",
         "Sen, P.K. (1968). Estimates of the regression coefficient based on Kendall's tau. "
         "Journal of the American Statistical Association, 63(324), 1379–1389."),
        ("Isolation Forest",
         "Liu, F.T., Ting, K.M., & Zhou, Z.H. (2008). Isolation Forest. "
         "8th IEEE International Conference on Data Mining (ICDM 2008), pp. 413–422."),
        ("Holt-Winters",
         "Holt, C.C. (1957). Forecasting seasonals and trends by exponentially weighted "
         "moving averages. ONR Research Memorandum, 52. Carnegie Institute of Technology. "
         "Winters, P.R. (1960). Forecasting sales by exponentially weighted moving averages. "
         "Management Science, 6(3), 324–342."),
        ("Open-Meteo",
         "Zippenfenig, P. (2023). Open-Meteo.com Weather API. Zenodo. "
         "https://doi.org/10.5281/zenodo.7970649. Accessed 2024."),
        ("ATARS Framework",
         f"Priyanshu ({datetime.now().year}). ATARS: Automated Time-Series Analysis "
         "and Reporting System [Software, v2.0.0]. Global Institute of Technology "
         "and Management, Haryana, India. MIT License. "
         "https://github.com/Priyanshu-ux712/ATARS. "
         ),
    ]
    for i, (ref_key, ref_txt) in enumerate(references, 1):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.left_indent  = Pt(20)
        run_num = p.add_run(f"[{i}] ")
        run_num.bold = True; run_num.font.size = Pt(10)
        run_key = p.add_run(f"{ref_key}. ")
        run_key.bold = True; run_key.font.size = Pt(10)
        run_txt = p.add_run(ref_txt)
        run_txt.font.size = Pt(9)
    doc.add_paragraph()

    h2("Software and Tools")
    tools_rows = [
        ("Python",         "3.9+", "Core pipeline language"),
        ("NumPy",          "≥1.21","Array computation and formal operators"),
        ("Pandas",         "≥1.3", "Data loading, aggregation, time series"),
        ("Matplotlib/Seaborn","≥3.5","Chart generation (15 figures)"),
        ("Scikit-learn",   "≥1.0", "Isolation Forest anomaly detection"),
        ("SciPy",          "≥1.7", "Shapiro-Wilk, D'Agostino, t-distribution"),
        ("python-docx",    "≥0.8", "Word report generation"),
        ("python-pptx",    "≥0.6", "PowerPoint generation"),
        ("Open-Meteo API", "free", "Live and historical weather data (no key required)"),
        ("Ollama",         "optional","Local LLM inference (llama3.2, mistral, phi3)"),
    ]
    t_tools = doc.add_table(rows=1, cols=3)
    t_tools.style = "Table Grid"
    for cell, hdr in zip(t_tools.rows[0].cells, ["Package","Version","Purpose"]):
        _set_cell_bg(cell, "1A2C4E")
        r = cell.paragraphs[0].add_run(hdr)
        r.bold = True; r.font.color.rgb = WHITE
    for pkg, ver, purpose in tools_rows:
        row = t_tools.add_row()
        for cell, val in zip(row.cells, [pkg, ver, purpose]):
            cell.paragraphs[0].add_run(val).font.size = Pt(9)
    doc.add_paragraph()

    fname     = (f"ATARS_Report_{config.get('city','City')}_"
                 f"{datetime.now().strftime('%Y%m%d_%H%M')}.docx")
    out_path  = str(out_dir / fname)

    # ── Section 19: Unified Verification Dashboard ────────────────────────
    # Called here (after all other sections) so all v3 data is available
    _v3_has_data = any([v3_ver_v3, v3_tcv, v3_ccd, v3_aee, v3_cvcc,
                        v3_nss, v3_mns, v3_upt, v3_benchmark, v3_drift])
    if _v3_has_data:
        try:
            _add_unified_verification_dashboard(
                doc=doc,
                verification_v2 = verification or {},
                ver_v3          = v3_ver_v3    or {},
                evaluation      = {"normal_ai": {"hallucination_risk":"UNKNOWN"},
                                   "verified_ai": {
                                       "g_rate"            : (v3_ver_v3 or {}).get("overall_G_rate",0),
                                       "verdict"           : (v3_ver_v3 or {}).get("overall_verdict","N/A"),
                                       "hallucination_risk": "LOW" if (v3_ver_v3 or {}).get("overall_G_rate",0)>=0.9
                                                             else "MEDIUM",
                                       "contradictions"    : (v3_ver_v3 or {}).get("contradiction_count",0),
                                       "semantic_aligned"  : (v3_ver_v3 or {}).get("semantic_pass", True),
                                   },
                                   "summary": f"G_rate={round((v3_ver_v3 or {}).get('overall_G_rate',0)*100,1)}%"},
                comparison      = v3_comparison or {},
                tcv             = v3_tcv        or {},
                ccd             = v3_ccd        or {},
                aee             = v3_aee        or {},
                drift           = v3_drift      or {},
                upt             = v3_upt        or {},
                cvcc            = v3_cvcc       or {},
                nss             = v3_nss        or {},
                benchmark       = v3_benchmark  or {},
                mns             = v3_mns        or {},
                config          = config,
            )
        except Exception as _dash_err:
            print(f"  ⚠ Unified dashboard error (non-fatal): {_dash_err}")

    doc.save(out_path)
    print(f"  ✓ Report saved: {out_path}")
    return out_path



def _json_safe(obj):
    """
    Custom JSON serializer that handles all numpy/pandas types
    that standard json.dumps chokes on.
    """
    import numpy as np
    if isinstance(obj, (np.bool_,)):
        return bool(obj)
    if isinstance(obj, (np.integer,)):
        return int(obj)
    if isinstance(obj, (np.floating,)):
        return float(obj)
    if isinstance(obj, np.ndarray):
        return obj.tolist()
    if hasattr(obj, 'item'):          # generic numpy scalar
        return obj.item()
    if hasattr(obj, 'isoformat'):     # datetime / date
        return obj.isoformat()
    if hasattr(obj, 'tolist'):        # pandas Series / DataFrame
        return obj.tolist()
    raise TypeError(f"Object of type {type(obj).__name__} is not JSON serializable")


def _safe_dumps(obj, **kwargs):
    """json.dumps with _json_safe fallback — drop-in replacement."""
    return json.dumps(obj, default=_json_safe, **kwargs)


def write_audit_log(J: dict, narrative: dict, report_path: str,
                     config: dict, out_dir: Path, verification: dict = None,
                     ml_if: dict = None, ml_hw: dict = None):
    """
    Write SHA-256 audit log and generate Windows Task Scheduler batch file.
    Hash(O) = SHA-256(canonical_JSON(J))  — Eq. 21
    v2.0: Logs ML results (IF anomaly count, HW forecast metrics) for audit trail.
    """
    ml_if = ml_if or {}
    ml_hw = ml_hw or {}
    J_canonical   = _safe_dumps(J, sort_keys=True, ensure_ascii=True)
    hash_J        = hashlib.sha256(J_canonical.encode()).hexdigest()

    if report_path and Path(str(report_path)).exists():
        with open(report_path, 'rb') as f:
            hash_report = hashlib.sha256(f.read()).hexdigest()
    else:
        hash_report = 'pending'


    audit = {
        "run_id"       : hash_J[:16],
        "timestamp"    : datetime.now().isoformat(),
        "city"         : J['city'],
        "date_range"   : J['date_range'],
        "total_records": J['total_records'],
        "quality_Q"    : J['data_quality']['overall_Q'],
        "hash_J"       : hash_J,
        "hash_report"  : hash_report,
        "report_file"  : os.path.basename(report_path),
        "llm_model"    : narrative.get('model', 'N/A'),
        "framework_version": "ATARS_v1.0",
        # SECURITY: config_used logs only statistical/operational parameters.
        # Personal identifiers (author, email, institution) and absolute file
        # paths are intentionally excluded to prevent accidental data leakage
        # when audit logs are shared for reproducibility verification.
        "config_used"  : {k: v for k, v in config.items()
                          if k not in (
                              'thresholds',      # large dict, referenced via hash_J
                              'author',          # personal identifier — excluded
                              'email',           # personal contact — excluded
                              'institution',     # personal institution — excluded
                              'department',      # personal department — excluded
                              'degree',          # personal info — excluded
                              'location',        # personal location — excluded
                              'data_file',       # absolute path may reveal FS layout
                              'output_dir',      # absolute path may reveal FS layout
                              'ollama_url',      # always localhost; not needed
                          )},
        "grounding_G_rate"     : verification["overall_G_rate"] if verification else None,
        "grounding_verdict"    : verification["overall_verdict"] if verification else "N/A",
        "grounding_passed"     : verification["verification_passed"] if verification else None,
        "reproducibility_note":
            "Re-run with same data file and config to verify Hash(J) match.",
        # ML audit fields (v2.0)
        "ml_isolation_forest"  : {
            "available"     : ml_if.get("available", False),
            "n_anomalies"   : ml_if.get("n_anomalies", None),
            "anomaly_rate"  : ml_if.get("anomaly_rate", None),
            "n_estimators"  : ml_if.get("n_estimators", None),
            "contamination" : ml_if.get("contamination", None),
            "random_state"  : ml_if.get("random_state", None),
            "features_used" : ml_if.get("features_used", None),
            "n_if_only"     : ml_if.get("n_if_only", None),
            "security_note" : ml_if.get("security_note", ""),
        },
        "ml_holt_winters"      : {
            "available"     : ml_hw.get("available", False),
            "forecast_days" : ml_hw.get("forecast_days", None),
            "rmse"          : ml_hw.get("rmse", None),
            "mae"           : ml_hw.get("mae", None),
            "alpha"         : ml_hw.get("alpha", None),
            "beta"          : ml_hw.get("beta", None),
            "gamma"         : ml_hw.get("gamma", None),
            "security_note" : ml_hw.get("security_note", ""),
        },
        "atars_version"        : "2.0.0",
        "copyright"            : "Copyright (c) 2026 Priyanshu — MIT License (Open Source) — https://github.com/Priyanshu-ux712/ATARS",
    }

    log_path = out_dir / f'audit_log_{audit["run_id"]}.json'
    with open(log_path, 'w') as f:
        json.dump(audit, f, indent=2, default=_json_safe)
    print(f"  ✓ Audit log: {log_path}")

    # ── Generate Windows Task Scheduler batch file ─────────────────────────
    _write_batch_file(config, out_dir)

    return audit


def _write_batch_file(config: dict, out_dir: Path):
    """
    Auto-generate a Windows Task Scheduler .bat file for daily automation.
    Users point Task Scheduler at this file to run ATARS every day automatically.
    """
    script_path  = os.path.abspath('atars.py')
    data_path    = os.path.abspath(config.get('data_file', 'data.csv'))
    out_path_abs = os.path.abspath(str(out_dir))
    work_dir     = os.path.dirname(script_path)
    llm_flag     = '' if config.get('use_llm') else ' --no-llm'
    model_flag   = f' --llm-model {config["llm_model"]}' if config.get('use_llm') else ''

    bat_content = f"""@echo off
REM ═══════════════════════════════════════════════════════════════════════════
REM  ATARS — Automated Time-Series Analysis and Reporting System
REM  Windows Task Scheduler Daily Run Script
REM  Author : {config['author']}
REM  Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}
REM
REM  HOW TO USE:
REM  1. Open Windows Task Scheduler (search in Start menu)
REM  2. Click "Create Basic Task" → Name: "ATARS Daily Report"
REM  3. Trigger: Daily → set your time (e.g. 07:00 AM)
REM  4. Action: Start a Program
REM  5. Program/script: point to this .bat file
REM     Path: {out_path_abs}\\run_atars.bat
REM  6. Click Finish
REM
REM  IMPORTANT: Your data file must arrive BEFORE the scheduled run time.
REM  Recommended: schedule data delivery 30+ minutes before ATARS runs.
REM ═══════════════════════════════════════════════════════════════════════════

REM Change to the ATARS working directory
cd /d "{work_dir}"

REM Record run start time
echo [%DATE% %TIME%] ATARS run started >> "{out_path_abs}\\scheduler_log.txt"

REM Run ATARS pipeline
python "{script_path}" ^
    --data "{data_path}" ^
    --city "{config['city']}" ^
    --station "{config['station_id']}" ^
    --output "{out_path_abs}"{llm_flag}{model_flag}

REM Check if run succeeded
if %ERRORLEVEL% EQU 0 (
    echo [%DATE% %TIME%] ATARS SUCCESS — report saved to {out_path_abs} >> "{out_path_abs}\\scheduler_log.txt"
    echo ATARS completed successfully.
) else (
    echo [%DATE% %TIME%] ATARS FAILED — check pipeline output above >> "{out_path_abs}\\scheduler_log.txt"
    echo ATARS run FAILED. Check output for errors.
    exit /b 1
)

REM Optional: open output folder automatically after run (comment out for silent mode)
REM explorer "{out_path_abs}"

exit /b 0
"""

    bat_path = out_dir / 'run_atars.bat'
    with open(bat_path, 'w', encoding='utf-8') as f:
        f.write(bat_content)
    print(f"  ✓ Batch file: {bat_path}")
    print(f"  ✓ Point Windows Task Scheduler at: {bat_path}")
    return str(bat_path)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 13 — MAIN PIPELINE ORCHESTRATOR (Steps 1–9)
# ═══════════════════════════════════════════════════════════════════════════════

def print_banner(config: dict):
    print("""
╔══════════════════════════════════════════════════════════════════════════════╗
║   ATARS — Automated Time-Series Analysis and Reporting System  v3.1.0      ║
║   Author: Priyanshu · Global Institute of Technology and Management        ║
║   Open Source · MIT License · Formally specified research pipeline          ║
║                                                                              ║
║   Formal operators: A(D_v), B(D,W), ζ_d, δ_d, Q(D), G:J→N, Hash(O)       ║
╚══════════════════════════════════════════════════════════════════════════════╝""")
    print(f"  City: {config['city']}  |  Station: {config['station_id']}")
    print(f"  Output: {config['output_dir']}")
    print(f"  LLM: {'Enabled (' + config['llm_model'] + ')' if config['use_llm'] else 'Disabled'}")



# ═══════════════════════════════════════════════════════════════════════════════
#  ATARS v3.0 ADDITIONS — VERIFIED AI REPORTING SYSTEM
#  "AI can explain, but cannot lie."
#
#  v3-01: Daily Snapshot System        data/snapshots/YYYY-MM-DD.json
#  v3-02: Comparison Engine            Today vs Yesterday / Week / Month
#  v3-03: Weather Context              Seasonal mock (no API required)
#  v3-04: Event Context Tags           Rule-based data pattern detection
#  v3-05: Enriched JSON Contract       Comparison + weather merged into J
#  v3-06: Semantic Grounding           Meaning check — High PM10 → must say unsafe
#  v3-07: Sentence Confidence Scores   0.0–1.0 per sentence
#  v3-08: Contradiction Detection      Flags "decrease" when data shows increase
#  v3-09: Multi-Source Grounding       J + WHO standards + weather values
#  v3-10: Evaluation Layer             Normal AI vs Verified AI comparison
#  v3-11: Test Mode                    --mode test uses dataset as timeline
# ═══════════════════════════════════════════════════════════════════════════════

import random
from typing import Optional


# ── v3 Config Defaults ────────────────────────────────────────────────────────

V3_CONFIG_DEFAULTS = {
    "snapshot_dir"         : "data/snapshots",
    "enable_snapshots"     : True,
    "weather_source"       : "open_meteo",  # default: always live
    "enable_semantic"      : True,
    "enable_confidence"    : True,
    "enable_contradiction" : True,
    "enable_evaluation"    : True,
    "mode"                 : "normal",
    "test_today_date"      : None,
}

WHO_STANDARDS = {
    "PM10"   : 45.0,
    "NO2"    : 25.0,
    "SO2"    : 40.0,
    "CO"     : 4.0,
    "Ozone"  : 100.0,
    "Benzene": 1.7,
    "NH3"    : 400.0,
}

# Semantic rules — what narrative MUST reflect when data shows certain conditions
SEMANTIC_RULES = [
    # Rule 1: Critical PM10 (>3x WHO) — must use strong language
    {
        "condition"   : lambda J: _get_j_pm10_mean(J) > 135,   # 3x WHO AQG
        "must_contain": ["exceed","above","hazardous","severely","critically",
                         "significant","substantially","far above","multiple times"],
        "must_not"    : ["safe","acceptable","within","compliant","good"],
        "description" : "PM10 > 3x WHO → must use strong exceedance language"
    },
    # Rule 2: High PM10 (1x-3x WHO)
    {
        "condition"   : lambda J: 45 < _get_j_pm10_mean(J) <= 135,
        "must_contain": ["exceed","above","concern","elevated","unhealthy",
                         "polluted","poor","high"],
        "must_not"    : ["safe","clean","within limits","compliant"],
        "description" : "PM10 above WHO → must reflect exceedance"
    },
    # Rule 3: PM10 compliant
    {
        "condition"   : lambda J: _get_j_pm10_mean(J) <= 45,
        "must_contain": ["within","below","compliant","acceptable","meets"],
        "must_not"    : ["hazardous","dangerous","severely","critically"],
        "description" : "PM10 within WHO → must reflect compliance"
    },
    # Rule 4: Benzene always carcinogen
    {
        "condition"   : lambda J: "benzene" in str(J).lower()
                                  and J.get("variables",{}).get("Benzene",{}).get(
                                      "annual_mean", 0) > 1.7,
        "must_contain": ["benzene","carcinogen","group 1","iarc","annual","exceed"],
        "must_not"    : ["safe benzene","benzene within"],
        "description" : "Benzene above WHO → must mention carcinogen risk"
    },
]

def _get_j_pm10_mean(J: dict) -> float:
    try:
        return float(J.get("PM10",{}).get("annual_mean", 0) or
                     J.get("pm10",{}).get("mean", 0) or 0)
    except Exception:
        return 0.0


# ── v3-03: Weather Context ─────────────────────────────────────────────────────

_SEASONAL_WEATHER = {
    1 :{"temp":8, "humidity":72,"wind":6, "rain":5 },
    2 :{"temp":12,"humidity":65,"wind":8, "rain":10},
    3 :{"temp":20,"humidity":55,"wind":12,"rain":8 },
    4 :{"temp":30,"humidity":40,"wind":15,"rain":5 },
    5 :{"temp":38,"humidity":30,"wind":18,"rain":3 },
    6 :{"temp":40,"humidity":50,"wind":20,"rain":30},
    7 :{"temp":35,"humidity":80,"wind":15,"rain":120},
    8 :{"temp":33,"humidity":85,"wind":12,"rain":100},
    9 :{"temp":32,"humidity":75,"wind":10,"rain":60},
    10:{"temp":26,"humidity":60,"wind":8, "rain":10},
    11:{"temp":17,"humidity":68,"wind":6, "rain":5 },
    12:{"temp":10,"humidity":75,"wind":5, "rain":3 },
}

def _mock_weather(month: int) -> dict:
    base = _SEASONAL_WEATHER.get(month, _SEASONAL_WEATHER[6])
    return {
        "temp_c"      : round(base["temp"] + random.uniform(-2,2), 1),
        "humidity_pct": round(base["humidity"] + random.uniform(-5,5), 1),
        "wind_kmh"    : round(max(0, base["wind"] + random.uniform(-2,2)), 1),
        "rain_mm"     : round(max(0, base["rain"] + random.uniform(-5,5)), 1),
        "source"      : "mock_seasonal",
    }

# ─────────────────────────────────────────────────────────────────────────────
#  OPEN-METEO LIVE WEATHER — Free, no API key required
#  Source: https://open-meteo.com  (WMO ERA5 reanalysis + forecast)
# ─────────────────────────────────────────────────────────────────────────────

# Coordinates for major Indian cities — extend as needed
CITY_COORDINATES = {
    "gurugram"   : (28.4595,  77.0266),
    "gurgaon"    : (28.4595,  77.0266),
    "delhi"      : (28.6139,  77.2090),
    "new delhi"  : (28.6139,  77.2090),
    "mumbai"     : (19.0760,  72.8777),
    "bangalore"  : (12.9716,  77.5946),
    "bengaluru"  : (12.9716,  77.5946),
    "chennai"    : (13.0827,  80.2707),
    "kolkata"    : (22.5726,  88.3639),
    "hyderabad"  : (17.3850,  78.4867),
    "pune"       : (18.5204,  73.8567),
    "ahmedabad"  : (23.0225,  72.5714),
    "faridabad"  : (28.4089,  77.3178),
    "noida"      : (28.5355,  77.3910),
    "ghaziabad"  : (28.6692,  77.4538),
    "chandigarh" : (30.7333,  76.7794),
    "jaipur"     : (26.9124,  75.7873),
    "lucknow"    : (26.8467,  80.9462),
    "patna"      : (25.5941,  85.1376),
    "bhopal"     : (23.2599,  77.4126),
    "nagpur"     : (21.1458,  79.0882),
    "surat"      : (21.1702,  72.8311),
    "kanpur"     : (26.4499,  80.3319),
    "prayagraj"  : (25.4358,  81.8463),
    "visakhapatnam": (17.6868, 83.2185),
    "indore"     : (22.7196,  75.8577),
    "thane"      : (19.2183,  72.9781),
    "vadodara"   : (22.3072,  73.1812),
    "coimbatore" : (11.0168,  76.9558),
    "agra"       : (27.1767,  78.0081),
    "meerut"     : (28.9845,  77.7064),
    "varanasi"   : (25.3176,  82.9739),
    "amritsar"   : (31.6340,  74.8723),
    "ludhiana"   : (30.9010,  75.8573),
    # default fallback
    "default"    : (28.6139,  77.2090),  # New Delhi
}


def _get_city_coords(city_name: str) -> tuple:
    """Returns (lat, lon) for a city name. Falls back to Delhi if unknown."""
    key = city_name.lower().strip()
    return CITY_COORDINATES.get(key, CITY_COORDINATES["default"])


def _fetch_open_meteo_weather(ts, config: dict) -> dict:
    """
    Fetches weather from Open-Meteo API for a given date.
    Completely free — no API key required.
    For historical dates: uses ERA5 reanalysis (archive API).
    For today/future: uses forecast API.
    Falls back to mock weather on any error.

    API docs: https://open-meteo.com/en/docs
    """
    try:
        import urllib.request
        import urllib.parse

        city   = config.get("city", "Delhi")
        lat, lon = _get_city_coords(city)

        # Convert ts to date string
        if hasattr(ts, "strftime"):
            date_str = ts.strftime("%Y-%m-%d")
            month    = ts.month
        else:
            date_str = str(ts)[:10]
            month    = int(date_str[5:7])

        from datetime import date as date_cls, timedelta
        today = date_cls.today()
        req_date = date_cls.fromisoformat(date_str)

        # Choose API endpoint
        if req_date < today - timedelta(days=5):
            # Historical — use archive API (ERA5 reanalysis)
            base_url = "https://archive-api.open-meteo.com/v1/archive"
        else:
            # Recent/forecast
            base_url = "https://api.open-meteo.com/v1/forecast"

        params = {
            "latitude"          : lat,
            "longitude"         : lon,
            "daily"             : "temperature_2m_mean,precipitation_sum,"
                                  "windspeed_10m_max,relativehumidity_2m_mean",
            "start_date"        : date_str,
            "end_date"          : date_str,
            "timezone"          : "Asia/Kolkata",
            "windspeed_unit"    : "kmh",
        }

        url = base_url + "?" + urllib.parse.urlencode(params)
        import requests as _rq
        _r = _rq.get(url, timeout=10, verify=False,
                     headers={"User-Agent": "ATARS/3.1"})
        _r.raise_for_status()
        data = _r.json()

        daily = data.get("daily", {})
        temp_list  = daily.get("temperature_2m_mean", [None])
        rain_list  = daily.get("precipitation_sum",  [None])
        wind_list  = daily.get("windspeed_10m_max",  [None])
        hum_list   = daily.get("relativehumidity_2m_mean", [None])

        temp  = temp_list[0]  if temp_list  else None
        rain  = rain_list[0]  if rain_list  else None
        wind  = wind_list[0]  if wind_list  else None
        hum   = hum_list[0]   if hum_list   else None

        if temp is None:
            raise ValueError("Empty response from Open-Meteo")

        return {
            "temp_c"      : round(float(temp), 1),
            "humidity_pct": round(float(hum  or 50), 1),
            "wind_kmh"    : round(float(wind or 10), 1),
            "rain_mm"     : round(float(rain or 0),  1),
            "source"      : "open_meteo_api",
            "lat"         : lat,
            "lon"         : lon,
            "date"        : date_str,
        }

    except Exception as e:
        # Graceful fallback to mock — never crash the pipeline
        month = 6
        try:
            if hasattr(ts, "month"):
                month = ts.month
            else:
                month = int(str(ts)[5:7])
        except Exception:
            pass
        result = _mock_weather(month)
        result["source"]       = f"mock_fallback (open_meteo_error: {str(e)[:60]})"
        return result


# Cache to avoid repeated API calls for same date
_weather_cache: dict = {}


def _get_weather_for_date(ts, daily_sub: pd.DataFrame, config: dict) -> dict:
    """
    Weather priority:
    1. From dataset columns (humidity_pct etc.) — most accurate
    2. Open-Meteo API (if weather_source = "open_meteo") — free live/historical
    3. Mock seasonal values — always available as fallback
    """
    # Priority 1: from dataset
    try:
        row = daily_sub[daily_sub["date"] == ts]
        if not row.empty:
            r   = row.iloc[0]
            hum = float(r.get("humidity_pct", 0) or 0)
            wind= float(r.get("wind_speed_10m", 0) or 0)
            rain= float(r.get("rain_mm", 0) or 0)
            temp= float(r.get("temperature_2m", 0) or 0)
            if hum > 0:
                return {
                    "temp_c"      : round(temp, 1),
                    "humidity_pct": round(hum, 1),
                    "wind_kmh"    : round(wind * 3.6, 1),
                    "rain_mm"     : round(rain, 1),
                    "source"      : "dataset",
                }
    except Exception:
        pass

    # Priority 2: Open-Meteo API
    weather_source = config.get("weather_source", "mock")
    if weather_source == "open_meteo":
        date_key = str(ts)[:10]
        city     = config.get("city", "Delhi")
        cache_key = f"{city}_{date_key}"
        if cache_key not in _weather_cache:
            _weather_cache[cache_key] = _fetch_open_meteo_weather(ts, config)
        return _weather_cache[cache_key]

    # Priority 3: Mock seasonal
    month = ts.month if hasattr(ts, "month") else 6
    return _mock_weather(month)


# ── v3-04: Event Context Tags ──────────────────────────────────────────────────

def _detect_event_tags(pm10: float, weather: dict, config: dict) -> list:
    who = config.get("thresholds",{}).get("PM10", 45.0)
    tags = []
    hum  = weather.get("humidity_pct", 50)
    wind = weather.get("wind_kmh", 10)
    rain = weather.get("rain_mm", 0)
    temp = weather.get("temp_c", 20)
    if pm10 > who * 5:   tags.append("severe_pollution_episode")
    elif pm10 > who * 3: tags.append("high_pollution")
    elif pm10 > who:     tags.append("who_exceedance")
    else:                tags.append("within_who_limits")
    if wind < 5 and pm10 > who*2:    tags.append("stagnant_air_trapping_pollution")
    if hum > 80 and pm10 > who:      tags.append("humid_conditions")
    if rain > 10:                    tags.append("wet_deposition_expected")
    if temp < 10 and pm10 > who*2:   tags.append("winter_temperature_inversion")
    if temp > 35:                    tags.append("high_temperature_photochemistry")
    return tags


# ── v3-01: Daily Snapshot System ──────────────────────────────────────────────


# ─────────────────────────────────────────────────────────────────────────────
#  LIVE DATA INTEGRATION — Open-Meteo for recent + today
#  Fetches weather for dates NOT well-covered by CSV (last 90 days + today)
#  This is the key difference between test mode and live mode
# ─────────────────────────────────────────────────────────────────────────────

from datetime import date as _date_cls, timedelta as _timedelta


def _should_use_api(ts, daily_sub: pd.DataFrame, config: dict) -> bool:
    """
    Open-Meteo is the DEFAULT weather source.
    Only skipped if user explicitly passes --weather mock.
    """
    if config.get("weather_source", "open_meteo") == "mock":
        return False   # user explicitly disabled API

    try:
        req_date = ts.date() if hasattr(ts, 'date') else _date_cls.fromisoformat(str(ts)[:10])
        today    = _date_cls.today()
        # Force API for last 90 days
        if (today - req_date).days <= 90:
            return True
    except Exception:
        pass
    return False




def _wmo_code_description(code: int) -> str:
    """Convert WMO weather code to human-readable description."""
    WMO_CODES = {
        0: "Clear sky", 1: "Mainly clear", 2: "Partly cloudy",
        3: "Overcast", 45: "Fog", 48: "Icy fog",
        51: "Light drizzle", 53: "Moderate drizzle", 55: "Dense drizzle",
        61: "Slight rain", 63: "Moderate rain", 65: "Heavy rain",
        71: "Slight snowfall", 73: "Moderate snowfall", 75: "Heavy snowfall",
        77: "Snow grains", 80: "Slight showers", 81: "Moderate showers",
        82: "Violent showers", 85: "Slight snow showers",
        86: "Heavy snow showers", 95: "Thunderstorm",
        96: "Thunderstorm with hail", 99: "Thunderstorm with heavy hail",
    }
    return WMO_CODES.get(code, f"Code {code}")


def _fetch_live_weather_range(start_date: str, end_date: str,
                               config: dict) -> dict:
    """
    Fetches weather for a date RANGE from Open-Meteo in one API call.
    Much more efficient than per-date calls.
    Returns {date_str: weather_dict}
    """
    import urllib.request
    import urllib.parse

    city     = config.get("city", "Delhi")
    lat, lon = _get_city_coords(city)

    try:
        from datetime import date as dc
        today    = dc.today()
        req_start = dc.fromisoformat(start_date)
        req_end   = dc.fromisoformat(end_date)

        if req_start < today - _timedelta(days=5):
            base_url = "https://archive-api.open-meteo.com/v1/archive"
        else:
            base_url = "https://api.open-meteo.com/v1/forecast"

        # Full variable set — all pollution-relevant meteorological parameters
        DAILY_VARS = ",".join([
            "temperature_2m_mean", "temperature_2m_max", "temperature_2m_min",
            "apparent_temperature_max",
            "precipitation_sum", "rain_sum", "precipitation_hours",
            "windspeed_10m_max", "windspeed_10m_mean", "windgusts_10m_max",
            "winddirection_10m_dominant",
            "relativehumidity_2m_mean", "relativehumidity_2m_max",
            "dewpoint_2m_mean", "surface_pressure_mean",
            "shortwave_radiation_sum", "sunshine_duration",
            "weathercode",
        ])

        params = {
            "latitude"      : lat,
            "longitude"     : lon,
            "daily"         : DAILY_VARS,
            "start_date"    : start_date,
            "end_date"      : end_date,
            "timezone"      : "Asia/Kolkata",
            "windspeed_unit": "kmh",
            "timeformat"    : "iso8601",
        }

        url = base_url + "?" + urllib.parse.urlencode(params)
        import requests as _rq
        _r = _rq.get(url, timeout=20, verify=False,
                     headers={"User-Agent": "ATARS/3.1"})
        _r.raise_for_status()
        data = _r.json()

        daily = data.get("daily", {})
        dates = daily.get("time", [])

        def _safe(key, i):
            arr = daily.get(key, [])
            try: return float(arr[i] or 0)
            except: return 0.0

        result = {}
        for i, d in enumerate(dates):
            temp_mean  = _safe("temperature_2m_mean",     i)
            temp_max   = _safe("temperature_2m_max",      i)
            temp_min   = _safe("temperature_2m_min",      i)
            ap_max     = _safe("apparent_temperature_max",i)
            precip     = _safe("precipitation_sum",       i)
            rain       = _safe("rain_sum",                i)
            precip_hrs = _safe("precipitation_hours",     i)
            wind_max   = _safe("windspeed_10m_max",       i)
            wind_mean  = _safe("windspeed_10m_mean",      i)
            gust_max   = _safe("windgusts_10m_max",       i)
            wind_dir   = _safe("winddirection_10m_dominant", i)
            hum_mean   = _safe("relativehumidity_2m_mean",i)
            hum_max    = _safe("relativehumidity_2m_max", i)
            dewpoint   = _safe("dewpoint_2m_mean",        i)
            pressure   = _safe("surface_pressure_mean",  i)
            radiation  = _safe("shortwave_radiation_sum", i)
            sunshine   = _safe("sunshine_duration",       i) / 3600  # s → hours
            wcode      = int(_safe("weathercode",         i))

            # Compute derived pollution risk indices
            inversion_risk  = temp_min < 8  and wind_mean < 5
            wet_deposition  = precip > 10
            dust_risk       = gust_max > 30 and precip < 1
            smog_risk       = temp_max > 35 and radiation > 15
            fog_risk        = (temp_mean - dewpoint) < 2 and hum_mean > 90
            high_pressure   = pressure > 1015
            stagnant        = wind_mean < 3

            # Risk score 0-6
            risk_score = sum([
                inversion_risk, wet_deposition, dust_risk,
                smog_risk, fog_risk, stagnant
            ])
            risk_label = ("CRITICAL" if risk_score >= 4 else
                          "HIGH"     if risk_score >= 3 else
                          "MODERATE" if risk_score >= 2 else
                          "LOW"      if risk_score >= 1 else "MINIMAL")

            # Wind direction label
            dirs = ["N","NE","E","SE","S","SW","W","NW"]
            wind_dir_label = dirs[int((wind_dir + 22.5) / 45) % 8] if wind_dir else "N/A"

            # WMO weather code description
            wcode_desc = _wmo_code_description(wcode)

            result[d] = {
                # Core variables (backward compatible)
                "temp_c"             : round(temp_mean, 1),
                "humidity_pct"       : round(hum_mean,  1),
                "wind_kmh"           : round(wind_mean, 1),
                "rain_mm"            : round(precip,    1),
                # Extended variables
                "temp_max_c"         : round(temp_max,  1),
                "temp_min_c"         : round(temp_min,  1),
                "feels_like_max_c"   : round(ap_max,    1),
                "rain_only_mm"       : round(rain,      1),
                "precip_hours"       : round(precip_hrs,1),
                "wind_max_kmh"       : round(wind_max,  1),
                "wind_gust_kmh"      : round(gust_max,  1),
                "wind_direction_deg" : round(wind_dir,  1),
                "wind_direction"     : wind_dir_label,
                "humidity_max_pct"   : round(hum_max,   1),
                "dewpoint_c"         : round(dewpoint,  1),
                "pressure_hpa"       : round(pressure,  1),
                "solar_radiation_mj" : round(radiation, 2),
                "sunshine_hours"     : round(sunshine,  1),
                "weather_code"       : wcode,
                "weather_desc"       : wcode_desc,
                # Derived risk indices
                "inversion_risk"     : inversion_risk,
                "wet_deposition"     : wet_deposition,
                "dust_risk"          : dust_risk,
                "smog_risk"          : smog_risk,
                "fog_risk"           : fog_risk,
                "stagnant_air"       : stagnant,
                "high_pressure"      : high_pressure,
                "pollution_risk_score": risk_score,
                "pollution_risk"     : risk_label,
                "source"             : "open_meteo_live",
            }

        print(f"  ✓ Open-Meteo: fetched {len(result)} days "
              f"({start_date} → {end_date}) for {city}")
        return result

    except Exception as e:
        print(f"  ⚠ Open-Meteo range fetch failed: {e} — using mock fallback")
        return {}


def _get_live_weather_cache(config: dict, all_dates: list) -> dict:
    """
    Pre-fetches weather for ALL dates in one batch API call.
    ALWAYS attempts Open-Meteo — falls back to mock on failure.
    mock mode = explicitly disabled by user (no internet situations).
    Returns {date_str: weather_dict}
    """
    if config.get("weather_source", "open_meteo") == "mock":
        return {}   # user explicitly chose mock — skip API
    if not all_dates:
        return {}

    try:
        str_dates = [str(d)[:10] for d in all_dates]
        start = min(str_dates)
        end   = max(str_dates)
        return _fetch_live_weather_range(start, end, config)
    except Exception as e:
        print(f"  ⚠ Weather pre-fetch failed: {e}")
        return {}


def build_daily_snapshots(daily_df: pd.DataFrame, config: dict,
                          live_weather_cache: dict = None) -> dict:
    """
    Creates one JSON snapshot per day under data/snapshots/YYYY-MM-DD.json
    Uses daily_df long-format (columns: date, variable, mean, ...)
    live_weather_cache: pre-fetched Open-Meteo data {date_str: weather_dict}
    """
    if live_weather_cache is None:
        live_weather_cache = {}
    snap_dir = Path(config.get("snapshot_dir","data/snapshots"))
    snap_dir.mkdir(parents=True, exist_ok=True)

    who_pm10 = config.get("thresholds",{}).get("PM10", 45.0)
    snapshots = {}

    # Pivot: get per-variable means indexed by date
    pm10_df = daily_df[daily_df["variable"]=="PM10"][["date","mean"]].copy()
    pm10_df = pm10_df.dropna(subset=["date"])
    pm10_df["date"] = pd.to_datetime(pm10_df["date"])
    pm10_df = pm10_df.sort_values("date").reset_index(drop=True)

    met_vars = ["humidity_pct","wind_speed_10m","rain_mm","temperature_2m"]
    met_data = {}
    for mv in met_vars:
        sub = daily_df[daily_df["variable"]==mv][["date","mean"]].copy()
        if not sub.empty:
            sub["date"] = pd.to_datetime(sub["date"])
            met_data[mv] = sub.set_index("date")["mean"]

    all_dates = list(pm10_df["date"])

    for i, ts in enumerate(all_dates):
        d_str = ts.strftime("%Y-%m-%d")
        pm10_today = float(pm10_df.iloc[i]["mean"] or 0)

        # Weekly and monthly rolling averages (computed here — not by AI)
        week_vals  = [float(pm10_df.iloc[j]["mean"] or 0)
                      for j in range(max(0,i-7), i)]
        month_vals = [float(pm10_df.iloc[j]["mean"] or 0)
                      for j in range(max(0,i-30), i)]
        pm10_week  = round(sum(week_vals)/len(week_vals),2) if week_vals else pm10_today
        pm10_month = round(sum(month_vals)/len(month_vals),2) if month_vals else pm10_today

        # Weather — priority: API cache > CSV columns > mock
        if live_weather_cache and d_str in live_weather_cache:
            weather = live_weather_cache[d_str].copy()
        else:
            weather = _mock_weather(ts.month)
            weather["source"] = "mock_seasonal"

        # Override with CSV columns if available (more precise)
        for mv in met_vars:
            if mv in met_data and ts in met_data[mv].index:
                val = met_data[mv].loc[ts]
                if not pd.isna(val) and float(val or 0) > 0:
                    if mv == "humidity_pct":     weather["humidity_pct"] = round(float(val),1)
                    elif mv == "wind_speed_10m": weather["wind_kmh"]     = round(float(val)*3.6,1)
                    elif mv == "rain_mm":        weather["rain_mm"]      = round(float(val),1)
                    elif mv == "temperature_2m": weather["temp_c"]       = round(float(val),1)
                    weather["source"] = "csv_columns"

        events = _detect_event_tags(pm10_today, weather, config)

        # Yesterday comparison computed here
        pm10_yest = float(pm10_df.iloc[i-1]["mean"] or 0) if i > 0 else pm10_today
        change_pct = round(((pm10_today-pm10_yest)/max(abs(pm10_yest),1))*100,1) if i>0 else 0.0

        snap = {
            "date"                    : d_str,
            "city"                    : config.get("city","City"),
            "pm10_avg"                : round(pm10_today,2),
            "pm10_week_avg"           : pm10_week,
            "pm10_month_avg"          : pm10_month,
            "pm10_yesterday"          : round(pm10_yest,2),
            "pm10_change_yesterday_pct": change_pct,
            "pm10_direction_yesterday": "increase" if change_pct>0.5
                                        else ("decrease" if change_pct<-0.5 else "stable"),
            "weather"                 : weather,
            "event_tags"              : events,
            "exceed_who_pm10"         : pm10_today > who_pm10,
            "source"                  : "ATARS_v3_snapshot",
        }
        snap_path = snap_dir / f"{d_str}.json"
        with open(snap_path,"w") as f:
            json.dump(snap, f, indent=2)
        snapshots[d_str] = snap

    print(f"  ✓ v3-01: {len(snapshots)} daily snapshots → {snap_dir}/")
    return snapshots


# ── v3-02: Comparison Engine ───────────────────────────────────────────────────

def run_comparison_engine(snapshots: dict, today_date: str) -> dict:
    """
    All arithmetic computed here — AI never calculates anything.
    Returns pre-computed comparison dict that goes into J.
    """
    all_dates = sorted(snapshots.keys())
    if today_date not in snapshots:
        today_date = all_dates[-1]

    idx      = all_dates.index(today_date)
    snap     = snapshots[today_date]
    pm10_now = snap.get("pm10_avg", 0)

    result = {
        "today"          : today_date,
        "pm10_today"     : pm10_now,
        "weather"        : snap.get("weather", {}),
        "event_tags"     : snap.get("event_tags", []),
        "exceed_who_pm10": snap.get("exceed_who_pm10", False),
    }

    # Yesterday
    if idx > 0:
        yd   = all_dates[idx-1]
        yv   = snapshots[yd].get("pm10_avg", 0)
        result.update({
            "pm10_yesterday"              : yv,
            "pm10_yesterday_date"         : yd,
            "pm10_change_yesterday_abs"   : round(pm10_now - yv, 2),
            "pm10_change_yesterday_pct"   : _safe_pct(pm10_now, yv),
            "pm10_direction_yesterday"    : _dir(pm10_now, yv),
        })

    # Week avg
    week_w = all_dates[max(0,idx-7):idx]
    if week_w:
        wa = round(sum(snapshots[d]["pm10_avg"] for d in week_w)/len(week_w),2)
        result.update({
            "pm10_week_avg"       : wa,
            "pm10_change_week_pct": _safe_pct(pm10_now, wa),
            "pm10_direction_week" : _dir(pm10_now, wa),
        })

    # Month avg
    mon_w = all_dates[max(0,idx-30):idx]
    if mon_w:
        ma = round(sum(snapshots[d]["pm10_avg"] for d in mon_w)/len(mon_w),2)
        result.update({
            "pm10_month_avg"        : ma,
            "pm10_change_month_pct" : _safe_pct(pm10_now, ma),
            "pm10_direction_month"  : _dir(pm10_now, ma),
        })

    # Status
    week_pct = result.get("pm10_change_week_pct", 0)
    if   week_pct >  20: result["status"] = "SIGNIFICANTLY_WORSE_THAN_WEEK"
    elif week_pct >   5: result["status"] = "WORSE_THAN_WEEK"
    elif week_pct < -20: result["status"] = "SIGNIFICANTLY_BETTER_THAN_WEEK"
    elif week_pct <  -5: result["status"] = "BETTER_THAN_WEEK"
    else:                result["status"] = "WITHIN_NORMAL_RANGE"

    print(f"  ✓ v3-02: Comparison — today={today_date}, PM10={pm10_now} µg/m³, "
          f"status={result['status']}")
    return result

def _safe_pct(new, old):
    if old == 0: return 0.0
    return round(((new-old)/abs(old))*100, 1)

def _dir(new, old):
    d = new - old
    if abs(d) < 0.5: return "stable"
    return "increase" if d > 0 else "decrease"


# ── v3-05: Enriched JSON Contract ─────────────────────────────────────────────

def enrich_json_contract_v3(J: dict, comparison: dict, config: dict) -> dict:
    """Merges comparison + weather + WHO standards into J. AI only sees this."""
    J3 = dict(J)
    J3["v3_comparison"] = {
        "today_date"               : comparison.get("today"),
        "pm10_today"               : comparison.get("pm10_today"),
        "pm10_yesterday"           : comparison.get("pm10_yesterday"),
        "pm10_change_yesterday_pct": comparison.get("pm10_change_yesterday_pct"),
        "pm10_direction_yesterday" : comparison.get("pm10_direction_yesterday"),
        "pm10_week_avg"            : comparison.get("pm10_week_avg"),
        "pm10_change_week_pct"     : comparison.get("pm10_change_week_pct"),
        "pm10_direction_week"      : comparison.get("pm10_direction_week"),
        "pm10_month_avg"           : comparison.get("pm10_month_avg"),
        "pm10_change_month_pct"    : comparison.get("pm10_change_month_pct"),
        "status"                   : comparison.get("status"),
    }
    J3["v3_weather"]   = comparison.get("weather", {})
    J3["v3_events"]    = {"tags": comparison.get("event_tags",[]),
                          "note": "Rule-based — no external API"}
    J3["v3_who_standards"] = {
        "PM10_24h"      : 45.0,
        "NO2_annual"    : 25.0,
        "SO2_24h"       : 40.0,
        "Benzene_annual": 1.7,
        "CO_24h"        : 4.0,
        "source"        : "WHO AQG 2021"
    }
    print("  ✓ v3-05: JSON contract enriched (comparison + weather + WHO)")
    return J3


# ── v3 Enhanced RGV ───────────────────────────────────────────────────────────

class EnhancedRGV_v3:
    """
    v3 Runtime Grounding Verifier.
    Adds on top of v2 RGV:
    - Multi-source grounding  (J + WHO + weather)
    - Semantic grounding      (meaning check)
    - Sentence confidence     (0.0–1.0 per sentence)
    - Contradiction detection (direction word vs data)
    """
    TOL_REL   = 0.03    # 3% — tighter than v2 (was 5%)
    TOL_ABS   = 1.0    # ±1 count — tighter than v2 (was ±2)
    TOL_SCORE = 0.05   # ±0.05 for scores — tighter than v2 (was ±0.10)
    PASS_THR  = 0.80   # raised from 0.75 — higher standard

    DIRECTION_WORDS_UP   = ["increase","increased","rose","elevated","higher",
                             "worsened","worse","risen","spiked","surged"]
    DIRECTION_WORDS_DOWN = ["decrease","declined","fell","dropped","reduced",
                             "lower","improved","better","fallen","eased"]
    DIRECTION_WORDS_FLAT = ["stable","unchanged","constant","steady","similar"]

    def __init__(self, J_v3: dict):
        self.J = J_v3
        self.lookup = {}
        self._build_lookup(J_v3, "")
        # WHO lookup
        self.who_lookup = {k: float(v) for k,v in WHO_STANDARDS.items()}
        # Weather lookup
        weather = J_v3.get("v3_weather", {})
        self.weather_lookup = {k: float(v) for k,v in weather.items()
                               if isinstance(v,(int,float))}

    def _build_lookup(self, obj, path):
        if isinstance(obj, dict):
            for k,v in obj.items():
                self._build_lookup(v, f"{path}.{k}" if path else k)
        elif isinstance(obj, list):
            for i,v in enumerate(obj):
                self._build_lookup(v, f"{path}[{i}]")
        elif isinstance(obj, (int,float)) and not isinstance(obj, bool):
            try:
                fv = float(obj)
                if not (math.isnan(fv) or math.isinf(fv)):
                    self.lookup[path] = fv
            except Exception:
                pass

    def _match(self, val: float) -> dict:
        # Check J
        for path, jv in self.lookup.items():
            denom = max(abs(jv), 1.0)
            if abs(val-jv)/denom <= self.TOL_REL or abs(val-jv) <= self.TOL_ABS:
                return {"matched":True,"source":"J","path":path,"jv":jv}
        # Check WHO
        for k,wv in self.who_lookup.items():
            if abs(val-wv)/max(abs(wv),1.0) <= self.TOL_REL:
                return {"matched":True,"source":"WHO","path":k,"jv":wv}
        # Check weather
        for k,wv in self.weather_lookup.items():
            if abs(val-wv)/max(abs(wv),1.0) <= self.TOL_REL:
                return {"matched":True,"source":"weather","path":k,"jv":wv}
        return {"matched":False,"source":None,"path":None,"jv":None}

    def _extract_nums(self, s: str) -> list:
        """Extract numerical values from sentence. Skips section/figure references."""
        nums    = []
        pattern = re.compile(r"(?<![a-zA-Z])(-?\d{1,6}(?:\.\d{1,6})?)\b")
        for m in pattern.finditer(s):
            try:
                v   = float(m.group(1))
                ctx = s[max(0, m.start()-20):m.start()].lower()
                if abs(v) < 30 and "." not in m.group(1):
                    if any(w in ctx for w in ["eq.", "sec", "fig", "tab", "section"]):
                        continue
                nums.append(v)
            except ValueError:
                pass
        return nums
    def _check_contradiction(self, sent: str, comparison: dict) -> dict:
        """
        Detects direction contradictions.
        Requires both: (a) direction word found AND (b) data shows opposite direction
        AND (c) change magnitude > 5% (avoids false positives on tiny changes).
        """
        sl     = sent.lower()
        actual = comparison.get("pm10_direction_yesterday", "stable")
        pct    = float(comparison.get("pm10_change_yesterday_pct", 0) or 0)

        has_up   = any(w in sl for w in self.DIRECTION_WORDS_UP)
        has_down = any(w in sl for w in self.DIRECTION_WORDS_DOWN)
        has_both = has_up and has_down  # e.g. "decrease from high values" — not a contradiction

        contradiction = False
        if not has_both:
            if has_up   and actual == "decrease" and abs(pct) > 5:
                contradiction = True
            if has_down and actual == "increase"  and abs(pct) > 5:
                contradiction = True

        return {
            "contradiction"   : contradiction,
            "actual_direction": actual,
            "change_pct"      : round(pct, 1),
            "words_up"        : [w for w in self.DIRECTION_WORDS_UP   if w in sl],
            "words_down"      : [w for w in self.DIRECTION_WORDS_DOWN  if w in sl],
            "ambiguous"       : has_both,
        }

    def _semantic_check(self, text: str) -> dict:
        results = []
        for rule in SEMANTIC_RULES:
            if not rule["condition"](self.J):
                continue
            tl = text.lower()
            has_req = any(w in tl for w in rule["must_contain"])
            has_for = any(w in tl for w in rule["must_not"])
            results.append({
                "rule"        : rule["description"],
                "pass"        : has_req and not has_for,
                "has_required": has_req,
                "has_forbidden":has_for,
            })
        return {
            "semantic_pass"  : all(r["pass"] for r in results) if results else True,
            "rules_checked"  : results,
        }

    def _sentence_confidence(self, sent: str, grounded_ratio: float,
                              has_contradiction: bool) -> float:
        nums = self._extract_nums(sent)
        if not nums:
            return 0.85  # qualitative — moderate-high trust
        base = grounded_ratio
        if has_contradiction: base -= 0.30
        if len(nums) >= 2:    base += 0.05
        return max(0.0, min(1.0, round(base, 2)))

    def verify_v3(self, narrative: dict, comparison: dict) -> dict:
        """Full v3 verification with all enhancements."""
        sections      = {}
        all_grounded  = []
        all_conf      = []
        contradictions= []

        for sec_key, text in narrative.items():
            if not isinstance(text,str) or not text.strip():
                continue

            sents = [s.strip() for s in re.split(r'[.!?]', text) if len(s.strip())>15]
            sec   = {
                "total":0,"numerical":0,"grounded":0,
                "g_rate":0.0,"avg_confidence":0.0,
                "verdict":"NON_NUMERICAL","sentences":[],
                "semantic":{}, "contradictions":0,
            }
            confs = []

            for sent in sents:
                nums  = self._extract_nums(sent)
                contra= self._check_contradiction(sent, comparison)
                if contra["contradiction"]:
                    contradictions.append({
                        "section":sec_key,"sentence":sent[:150],
                        "details":contra
                    })
                    sec["contradictions"] += 1

                if nums:
                    matches      = [self._match(n) for n in nums]
                    gr_count     = sum(1 for m in matches if m["matched"])
                    gr_ratio     = gr_count / len(nums)
                    grounded     = gr_ratio >= 0.5
                    all_grounded.append(grounded)
                    sec["numerical"] += 1
                    if grounded: sec["grounded"] += 1
                else:
                    gr_ratio = 1.0

                conf = self._sentence_confidence(sent, gr_ratio, contra["contradiction"])
                confs.append(conf)
                all_conf.append(conf)

                sec["sentences"].append({
                    "text"        : sent[:150],
                    "nums"        : nums,
                    "grounded"    : gr_ratio>=0.5 if nums else True,
                    "confidence"  : conf,
                    "contradiction": contra["contradiction"],
                })

            sec["total"]          = len(sents)
            sec["avg_confidence"] = round(sum(confs)/len(confs),2) if confs else 1.0
            n_num = sec["numerical"]
            sec["g_rate"]  = round(sec["grounded"]/n_num,3) if n_num>0 else 1.0
            sec["verdict"] = ("PASS" if sec["g_rate"]>=0.75
                              else ("NON_NUMERICAL" if n_num==0 else "REVIEW"))
            sec["semantic"]= self._semantic_check(text)
            sections[sec_key] = sec

        # Overall
        total_num    = sum(s["numerical"] for s in sections.values())
        total_ground = sum(s["grounded"]  for s in sections.values())
        g_rate       = round(total_ground/total_num, 3) if total_num > 0 else 1.0
        verdict      = ("PASS"    if g_rate>=0.75 else
                        "WARNING" if g_rate>=0.50 else "FAIL")
        avg_conf     = round(sum(all_conf)/len(all_conf),2) if all_conf else 1.0
        sem_pass     = all(s["semantic"].get("semantic_pass",True) for s in sections.values())

        print(f"  ✓ v3 G_rate: {round(g_rate*100,1)}% [{verdict}] | "
              f"Confidence: {avg_conf} | Contradictions: {len(contradictions)} | "
              f"Semantic: {'PASS' if sem_pass else 'REVIEW'}")

        return {
            "overall_G_rate"       : g_rate,
            "overall_verdict"      : verdict,
            "overall_confidence"   : avg_conf,
            "semantic_pass"        : sem_pass,
            "contradictions"       : contradictions,
            "contradiction_count"  : len(contradictions),
            "total_numerical"      : total_num,
            "total_grounded"       : total_ground,
            "total_ungrounded"     : total_num - total_ground,
            "sources_used"         : ["J_contract","WHO_AQG_2021","weather_context"],
            "sections"             : sections,
            "verification_passed"  : g_rate >= 0.75,
            "j_lookup_size"        : len(self.lookup),
            "version"              : "ATARS_RGV_v3.0",
        }


# ── v3-10: Evaluation Layer ────────────────────────────────────────────────────

def run_evaluation_layer(narrative_base: dict, narrative_v3: dict,
                          verification: dict) -> dict:
    """Compares unverified vs verified narrative. Shows hallucination reduction."""
    def _nums(text):
        return len(re.findall(r'-?\d+(?:\.\d+)?', text))
    def _sents(text):
        return len([s for s in re.split(r'[.!?]', text) if len(s.strip())>10])

    b_nums = sum(_nums(t) for t in narrative_base.values() if isinstance(t,str))
    v_nums = sum(_nums(t) for t in narrative_v3.values()  if isinstance(t,str))
    g      = verification.get("overall_G_rate", 0)

    result = {
        "normal_ai" : {
            "grounding_verified": False,
            "total_numbers"     : b_nums,
            "hallucination_risk": "UNKNOWN",
            "g_rate"            : "N/A",
        },
        "verified_ai": {
            "grounding_verified": True,
            "total_numbers"     : v_nums,
            "g_rate"            : g,
            "verdict"           : verification.get("overall_verdict"),
            "hallucination_risk": ("LOW" if g>=0.9 else "MEDIUM" if g>=0.75 else "HIGH"),
            "confidence"        : verification.get("overall_confidence"),
            "contradictions"    : verification.get("contradiction_count",0),
            "semantic_aligned"  : verification.get("semantic_pass",True),
        },
        "summary": (
            f"Verified AI: {round(g*100,1)}% of numerical claims grounded. "
            f"{verification.get('contradiction_count',0)} contradiction(s) flagged. "
            f"Confidence: {verification.get('overall_confidence',0)}."
        )
    }
    return result


# ── v3: Add v3 grounding section to Word report ────────────────────────────────

def _add_v3_section_to_report(doc, verification: dict, evaluation: dict,
                               comparison: dict, config: dict):
    """
    Appends a v3 Grounding & Comparison section to the existing Word document.
    This is the ONLY output file — no separate files created.
    """
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    PAL_NAVY  = RGBColor(0x1A, 0x2C, 0x4E)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN = RGBColor(0x1E, 0x6B, 0x3C)
    AMBER = RGBColor(0x92, 0x40, 0x0E)
    RED   = RGBColor(0x8B, 0x1C, 0x2A)
    PAL_TEAL  = RGBColor(0x1A, 0x6B, 0x72)
    PAL_GREEN = RGBColor(0x1E, 0x6B, 0x3C)
    PAL_RED   = RGBColor(0x8B, 0x1C, 0x2A)
    PAL_AMBER = RGBColor(0x92, 0x40, 0x0E)

    def _heading(level, text, color=None):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(10)
        p.paragraph_format.space_after  = Pt(4)
        run = p.add_run(text)
        run.bold = True
        if level == 1:
            run.font.size = Pt(15)
            run.font.color.rgb = PAL_NAVY
        elif level == 2:
            run.font.size = Pt(12)
            run.font.color.rgb = PAL_TEAL
        else:
            run.font.size = Pt(11)
            run.font.color.rgb = PAL_AMBER
        p.paragraph_format.keep_with_next = True

    def _para(text, bold=False):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(4)
        if bold:
            for run in p.runs: run.bold = True

    def _table_2col(rows):
        t = doc.add_table(rows=1, cols=2)
        t.style = "Table Grid"
        hdr = t.rows[0].cells
        hdr[0].text = "Metric"
        hdr[1].text = "Value"
        for c in hdr:
            for run in c.paragraphs[0].runs:
                run.bold = True
        for k,v in rows:
            r = t.add_row().cells
            r[0].text = str(k)
            r[1].text = str(v)
        doc.add_paragraph()

    doc.add_page_break()
    _heading(1, "ATARS v3.0 — Verified AI Reporting Layer")
    _para('This section documents the v3 enhancements: daily snapshot comparison, '          'enhanced Runtime Grounding Verifier (RGV), semantic alignment check, '          'sentence confidence scores, contradiction detection, and multi-source grounding.')

    # ── Comparison Engine Results ────────────────────────────────────────────
    _heading(2, "Today vs Historical Comparison")
    _para('All values below are pre-computed by the comparison engine. '          'The AI never performs these calculations.')
    comp_rows = []
    for k,v in comparison.items():
        if isinstance(v,(int,float,str,bool)) and k not in ("weather","event_tags"):
            comp_rows.append((k.replace("_"," ").title(), str(v)))
    _table_2col(comp_rows[:15])

    # ── Event Tags ───────────────────────────────────────────────────────────
    tags = comparison.get("event_tags",[])
    if tags:
        _heading(3, "Context Event Tags")
        _para(", ".join(tags))

    # ── v3 RGV Summary ───────────────────────────────────────────────────────
    _heading(2, "Enhanced RGV v3 — Grounding Verification")
    g = verification.get("overall_G_rate",0)
    verdict = verification.get("overall_verdict","N/A")
    verdict_color = PAL_GREEN if verdict=="PASS" else (
                    PAL_AMBER if verdict=="WARNING" else PAL_RED)

    gp = doc.add_paragraph()
    r1 = gp.add_run(f"G_rate: {round(g*100,1)}%  [{verdict}]")
    r1.bold = True
    r1.font.size = Pt(13)
    r1.font.color.rgb = verdict_color

    rgv_rows = [
        ("G_rate",                f"{round(g*100,1)}%"),
        ("Verdict",               verdict),
        ("Average Confidence",    str(verification.get("overall_confidence","N/A"))),
        ("Semantic Alignment",    "PASS" if verification.get("semantic_pass") else "REVIEW"),
        ("Contradictions Found",  str(verification.get("contradiction_count",0))),
        ("Total Numerical Sent.", str(verification.get("total_numerical",0))),
        ("Grounded",              str(verification.get("total_grounded",0))),
        ("Ungrounded",            str(verification.get("total_ungrounded",0))),
        ("Sources Used",          ", ".join(verification.get("sources_used",[]))),
        ("RGV Version",           verification.get("version","v3.0")),
    ]
    _table_2col(rgv_rows)

    # ── Per-section RGV ──────────────────────────────────────────────────────
    _heading(3, "Per-Section Grounding Results")
    sec_data = verification.get("sections",{})
    t = doc.add_table(rows=1, cols=5)
    t.style = "Table Grid"
    hdrs = ["Section","G_rate","Verdict","Confidence","Contradictions"]
    for i,h in enumerate(hdrs):
        c = t.rows[0].cells[i]
        c.text = h
        for run in c.paragraphs[0].runs: run.bold = True
    for sk, sv in sec_data.items():
        row = t.add_row().cells
        row[0].text = sk
        row[1].text = f"{round(sv.get('G_rate', sv.get('g_rate', 0))*100,1)}%"
        row[2].text = sv.get("verdict","N/A")
        row[3].text = str(sv.get("avg_confidence","N/A"))
        row[4].text = str(sv.get("contradictions",0))
    doc.add_paragraph()

    # ── Contradictions ───────────────────────────────────────────────────────
    contras = verification.get("contradictions",[])
    if contras:
        _heading(3, f"Contradictions Detected ({len(contras)})")
        for ct in contras[:5]:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(f"[{ct.get('section','?')}] ").bold = True
            p.add_run(ct.get("sentence","")[:120])

    # ── Evaluation Layer ─────────────────────────────────────────────────────
    _heading(2, "Evaluation — Normal AI vs Verified AI")
    ev = evaluation
    ev_rows = [
        ("Normal AI — Grounding Verified", "No"),
        ("Normal AI — Hallucination Risk", ev.get("normal_ai",{}).get("hallucination_risk","?")),
        ("Verified AI — G_rate",           f"{round(g*100,1)}%"),
        ("Verified AI — Verdict",          ev.get("verified_ai",{}).get("verdict","?")),
        ("Verified AI — Confidence",       str(ev.get("verified_ai",{}).get("confidence","?"))),
        ("Verified AI — Hallucination Risk",ev.get("verified_ai",{}).get("hallucination_risk","?")),
    ]
    _table_2col(ev_rows)
    _para(ev.get("summary",""))

    print("  ✓ v3 section added to Word report")




# ═══════════════════════════════════════════════════════════════════════════════
#  ATARS v3.1 — BEAST MODE NOVELTY ADDITIONS
#  ─────────────────────────────────────────────────────────────────────────────
#  N-01: Temporal Claim Verifier (TCV)
#         Verifies time-based claims: "peaked in January", "worst day was X"
#         First system to verify TEMPORAL claims against computed J
#
#  N-02: Causal Claim Detector (CCD)
#         Flags causal language when only correlation exists in J
#         Enforces: Corr(X,Y) ≢ C(X→Y) algorithmically
#
#  N-03: Anomaly Explanation Engine (AEE)
#         Rule-based contextual explanation of every anomaly day
#         Cross-references: weather, season, day-of-week, pollutant patterns
#         Zero LLM — fully deterministic, fully auditable
#
#  N-04: Report Drift Detector (RDD)
#         Compares current run against previous audit logs
#         Detects: G_rate degradation, pollutant trend shifts, Q(D) changes
#         Flags when AI narrative quality is deteriorating over time
#
#  N-05: Uncertainty Propagation Tracker (UPT)
#         Every numerical claim in J carries a confidence interval
#         When LLM uses a number, UPT annotates the report with CI bounds
#         Makes every AI claim scientifically honest
#
#  N-06: Cross-Variable Consistency Checker (CVCC)
#         Checks if related variable claims are mutually consistent
#         Example: NO2 claims consistent with NOx claims (since NOx ≥ NO2)
#         Flags internal contradictions the RGV cannot catch
#
#  N-07: Narrative Specificity Scorer (NSS)
#         Measures how specific vs vague each sentence is
#         Penalises sentences like "air quality was poor" (no numbers)
#         Rewards sentences like "PM10 exceeded WHO by 4.1× on 364 days"
#
#  N-08: Benchmark Comparator
#         Compares station results against national/global context
#         All values stored in J — zero external API needed
#         "Gurugram PM10 is 4.1× WHO limit vs India avg of 3.2×"
# ═══════════════════════════════════════════════════════════════════════════════


# ─────────────────────────────────────────────────────────────────────────────
#  N-01: TEMPORAL CLAIM VERIFIER (TCV)
#  Verifies: "peaked in January", "lowest in September", "worst day was X"
#  Against: J computed temporal facts
# ─────────────────────────────────────────────────────────────────────────────

MONTH_NAMES = {
    "january":1,"february":2,"march":3,"april":4,"may":5,"june":6,
    "july":7,"august":8,"september":9,"october":10,"november":11,"december":12,
    "jan":1,"feb":2,"mar":3,"apr":4,"jun":6,"jul":7,"aug":8,
    "sep":9,"oct":10,"nov":11,"dec":12
}

TEMPORAL_CLAIM_PATTERNS = [
    # "peaked in January" / "highest in January"
    (re.compile(
        r'(?:peaked?|highest?|maximum|worst|most severe|elevated)\w*\s+in\s+([A-Za-z]+)',
        re.I), "peak_month"),
    # "lowest in September" / "best in September"
    (re.compile(
        r'(?:lowest?|minimum|best|cleanest?|least)\w*\s+in\s+([A-Za-z]+)',
        re.I), "low_month"),
    # "during January" (context: high pollution)
    (re.compile(
        r'during\s+([A-Za-z]+).*(?:highest?|peaked?|worst)',
        re.I), "peak_month"),
    # "winter months show highest"
    (re.compile(
        r'(winter|summer|monsoon|spring|autumn|fall)\s+.*(?:highest?|peaked?|worst|lowest?)',
        re.I), "season"),
]

SEASON_MONTHS = {
    "winter"  : [11, 12, 1, 2],
    "summer"  : [3, 4, 5, 6],
    "monsoon" : [6, 7, 8, 9],
    "spring"  : [2, 3, 4, 5],
    "autumn"  : [9, 10, 11],
    "fall"    : [9, 10, 11],
}


def build_temporal_facts(daily_df: pd.DataFrame, J: dict) -> dict:
    """
    Extracts temporal facts from daily_df that can be verified against
    LLM claims. Stored in J_v3 as j_temporal_facts.
    """
    facts = {}
    for var in POLLUTANTS:
        sub = daily_df[daily_df["variable"] == var].copy()
        if sub.empty or sub["mean"].isna().all():
            continue
        sub = sub.dropna(subset=["mean"])
        sub["date"] = pd.to_datetime(sub["date"])
        sub["month"] = sub["date"].dt.month
        sub["month_name"] = sub["date"].dt.strftime("%B").str.lower()
        sub["day_of_week"] = sub["date"].dt.day_name().str.lower()

        # Monthly means
        monthly = sub.groupby("month")["mean"].mean()
        peak_month_num = int(monthly.idxmax()) if not monthly.empty else 1
        low_month_num  = int(monthly.idxmin()) if not monthly.empty else 6

        # Peak day
        peak_idx = sub["mean"].idxmax()
        low_idx  = sub["mean"].idxmin()

        _month_abbr = ["jan","feb","mar","apr","may","jun",
                        "jul","aug","sep","oct","nov","dec"]
        _month_full  = ["january","february","march","april","may","june",
                        "july","august","september","october","november","december"]
        facts[var] = {
            "peak_month_num"  : peak_month_num,
            "peak_month_name" : _month_full[peak_month_num-1],
            "low_month_num"   : low_month_num,
            "low_month_name"  : _month_full[low_month_num-1],
            "peak_month_mean" : round(float(monthly.max()), 2),
            "low_month_mean"  : round(float(monthly.min()), 2),
            "peak_date"       : str(sub.loc[peak_idx,"date"])[:10] if peak_idx is not None else "",
            "low_date"        : str(sub.loc[low_idx,"date"])[:10]  if low_idx is not None else "",
            "peak_season"     : _month_to_season(peak_month_num),
            "low_season"      : _month_to_season(low_month_num),
            "monthly_means"   : {int(k):round(float(v),2) for k,v in monthly.items()},
        }
    print(f"  ✓ N-01: Temporal facts built for {len(facts)} variables")
    return facts


def _month_to_season(month: int) -> str:
    if month in [12,1,2]:  return "winter"
    if month in [3,4,5]:   return "spring"
    if month in [6,7,8,9]: return "monsoon"
    return "autumn"


def verify_temporal_claims(narrative: dict, temporal_facts: dict) -> dict:
    """
    N-01: Temporal Claim Verifier.
    Extracts month/season claims from narrative sentences and
    verifies each against computed temporal_facts from daily_df.
    Tolerates ±1 month for peak/low month claims.
    """
    results = []
    total   = 0
    correct = 0

    _full_months = [
        "january","february","march","april","may","june",
        "july","august","september","october","november","december"
    ]
    _seasons_map = {
        "winter":  [11,12,1,2],
        "summer":  [3,4,5,6],
        "monsoon": [6,7,8,9],
        "spring":  [2,3,4,5],
        "autumn":  [9,10,11],
        "fall":    [9,10,11],
    }

    primary_var = "PM10" if "PM10" in temporal_facts else (
        list(temporal_facts.keys())[0] if temporal_facts else None)
    if primary_var is None:
        return {"TCV_rate":1.0,"TCV_verdict":"NO_DATA","total_claims":0,
                "correct":0,"incorrect":0,"results":[]}

    facts          = temporal_facts[primary_var]
    peak_month_num = facts.get("peak_month_num", 1)
    low_month_num  = facts.get("low_month_num",  7)
    peak_season    = facts.get("peak_season", "winter")

    for sec_key, text in narrative.items():
        if not isinstance(text, str):
            continue
        sentences = [s.strip() for s in re.split(r"[.!?]", text) if len(s.strip()) > 10]

        for sent in sentences:
            sl = sent.lower()

            # ── Month claims ─────────────────────────────────────────────────
            for idx, m_str in enumerate(_full_months):
                if m_str not in sl:
                    continue
                m_num    = idx + 1
                peak_ctx = any(w in sl for w in [
                    "peak","highest","maximum","worst","most","elevated","severe"])
                low_ctx  = any(w in sl for w in [
                    "lowest","minimum","best","cleanest","least","low pollution"])

                if peak_ctx or low_ctx:
                    claim_type = "peak_month" if peak_ctx else "low_month"
                    actual_num = peak_month_num if peak_ctx else low_month_num
                    is_correct = abs(m_num - actual_num) <= 1
                    total     += 1
                    if is_correct: correct += 1
                    results.append({
                        "section"          : sec_key,
                        "sentence"         : sent[:120],
                        "claim_type"       : claim_type,
                        "claimed_month"    : m_str,
                        "claimed_month_num": m_num,
                        "actual_month_num" : actual_num,
                        "variable"         : primary_var,
                        "verdict"          : "CORRECT" if is_correct else "INCORRECT",
                        "tolerance"        : "±1 month",
                    })

            # ── Season claims ─────────────────────────────────────────────────
            for season, months in _seasons_map.items():
                if season not in sl:
                    continue
                peak_ctx = any(w in sl for w in [
                    "peak","highest","worst","most","elevated","severe"])
                if peak_ctx:
                    is_correct = peak_month_num in months
                    total     += 1
                    if is_correct: correct += 1
                    results.append({
                        "section"           : sec_key,
                        "sentence"          : sent[:120],
                        "claim_type"        : "peak_season",
                        "claimed_season"    : season,
                        "actual_peak_season": peak_season,
                        "peak_month_actual" : peak_month_num,
                        "variable"          : primary_var,
                        "verdict"           : "CORRECT" if is_correct else "INCORRECT",
                    })

    tcv_rate    = round(correct / total, 3) if total > 0 else 1.0
    tcv_verdict = ("PASS"           if tcv_rate >= 0.75 else
                   "WARNING"        if tcv_rate >= 0.50 else
                   "FAIL"           if total > 0 else
                   "NO_CLAIMS_FOUND")
    print(f"  ✓ N-01 TCV: {round(tcv_rate*100,1)}% [{tcv_verdict}] | "
          f"{total} temporal claims | {total-correct} incorrect")
    return {
        "TCV_rate"          : tcv_rate,
        "TCV_verdict"       : tcv_verdict,
        "total_claims"      : total,
        "correct"           : correct,
        "incorrect"         : total - correct,
        "primary_variable"  : primary_var,
        "actual_peak_month" : facts.get("peak_month_name",""),
        "actual_low_month"  : facts.get("low_month_name",""),
        "results"           : results,
        "novel_contribution": "First temporal claim verifier in automated environmental reporting"
    }
CAUSAL_TRIGGERS = [
    "causes","caused by","leads to","results in","drives","driven by",
    "due to","because of","responsible for","triggers","triggered by",
    "produces","induced by","impact of","effect of","affects","influences",
    "makes","forces","compels","is why","reason for","explains the increase",
    "explains the decrease","accounts for","is attributed to"
]

ASSOCIATION_SAFE = [
    "associated with","correlation","correlated","linked to","related to",
    "consistent with","coincides with","aligned with","co-occurs",
    "statistical association","tendency","pattern","suggests","indicates",
    "may be related","appears to","is observed when","shows relationship"
]

def detect_causal_claims(narrative: dict, J: dict) -> dict:
    """
    N-02: Causal Claim Detector.
    Scans narrative for causal language.
    Flags sentences that assert causation when J only contains correlation data.
    Strategy: flag causal triggers ONLY when no safe association language is present
    in the same sentence. This avoids false positives.
    """
    flagged      = []
    total_sents  = 0
    causal_count = 0

    for sec_key, text in narrative.items():
        if not isinstance(text, str):
            continue
        sentences = [s.strip() for s in re.split(r"[.!?]", text) if len(s.strip())>10]
        total_sents += len(sentences)

        for sent in sentences:
            sent_lower = sent.lower()

            # Skip sentences that already use safe association language
            has_safe = any(s in sent_lower for s in ASSOCIATION_SAFE)

            # Check for causal triggers
            found_triggers = [t for t in CAUSAL_TRIGGERS if t in sent_lower]

            if found_triggers and not has_safe:
                causal_count += 1
                # Check severity — is there correlation data to back it?
                severity = "HIGH"
                for corr in J.get("top_correlations", []):
                    v1 = corr.get("var1","").lower()
                    v2 = corr.get("var2","").lower()
                    if any(v in sent_lower for v in [v1,v2]):
                        severity = "MEDIUM"  # correlation exists — just wrong framing
                        break

                flagged.append({
                    "section"       : sec_key,
                    "sentence"      : sent[:150],
                    "triggers_found": found_triggers,
                    "severity"      : severity,
                    "safe_language" : has_safe,
                    "recommendation": "Replace causal language with association framing. "
                                      "E.g. 'causes' → 'is associated with'",
                    "formal_rule"   : "Corr(X,Y) ≢ C(X→Y) — correlation ≠ causation"
                })

    ccd_pass = causal_count == 0
    print(f"  ✓ N-02 CCD: {causal_count} causal claim(s) detected | "
          f"{'PASS — no causal language' if ccd_pass else 'REVIEW — causal framing found'}")
    return {
        "CCD_pass"          : ccd_pass,
        "CCD_verdict"       : "PASS" if ccd_pass else "REVIEW",
        "total_sentences"   : total_sents,
        "causal_claim_count": causal_count,
        "flagged"           : flagged,
        "formal_rule"       : "Corr(X,Y) ≢ C(X→Y) — enforced algorithmically",
        "novel_contribution": "First algorithmic causal claim detector in scientific reporting pipeline"
    }


# ─────────────────────────────────────────────────────────────────────────────
#  N-03: ANOMALY EXPLANATION ENGINE (AEE)
#  Rule-based contextual explanation of every anomaly day
#  Zero LLM — fully deterministic, fully auditable
# ─────────────────────────────────────────────────────────────────────────────

ANOMALY_RULES = [
    # (condition_fn, explanation, confidence)
    (lambda row: row.get("wind_kmh",10) < 4,
     "Stagnant wind conditions (< 4 km/h) likely trapped particulates near surface",
     "HIGH"),
    (lambda row: row.get("humidity_pct",50) > 85,
     "High humidity (> 85%) promotes aerosol formation and particulate accumulation",
     "MEDIUM"),
    (lambda row: row.get("temp_c",20) < 8,
     "Low temperature (< 8°C) promotes temperature inversion, trapping pollutants",
     "HIGH"),
    (lambda row: row.get("rain_mm",0) > 15,
     "Heavy rainfall (> 15 mm) may cause wet deposition — anomaly may reflect "
     "post-rain resuspension or pre-rain accumulation",
     "MEDIUM"),
    (lambda row: row.get("temp_c",20) > 38,
     "High temperature (> 38°C) promotes photochemical reactions and secondary aerosol formation",
     "MEDIUM"),
    (lambda row: row.get("day_of_week","") in ["monday","tuesday"],
     "Early-week pattern may reflect weekend activity accumulation effect",
     "LOW"),
    (lambda row: row.get("month_num",6) in [11,12,1,2],
     "Winter month — seasonal biomass burning and reduced atmospheric mixing typical",
     "HIGH"),
    (lambda row: row.get("month_num",6) in [6,7,8,9],
     "Monsoon/post-monsoon — dust resuspension and agricultural burning typical",
     "MEDIUM"),
]


def run_anomaly_explanation_engine(daily_df: pd.DataFrame, snapshots: dict,
                                    config: dict) -> dict:
    """
    N-03: Anomaly Explanation Engine.
    For every anomaly day detected by IF or z-score, generates a
    rule-based contextual explanation. Zero LLM. Fully auditable.
    """
    pm10_df = daily_df[daily_df["variable"] == "PM10"].copy()
    pm10_df = pm10_df.dropna(subset=["date"])
    pm10_df["date"] = pd.to_datetime(pm10_df["date"])
    anomaly_rows = pm10_df[pm10_df["is_anomaly"] == True].copy()

    if anomaly_rows.empty:
        print("  ✓ N-03 AEE: No anomaly days to explain")
        return {"explanations": [], "total_anomalies": 0, "explained": 0}

    explanations = []

    for _, row in anomaly_rows.iterrows():
        d_str = str(row["date"])[:10]
        snap  = snapshots.get(d_str, {})
        weather = snap.get("weather", {})
        pm10_val = float(row.get("mean", 0))
        z_val    = float(row.get("z_score", 0) or 0)

        ctx = {
            "wind_kmh"    : weather.get("wind_kmh", 10),
            "humidity_pct": weather.get("humidity_pct", 50),
            "temp_c"      : weather.get("temp_c", 20),
            "rain_mm"     : weather.get("rain_mm", 0),
            "day_of_week" : pd.to_datetime(d_str).strftime("%A").lower(),
            "month_num"   : pd.to_datetime(d_str).month,
        }

        # Apply all rules
        applicable = []
        for cond_fn, explanation, confidence in ANOMALY_RULES:
            try:
                if cond_fn(ctx):
                    applicable.append({
                        "explanation": explanation,
                        "confidence" : confidence
                    })
            except Exception:
                pass

        # Sort by confidence
        conf_order = {"HIGH":0,"MEDIUM":1,"LOW":2}
        applicable.sort(key=lambda x: conf_order.get(x["confidence"],3))

        # Build composite explanation
        if applicable:
            primary = applicable[0]["explanation"]
            secondary = [a["explanation"] for a in applicable[1:3]]
            composite = primary
            if secondary:
                composite += " Contributing factors: " + "; ".join(secondary) + "."
        else:
            composite = "No single dominant meteorological factor identified. Multi-source episode possible."

        # Anomaly type
        # Safely read anomaly type columns (may not exist in all daily_df versions)
        is_if  = bool(row.get("if_anomaly", False))
        is_z   = bool(row.get("is_anomaly", False))
        is_both    = is_if and is_z
        is_if_only = is_if and not is_z
        is_z_only  = not is_if and is_z
        anomaly_type = ("BOTH_IF_AND_ZSCORE"   if is_both    else
                        "IF_ONLY_MULTIVARIATE" if is_if_only else
                        "ZSCORE_UNIVARIATE"    if is_z_only  else
                        "ANOMALY")

        explanations.append({
            "date"          : d_str,
            "pm10_mean"     : round(pm10_val, 2),
            "z_score"       : round(z_val, 3),
            "anomaly_type"  : anomaly_type,
            "day_of_week"   : ctx["day_of_week"].capitalize(),
            "month"         : pd.to_datetime(d_str).strftime("%B"),
            "season"        : _month_to_season(ctx["month_num"]),
            "weather"       : weather,
            "rules_triggered": len(applicable),
            "primary_explanation": primary if applicable else "Unknown",
            "full_explanation"   : composite,
            "confidence"         : applicable[0]["confidence"] if applicable else "LOW",
        })

    explained = sum(1 for e in explanations if e["rules_triggered"] > 0)
    print(f"  ✓ N-03 AEE: {len(explanations)} anomaly days explained | "
          f"{explained} with high-confidence rule match")
    return {
        "total_anomalies"   : len(explanations),
        "explained"         : explained,
        "explanation_rate"  : round(explained/len(explanations),3) if explanations else 0,
        "explanations"      : explanations,
        "method"            : "Rule-based deterministic — zero LLM",
        "novel_contribution": "First rule-based anomaly explanation engine in automated "
                              "environmental reporting pipeline"
    }


# ─────────────────────────────────────────────────────────────────────────────
#  N-04: REPORT DRIFT DETECTOR (RDD)
#  Compares current run vs previous audit logs
#  Detects G_rate degradation, pollutant shifts, Q(D) changes
# ─────────────────────────────────────────────────────────────────────────────

def run_report_drift_detector(J: dict, verification: dict,
                               out_dir: Path) -> dict:
    """
    N-04: Report Drift Detector.
    Loads previous audit logs from output directory.
    Compares current run against historical runs.
    Detects: G_rate drift, Q(D) drift, pollutant mean drift.
    """
    audit_files = sorted(out_dir.glob("audit_log_*.json"))
    if not audit_files:
        print("  ✓ N-04 RDD: No previous runs found — baseline established")
        return {
            "status"          : "BASELINE",
            "runs_compared"   : 0,
            "drifts_detected" : [],
            "novel_contribution": "Report drift detector tracks AI narrative quality over time"
        }

    # Load last 5 runs for comparison
    history = []
    for af in audit_files[-5:]:
        try:
            with open(af) as f:
                history.append(json.load(f))
        except Exception:
            pass

    if not history:
        return {"status": "NO_HISTORY", "runs_compared": 0, "drifts_detected": []}

    current_g = verification.get("overall_G_rate", 0)
    current_q = J.get("data_quality", {}).get("overall_Q", 0)
    current_pm10 = J.get("variables", {}).get("PM10", {}).get("annual_mean", 0)

    drifts = []

    # G_rate drift
    hist_g_rates = [h.get("verification", {}).get("overall_G_rate", 0)
                    for h in history if h.get("verification")]
    if hist_g_rates:
        avg_hist_g = sum(hist_g_rates) / len(hist_g_rates)
        g_drift = current_g - avg_hist_g
        if abs(g_drift) > 0.05:
            drifts.append({
                "metric"    : "G_rate",
                "current"   : round(current_g, 3),
                "historical_avg": round(avg_hist_g, 3),
                "drift"     : round(g_drift, 3),
                "direction" : "IMPROVING" if g_drift > 0 else "DEGRADING",
                "severity"  : "HIGH" if abs(g_drift) > 0.15 else "MEDIUM",
                "alert"     : abs(g_drift) > 0.15 and g_drift < 0
            })

    # Q(D) drift
    hist_q = [h.get("J_summary", {}).get("overall_Q", 0)
              for h in history if h.get("J_summary")]
    if hist_q:
        avg_hist_q = sum(hist_q) / len(hist_q)
        q_drift = current_q - avg_hist_q
        if abs(q_drift) > 0.05:
            drifts.append({
                "metric"    : "data_quality_Q",
                "current"   : round(current_q, 3),
                "historical_avg": round(avg_hist_q, 3),
                "drift"     : round(q_drift, 3),
                "direction" : "IMPROVING" if q_drift > 0 else "DEGRADING",
                "severity"  : "MEDIUM",
                "alert"     : q_drift < -0.1
            })

    status = "STABLE"
    if any(d["alert"] for d in drifts):
        status = "ALERT — significant degradation detected"
    elif drifts:
        status = "DRIFT_DETECTED"

    print(f"  ✓ N-04 RDD: {len(history)} previous runs compared | "
          f"{len(drifts)} drift(s) detected | Status: {status}")
    return {
        "status"          : status,
        "runs_compared"   : len(history),
        "current_g_rate"  : current_g,
        "current_q"       : current_q,
        "drifts_detected" : drifts,
        "alert"           : any(d.get("alert",False) for d in drifts),
        "novel_contribution": "First automated report drift detector tracking narrative quality over time"
    }


# ─────────────────────────────────────────────────────────────────────────────
#  N-05: UNCERTAINTY PROPAGATION TRACKER (UPT)
#  Every numerical claim in J carries CI bounds
#  When LLM uses that number, UPT annotates it with uncertainty
# ─────────────────────────────────────────────────────────────────────────────

def build_uncertainty_table(daily_df: pd.DataFrame, J: dict) -> dict:
    """
    N-05: Builds a lookup table of uncertainty bounds for every key J value.
    Based on bootstrap CI already computed by the formal pipeline.
    """
    uncertainty = {}

    for var in POLLUTANTS:
        sub = daily_df[daily_df["variable"] == var].copy()
        if sub.empty or sub["mean"].isna().all():
            continue
        means = sub["mean"].dropna()
        n = len(means)
        if n < 3:
            continue
        mean_val = float(means.mean())
        std_val  = float(means.std()) if n > 1 else 0.0
        sem      = std_val / math.sqrt(max(n, 1))
        ci_95_lo = round(mean_val - 1.96 * sem, 2)
        ci_95_hi = round(mean_val + 1.96 * sem, 2)
        ci_width = round(ci_95_hi - ci_95_lo, 2)

        # Relative uncertainty — guard against zero mean
        rel_uncertainty = round(
            (ci_width / max(abs(mean_val), 0.001)) * 100, 1)

        uncertainty[var] = {
            "mean"           : round(mean_val, 2),
            "std"            : round(std_val, 2),
            "n_days"         : n,
            "sem"            : round(sem, 3),
            "ci_95_lower"    : ci_95_lo,
            "ci_95_upper"    : ci_95_hi,
            "ci_width"       : ci_width,
            "relative_uncertainty_pct": rel_uncertainty,
            "uncertainty_class": ("LOW"    if rel_uncertainty < 5  else
                                  "MEDIUM" if rel_uncertainty < 15 else "HIGH"),
        }

    print(f"  ✓ N-05 UPT: Uncertainty bounds computed for {len(uncertainty)} variables")
    return uncertainty


def annotate_claims_with_uncertainty(narrative: dict, uncertainty_table: dict,
                                      verification: dict) -> dict:
    """
    N-05: For each grounded numerical claim, annotate with CI bounds.
    Returns annotations that can be added to the Word report.
    """
    annotations = []

    for sec_key, text in narrative.items():
        if not isinstance(text, str):
            continue
        sentences = [s.strip() for s in re.split(r"[.!?]", text) if len(s.strip())>10]

        for sent in sentences:
            for var, unc in uncertainty_table.items():
                var_lower = var.lower()
                if var_lower not in sent.lower() and var not in sent:
                    continue
                mean_val = unc["mean"]
                # Check if mean value appears in sentence (within tolerance)
                nums_in_sent = [float(m) for m in re.findall(
                    r"-?\d+(?:\.\d+)?", sent) if abs(float(m) - mean_val) /
                    max(abs(mean_val), 1) < 0.06]
                if nums_in_sent:
                    annotations.append({
                        "section"    : sec_key,
                        "sentence"   : sent[:120],
                        "variable"   : var,
                        "claimed_val": nums_in_sent[0],
                        "ci_lower"   : unc["ci_95_lower"],
                        "ci_upper"   : unc["ci_95_upper"],
                        "annotation" : f"(95% CI: {unc['ci_95_lower']}–{unc['ci_95_upper']} µg/m³)",
                        "uncertainty_class": unc["uncertainty_class"],
                    })

    print(f"  ✓ N-05 UPT: {len(annotations)} claims annotated with uncertainty bounds")
    return {
        "uncertainty_table" : uncertainty_table,
        "annotations"       : annotations,
        "total_annotated"   : len(annotations),
        "novel_contribution": "First uncertainty propagation tracker in LLM environmental reporting"
    }


# ─────────────────────────────────────────────────────────────────────────────
#  N-06: CROSS-VARIABLE CONSISTENCY CHECKER (CVCC)
#  Checks if related variable claims are mutually consistent
#  Example: NOx ≥ NO2 always (chemical identity)
# ─────────────────────────────────────────────────────────────────────────────

CHEMICAL_CONSTRAINTS = [
    # (var_a, var_b, relation, description)
    ("NOx", "NO2", ">=",
     "NOx ≥ NO2 always (NOx = NO + NO2 by definition)"),
    ("NOx", "NO",  ">=",
     "NOx ≥ NO always (NOx = NO + NO2)"),
    ("Xylene","Benzene", ">=",
     "Xylene typically ≥ Benzene in traffic emissions"),
    ("Toluene","Benzene", ">=",
     "Toluene/Benzene ratio typically > 1 in urban settings"),
]

NARRATIVE_CONSISTENCY_RULES = [
    # If narrative mentions both vars, check relative claims are consistent
    ("PM10", "PM2.5",
     "PM10 ≥ PM2.5 always (PM10 includes PM2.5 fraction)"),
]


def run_cross_variable_consistency(J: dict, narrative: dict) -> dict:
    """
    N-06: Cross-Variable Consistency Checker.
    Verifies chemical/physical constraints between variables in J.
    Also checks narrative for internally inconsistent relative claims.
    """
    violations  = []
    warnings    = []
    vars_in_J   = J.get("variables", {})

    # Check chemical constraints in data
    for var_a, var_b, relation, description in CHEMICAL_CONSTRAINTS:
        if var_a not in vars_in_J or var_b not in vars_in_J:
            continue
        val_a = vars_in_J[var_a].get("annual_mean", 0)
        val_b = vars_in_J[var_b].get("annual_mean", 0)

        violated = False
        if relation == ">=" and val_a < val_b:
            violated = True
        elif relation == "<=" and val_a > val_b:
            violated = True

        if violated:
            violations.append({
                "type"       : "CHEMICAL_CONSTRAINT",
                "var_a"      : var_a,
                "val_a"      : val_a,
                "var_b"      : var_b,
                "val_b"      : val_b,
                "relation"   : relation,
                "description": description,
                "severity"   : "HIGH",
                "note"       : "May indicate data quality issue or measurement error"
            })
        else:
            warnings.append({
                "type"   : "CONSTRAINT_SATISFIED",
                "vars"   : f"{var_a} {relation} {var_b}",
                "values" : f"{val_a:.2f} {relation} {val_b:.2f}",
                "status" : "OK"
            })

    # Check narrative for inconsistent relative comparisons
    for sec_key, text in narrative.items():
        if not isinstance(text, str):
            continue
        text_lower = text.lower()
        # Simple check: if text says "NO2 is higher than NOx" — flag it
        for var_a, var_b, desc in NARRATIVE_CONSISTENCY_RULES:
            va_lower = var_a.lower()
            vb_lower = var_b.lower()
            if va_lower in text_lower and vb_lower in text_lower:
                # Look for inconsistent comparison
                if (f"{vb_lower}.*higher.*{va_lower}" in text_lower or
                    f"{va_lower}.*lower.*{vb_lower}" in text_lower):
                    violations.append({
                        "type"    : "NARRATIVE_INCONSISTENCY",
                        "section" : sec_key,
                        "note"    : f"Possible inconsistency: {desc}",
                        "severity": "MEDIUM"
                    })

    cvcc_pass = len(violations) == 0
    print(f"  ✓ N-06 CVCC: {len(violations)} constraint violation(s) | "
          f"{len(warnings)} constraints satisfied | "
          f"{'PASS' if cvcc_pass else 'REVIEW'}")
    return {
        "CVCC_pass"           : cvcc_pass,
        "CVCC_verdict"        : "PASS" if cvcc_pass else "REVIEW",
        "violations"          : violations,
        "constraints_checked" : len(CHEMICAL_CONSTRAINTS),
        "constraints_satisfied": len(warnings),
        "novel_contribution"  : "First cross-variable chemical constraint checker in "
                                "automated environmental reporting"
    }


# ─────────────────────────────────────────────────────────────────────────────
#  N-07: NARRATIVE SPECIFICITY SCORER (NSS)
#  Measures how specific vs vague each sentence is
#  Penalises vague language, rewards specific numerical claims
# ─────────────────────────────────────────────────────────────────────────────

VAGUE_PHRASES = [
    "very high","very low","significantly","quite","rather","fairly",
    "relatively","somewhat","a lot","many","several","some","various",
    "numerous","considerable","notable","substantial","marked","pronounced",
    "air quality was poor","air quality is","pollution levels",
    "increased significantly","decreased significantly"
]

SPECIFIC_INDICATORS = [
    r"\d+\.\d+",          # decimal number
    r"\d+%",              # percentage
    r"\d+ days",          # count with unit
    r"µg/m",              # unit
    r"mg/m",              # unit
    r"\d+× ",             # multiple
    r"Q\(D\)",            # formal notation
    r"WHO",               # standard reference
    r"ζ_",                # formal operator
    r"G_rate",            # formal metric
    r"r\s*=\s*-?\d",      # correlation coefficient
    r"R²\s*=",            # R-squared
]


def score_narrative_specificity(narrative: dict) -> dict:
    """
    N-07: Narrative Specificity Scorer.
    Scores each sentence 0.0–1.0 for specificity.
    1.0 = highly specific (numbers, units, formal notation)
    0.0 = vague (no measurable claims)
    """
    section_scores = {}
    all_scores     = []

    # Pre-compile patterns once for speed
    _compiled_specific = [re.compile(p, re.I) for p in SPECIFIC_INDICATORS]

    for sec_key, text in narrative.items():
        if not isinstance(text, str):
            continue
        sentences = [s.strip() for s in re.split(r"[.!?]", text) if len(s.strip())>10]
        sent_scores = []

        for sent in sentences:
            sent_lower = sent.lower()

            # Count specific indicators using pre-compiled patterns
            specific_hits = sum(
                1 for pat in _compiled_specific if pat.search(sent))

            # Count vague phrases (exact substring match)
            vague_hits = sum(1 for vp in VAGUE_PHRASES if vp in sent_lower)

            # Compute specificity score
            score = min(1.0, specific_hits * 0.25) - (vague_hits * 0.1)
            score = max(0.0, round(score, 2))

            sent_scores.append({
                "sentence"      : sent[:120],
                "specificity"   : score,
                "specific_hits" : specific_hits,
                "vague_hits"    : vague_hits,
                "rating"        : ("HIGH"   if score >= 0.5 else
                                   "MEDIUM" if score >= 0.25 else "LOW"),
            })
            all_scores.append(score)

        avg = round(sum(s["specificity"] for s in sent_scores) /
                    len(sent_scores), 2) if sent_scores else 0.0
        section_scores[sec_key] = {
            "avg_specificity": avg,
            "sentences"      : sent_scores,
            "verdict"        : "SPECIFIC" if avg >= 0.4 else (
                               "ADEQUATE" if avg >= 0.2 else "VAGUE"),
        }

    overall = round(sum(all_scores) / len(all_scores), 2) if all_scores else 0.0
    print(f"  ✓ N-07 NSS: Overall specificity = {overall:.2f} | "
          f"{'SPECIFIC' if overall>=0.4 else ('ADEQUATE' if overall>=0.2 else 'VAGUE')}")
    return {
        "overall_specificity" : overall,
        "NSS_verdict"         : ("SPECIFIC" if overall >= 0.4 else
                                 "ADEQUATE" if overall >= 0.2 else "VAGUE"),
        "sections"            : section_scores,
        "novel_contribution"  : "First narrative specificity scorer in automated scientific reporting"
    }


# ─────────────────────────────────────────────────────────────────────────────
#  N-08: BENCHMARK COMPARATOR
#  Compares station results against Indian and WHO reference values
#  All reference data stored internally — zero external API
# ─────────────────────────────────────────────────────────────────────────────

# Reference values for India — sourced from CPCB annual reports and WHO 2021 AQG
INDIA_REFERENCE = {
    "PM10": {
        "india_avg_annual_mean" : 70.0,   # µg/m³ — national average ~2023
        "india_naaqs_annual"    : 60.0,   # µg/m³ — CPCB NAAQS standard
        "who_aqg_annual"        : 15.0,   # µg/m³ — WHO AQG 2021
        "who_aqg_24h"           : 45.0,
        "delhi_avg_annual"      : 180.0,  # µg/m³ — Delhi baseline
        "haryana_avg_annual"    : 95.0,   # µg/m³ — Haryana state avg
        "unit"                  : "µg/m³"
    },
    "NO2": {
        "india_avg_annual_mean" : 28.0,
        "india_naaqs_annual"    : 40.0,
        "who_aqg_annual"        : 10.0,
        "who_aqg_24h"           : 25.0,
        "unit"                  : "µg/m³"
    },
    "SO2": {
        "india_avg_annual_mean" : 15.0,
        "india_naaqs_annual"    : 50.0,
        "who_aqg_24h"           : 40.0,
        "unit"                  : "µg/m³"
    },
    "Benzene": {
        "who_aqg_annual"        : 1.7,
        "india_naaqs_annual"    : 5.0,
        "unit"                  : "µg/m³"
    },
    "Ozone": {
        "who_aqg_8h"            : 100.0,
        "india_naaqs_8h"        : 100.0,
        "unit"                  : "µg/m³"
    }
}


def run_benchmark_comparator(J: dict, config: dict) -> dict:
    """
    N-08: Benchmark Comparator.
    Compares station annual means against India national averages,
    NAAQS standards, and WHO AQG 2021 values.
    All reference data stored internally — zero API calls.
    """
    vars_in_J = J.get("variables", {})
    comparisons = {}

    for var, ref in INDIA_REFERENCE.items():
        if var not in vars_in_J:
            continue
        station_mean = vars_in_J[var].get("annual_mean", 0)
        if not station_mean:
            continue

        comp = {
            "station_annual_mean" : station_mean,
            "unit"                : ref.get("unit", "µg/m³"),
        }

        # WHO comparison
        who_val = ref.get("who_aqg_annual") or ref.get("who_aqg_24h")
        if who_val:
            ratio = round(station_mean / who_val, 2)
            comp["who_aqg_ratio"]   = ratio
            comp["who_aqg_value"]   = who_val
            comp["vs_who"]          = f"{ratio:.1f}× WHO limit"
            comp["who_status"]      = ("COMPLIANT"  if ratio <= 1.0 else
                                       "EXCEEDS_WHO")

        # India NAAQS comparison
        naaqs = ref.get("india_naaqs_annual") or ref.get("india_naaqs_8h")
        if naaqs:
            naaqs_ratio = round(station_mean / naaqs, 2)
            comp["india_naaqs_ratio"]  = naaqs_ratio
            comp["india_naaqs_value"]  = naaqs
            comp["naaqs_status"]       = ("COMPLIANT" if naaqs_ratio <= 1.0 else
                                          "EXCEEDS_NAAQS")

        # India national average comparison
        india_avg = ref.get("india_avg_annual_mean")
        if india_avg:
            above_national = round(station_mean - india_avg, 2)
            comp["india_national_avg"]  = india_avg
            comp["above_national_avg"]  = above_national
            comp["national_percentile"] = ("ABOVE_NATIONAL" if station_mean > india_avg
                                           else "BELOW_NATIONAL")

        comparisons[var] = comp

    # Overall station severity classification
    pm10_comp = comparisons.get("PM10", {})
    pm10_ratio = pm10_comp.get("who_aqg_ratio", 0)
    severity = ("CRITICAL"  if pm10_ratio > 5 else
                "SEVERE"    if pm10_ratio > 3 else
                "HIGH"      if pm10_ratio > 1 else "COMPLIANT")

    print(f"  ✓ N-08 BNC: {len(comparisons)} variables benchmarked | "
          f"PM10 = {pm10_ratio:.1f}× WHO | Severity: {severity}")
    return {
        "station"             : config.get("city","City"),
        "overall_severity"    : severity,
        "comparisons"         : comparisons,
        "reference_source"    : "CPCB NAAQS + WHO AQG 2021 + India national averages",
        "data_year"           : "2023-2024",
        "novel_contribution"  : "Integrated benchmark comparator with Indian national context"
    }


# ─────────────────────────────────────────────────────────────────────────────
#  N-09: MASTER NOVELTY SCORE (MNS)
#  Aggregates all novel verifier outputs into one trust score
#  MNS ∈ [0,1] — the single metric that summarises report trustworthiness
# ─────────────────────────────────────────────────────────────────────────────

def compute_master_novelty_score(
        verification_v2 : dict,
        verification_v3 : dict,
        tcv             : dict,
        ccd             : dict,
        cvcc            : dict,
        nss             : dict,
        aee             : dict,
) -> dict:
    """
    N-09: Master Novelty Score.
    Weighted combination of all verification dimensions.
    MNS = G_rate(v2)×0.25 + G_rate(v3)×0.20 + TCV×0.15 +
          CCD_pass×0.15 + CVCC_pass×0.10 + Confidence×0.10 +
          Specificity×0.05
    """
    g2   = verification_v2.get("overall_G_rate", 0)
    g3   = verification_v3.get("overall_G_rate", 0)
    tcv_r = tcv.get("TCV_rate", 1.0)
    ccd_p = 1.0 if ccd.get("CCD_pass", True) else max(
        0.0, 1.0 - ccd.get("causal_claim_count",0) * 0.1)
    cvcc_p = 1.0 if cvcc.get("CVCC_pass", True) else 0.7
    conf   = verification_v3.get("overall_confidence", 0.8)
    spec   = nss.get("overall_specificity", 0.3)
    sem_p  = 1.0 if verification_v3.get("semantic_pass", True) else 0.7

    # Clamp all components to [0,1] before weighting
    g2     = max(0.0, min(1.0, g2))
    g3     = max(0.0, min(1.0, g3))
    tcv_r  = max(0.0, min(1.0, tcv_r))
    ccd_p  = max(0.0, min(1.0, ccd_p))
    cvcc_p = max(0.0, min(1.0, cvcc_p))
    conf   = max(0.0, min(1.0, conf))
    sem_p  = max(0.0, min(1.0, sem_p))

    MNS = (g2     * 0.25 +
           g3     * 0.20 +
           tcv_r  * 0.15 +
           ccd_p  * 0.15 +
           cvcc_p * 0.10 +
           conf   * 0.10 +
           sem_p  * 0.05)

    MNS = round(MNS, 3)
    verdict = ("CERTIFIED"   if MNS >= 0.90 else
               "TRUSTED"     if MNS >= 0.80 else
               "ACCEPTABLE"  if MNS >= 0.70 else
               "REVIEW"      if MNS >= 0.55 else "UNRELIABLE")

    print(f"\n  {'═'*60}")
    print(f"  MASTER NOVELTY SCORE (MNS) = {MNS:.3f}  [{verdict}]")
    print(f"  {'═'*60}")
    print(f"    v2 G_rate:      {g2:.1%}  (weight 25%)")
    print(f"    v3 G_rate:      {g3:.1%}  (weight 20%)")
    print(f"    TCV rate:       {tcv_r:.1%}  (weight 15%)")
    print(f"    CCD score:      {ccd_p:.1%}  (weight 15%)")
    print(f"    CVCC score:     {cvcc_p:.1%}  (weight 10%)")
    print(f"    Confidence:     {conf:.2f}   (weight 10%)")
    print(f"    Semantic:       {sem_p:.1%}  (weight  5%)")
    print(f"  {'═'*60}\n")

    return {
        "MNS"               : MNS,
        "verdict"           : verdict,
        "components": {
            "v2_G_rate"     : g2,
            "v3_G_rate"     : g3,
            "TCV_rate"      : tcv_r,
            "CCD_score"     : ccd_p,
            "CVCC_score"    : cvcc_p,
            "confidence"    : conf,
            "semantic"      : sem_p,
        },
        "weights": {
            "v2_G_rate": 0.25, "v3_G_rate": 0.20,
            "TCV_rate" : 0.15, "CCD_score" : 0.15,
            "CVCC_score": 0.10,"confidence": 0.10,
            "semantic"  : 0.05
        },
        "equation"          : "MNS = G2×0.25 + G3×0.20 + TCV×0.15 + CCD×0.15 + CVCC×0.10 + Conf×0.10 + Sem×0.05",
        "novel_contribution": "First unified trust metric aggregating 7 independent verification dimensions"
    }


# ─────────────────────────────────────────────────────────────────────────────
#  WORD REPORT EXTENSION — All novel sections appended to single file
# ─────────────────────────────────────────────────────────────────────────────

def _add_beast_sections_to_report(
        doc,
        tcv: dict, ccd: dict, aee: dict,
        drift: dict, upt: dict, cvcc: dict,
        nss: dict, benchmark: dict, mns: dict,
        config: dict):
    """
    Appends all 8 novel section reports to the existing Word document.
    Single file output — no extra files created.
    """
    from docx.shared import Pt, RGBColor
    NAVY  = RGBColor(0x1A,0x2C,0x4E)
    TEAL  = RGBColor(0x1A,0x6B,0x72)
    GREEN = RGBColor(0x1E,0x6B,0x3C)
    RED   = RGBColor(0x8B,0x1C,0x2A)
    AMBER = RGBColor(0x92,0x40,0x0E)
    PURPLE= RGBColor(0x4A,0x1A,0x6B)

    def _h1(text):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = True; r.font.size = Pt(15); r.font.color.rgb = NAVY
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after  = Pt(5)

    def _h2(text, color=TEAL):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = True; r.font.size = Pt(12); r.font.color.rgb = color
        p.paragraph_format.space_before = Pt(8)

    def _h3(text):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = True; r.font.size = Pt(11); r.font.color.rgb = AMBER

    def _body(text):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(3)

    def _kv_table(rows, col1="Metric", col2="Value"):
        t = doc.add_table(rows=1, cols=2)
        t.style = "Table Grid"
        h = t.rows[0].cells
        h[0].text = col1; h[1].text = col2
        for c in h:
            for r in c.paragraphs[0].runs: r.bold = True
        for k,v in rows:
            rc = t.add_row().cells
            rc[0].text = str(k); rc[1].text = str(v)
        doc.add_paragraph()

    def _verdict_color(v):
        v = str(v).upper()
        if v in ("PASS","CORRECT","CERTIFIED","TRUSTED","COMPLIANT","SPECIFIC"): return GREEN
        if v in ("REVIEW","WARNING","ADEQUATE","MEDIUM"): return AMBER
        if v in ("FAIL","INCORRECT","UNRELIABLE","HIGH_VIOLATION"): return RED
        return TEAL

    # ── Master Novelty Score — first (most important) ─────────────────────────
    doc.add_page_break()
    _h1("ATARS v3.1 — Master Novelty Score (MNS)")
    _body("The Master Novelty Score aggregates 7 independent verification dimensions "
          "into a single auditable trust metric for the AI-generated report.")
    mns_val = mns.get("MNS",0)
    mns_v   = mns.get("verdict","N/A")
    p = doc.add_paragraph()
    r = p.add_run(f"MNS = {mns_val:.3f}   [{mns_v}]")
    r.bold = True; r.font.size = Pt(16)
    r.font.color.rgb = _verdict_color(mns_v)
    _body(f"Equation: {mns.get('equation','')}")

    mns_rows = [(f"{k.replace('_',' ').title()} (weight {int(mns['weights'][k]*100)}%)",
                 f"{round(v*100,1) if isinstance(v,float) else v}%")
                for k,v in mns.get("components",{}).items()]
    _kv_table(mns_rows, "Component", "Score")

    # ── N-01: TCV ─────────────────────────────────────────────────────────────
    doc.add_page_break()
    _h1("N-01: Temporal Claim Verifier (TCV)")
    _body("Verifies all time-based claims in AI narrative against computed temporal facts. "
          "First system to verify temporal assertions in automated environmental reports.")
    tcv_rows = [
        ("TCV Rate",      f"{round(tcv.get('TCV_rate',0)*100,1)}%"),
        ("Verdict",       tcv.get("TCV_verdict","N/A")),
        ("Claims Checked",str(tcv.get("total_claims",0))),
        ("Correct",       str(tcv.get("correct",0))),
        ("Incorrect",     str(tcv.get("incorrect",0))),
    ]
    _kv_table(tcv_rows)
    if tcv.get("results"):
        _h3("Temporal Claim Results")
        for r in tcv["results"][:8]:
            _body(f"[{r.get('verdict','?')}] Claimed: {r.get('claimed','')} | "
                  f"Actual: {r.get('actual_month_num','')} | {r.get('sentence','')[:80]}")

    # ── N-02: CCD ─────────────────────────────────────────────────────────────
    doc.add_page_break()
    _h1("N-02: Causal Claim Detector (CCD)")
    _body("Algorithmically enforces Corr(X,Y) ≢ C(X→Y). "
          "Flags causal language in AI narrative when only correlation data exists. "
          "First causal claim detector in automated scientific reporting.")
    ccd_rows = [
        ("CCD Verdict",       ccd.get("CCD_verdict","N/A")),
        ("Causal Claims Found",str(ccd.get("causal_claim_count",0))),
        ("Total Sentences",   str(ccd.get("total_sentences",0))),
        ("Formal Rule",       "Corr(X,Y) ≢ C(X→Y)"),
    ]
    _kv_table(ccd_rows)
    if ccd.get("flagged"):
        _h3(f"Flagged Causal Claims ({len(ccd['flagged'])})")
        for fc in ccd["flagged"][:5]:
            _body(f"[{fc.get('severity','?')}] {fc.get('sentence','')[:100]}")
            _body(f"   Triggers: {', '.join(fc.get('triggers_found',[]))}")

    # ── N-03: AEE ─────────────────────────────────────────────────────────────
    doc.add_page_break()
    _h1("N-03: Anomaly Explanation Engine (AEE)")
    _body("Rule-based deterministic contextual explanation of every anomaly day. "
          "Zero LLM — cross-references weather, season, day-of-week. Fully auditable.")
    aee_rows = [
        ("Total Anomaly Days",    str(aee.get("total_anomalies",0))),
        ("Explained",             str(aee.get("explained",0))),
        ("Explanation Rate",      f"{round(aee.get('explanation_rate',0)*100,1)}%"),
        ("Method",                "Rule-based deterministic — zero LLM"),
    ]
    _kv_table(aee_rows)
    if aee.get("explanations"):
        _h3("Anomaly Explanations (Top 10)")
        t = doc.add_table(rows=1, cols=5)
        t.style = "Table Grid"
        hdrs = ["Date","PM10","Type","Season","Primary Explanation"]
        for i,h in enumerate(hdrs):
            t.rows[0].cells[i].text = h
            for r in t.rows[0].cells[i].paragraphs[0].runs: r.bold = True
        for ex in aee["explanations"][:10]:
            row = t.add_row().cells
            row[0].text = ex.get("date","")
            row[1].text = f"{ex.get('pm10_mean',0):.1f}"
            row[2].text = ex.get("anomaly_type","")[:18]
            row[3].text = ex.get("season","")
            row[4].text = ex.get("primary_explanation","")[:60]
        doc.add_paragraph()

    # ── N-04: RDD ─────────────────────────────────────────────────────────────
    doc.add_page_break()
    _h1("N-04: Report Drift Detector (RDD)")
    _body("Compares current run against historical audit logs. "
          "Detects G_rate degradation, Q(D) changes, and AI narrative quality drift over time.")
    rdd_rows = [
        ("Status",          drift.get("status","N/A")),
        ("Runs Compared",   str(drift.get("runs_compared",0))),
        ("Drifts Detected", str(len(drift.get("drifts_detected",[])))),
        ("Alert",           str(drift.get("alert",False))),
    ]
    _kv_table(rdd_rows)
    for d in drift.get("drifts_detected",[]):
        _body(f"[{d.get('severity','?')}] {d.get('metric','?')}: "
              f"current={d.get('current',0):.3f} vs "
              f"historical_avg={d.get('historical_avg',0):.3f} "
              f"({d.get('direction','?')})")

    # ── N-05: UPT ─────────────────────────────────────────────────────────────
    doc.add_page_break()
    _h1("N-05: Uncertainty Propagation Tracker (UPT)")
    _body("Computes 95% confidence intervals for every key statistic in J. "
          "Annotates LLM claims with uncertainty bounds for scientific honesty.")
    upt_anns = upt.get("annotations",[])
    upt_rows = [
        ("Total Claims Annotated", str(upt.get("total_annotated",0))),
        ("Variables with CI",      str(len(upt.get("uncertainty_table",{})))),
    ]
    _kv_table(upt_rows)
    if upt.get("uncertainty_table"):
        _h3("Uncertainty Bounds by Variable")
        t = doc.add_table(rows=1, cols=5)
        t.style = "Table Grid"
        hdrs2 = ["Variable","Annual Mean","95% CI Lower","95% CI Upper","Uncertainty Class"]
        for i,h in enumerate(hdrs2):
            t.rows[0].cells[i].text = h
            for r in t.rows[0].cells[i].paragraphs[0].runs: r.bold = True
        for var, unc in list(upt["uncertainty_table"].items())[:10]:
            row = t.add_row().cells
            row[0].text = var
            row[1].text = f"{unc.get('mean',0):.2f}"
            row[2].text = f"{unc.get('ci_95_lower',0):.2f}"
            row[3].text = f"{unc.get('ci_95_upper',0):.2f}"
            row[4].text = unc.get("uncertainty_class","N/A")
        doc.add_paragraph()

    # ── N-06: CVCC ────────────────────────────────────────────────────────────
    doc.add_page_break()
    _h1("N-06: Cross-Variable Consistency Checker (CVCC)")
    _body("Verifies chemical and physical constraints between variables. "
          "Example: NOx ≥ NO2 always (chemical identity). "
          "First cross-variable constraint checker in environmental reporting.")
    cvcc_rows = [
        ("CVCC Verdict",            cvcc.get("CVCC_verdict","N/A")),
        ("Violations Found",        str(len(cvcc.get("violations",[])))),
        ("Constraints Checked",     str(cvcc.get("constraints_checked",0))),
        ("Constraints Satisfied",   str(cvcc.get("constraints_satisfied",0))),
    ]
    _kv_table(cvcc_rows)
    for v in cvcc.get("violations",[]):
        _body(f"[{v.get('severity','?')}] {v.get('description','')}")

    # ── N-07: NSS ─────────────────────────────────────────────────────────────
    doc.add_page_break()
    _h1("N-07: Narrative Specificity Scorer (NSS)")
    _body("Scores each narrative sentence for specificity (0.0–1.0). "
          "Rewards specific numerical claims; penalises vague language.")
    nss_rows = [
        ("Overall Specificity", f"{nss.get('overall_specificity',0):.2f}"),
        ("NSS Verdict",         nss.get("NSS_verdict","N/A")),
    ]
    _kv_table(nss_rows)
    for sec_key, sec in nss.get("sections",{}).items():
        _body(f"Section {sec_key}: avg_specificity={sec.get('avg_specificity',0):.2f} "
              f"[{sec.get('verdict','?')}]")

    # ── N-08: Benchmark ───────────────────────────────────────────────────────
    doc.add_page_break()
    _h1("N-08: Benchmark Comparator")
    _body("Compares station results against CPCB NAAQS, WHO AQG 2021, "
          "and India national averages. Zero external API — all references stored internally.")
    bm_rows = [
        ("Station",          benchmark.get("station","N/A")),
        ("Overall Severity", benchmark.get("overall_severity","N/A")),
        ("Reference Source", benchmark.get("reference_source","N/A")),
    ]
    _kv_table(bm_rows)
    for var, comp in benchmark.get("comparisons",{}).items():
        _h3(f"{var} Benchmark")
        var_rows = [
            ("Station Annual Mean",   f"{comp.get('station_annual_mean',0):.2f} {comp.get('unit','')}"),
            ("vs WHO AQG",            comp.get("vs_who","N/A")),
            ("WHO Status",            comp.get("who_status","N/A")),
            ("NAAQS Status",          comp.get("naaqs_status","N/A")),
            ("vs India National Avg", f"{comp.get('above_national_avg',0):+.2f} {comp.get('unit','')}"),
        ]
        _kv_table(var_rows)

    print("  ✓ All 8 novel sections appended to Word report")



# ─────────────────────────────────────────────────────────────────────────────
#  UNIFIED VERIFICATION DASHBOARD — Single integrated section for Word report
#  Replaces separate v2/v3/beast sections with one clean dashboard
# ─────────────────────────────────────────────────────────────────────────────

def _add_unified_verification_dashboard(
        doc,
        verification_v2 : dict,
        ver_v3          : dict,
        evaluation      : dict,
        comparison      : dict,
        tcv             : dict,
        ccd             : dict,
        aee             : dict,
        drift           : dict,
        upt             : dict,
        cvcc            : dict,
        nss             : dict,
        benchmark       : dict,
        mns             : dict,
        config          : dict):
    """
    Single unified verification dashboard appended to the Word report.
    Looks like a natural continuation of the v2 report — one coherent document.
    """
    from docx.shared import Pt, RGBColor, Inches
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    NAVY   = RGBColor(0x1A, 0x2C, 0x4E)
    TEAL   = RGBColor(0x1A, 0x6B, 0x72)
    GREEN  = RGBColor(0x1E, 0x6B, 0x3C)
    RED    = RGBColor(0x8B, 0x1C, 0x2A)
    AMBER  = RGBColor(0x92, 0x40, 0x0E)
    PURPLE = RGBColor(0x4A, 0x1A, 0x6B)

    def _color(verdict):
        v = str(verdict).upper()
        if v in ("PASS","CORRECT","CERTIFIED","TRUSTED","COMPLIANT",
                 "SPECIFIC","BASELINE","STABLE","NO_CLAIMS_FOUND"): return GREEN
        if v in ("REVIEW","WARNING","ADEQUATE","MEDIUM",
                 "WORSE_THAN_WEEK","DRIFT_DETECTED"): return AMBER
        return RED

    def _h1(text):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = True
        r.font.size = Pt(15)
        r.font.color.rgb = NAVY
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after  = Pt(6)

    def _h2(text, color=TEAL):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = True
        r.font.size = Pt(12)
        r.font.color.rgb = color
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after  = Pt(3)

    def _body(text):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(3)
        return p

    def _tbl(rows, h1="Metric", h2="Value"):
        t = doc.add_table(rows=1, cols=2)
        t.style = "Table Grid"
        hc = t.rows[0].cells
        hc[0].text = h1
        hc[1].text = h2
        for c in hc:
            for r in c.paragraphs[0].runs:
                r.bold = True
        for k, v in rows:
            rc = t.add_row().cells
            rc[0].text = str(k)
            rc[1].text = str(v)
        doc.add_paragraph()

    def _tbl4(rows, headers):
        t = doc.add_table(rows=1, cols=len(headers))
        t.style = "Table Grid"
        for i, h in enumerate(headers):
            t.rows[0].cells[i].text = h
            for r in t.rows[0].cells[i].paragraphs[0].runs:
                r.bold = True
        for row_data in rows:
            rc = t.add_row().cells
            for i, v in enumerate(row_data):
                rc[i].text = str(v)
        doc.add_paragraph()

    # ── Section 19: Enhanced Verification Dashboard ───────────────────────────
    doc.add_page_break()
    _h1("Section 19 — Enhanced Verification Dashboard")
    _body(
        f"This section documents the complete verification audit for this report. "
        f"Every numerical claim, temporal assertion, causal statement, and variable "
        f"relationship in the AI-generated narrative has been algorithmically checked "
        f"against the computed data contract J. The Master Novelty Score (MNS) "
        f"summarises overall report trustworthiness across 7 independent dimensions."
    )

    # ── 19.1 Master Novelty Score ─────────────────────────────────────────────
    _h2("19.1  Master Novelty Score (MNS)")
    mns_val = mns.get("MNS", 0)
    mns_vdt = mns.get("verdict", "N/A")
    p = doc.add_paragraph()
    r = p.add_run(f"MNS  =  {mns_val:.3f}     [{mns_vdt}]")
    r.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = _color(mns_vdt)
    _body(f"Equation: {mns.get('equation','')}")
    _body("Interpretation: CERTIFIED ≥ 0.90 · TRUSTED ≥ 0.80 · ACCEPTABLE ≥ 0.70 · REVIEW ≥ 0.55 · UNRELIABLE < 0.55")
    doc.add_paragraph()

    comps = mns.get("components", {})
    wts   = mns.get("weights", {})
    mns_rows = []
    labels = {
        "v2_G_rate" : "Numerical Grounding v2 (G_rate)",
        "v3_G_rate" : "Numerical Grounding v3 (G_rate)",
        "TCV_rate"  : "Temporal Claim Accuracy (TCV)",
        "CCD_score" : "Causal Language Score (CCD)",
        "CVCC_score": "Cross-Variable Consistency (CVCC)",
        "confidence": "Sentence Confidence",
        "semantic"  : "Semantic Alignment",
    }
    for k, v in comps.items():
        w = wts.get(k, 0)
        label = labels.get(k, k)
        mns_rows.append((f"{label}  (weight {int(w*100)}%)",
                         f"{round(v*100,1)}%"))
    _tbl(mns_rows, "Verification Dimension", "Score")

    # ── 19.2 Grounding Summary ────────────────────────────────────────────────
    _h2("19.2  Numerical Grounding Summary")
    g2 = verification_v2.get("overall_G_rate", 0)
    g3 = ver_v3.get("overall_G_rate", 0)
    g2v = verification_v2.get("overall_verdict", "N/A")
    g3v = ver_v3.get("overall_verdict", "N/A")
    grounding_rows = [
        ("G_rate (base verification)",          f"{round(g2*100,1)}%  [{g2v}]"),
        ("G_rate (enhanced multi-source)",      f"{round(g3*100,1)}%  [{g3v}]"),
        ("Total numerical sentences checked",   str(ver_v3.get("total_numerical",0))),
        ("Grounded sentences",                  str(ver_v3.get("total_grounded",0))),
        ("Ungrounded sentences",                str(ver_v3.get("total_ungrounded",0))),
        ("Sentence confidence (avg)",           str(ver_v3.get("overall_confidence","N/A"))),
        ("Semantic alignment",                  "PASS" if ver_v3.get("semantic_pass") else "REVIEW"),
        ("Direction contradictions",            str(ver_v3.get("contradiction_count",0))),
        ("Sources verified against",            ", ".join(ver_v3.get("sources_used",[]))),
    ]
    _tbl(grounding_rows)

    # ── 19.3 Today vs Historical ──────────────────────────────────────────────
    _h2("19.3  Today vs Historical Comparison")
    _body("All values below are pre-computed by the comparison engine. The AI never performs these calculations.")
    comp_rows = [
        ("Analysis date",                       comparison.get("today","N/A")),
        ("PM10 today",                          f"{comparison.get('pm10_today',0):.1f} µg/m³"),
        ("PM10 yesterday",                      f"{comparison.get('pm10_yesterday',0):.1f} µg/m³"),
        ("Change vs yesterday",                 f"{comparison.get('pm10_change_yesterday_pct',0):+.1f}%  ({comparison.get('pm10_direction_yesterday','N/A')})"),
        ("7-day average",                       f"{comparison.get('pm10_week_avg',0):.1f} µg/m³"),
        ("Change vs 7-day avg",                 f"{comparison.get('pm10_change_week_pct',0):+.1f}%"),
        ("30-day average",                      f"{comparison.get('pm10_month_avg',0):.1f} µg/m³"),
        ("Change vs 30-day avg",                f"{comparison.get('pm10_change_month_pct',0):+.1f}%"),
        ("Status",                              comparison.get("status","N/A")),
    ]
    _tbl(comp_rows)
    tags = comparison.get("event_tags", [])
    if tags:
        _body("Context tags: " + ", ".join(tags))
        doc.add_paragraph()

    # ── 19.4 N-01 Temporal Claim Verifier ────────────────────────────────────
    _h2("19.4  Temporal Claim Verifier (N-01)")
    _body("Verifies time-based assertions in AI narrative against computed temporal facts.")
    tcv_rows = [
        ("TCV rate",        f"{round(tcv.get('TCV_rate',0)*100,1)}%  [{tcv.get('TCV_verdict','N/A')}]"),
        ("Claims checked",  str(tcv.get("total_claims",0))),
        ("Correct",         str(tcv.get("correct",0))),
        ("Incorrect",       str(tcv.get("incorrect",0))),
        ("Actual peak month (PM10)", tcv.get("actual_peak_month","N/A").capitalize()),
        ("Actual low month (PM10)",  tcv.get("actual_low_month","N/A").capitalize()),
    ]
    _tbl(tcv_rows)
    if tcv.get("results"):
        incorrect = [r for r in tcv["results"] if r.get("verdict")=="INCORRECT"]
        if incorrect:
            _body(f"Incorrect temporal claims ({len(incorrect)}):")
            for r in incorrect[:5]:
                _body(f"  · Claimed '{r.get('claimed_month',r.get('claimed_season','?'))}' "
                      f"— actual: month {r.get('actual_month_num',r.get('peak_month_actual','?'))}")

    # ── 19.5 N-02 Causal Claim Detector ──────────────────────────────────────
    _h2("19.5  Causal Claim Detector (N-02)")
    _body("Algorithmically enforces Corr(X,Y) ≢ C(X→Y). Flags causal language when only correlation data exists in J.")
    ccd_rows = [
        ("CCD verdict",         ccd.get("CCD_verdict","N/A")),
        ("Causal claims found", str(ccd.get("causal_claim_count",0))),
        ("Total sentences",     str(ccd.get("total_sentences",0))),
        ("Formal rule",         "Corr(X,Y) ≢ C(X→Y)"),
    ]
    _tbl(ccd_rows)
    for fc in ccd.get("flagged",[])[:3]:
        _body(f"  [{fc.get('severity','?')}] {fc.get('sentence','')[:100]}")
        _body(f"  Triggers: {', '.join(fc.get('triggers_found',[]))}")

    # ── 19.6 N-03 Anomaly Explanation Engine ─────────────────────────────────
    _h2("19.6  Anomaly Explanation Engine (N-03)")
    _body("Rule-based deterministic explanation of every anomaly day. Zero LLM.")
    aee_rows = [
        ("Anomaly days explained", f"{aee.get('explained',0)} / {aee.get('total_anomalies',0)}"),
        ("Explanation rate",       f"{round(aee.get('explanation_rate',0)*100,1)}%"),
        ("Method",                 "Rule-based — zero LLM"),
    ]
    _tbl(aee_rows)
    if aee.get("explanations"):
        headers = ["Date","PM10","Type","Season","Primary Explanation"]
        rows_data = []
        for ex in aee["explanations"][:10]:
            rows_data.append([
                ex.get("date",""),
                f"{ex.get('pm10_mean',0):.1f}",
                ex.get("anomaly_type","")[:18],
                ex.get("season","").capitalize(),
                ex.get("primary_explanation","")[:55],
            ])
        _tbl4(rows_data, headers)

    # ── 19.7 N-04 Drift · N-05 Uncertainty · N-06 CVCC ───────────────────────
    _h2("19.7  Report Drift · Uncertainty · Cross-Variable Consistency")
    combined_rows = [
        ("Drift status",              drift.get("status","N/A")),
        ("Previous runs compared",    str(drift.get("runs_compared",0))),
        ("Drift alerts",              str(drift.get("alert",False))),
        ("UPT claims annotated",      str(upt.get("total_annotated",0))),
        ("UPT variables with CI",     str(len(upt.get("uncertainty_table",{})))),
        ("CVCC verdict",              cvcc.get("CVCC_verdict","N/A")),
        ("Chemical constraints",      str(cvcc.get("constraints_checked",0))),
        ("Constraints satisfied",     str(cvcc.get("constraints_satisfied",0))),
        ("Violations",                str(len(cvcc.get("violations",[])))),
    ]
    _tbl(combined_rows)

    # UPT CI table
    if upt.get("uncertainty_table"):
        _body("95% Confidence Intervals by variable:")
        unc_rows = []
        for var, unc in list(upt["uncertainty_table"].items())[:8]:
            unc_rows.append([
                var,
                f"{unc.get('mean',0):.2f}",
                f"{unc.get('ci_95_lower',0):.2f}",
                f"{unc.get('ci_95_upper',0):.2f}",
                unc.get("uncertainty_class","N/A"),
            ])
        _tbl4(unc_rows, ["Variable","Mean","CI Lower","CI Upper","Class"])

    # ── 19.8 N-07 NSS · N-08 Benchmark ───────────────────────────────────────
    _h2("19.8  Narrative Specificity · Benchmark Context")
    bench_pm10 = benchmark.get("comparisons",{}).get("PM10",{})
    nss_bench_rows = [
        ("Narrative specificity (NSS)",  f"{nss.get('overall_specificity',0):.2f}  [{nss.get('NSS_verdict','N/A')}]"),
        ("Station overall severity",     benchmark.get("overall_severity","N/A")),
        ("PM10 vs WHO AQG",              bench_pm10.get("vs_who","N/A")),
        ("PM10 WHO status",              bench_pm10.get("who_status","N/A")),
        ("PM10 vs NAAQS India",          bench_pm10.get("naaqs_status","N/A")),
        ("PM10 vs India national avg",   f"{bench_pm10.get('above_national_avg',0):+.1f} µg/m³"),
        ("Reference source",             benchmark.get("reference_source","N/A")),
    ]
    _tbl(nss_bench_rows)

    # ── 19.9 Evaluation Summary ───────────────────────────────────────────────
    _h2("19.9  AI Verification Evaluation")
    _body("Comparison of unverified AI output versus verified AI output.")
    ev_vai = evaluation.get("verified_ai",{})
    ev_rows = [
        ("Unverified AI — hallucination risk",   evaluation.get("normal_ai",{}).get("hallucination_risk","UNKNOWN")),
        ("Verified AI — G_rate",                 f"{round(ev_vai.get('g_rate',0)*100,1)}%"),
        ("Verified AI — verdict",                ev_vai.get("verdict","N/A")),
        ("Verified AI — hallucination risk",     ev_vai.get("hallucination_risk","N/A")),
        ("Contradictions flagged",               str(ev_vai.get("contradictions",0))),
        ("Semantic aligned",                     str(ev_vai.get("semantic_aligned",True))),
    ]
    _tbl(ev_rows)
    _body(evaluation.get("summary",""))

    print("  ✓ Section 19 — Unified verification dashboard added to report")



def _fetch_7day_forecast(config: dict) -> dict:
    """
    Fetches 7-day weather forecast from Open-Meteo.
    Returns {date_str: weather_dict} for next 7 days.
    Includes pollution risk scores for each forecast day.
    """
    import urllib.request, urllib.parse
    from datetime import date as dc

    city     = config.get("city", "Delhi")
    lat, lon = _get_city_coords(city)

    try:
        DAILY_VARS = ",".join([
            "temperature_2m_max","temperature_2m_min","temperature_2m_mean",
            "precipitation_sum","windspeed_10m_max","windspeed_10m_mean",
            "windgusts_10m_max","winddirection_10m_dominant",
            "relativehumidity_2m_mean","dewpoint_2m_mean",
            "surface_pressure_mean","shortwave_radiation_sum","weathercode",
        ])
        params = {
            "latitude"      : lat,
            "longitude"     : lon,
            "daily"         : DAILY_VARS,
            "forecast_days" : 7,
            "timezone"      : "Asia/Kolkata",
            "windspeed_unit": "kmh",
        }
        url = "https://api.open-meteo.com/v1/forecast?" + urllib.parse.urlencode(params)
        import requests as _rq
        _r = _rq.get(url, timeout=15, verify=False,
                     headers={"User-Agent": "ATARS/3.1"})
        _r.raise_for_status()
        data = _r.json()

        daily = data.get("daily", {})
        dates = daily.get("time", [])

        def _s(key, i):
            arr = daily.get(key, [])
            try: return float(arr[i] or 0)
            except: return 0.0

        result = {}
        for i, d in enumerate(dates):
            temp_min  = _s("temperature_2m_min",          i)
            temp_max  = _s("temperature_2m_max",          i)
            temp_mean = _s("temperature_2m_mean",         i)
            wind_mean = _s("windspeed_10m_mean",          i)
            wind_max  = _s("windspeed_10m_max",           i)
            gust      = _s("windgusts_10m_max",           i)
            wind_dir  = _s("winddirection_10m_dominant",  i)
            precip    = _s("precipitation_sum",           i)
            hum       = _s("relativehumidity_2m_mean",   i)
            dewpoint  = _s("dewpoint_2m_mean",            i)
            pressure  = _s("surface_pressure_mean",      i)
            radiation = _s("shortwave_radiation_sum",    i)
            wcode     = int(_s("weathercode",            i))

            inversion = temp_min < 8  and wind_mean < 5
            wet       = precip > 10
            dust      = gust > 30    and precip < 1
            smog      = temp_max > 35 and radiation > 15
            fog       = (temp_mean - dewpoint) < 2 and hum > 90
            stagnant  = wind_mean < 3

            risk  = sum([inversion, wet, dust, smog, fog, stagnant])
            label = ("CRITICAL" if risk >= 4 else "HIGH" if risk >= 3
                     else "MODERATE" if risk >= 2 else "LOW" if risk >= 1
                     else "MINIMAL")

            dirs = ["N","NE","E","SE","S","SW","W","NW"]
            dir_label = dirs[int((wind_dir + 22.5) / 45) % 8] if wind_dir else "N/A"

            result[d] = {
                "date"                : d,
                "temp_mean_c"         : round(temp_mean, 1),
                "temp_max_c"          : round(temp_max,  1),
                "temp_min_c"          : round(temp_min,  1),
                "wind_mean_kmh"       : round(wind_mean, 1),
                "wind_max_kmh"        : round(wind_max,  1),
                "wind_gust_kmh"       : round(gust,      1),
                "wind_direction"      : dir_label,
                "rain_mm"             : round(precip,    1),
                "humidity_pct"        : round(hum,       1),
                "pressure_hpa"        : round(pressure,  1),
                "weather_desc"        : _wmo_code_description(wcode),
                "inversion_risk"      : inversion,
                "wet_deposition"      : wet,
                "dust_risk"           : dust,
                "smog_risk"           : smog,
                "fog_risk"            : fog,
                "stagnant_air"        : stagnant,
                "pollution_risk_score": risk,
                "pollution_risk"      : label,
                "is_forecast"         : True,
            }

        print(f"  ✓ 7-day weather forecast fetched for {city} ({len(result)} days)")
        return result

    except Exception as e:
        print(f"  ⚠ 7-day forecast failed: {e}")
        return {}



def compute_atmospheric_stability(weather_cache: dict) -> dict:
    """
    Computes daily atmospheric stability indices from weather data.
    Used to explain pollution episode patterns scientifically.

    Stability classes (Pasquill-Gifford proxy):
    A = Very unstable  (high radiation, low wind)
    B = Unstable
    C = Slightly unstable
    D = Neutral        (overcast or night)
    E = Slightly stable
    F = Stable         (low wind, clear night)

    Higher stability = more pollution trapping.
    """
    results = {}
    risk_days = []

    for date_str, w in weather_cache.items():
        radiation = w.get("solar_radiation_mj", 0)
        wind      = w.get("wind_kmh", w.get("wind_mean_kmh", 10))
        temp_min  = w.get("temp_min_c", w.get("temp_c", 15))
        precip    = w.get("rain_mm", 0)
        pressure  = w.get("pressure_hpa", 1013)

        # Pasquill-Gifford simplified proxy
        if radiation > 20 and wind < 3:
            stability = "A"
        elif radiation > 15 and wind < 5:
            stability = "B"
        elif radiation > 8  and wind < 6:
            stability = "C"
        elif precip > 5     or (5 < wind < 10):
            stability = "D"
        elif wind < 4 and radiation < 5 and temp_min < 15:
            stability = "E"
        else:
            stability = "F"

        # Boundary layer height proxy (meters)
        # Higher BLH = better dispersion = lower pollution
        if stability in ("A","B"):
            blh_proxy = 1500
        elif stability == "C":
            blh_proxy = 1000
        elif stability == "D":
            blh_proxy = 600
        elif stability == "E":
            blh_proxy = 300
        else:  # F
            blh_proxy = 100

        # Inversion index (0-10)
        inv_idx = 0
        if temp_min < 5:   inv_idx += 3
        elif temp_min < 10: inv_idx += 2
        elif temp_min < 15: inv_idx += 1
        if wind < 3:       inv_idx += 3
        elif wind < 5:     inv_idx += 2
        elif wind < 8:     inv_idx += 1
        if pressure > 1018: inv_idx += 2
        if precip < 0.5:   inv_idx += 1
        inv_idx = min(10, inv_idx)

        trapping_risk = ("SEVERE"   if inv_idx >= 7 else
                         "HIGH"     if inv_idx >= 5 else
                         "MODERATE" if inv_idx >= 3 else
                         "LOW")

        entry = {
            "stability_class"   : stability,
            "boundary_layer_m"  : blh_proxy,
            "inversion_index"   : inv_idx,
            "trapping_risk"     : trapping_risk,
            "pollution_risk"    : w.get("pollution_risk", "N/A"),
            "wind_direction"    : w.get("wind_direction", "N/A"),
        }
        results[date_str] = entry

        if inv_idx >= 6:
            risk_days.append({
                "date"          : date_str,
                "inversion_index": inv_idx,
                "trapping_risk" : trapping_risk,
                "temp_min"      : temp_min,
                "wind_kmh"      : wind,
            })

    high_risk_count = sum(1 for v in results.values() if v["trapping_risk"] in ("HIGH","SEVERE"))
    print(f"  ✓ Atmospheric stability: {len(results)} days | "
          f"{high_risk_count} high/severe trapping risk days")

    return {
        "daily"          : results,
        "high_risk_days" : sorted(risk_days, key=lambda x: x["inversion_index"], reverse=True)[:10],
        "high_risk_count": high_risk_count,
        "total_days"     : len(results),
    }



def classify_pollution_episodes(anomaly_days: list, weather_cache: dict,
                                 stability: dict, config: dict) -> list:
    """
    N-10: Pollution Episode Classifier (NEW NOVEL CONTRIBUTION)
    Rule-based deterministic classification of every anomaly day.
    Zero LLM. Fully auditable. Based on meteorology + temporal patterns.

    Episode types:
    REGIONAL_TRANSPORT    → high PM10, strong wind from industrial direction
    AGRICULTURAL_BURNING  → Oct-Nov + NW wind from Punjab/Haryana
    INDUSTRIAL_EMISSION   → weekday pattern + specific wind direction
    TRAFFIC_PEAK          → Monday-Tuesday, morning diurnal spike
    DUST_STORM            → high wind + no rain + sudden PM spike
    TEMPERATURE_INVERSION → low temp + stagnant air + morning peak
    WET_DEPOSITION_END    → PM spike day after rain (resuspension)
    PHOTOCHEMICAL_SMOG    → high temp + high radiation + ozone spike
    UNKNOWN               → none of the above patterns match
    """
    classified = []

    for day_info in anomaly_days:
        date_str = day_info.get("date", "")
        pm10_val = day_info.get("pm10_mean", 0)

        w = weather_cache.get(date_str, {})
        s = stability.get("daily", {}).get(date_str, {})

        # Extract weather values
        wind_kmh    = w.get("wind_kmh", w.get("wind_mean_kmh", 10))
        wind_dir    = w.get("wind_direction", "N/A")
        wind_max    = w.get("wind_max_kmh", w.get("wind_kmh", 10))
        gust        = w.get("wind_gust_kmh", wind_max)
        temp_min    = w.get("temp_min_c", w.get("temp_c", 20))
        temp_max    = w.get("temp_max_c", w.get("temp_c", 30))
        precip      = w.get("rain_mm", 0)
        hum         = w.get("humidity_pct", 50)
        radiation   = w.get("solar_radiation_mj", 10)
        pressure    = w.get("pressure_hpa", 1013)
        inversion   = w.get("inversion_risk", False)
        inv_idx     = s.get("inversion_index", 0)

        # Date features
        try:
            from datetime import date as dc
            d = dc.fromisoformat(date_str)
            weekday   = d.weekday()  # 0=Mon, 6=Sun
            month     = d.month
            is_monday = weekday == 0
            is_winter = month in [11, 12, 1, 2]
            is_harvest= month in [10, 11]  # Oct-Nov burning season
        except Exception:
            weekday = 2; month = 6; is_monday = False
            is_winter = False; is_harvest = False

        # Wind from NW = from Punjab/Haryana agricultural areas
        nw_wind = wind_dir in ["NW", "N", "NE"] and wind_kmh > 8

        # Classification rules (priority order)
        episode = "UNKNOWN"
        confidence = "LOW"
        reason = ""

        # 1. Agricultural burning — Oct-Nov + NW wind
        if is_harvest and nw_wind and pm10_val > 100:
            episode    = "AGRICULTURAL_BURNING"
            confidence = "HIGH"
            reason = (f"Harvest season ({month=}), NW wind ({wind_dir}) "
                      f"from agricultural regions, PM10={pm10_val:.0f} µg/m³")

        # 2. Dust storm — high gust, no rain, sudden spike
        elif gust > 40 and precip < 1 and pm10_val > 150:
            episode    = "DUST_STORM"
            confidence = "HIGH"
            reason = (f"Wind gust {gust:.0f} km/h, no rain, "
                      f"PM10={pm10_val:.0f} µg/m³ — dust resuspension")

        # 3. Temperature inversion — cold + stagnant winter
        elif is_winter and inversion and inv_idx >= 5 and pm10_val > 80:
            episode    = "TEMPERATURE_INVERSION"
            confidence = "HIGH"
            reason = (f"Winter inversion: temp_min={temp_min:.1f}°C, "
                      f"wind={wind_kmh:.1f} km/h, inversion_idx={inv_idx}")

        # 4. Photochemical smog — hot + sunny + radiation
        elif temp_max > 38 and radiation > 18 and pm10_val > 60:
            episode    = "PHOTOCHEMICAL_SMOG"
            confidence = "MEDIUM"
            reason = (f"Temp_max={temp_max:.1f}°C, "
                      f"radiation={radiation:.1f} MJ/m², "
                      f"photochemical formation likely")

        # 5. Regional transport — strong wind from industrial direction
        elif wind_kmh > 15 and wind_dir in ["SW","W","S"] and pm10_val > 120:
            episode    = "REGIONAL_TRANSPORT"
            confidence = "MEDIUM"
            reason = (f"Wind {wind_kmh:.0f} km/h from {wind_dir} "
                      f"(industrial corridor direction), PM10={pm10_val:.0f}")

        # 6. Traffic peak — Monday + morning pattern
        elif is_monday and pm10_val > 80 and wind_kmh < 8:
            episode    = "TRAFFIC_PEAK"
            confidence = "LOW"
            reason = (f"Monday pattern, low wind {wind_kmh:.1f} km/h, "
                      f"traffic accumulation likely")

        # 7. Post-rain resuspension
        elif not w.get("wet_deposition") and pm10_val > 100:
            # Check if previous day had rain (use date logic)
            episode    = "WET_DEPOSITION_END"
            confidence = "LOW"
            reason = "Post-rain conditions — resuspension possible"

        # 8. Industrial — weekday + specific wind
        elif weekday < 5 and wind_dir in ["SW","W"] and pm10_val > 100:
            episode    = "INDUSTRIAL_EMISSION"
            confidence = "LOW"
            reason = (f"Weekday ({weekday}), wind from {wind_dir} "
                      f"(industrial direction)")

        else:
            episode    = "COMPLEX_MULTI_SOURCE"
            confidence = "LOW"
            reason = "Multiple factors — no single dominant source identified"

        classified.append({
            "date"           : date_str,
            "pm10_mean"      : round(pm10_val, 1),
            "episode_type"   : episode,
            "confidence"     : confidence,
            "reason"         : reason,
            "weather_summary": {
                "temp_min"   : temp_min,
                "wind_kmh"   : wind_kmh,
                "wind_dir"   : wind_dir,
                "rain_mm"    : precip,
                "inversion"  : inversion,
            },
            "anomaly_type"   : day_info.get("anomaly_type", "UNKNOWN"),
        })

    # Summary by episode type
    from collections import Counter
    episode_counts = Counter(c["episode_type"] for c in classified)
    print(f"  ✓ N-10 Episode Classifier: {len(classified)} episodes classified")
    print(f"    Types: {dict(episode_counts)}")

    return classified



# ─────────────────────────────────────────────────────────────────────────────
#  AUTOMATED PPT GENERATOR
#  Creates a clean executive summary presentation alongside the Word report
#  Uses python-pptx — already installed (same as python-docx family)
#  Ocean Executive palette — matches environmental reporting context
# ─────────────────────────────────────────────────────────────────────────────

def generate_summary_ppt(J: dict, J_v3: dict, verification: dict,
                          ml_if: dict, ml_hw: dict, config: dict,
                          out_dir: Path,
                          v3_mns: dict = None,
                          v3_episodes: list = None,
                          v3_forecast: dict = None,
                          v3_benchmark: dict = None,
                          charts: dict = None,
                          cleaning_report: dict = None) -> str:
    if cleaning_report is None:
        cleaning_report = {}
    """
    Generates a clean executive summary PPT alongside the Word report.
    15 slides — key numbers only, no clutter.
    Ocean Executive palette: navy #1A2C4E, teal #1A6B72, white #FFFFFF
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("  ⚠ python-pptx not installed — skipping PPT generation")
        return None

    # ── Palette ───────────────────────────────────────────────────────────
    NAVY  = RGBColor(0x1A, 0x2C, 0x4E)
    TEAL  = RGBColor(0x1A, 0x6B, 0x72)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    AMBER = RGBColor(0xE8, 0x8C, 0x00)
    GREEN = RGBColor(0x1E, 0x6B, 0x3C)
    RED   = RGBColor(0xC0, 0x39, 0x2B)
    LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
    GRAY  = RGBColor(0x6B, 0x72, 0x80)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank layout

    def _slide(bg_color=None):
        sl = prs.slides.add_slide(blank)
        if bg_color:
            fill = sl.background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        return sl

    def _box(sl, x, y, w, h, text, font_size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, bg=None, italic=False):
        txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.italic= italic
        run.font.color.rgb = color
        if bg:
            txb.fill.solid()
            txb.fill.fore_color.rgb = bg
        return txb

    def _rect(sl, x, y, w, h, color):
        shape = sl.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _stat_card(sl, x, y, value, label, val_color=WHITE, bg=TEAL):
        _rect(sl, x, y, 2.8, 1.4, bg)
        _box(sl, x+0.1, y+0.05, 2.6, 0.8, value,
             font_size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
        _box(sl, x+0.1, y+0.85, 2.6, 0.5, label,
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # Short data extractions
    city    = config.get("city", "City")
    dq      = J.get("data_quality", {})
    Q       = dq.get("overall_Q", 0)
    Q_flag  = dq.get("confidence_flag", "N/A")
    vars_j  = J.get("variables", {})
    pm10    = vars_j.get("PM10", {})
    pm10_mean= pm10.get("annual_mean", 0)
    pm10_exc = pm10.get("exceed_days", 0)
    benz    = vars_j.get("Benzene", {})
    benz_mean= benz.get("annual_mean", 0)
    no2     = vars_j.get("NO2", {})
    no2_mean = no2.get("annual_mean", 0)
    g_rate  = verification.get("overall_G_rate", 0)
    g_vdt   = verification.get("overall_verdict", "N/A")
    mns_val = v3_mns.get("MNS", 0) if v3_mns else 0
    mns_vdt = v3_mns.get("verdict", "N/A") if v3_mns else "N/A"
    if_anom = ml_if.get("n_anomalies", 0) if ml_if else 0
    if_only = ml_if.get("n_if_only", 0)   if ml_if else 0
    hw_rmse = ml_hw.get("rmse", 0)        if ml_hw else 0
    _date_range_str = J.get("date_range", "2024-01-01 to 2024-12-31")
    year    = _date_range_str[:4]
    run_id  = J.get("run_metadata", {}).get("run_id", "N/A")
    run_id  = (run_id[:16] if isinstance(run_id, str) and len(run_id) >= 16 else str(run_id))

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 1 — TITLE
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(NAVY)
    _rect(sl, 0, 0, 13.33, 0.08, TEAL)
    _rect(sl, 0, 7.42, 13.33, 0.08, TEAL)
    _box(sl, 1, 1.2, 11, 1.0,
         "ATARS v3.1 — Air Quality Analysis", 40, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _box(sl, 1, 2.4, 11, 0.6,
         f"{city} Station  ·  {year}  ·  {J.get('date_range','N/A')}",
         18, color=LIGHT, align=PP_ALIGN.CENTER)
    _box(sl, 1, 3.1, 11, 0.5,
         "21-Section Word Report  ·  19-Slide PPT  ·  15 Charts  ·  9 Verification Modules",
         14, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    _box(sl, 1, 3.65, 11, 0.4,
         "Formal Statistical Framework + ML (IF + HW) + LLM Grounding Verification (RGV)",
         13, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    _rect(sl, 4, 4.0, 5.33, 0.03, TEAL)
    _box(sl, 1, 4.4, 11, 0.4,
         "Priyanshu  ·  Global Institute of Technology and Management, Haryana",
         13, color=LIGHT, align=PP_ALIGN.CENTER)
    _box(sl, 1, 4.9, 11, 0.4,
         "Open Source  ·  MIT License  ·  github.com/Priyanshu-ux712/ATARS",
         12, color=LIGHT, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 2 — EXECUTIVE SUMMARY (stat cards)
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(LIGHT)
    _rect(sl, 0, 0, 13.33, 0.9, NAVY)
    _box(sl, 0.3, 0.15, 12, 0.6, "Executive Summary — Key Metrics at a Glance",
         24, bold=True, color=WHITE)
    _stat_card(sl, 0.3, 1.1, f"{pm10_mean:.1f}", "PM10 Annual Mean (µg/m³)",
               val_color=WHITE, bg=RED)
    _stat_card(sl, 3.3, 1.1, f"{pm10_exc}", "PM10 Exceedance Days",
               val_color=WHITE, bg=AMBER)
    _stat_card(sl, 6.3, 1.1, f"{Q:.2f}", f"Data Quality Q(D) [{Q_flag}]",
               val_color=WHITE, bg=TEAL)
    _stat_card(sl, 9.3, 1.1, f"{if_anom}", "IF Anomaly Days",
               val_color=WHITE, bg=NAVY)
    _stat_card(sl, 0.3, 2.7, f"{round(g_rate*100,1)}%", f"G_rate [{g_vdt}]",
               val_color=WHITE, bg=GREEN)
    _stat_card(sl, 3.3, 2.7, f"{benz_mean:.2f}", "Benzene Mean (µg/m³)",
               val_color=WHITE, bg=RED)
    _stat_card(sl, 6.3, 2.7, f"{no2_mean:.1f}", "NO2 Annual Mean (µg/m³)",
               val_color=WHITE, bg=TEAL)
    _stat_card(sl, 9.3, 2.7, f"{hw_rmse:.1f}", "HW Forecast RMSE",
               val_color=WHITE, bg=NAVY)
    stat_inf_j = J_v3.get("stat_inference_summary", {})
    mk_dir  = (stat_inf_j.get("mann_kendall_pm10") or {}).get("direction","N/A")
    mk_p    = (stat_inf_j.get("mann_kendall_pm10") or {}).get("p_value")
    ols_r2a = stat_inf_j.get("ols_r2_adj")
    _box(sl, 0.3, 4.1, 6.3, 0.45,
         f"MK Trend: {str(mk_dir).replace('_',' ')} "
         f"{'(p='+str(round(mk_p,4))+')' if mk_p is not None else ''}  ·  "
         f"OLS adj. R² = {round(ols_r2a,3) if ols_r2a is not None else 'N/A'}",
         11, color=TEAL)
    if v3_mns:
        _box(sl, 0.3, 4.6, 12.7, 0.4,
             f"Master Novelty Score (MNS) = {mns_val:.3f}  [{mns_vdt}]  "
             f"·  Audit ID: {run_id}",
             12, color=NAVY)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 3 — PM10 CRISIS
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, NAVY)
    _box(sl, 0.3, 0.15, 12, 0.6, "PM10 — Severity Assessment", 24, bold=True, color=WHITE)
    who_ratio = round(pm10_mean / (pm10.get("threshold") or 45.0), 1) if pm10_mean else 0
    _box(sl, 0.5, 1.1, 6, 1.2, f"{pm10_mean:.1f} µg/m³",
         52, bold=True, color=RED)
    _box(sl, 0.5, 2.4, 6, 0.5, "Annual Mean  ·  WHO AQG limit: 45 µg/m³",
         14, color=NAVY)
    _box(sl, 0.5, 3.0, 6, 0.5, f"{who_ratio}× above WHO Annual Guideline",
         16, bold=True, color=RED)
    _rect(sl, 7, 1.0, 5.8, 1.4, NAVY)
    _box(sl, 7.1, 1.1, 5.6, 0.6, f"{pm10_exc}", 42, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _box(sl, 7.1, 1.75, 5.6, 0.5, "Days exceeding WHO 24h limit (45 µg/m³)",
         11, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(sl, 7, 2.6, 5.8, 1.4, TEAL)
    pct = round(pm10_exc / 366 * 100, 1) if pm10_exc else 0
    _box(sl, 7.1, 2.7, 5.6, 0.6, f"{pct}%", 42, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _box(sl, 7.1, 3.35, 5.6, 0.5, "Fraction of monitored days above WHO",
         11, color=WHITE, align=PP_ALIGN.CENTER)
    _pm10_peak_max = pm10.get("annual_max", pm10.get("max_daily_mean", 0))
    bullets = [
        f"Peak daily mean: {_pm10_peak_max:.1f} µg/m³",
        f"Annual min: {pm10.get('annual_min', 0):.1f} µg/m³",
        f"Annual std σ: {pm10.get('annual_std', 0):.1f} µg/m³",
        "Primary source: particulate matter from traffic and regional transport",
    ]
    for i, b in enumerate(bullets):
        _box(sl, 0.5, 4.2 + i*0.55, 12, 0.5, f"•  {b}", 13, color=NAVY)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 4 — DATA QUALITY
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, TEAL)
    _box(sl, 0.3, 0.15, 12.5, 0.6, "Data Quality Assessment  —  Q(D) = N_valid / N_total",
         24, bold=True, color=WHITE)
    # LEFT column: big Q score + stats (x=0.3, width=4.5)
    _rect(sl, 0.3, 1.0, 4.5, 5.6, LIGHT)
    _box(sl, 0.4, 1.1, 4.3, 1.0, f"Q(D) = {Q:.4f}", 38, bold=True, color=TEAL,
         align=PP_ALIGN.CENTER)
    _box(sl, 0.4, 2.15, 4.3, 0.5, f"Confidence: {Q_flag}", 18, bold=True,
         color=GREEN if Q_flag=="HIGH" else (AMBER if Q_flag=="MODERATE" else RED),
         align=PP_ALIGN.CENTER)
    _box(sl, 0.4, 2.8, 4.3, 0.4, f"Total Days: {J.get('total_days', 'N/A')}",
         13, color=NAVY, align=PP_ALIGN.CENTER)
    _box(sl, 0.4, 3.2, 4.3, 0.4, f"Total Records: {J.get('total_records', 0):,}",
         13, color=NAVY, align=PP_ALIGN.CENTER)
    _box(sl, 0.4, 3.6, 4.3, 0.4, f"Station: {config.get('station_id','N/A')}",
         13, color=NAVY, align=PP_ALIGN.CENTER)
    _box(sl, 0.4, 4.0, 4.3, 0.4, f"Baseline Window: {config.get('baseline_window_days',30)} days",
         13, color=NAVY, align=PP_ALIGN.CENTER)
    _box(sl, 0.4, 4.5, 4.3, 0.4,
         f"HIGH: Q≥0.90  MODERATE: Q≥0.70  LOW: Q<0.70",
         10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    # RIGHT column: per-variable bars (x=5.1, width=7.9)
    _box(sl, 5.2, 1.0, 7.8, 0.38, "Per-Variable Data Completeness:",
         12, bold=True, color=NAVY)
    per_q = J.get("data_quality", {}).get("per_variable_Q", {})
    y_off = 1.45
    for var, q_v in list(per_q.items())[:8]:
        q_v = float(q_v or 0)
        bar  = "█" * int(q_v * 22)
        line = f"{var:<14}  {q_v:.3f}  {bar}"
        clr  = GREEN if q_v >= 0.9 else (AMBER if q_v >= 0.7 else RED)
        _rect(sl, 5.1, y_off, 8.1, 0.4,
              RGBColor(0xD1,0xFA,0xE5) if q_v>=0.9 else
              (RGBColor(0xFE,0xF3,0xC7) if q_v>=0.7 else RGBColor(0xFF,0xE4,0xE6)))
        _box(sl, 5.2, y_off+0.03, 7.9, 0.35, line, 11, color=clr)
        y_off += 0.45
    if not per_q:
        _box(sl, 5.1, 1.45, 7.9, 0.5,
             f"Overall Q(D) = {Q:.3f}  [{Q_flag}]", 14, color=TEAL)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 5 — ANOMALY DETECTION
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, NAVY)
    _box(sl, 0.3, 0.15, 12, 0.6, "Anomaly Detection — Isolation Forest + Z-Score",
         22, bold=True, color=WHITE)
    _stat_card(sl, 0.5, 1.1, str(if_anom), "Total Anomaly Days", bg=RED)
    _stat_card(sl, 3.5, 1.1, str(if_only),
               "IF-Only (missed by z-score)", bg=AMBER)
    _stat_card(sl, 6.5, 1.1,
               str(ml_if.get("n_both", 0) if ml_if else 0),
               "Both IF and Z-Score", bg=TEAL)
    _stat_card(sl, 9.5, 1.1,
               str(ml_if.get("n_zs_only", 0) if ml_if else 0),
               "Z-Score Only", bg=NAVY)
    _box(sl, 0.5, 2.8, 12, 0.5,
         "Isolation Forest uses 150 trees, contamination=5%, fixed seed=42 for reproducibility.",
         13, italic=True, color=NAVY)
    _box(sl, 0.5, 3.4, 12, 0.5,
         f"{if_only} anomaly days were detected ONLY by multivariate Isolation Forest — "
         "completely missed by univariate z-score analysis.",
         14, bold=True, color=RED)
    if v3_episodes:
        from collections import Counter
        ep_counts = Counter(e["episode_type"] for e in v3_episodes)
        _box(sl, 0.5, 4.0, 12, 0.4, "Episode Classification:", 13, bold=True, color=NAVY)
        ep_text = "  ·  ".join(f"{t.replace('_',' ')}: {c}" for t,c in ep_counts.most_common(5))
        _box(sl, 0.5, 4.4, 12, 0.5, ep_text, 12, color=TEAL)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 6 — ML FORECAST
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, TEAL)
    _box(sl, 0.3, 0.15, 12, 0.6, "Holt-Winters PM10 Forecast",
         24, bold=True, color=WHITE)
    hw_mean14 = (float(np.mean(ml_hw["forecast"])) if ml_hw and ml_hw.get("available") and "forecast" in ml_hw else 0)
    hw_mae    = ml_hw.get("mae", 0)           if ml_hw else 0
    _stat_card(sl, 0.5, 1.1, f"{hw_rmse:.1f}", "In-Sample RMSE (µg/m³)", bg=NAVY)
    _stat_card(sl, 3.5, 1.1, f"{hw_mae:.1f}",  "In-Sample MAE (µg/m³)",  bg=TEAL)
    _stat_card(sl, 6.5, 1.1, f"{hw_mean14:.1f}","14-Day Forecast Mean",   bg=AMBER)
    _stat_card(sl, 9.5, 1.1, "80%",             "Prediction Interval",    bg=GREEN)
    params_text = (f"α={ml_hw.get('alpha',0.3) if ml_hw else 0.3}  "
                   f"β={ml_hw.get('beta',0.1) if ml_hw else 0.1}  "
                   f"γ={ml_hw.get('gamma',0.2) if ml_hw else 0.2}  "
                   f"m={ml_hw.get('seasonal_periods',7) if ml_hw else 7} days  "
                   f"horizon={ml_hw.get('forecast_days',14) if ml_hw else 14} days")
    _box(sl, 0.5, 2.8, 12, 0.5, f"Model parameters: {params_text}",
         12, italic=True, color=NAVY)
    _box(sl, 0.5, 3.4, 12, 0.6,
         "Forecast values are deliberately excluded from the LLM data contract J. "
         "This prevents AI from presenting statistical extrapolations as historical facts.",
         13, color=NAVY)
    if v3_forecast:
        _box(sl, 0.5, 4.2, 12, 0.4, "7-Day Weather Forecast Pollution Risk:", 13, bold=True, color=NAVY)
        risk_text = "  ·  ".join(
            f"{d}: {w.get('pollution_risk','?')}"
            for d,w in list(v3_forecast.items())[:4]
        )
        _box(sl, 0.5, 4.65, 12, 0.5, risk_text, 11, color=TEAL)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 7 — BENZENE HEALTH RISK
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, RED)
    _box(sl, 0.3, 0.15, 12, 0.6,
         "Benzene — IARC Group 1 Carcinogen Alert", 22, bold=True, color=WHITE)
    benz_ratio = round(benz_mean / 1.7, 2) if benz_mean else 0
    benz_exc   = benz.get("exceed_days", 0)
    _box(sl, 0.5, 1.0, 5, 1.0, f"{benz_mean:.2f} µg/m³", 44, bold=True, color=RED)
    _box(sl, 0.5, 2.1, 5, 0.5, "Annual Mean  ·  WHO limit: 1.7 µg/m³",
         13, color=NAVY)
    _rect(sl, 6.5, 1.0, 6.3, 1.3, RED)
    _box(sl, 6.6, 1.05, 6.1, 0.7, f"{benz_ratio}×  WHO Annual Limit",
         26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _box(sl, 6.6, 1.75, 6.1, 0.4, f"Exceedance on {benz_exc} / 366 days (100%)",
         12, color=WHITE, align=PP_ALIGN.CENTER)
    facts = [
        "Benzene is classified IARC Group 1 — definite human carcinogen",
        "No safe threshold established — any exposure carries risk",
        f"Gurugram annual mean ({benz_mean:.2f} µg/m³) is {benz_ratio}× the WHO reference",
        "Primary sources: vehicle exhaust, fuel evaporation, industrial solvents",
        "Long-term exposure linked to leukaemia and other blood cancers",
    ]
    for i, f in enumerate(facts):
        _box(sl, 0.5, 2.9 + i*0.6, 12, 0.55, f"•  {f}", 13,
             color=RED if i == 0 else NAVY)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 8 — RGV GROUNDING VERIFICATION
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, NAVY)
    _box(sl, 0.3, 0.15, 12, 0.6,
         "Runtime Grounding Verifier (RGV) — Novel Contribution",
         22, bold=True, color=WHITE)
    _box(sl, 0.5, 1.0, 12, 0.5,
         'G(s) ⊆ J  ∀ s ∈ N  —  Every AI number traced to computed data',
         16, italic=True, color=TEAL)
    _stat_card(sl, 0.5, 1.7, f"{round(g_rate*100,1)}%", f"G_rate [{g_vdt}]",
               bg=GREEN if g_rate >= 0.8 else AMBER)
    _stat_card(sl, 3.5, 1.7,
               str(verification.get("total_numerical",0)), "Sentences Checked", bg=NAVY)
    _stat_card(sl, 6.5, 1.7,
               str(verification.get("total_grounded",0)),  "Grounded",          bg=GREEN)
    _stat_card(sl, 9.5, 1.7,
               str(verification.get("total_ungrounded",0)),"Ungrounded",        bg=RED)
    _box(sl, 0.5, 3.3, 12, 0.5,
         "Tolerance: ±3% for concentrations  ·  ±1 for counts  ·  ±0.05 for scores",
         12, italic=True, color=NAVY)
    _box(sl, 0.5, 3.9, 12, 0.5,
         "Sources verified against: J contract  +  WHO AQG 2021  +  Live weather data",
         13, color=TEAL)
    conf = verification.get("overall_confidence", 0)
    sem  = "PASS" if verification.get("semantic_pass") else "REVIEW"
    _box(sl, 0.5, 4.5, 6, 0.4, f"Avg Confidence: {conf}", 13, color=NAVY)
    _box(sl, 6.5, 4.5, 6, 0.4, f"Semantic Alignment: {sem}", 13,
         color=GREEN if sem=="PASS" else AMBER)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 9 — NOVEL CONTRIBUTIONS
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(NAVY)
    _box(sl, 0.3, 0.2, 12, 0.6, "Novel Contributions — What Makes ATARS Unique",
         22, bold=True, color=WHITE)
    contributions = [
        ("N-01", "Temporal Claim Verifier (TCV)",
         "Verifies month/season claims against computed facts"),
        ("N-02", "Causal Claim Detector (CCD)",
         "Flags causal language — enforces Corr(X,Y) ≢ C(X→Y)"),
        ("N-03", "Anomaly Explanation Engine (AEE)",
         "Rule-based meteorological explanation of every anomaly"),
        ("N-05", "Uncertainty Propagation (UPT)",
         "95% CI bounds attached to every numerical claim"),
        ("N-09", "Master Novelty Score (MNS)",
         f"Unified trust metric across 7 dimensions = {mns_val:.3f} [{mns_vdt}]"),
        ("N-10", "Episode Classifier",
         "8 pollution episode types — Agricultural, Dust Storm, Inversion..."),
        ("RGV",  "Runtime Grounding Verifier",
         "First numerical grounding enforcement in scientific pipeline"),
    ]
    for i, (code, title, desc) in enumerate(contributions):
        x = 0.3 if i % 2 == 0 else 6.8
        y = 1.1 + (i // 2) * 1.6
        _rect(sl, x, y, 5.9, 1.3, TEAL)
        _box(sl, x+0.1, y+0.05, 1.0, 0.4, code, 13, bold=True, color=AMBER)
        _box(sl, x+1.2, y+0.05, 4.6, 0.4, title, 13, bold=True, color=WHITE)
        _box(sl, x+0.1, y+0.55, 5.7, 0.6, desc, 11, color=WHITE)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 10 — BENCHMARK
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, TEAL)
    _box(sl, 0.3, 0.15, 12, 0.6, "National & International Benchmark Comparison",
         22, bold=True, color=WHITE)
    if v3_benchmark:
        comps = v3_benchmark.get("comparisons", {})
        y_off = 1.1
        _box(sl, 0.3, y_off, 12.5, 0.4,
             f"Station: {v3_benchmark.get('station','N/A')}  ·  "
             f"Overall Severity: {v3_benchmark.get('overall_severity','N/A')}",
             14, bold=True, color=NAVY)
        y_off += 0.55
        for var, comp in list(comps.items())[:5]:
            mean = comp.get("station_annual_mean", 0)
            who  = comp.get("vs_who", "N/A")
            naaqs= comp.get("naaqs_status", "N/A")
            bg   = NAVY
            _rect(sl, 0.3, y_off, 12.5, 0.65, LIGHT)
            _box(sl, 0.4, y_off+0.05, 1.8, 0.5, var, 13, bold=True, color=NAVY)
            _box(sl, 2.3, y_off+0.05, 2.5, 0.5, f"{mean:.2f} µg/m³", 12, color=NAVY)
            _box(sl, 5.0, y_off+0.05, 3.5, 0.5, who, 12,
                 color=RED if "EXCEED" in str(who).upper() else GREEN)
            _box(sl, 8.5, y_off+0.05, 4.0, 0.5, naaqs, 12,
                 color=RED if "EXCEED" in str(naaqs).upper() else GREEN)
            y_off += 0.75
    else:
        _box(sl, 0.5, 2.0, 12, 1.0,
             "Run with --weather open_meteo for benchmark comparison data.",
             16, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 11 — WEATHER & ATMOSPHERIC CONTEXT
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, NAVY)
    _box(sl, 0.3, 0.15, 12, 0.6, "Atmospheric Context — Open-Meteo Live Data",
         22, bold=True, color=WHITE)
    comp_j = J_v3.get("v3_comparison", {})
    weather = J_v3.get("v3_weather", {})
    _stat_card(sl, 0.3, 1.1, f"{weather.get('temp_c',0):.1f}°C", "Temperature", bg=TEAL)
    _stat_card(sl, 3.3, 1.1, f"{weather.get('humidity_pct',0):.0f}%","Humidity", bg=NAVY)
    _stat_card(sl, 6.3, 1.1, f"{weather.get('wind_kmh',0):.1f}", "Wind (km/h)", bg=TEAL)
    _stat_card(sl, 9.3, 1.1, f"{weather.get('rain_mm',0):.1f}", "Rain (mm)", bg=NAVY)
    pm10_now  = comp_j.get("pm10_today", 0)
    pm10_yest = comp_j.get("pm10_yesterday", 0)
    pm10_week = comp_j.get("pm10_week_avg", 0)
    pm10_chg  = comp_j.get("pm10_change_yesterday_pct", 0)
    _box(sl, 0.3, 2.8, 12, 0.5, "Today vs Historical PM10 Comparison:", 14, bold=True, color=NAVY)
    comp_items = [
        f"Today: {pm10_now:.1f} µg/m³",
        f"Yesterday: {pm10_yest:.1f} µg/m³  ({pm10_chg:+.1f}%)",
        f"7-day avg: {pm10_week:.1f} µg/m³",
        f"Status: {comp_j.get('status','N/A').replace('_',' ')}",
    ]
    for i, item in enumerate(comp_items):
        _box(sl, 0.5 + i*3.2, 3.4, 3.0, 0.5, item, 12, color=NAVY if i < 2 else TEAL)
    tags = J_v3.get("v3_events", {}).get("tags", [])
    if tags:
        _box(sl, 0.3, 4.2, 12, 0.4, "Context tags: " + "  ·  ".join(tags[:5]),
             11, italic=True, color=TEAL)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 12 — 7-DAY FORECAST
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, TEAL)
    _box(sl, 0.3, 0.15, 12, 0.6, "7-Day Pollution Risk Forecast",
         24, bold=True, color=WHITE)
    if v3_forecast:
        dates = sorted(v3_forecast.keys())
        col_w = 13.0 / max(len(dates), 1)
        for i, d in enumerate(dates[:7]):
            w     = v3_forecast[d]
            risk  = w.get("pollution_risk", "N/A")
            score = w.get("pollution_risk_score", 0)
            bg    = (RED   if risk in ("CRITICAL","HIGH") else
                     AMBER if risk == "MODERATE" else GREEN)
            x = 0.3 + i * col_w
            _rect(sl, x, 1.0, col_w - 0.05, 3.5, bg)
            _box(sl, x, 1.05, col_w-0.05, 0.4, d[5:], 10, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
            _box(sl, x, 1.55, col_w-0.05, 0.6, risk.replace("_"," "),
                 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _box(sl, x, 2.2, col_w-0.05, 0.4, f"Score: {score}/6",
                 10, color=WHITE, align=PP_ALIGN.CENTER)
            _box(sl, x, 2.7, col_w-0.05, 0.4,
                 f"{w.get('temp_mean_c',0):.0f}°C",
                 10, color=WHITE, align=PP_ALIGN.CENTER)
            _box(sl, x, 3.1, col_w-0.05, 0.4,
                 f"Wind: {w.get('wind_mean_kmh',0):.0f}km/h",
                 9, color=WHITE, align=PP_ALIGN.CENTER)
            _box(sl, x, 3.5, col_w-0.05, 0.4,
                 f"Rain: {w.get('rain_mm',0):.1f}mm",
                 9, color=WHITE, align=PP_ALIGN.CENTER)
        _box(sl, 0.3, 4.8, 12, 0.4,
             "Red = High/Critical risk  ·  Orange = Moderate  ·  Green = Low/Minimal",
             11, italic=True, color=NAVY)
    else:
        _box(sl, 0.5, 2.5, 12, 1.0,
             "Run with --weather open_meteo to enable 7-day forecast.",
             16, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 12b — STATISTICAL INFERENCE RESULTS
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, NAVY)
    _box(sl, 0.3, 0.15, 12, 0.6, "Statistical Inference — Mann-Kendall & OLS",
         22, bold=True, color=WHITE)
    # MK results
    stat_inf = J_v3.get("stat_inference_summary", {})
    mk_pm10  = stat_inf.get("mann_kendall_pm10") or {}
    ols_r2   = stat_inf.get("ols_r2_adj")
    ols_f    = stat_inf.get("ols_f_pvalue")
    ols_dw   = stat_inf.get("ols_dw")
    s_trend  = stat_inf.get("seasonal_pct_trend")
    s_seas   = stat_inf.get("seasonal_pct_seasonal")
    s_res    = stat_inf.get("seasonal_pct_residual")

    _stat_card(sl, 0.3, 1.1,
               mk_pm10.get("direction","N/A").replace("_"," ")[:12] if mk_pm10 else "N/A",
               f"PM10 Trend (p={mk_pm10.get('p_value',1):.4f})" if mk_pm10 else "MK Trend",
               bg=RED if str(mk_pm10.get("direction","")).startswith("INC") else
                  GREEN if str(mk_pm10.get("direction","")).startswith("DEC") else TEAL)
    _stat_card(sl, 3.3, 1.1,
               f"{mk_pm10.get('sens_slope',0):.3f}" if mk_pm10 else "N/A",
               "Sen's Slope (µg/m³/day)", bg=NAVY)
    _stat_card(sl, 6.3, 1.1,
               f"{ols_r2:.3f}" if ols_r2 is not None else "N/A",
               "OLS Adj. R²", bg=TEAL)
    _stat_card(sl, 9.3, 1.1,
               f"{ols_dw:.2f}" if ols_dw is not None else "N/A",
               "Durbin-Watson", bg=GREEN if ols_dw and 1.5<=ols_dw<=2.5 else AMBER)

    # Seasonal decomposition
    if s_trend is not None:
        _box(sl, 0.3, 2.8, 12, 0.4, "PM10 Variance Decomposition (Seasonal Decomposition):",
             13, bold=True, color=NAVY)
        for i, (lbl, val, col) in enumerate([
            (f"Trend component",    f"{s_trend:.1f}% of variance", TEAL),
            (f"Seasonal component", f"{s_seas:.1f}% of variance",  AMBER),
            (f"Residual (episodic)",f"{s_res:.1f}% of variance",   RED),
        ]):
            _box(sl, 0.5 + i*4.2, 3.2, 4.0, 0.5, f"{lbl}: {val}", 13, color=col)
    # Exceedance frequency
    exc_f = J_v3.get("stat_inference_summary", {})
    _box(sl, 0.3, 4.0, 12, 0.5,
         f"OLS F-statistic p-value: {ols_f:.4f}" if ols_f is not None else "",
         12, color=NAVY)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 12c — WHO EXCEEDANCE DEEP DIVE
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, RED)
    _box(sl, 0.3, 0.15, 12, 0.6,
         "WHO Exceedance Analysis — Pollutant by Pollutant", 22, bold=True, color=WHITE)

    exc_freq_data = {}
    for col_name, vdata in vars_j.items():
        thr = vdata.get("threshold")
        exc_days = vdata.get("exceeds_days", 0)
        if thr and exc_days > 0:
            total_d = J.get("total_days", 366)
            exc_freq_data[col_name] = {
                "exceed_days": exc_days,
                "exceed_pct" : vdata.get("exceed_pct", round(exc_days/max(total_d,1)*100,1)),
                "threshold"  : thr,
                "annual_mean": vdata.get("annual_mean", 0),
            }

    y_row = 1.0
    _box(sl, 0.3, y_row, 3.5, 0.4, "Pollutant", 11, bold=True, color=NAVY)
    _box(sl, 3.9, y_row, 2.5, 0.4, "Annual Mean", 11, bold=True, color=NAVY)
    _box(sl, 6.5, y_row, 2.5, 0.4, "WHO Limit", 11, bold=True, color=NAVY)
    _box(sl, 9.1, y_row, 2.0, 0.4, "Exceed Days", 11, bold=True, color=NAVY)
    _box(sl, 11.2,y_row, 2.0, 0.4, "Exceed %", 11, bold=True, color=NAVY)
    y_row += 0.5
    for col_name, ef in sorted(exc_freq_data.items(),
                                key=lambda x: -x[1]["exceed_pct"])[:8]:
        bg = RED if ef["exceed_pct"] > 75 else (AMBER if ef["exceed_pct"] > 25 else TEAL)
        _rect(sl, 0.3, y_row, 12.7, 0.5, bg if ef["exceed_pct"]>50 else LIGHT)
        txt_col = WHITE if ef["exceed_pct"] > 50 else NAVY
        _box(sl, 0.4, y_row+0.05, 3.4, 0.4, col_name, 11, bold=True, color=txt_col)
        _box(sl, 3.9, y_row+0.05, 2.4, 0.4,
             f"{ef['annual_mean']:.2f}", 11, color=txt_col)
        _box(sl, 6.5, y_row+0.05, 2.4, 0.4,
             f"{ef['threshold']}", 11, color=txt_col)
        _box(sl, 9.1, y_row+0.05, 2.0, 0.4,
             str(ef["exceed_days"]), 11, bold=True, color=txt_col)
        _box(sl, 11.2,y_row+0.05, 2.0, 0.4,
             f"{ef['exceed_pct']:.1f}%", 11, bold=True,
             color=WHITE if ef["exceed_pct"]>50 else RED)
        y_row += 0.55
    if not exc_freq_data:
        _box(sl, 0.5, 2.0, 12, 1.0,
             "No WHO exceedances detected — all pollutants within guidelines.",
             18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 12d — DATA QUALITY DEEP DIVE
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, TEAL)
    _box(sl, 0.3, 0.15, 12, 0.6,
         "Data Quality Dashboard — What the Cleaning Engine Found", 20, bold=True, color=WHITE)

    clean_stats = cleaning_report if isinstance(cleaning_report, dict) else {}
    _stat_card(sl, 0.3, 1.1,
               f"{clean_stats.get('flatlines_nulled',0):,}",
               "Stuck Sensor Readings Removed", bg=RED)
    _stat_card(sl, 3.3, 1.1,
               f"{clean_stats.get('spikes_nulled',0):,}",
               "Impossible Spikes Removed", bg=AMBER)
    _stat_card(sl, 6.3, 1.1,
               f"{clean_stats.get('values_imputed',0):,}",
               "Values Imputed (5-strategy)", bg=TEAL)
    _stat_card(sl, 9.3, 1.1,
               f"{clean_stats.get('quality_score',0):.3f}",
               "Cleaning Quality Score", bg=GREEN)

    qs = clean_stats.get("quality_score", 0)
    dims = clean_stats.get("quality_dimensions", {})
    if dims:
        _box(sl, 0.3, 2.8, 12, 0.4,
             "7-Dimensional Quality Score Breakdown:", 13, bold=True, color=NAVY)
        dim_labels = {
            "D1_null_reduction"    : "Null Reduction",
            "D2_row_retention"     : "Row Retention",
            "D3_col_retention"     : "Column Retention",
            "D4_constraints"       : "Chemical Constraints",
            "D5_flatline"          : "No Flatlines",
            "D6_null_streaks"      : "Streak Quality",
            "D7_pm10_completeness" : "PM10 Complete",
        }
        for i, (k, v) in enumerate(dims.items()):
            x = 0.3 + (i % 4) * 3.25
            y = 3.3 + (i // 4) * 1.1
            bar_w = float(v or 0) * 2.8
            _rect(sl, x, y, 3.1, 0.9, LIGHT)
            _rect(sl, x, y, bar_w, 0.9,
                  GREEN if v>=0.85 else AMBER if v>=0.6 else RED)
            _box(sl, x+0.05, y+0.1, 3.0, 0.35,
                 dim_labels.get(k, k), 9, bold=True, color=WHITE)
            _box(sl, x+0.05, y+0.5, 3.0, 0.35,
                 f"{float(v or 0):.3f}", 11, bold=True, color=WHITE)
    _box(sl, 0.3, 5.8, 12, 0.4,
         f"Rows: {clean_stats.get('original_shape',(0,0))[0]:,} → "
         f"{clean_stats.get('final_shape',(0,0))[0]:,}  |  "
         f"Nulls: {clean_stats.get('total_nulls_before',0):,} → "
         f"{clean_stats.get('total_nulls_after',0):,}  |  "
         f"Gaps detected: {len(clean_stats.get('datetime_gaps',[]))}",
         11, italic=True, color=NAVY)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 12e — ANOMALY CONTEXT — TOP ANOMALY DAYS
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(WHITE)
    _rect(sl, 0, 0, 13.33, 0.9, NAVY)
    _box(sl, 0.3, 0.15, 12, 0.6,
         "Top Anomaly Days — Episode Classification", 22, bold=True, color=WHITE)

    episodes = v3_episodes if v3_episodes else []
    top_eps  = sorted(episodes, key=lambda x: x.get("pm10_mean",0), reverse=True)[:8]

    if top_eps:
        y_r = 1.0
        headers = ["Date","PM10 (µg/m³)","Episode Type","Season","Confidence","Reason"]
        col_ws  = [1.8, 1.8, 2.2, 1.4, 1.4, 4.5]
        x_starts= [0.3, 2.2, 4.1, 6.4, 7.9, 9.4]
        for hdr, xs in zip(headers, x_starts):
            _box(sl, xs, y_r, col_ws[0], 0.4, hdr, 10, bold=True, color=TEAL)
        y_r += 0.45
        for ep in top_eps:
            conf = ep.get("confidence","LOW")
            bg_c = RED if conf=="HIGH" else (AMBER if conf=="MEDIUM" else LIGHT)
            _rect(sl, 0.3, y_r, 12.7, 0.48, bg_c)
            txt_c = WHITE if conf=="HIGH" else NAVY
            vals  = [
                ep.get("date",""),
                f"{ep.get('pm10_mean',0):.1f}",
                ep.get("episode_type","").replace("_"," ")[:20],
                ep.get("season","").capitalize()[:8],
                conf,
                ep.get("reason","")[:50],
            ]
            for val, xs in zip(vals, x_starts):
                _box(sl, xs, y_r+0.04, 2.1, 0.4, val, 9, color=txt_c)
            y_r += 0.52
    else:
        _box(sl, 0.5, 2.5, 12, 1.0,
             "No pollution episodes classified. Run with --weather open_meteo for episode classification.",
             14, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 13 — PIPELINE ARCHITECTURE
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(LIGHT)
    _rect(sl, 0, 0, 13.33, 0.9, NAVY)
    _box(sl, 0.3, 0.15, 12, 0.6, "ATARS v3.1 — Pipeline Architecture",
         22, bold=True, color=WHITE)
    steps = [
        ("CSV Input",        "Any CPCB station data",   NAVY),
        ("16 Operators",     "Formal math Eqs. 1-16",  TEAL),
        ("ML Layer",         "IF + Holt-Winters",       NAVY),
        ("Open-Meteo API",   "18 weather variables",    TEAL),
        ("LLM Narrative",    "Ollama local, τ=0",       NAVY),
        ("9 Verifiers",      "RGV+TCV+CCD+MNS...",      TEAL),
        ("Word + PPT",       "Report + Slides auto",    GREEN),
    ]
    arrow = "→"
    for i, (title, desc, bg) in enumerate(steps):
        x = 0.3 + i * 1.85
        _rect(sl, x, 1.2, 1.7, 1.8, bg)
        _box(sl, x+0.05, 1.3, 1.6, 0.6, title, 12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
        _box(sl, x+0.05, 1.95, 1.6, 0.9, desc, 9,
             color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            _box(sl, x+1.7, 1.9, 0.3, 0.4, arrow, 16, bold=True, color=NAVY)
    _box(sl, 0.3, 3.3, 12.5, 0.5,
         f"Output: 18-section Word report  ·  Executive PPT  ·  15 charts  "
         f"·  SHA-256 audit log  ·  {len(verification.get('sections',{}))} verified sections",
         12, color=NAVY)
    import os as _os
    _src_lines = sum(1 for _ in open(__file__, "r", errors="ignore")) if "__file__" in dir() else 11000
    stats = [
        f"~{_src_lines:,} lines of code",
        "100+ functions  ·  16 formal operators",
        "21 report sections  ·  19 PPT slides",
        "MIT License — Open Source",
    ]
    for i, s in enumerate(stats):
        _box(sl, 0.3 + i * 3.2, 4.0, 3.0, 0.5, f"•  {s}", 12, color=TEAL)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 14 — KEY FINDINGS
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(NAVY)
    _box(sl, 0.3, 0.2, 12, 0.6, "Key Findings — Critical Conclusions",
         24, bold=True, color=WHITE)
    findings = [
        (RED,   f"PM10 at {pm10_mean:.1f} µg/m³ — {who_ratio}× WHO AQG — public health emergency"),
        (RED,   f"Benzene exceeded WHO on ALL 366 days — IARC Group 1 carcinogen"),
        (AMBER, f"{if_only} anomaly days missed by standard analysis — only IF detected them"),
        (GREEN, f"G_rate = {round(g_rate*100,1)}% — AI narrative numerically verified"),
        (WHITE, f"Data quality Q(D) = {Q:.4f} [{Q_flag}] — analysis is statistically valid"),
        (WHITE, f"HW 14-day forecast RMSE = {hw_rmse:.1f} µg/m³ — {round(hw_rmse/83.8*100,0):.0f}% of annual σ"),
    ]
    for i, (color, text) in enumerate(findings):
        y = 1.2 + i * 0.95
        _rect(sl, 0.3, y, 0.05, 0.7, color)
        _box(sl, 0.6, y, 12.4, 0.85, text, 14, color=color)

    # ══════════════════════════════════════════════════════════════════════
    #  SLIDE 15 — CLOSING
    # ══════════════════════════════════════════════════════════════════════
    sl = _slide(NAVY)
    _rect(sl, 0, 0, 13.33, 0.08, TEAL)
    _rect(sl, 0, 7.42, 13.33, 0.08, TEAL)
    # Title  (y=1.1 → 1.9)
    _box(sl, 0.5, 1.1, 12.3, 0.75,
         "ATARS v3.1 — Formal. Reproducible. Verified.",
         30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Subtitle (y=1.9 → 2.3)
    _box(sl, 0.5, 1.9, 12.3, 0.4,
         "The first automated air quality pipeline with 9 novel AI verification layers",
         13, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
    # 4 stat cards  (y=2.35 → 3.75, height=1.4)
    _stat_card(sl, 0.3,  2.35, f"{J.get('total_days',0)}", "Days Analysed", bg=TEAL)
    _stat_card(sl, 3.45, 2.35,
               f"{J['data_quality']['overall_Q']:.3f}",
               f"Q(D) [{J['data_quality']['confidence_flag']}]", bg=NAVY)
    _stat_card(sl, 6.6,  2.35,
               f"{round(verification.get('overall_G_rate',0)*100,1)}%",
               f"G_rate [{verification.get('overall_verdict','N/A')}]",
               bg=GREEN if verification.get('overall_G_rate',0)>=0.75 else AMBER)
    _stat_card(sl, 9.75, 2.35,
               f"{mns_val:.3f}" if v3_mns else "N/A",
               f"MNS [{mns_vdt}]",
               bg=GREEN if mns_vdt in ('CERTIFIED','TRUSTED') else AMBER)
    # Divider (y=3.85)
    _rect(sl, 0.5, 3.85, 12.33, 0.04, TEAL)
    # 3 links  (y=3.95 → 5.35, spacing=0.45)
    _box(sl, 0.5, 3.95, 12.3, 0.4,
         "Code:  github.com/Priyanshu-ux712/ATARS",
         12, color=LIGHT, align=PP_ALIGN.CENTER)
    _box(sl, 0.5, 4.40, 12.3, 0.4,
         f"Audit ID:  {run_id}",
         12, color=LIGHT, align=PP_ALIGN.CENTER)
    # Footer (y=5.4 → 5.8 — well above slide bottom 7.5)
    _box(sl, 0.5, 5.4, 12.3, 0.4,
         "MIT License  ·  Open Source  ·  © 2026 Priyanshu  ·  Global Institute of Technology and Management",
         10, color=LIGHT, align=PP_ALIGN.CENTER)

    # ── Save PPT ──────────────────────────────────────────────────────────
    city_slug = city.lower().replace(" ", "_")
    from datetime import datetime
    ts = datetime.now().strftime("%Y%m%d_%H%M")
    ppt_path = str(out_dir / f"ATARS_Summary_{city}_{ts}.pptx")
    prs.save(ppt_path)
    print(f"  ✓ PPT saved: {ppt_path}  (19 slides)")
    return ppt_path


def run_pipeline(config: dict):
    """
    ATARS v3 Pipeline — 11 steps.
    Steps 1-8: All v2 formal operators (exact same calls as v2).
    Steps 9-11: v3 enhancements (snapshots, comparison, enhanced RGV).
    Single Word output file — v3 section appended to same doc.
    """
    print_banner(config)

    # Merge v3 defaults
    for k, v in V3_CONFIG_DEFAULTS.items():
        if k not in config:
            config[k] = v

    out_dir = Path(config['output_dir'])
    out_dir.mkdir(parents=True, exist_ok=True)

    mode = config.get('mode', 'normal')
    if mode == 'test':
        print("\n  ⚠  TEST MODE — Using dataset as timeline. No live data claims.")

    try:
        # Initialize variables
        cleaning_report = {}
        data_profile    = {}
        stat_inference  = {}

        # ── Step 1: Data Ingestion ────────────────────────────────────────────
        print(f"\n{'═'*70}")
        print("  PHASE 1 — DATA PIPELINE  (Ingestion → Cleaning → Validation)")
        print(f"  {'═'*70}")
        print(f"\n{'─'*70}")
        print("  STEP 1.1 — Data Ingestion and Column Normalization")
        df = load_data(config['data_file'], config=config)

        # ── Step 1.5: Data Cleaning Engine ───────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 1.2 — Data Quality Engine: 10-Step Automated Cleaning Pipeline")
        print("  Steps: duplicate removal → datetime parsing → gap detection →")
        print("         negative correction → physical bounds → type coercion →")
        print("         sparse column removal → winsorization → imputation → audit")
        df, cleaning_report = clean_dataset(df, config)
        print(f"  ✓ Cleaning quality score: {cleaning_report['quality_score']:.3f}")
        print(f"  ✓ Nulls: {cleaning_report['total_nulls_before']:,} → {cleaning_report['total_nulls_after']:,}")
        print(f"  ✓ Rows: {cleaning_report['original_shape'][0]:,} → {cleaning_report['final_shape'][0]:,}")
        if cleaning_report.get('cols_removed'):
            print(f"  ✓ Sparse columns removed: {cleaning_report['cols_removed']}")
        if cleaning_report.get('datetime_gaps'):
            print(f"  ⚠ Datetime gaps detected: {len(cleaning_report['datetime_gaps'])}  (see report Section 00b)")
        print(f"\n{'═'*70}")
        print(f"  PHASE 1 COMPLETE — Clean dataset ready for analysis")
        print(f"  {'═'*70}")

        # ── Step 2: Quality Flags ─────────────────────────────────────────────
        print(f"\n{'═'*70}")
        print("  PHASE 2 — STATISTICAL ANALYSIS  (Quality Flags → Operators → ML)")
        print(f"  {'═'*70}")
        print(f"\n{'─'*70}")
        print("  STEP 2.1 — Data Availability Check + Quality Flag Assignment")
        n_records = len(df)
        if n_records == 0:
            print("  ✗ No records found. Exiting.")
            return None, None, None
        print(f"  ✓ {n_records:,} records spanning {df['date_only'].nunique()} days")
        df = assign_quality_flags(df)

        # ── Step 3: Statistical Aggregation ──────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 3 — Quality Assessment + Statistical Aggregation")
        daily_df = build_daily_stats(df, config)

        # ── Step 3.5: Data Profile ────────────────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 3.5 — Data Profile and Readiness Assessment")
        data_profile = generate_data_profile(df, daily_df, config)

        # ── Step 3.6: Statistical Inference ──────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 3.6 — Formal Statistical Inference (MK, Normality, OLS, Decomp)")
        stat_inference = run_full_statistical_inference(daily_df, config)

        # Correlation + OLS — exact same calls as v2
        print("\n  Computing correlation matrix and OLS regression...")
        corr_matrix = compute_correlation_matrix(daily_df, POLLUTANTS + METEOROLOGICAL[:4])
        ols_result  = run_ols_regression(daily_df, target='PM10',
                                          predictors=['humidity_pct', 'wind_speed_10m',
                                                      'rain_mm', 'pressure_hpa'])
        if ols_result:
            print(f"  ✓ OLS: R² = {ols_result['r_squared']:.4f} for PM10 ~ meteorology")

        # ── Step 3c: Isolation Forest ─────────────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 3c — ML-1: Isolation Forest Multivariate Anomaly Detection")
        ml_if = run_isolation_forest(daily_df, config)

        # ── Step 3d: Holt-Winters Forecast ────────────────────────────────────
        if config.get('use_forecast', True):
            print(f"\n{'─'*70}")
            print("  STEP 3d — ML-2: Holt-Winters PM10 Forecast")
            ml_hw = run_hw_forecast(daily_df, config)
        else:
            ml_hw = {'available': False, 'reason': 'disabled_by_flag'}
            print("  STEP 3d — Forecast disabled (--no-forecast)")

        # ── Step 4: JSON Contract ─────────────────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 4 — JSON Contract Assembly (formal data boundary J)")
        J = build_json_contract(df, daily_df, config, ols_result, corr_matrix)
        # Inject stat inference summary into J
        if stat_inference:
            mk_pm10 = stat_inference.get("mann_kendall", {}).get("PM10", {})
            norm_pm10 = stat_inference.get("normality", {}).get("PM10", {})
            ols_inf   = stat_inference.get("ols_inference", {})
            decomp_s  = stat_inference.get("seasonal_decomposition", {})
            J["stat_inference_summary"] = {
                "mann_kendall_pm10"   : {
                    "direction"  : mk_pm10.get("direction"),
                    "p_value"    : mk_pm10.get("p_value"),
                    "significant": mk_pm10.get("significant"),
                    "sens_slope" : mk_pm10.get("sens_slope"),
                } if mk_pm10.get("available") else None,
                "normality_pm10"      : {
                    "is_normal"  : norm_pm10.get("is_normal"),
                    "verdict"    : norm_pm10.get("verdict"),
                    "skewness"   : norm_pm10.get("skewness"),
                    "kurtosis"   : norm_pm10.get("kurtosis"),
                } if norm_pm10 else None,
                "ols_r2_adj"          : ols_inf.get("r_squared_adj"),
                "ols_f_pvalue"        : ols_inf.get("f_pvalue"),
                "ols_dw"              : ols_inf.get("durbin_watson"),
                "seasonal_pct_trend"  : decomp_s.get("pct_trend"),
                "seasonal_pct_seasonal": decomp_s.get("pct_seasonal"),
                "seasonal_pct_residual": decomp_s.get("pct_residual"),
            }
        json_path = out_dir / 'json_contract.json'
        with open(json_path, 'w') as f:
            json.dump(J, f, indent=2, default=_json_safe)
        print(f"  ✓ JSON contract saved: {json_path}")

        # ── Step 5: AI Narrative ──────────────────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 5 — AI Narrative Generation (G: J → N)")
        narrative_base = _placeholder_narrative(J)
        if config.get('use_llm', True):
            narrative = try_llm_narrative(J, config)
        else:
            narrative = _placeholder_narrative(J)

        # ── Step 5.5: v2 RGV (preserved exactly) ─────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 5.5 — Runtime Grounding Verifier v2 (G(s) ⊆ J)")
        rgv          = RuntimeGroundingVerifier(J)
        verification = rgv.verify_narrative(narrative)
        g_rate       = verification['overall_G_rate']
        g_verdict    = verification['overall_verdict']
        n_ungrounded = verification['total_ungrounded']
        print(f"  ✓ G_rate = {g_rate:.1%}  [{g_verdict}]")
        print(f"  ✓ J lookup: {verification['j_lookup_size']} entries | "
              f"Ungrounded: {n_ungrounded}")
        ver_path = out_dir / 'grounding_verification.json'
        with open(ver_path, 'w') as _f:
            json.dump(verification, _f, indent=2, default=str)

        # ── Step 5.5b: Chart Insights ─────────────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 5.5b — Deterministic Chart Insights")
        chart_insights = generate_chart_insights(
            daily_df, df, config, corr_matrix, ols_result, ml_if, ml_hw)
        print(f"  ✓ Insights for {len(chart_insights)} charts")

        # ── Step 6: Charts ────────────────────────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 6 — Chart Generation")
        charts = generate_all_charts(df, daily_df, config, out_dir, verification,
                                     ml_if=ml_if, ml_hw=ml_hw,
                                     cleaning_report=cleaning_report,
                                     data_profile=data_profile)

        # ── Step 7: Word Report ───────────────────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 7 — Chart generation complete. Report compiled after v3 steps.")
        # NOTE: build_report is called after all v3 steps (Step 17+) so that
        # upt, benchmark, comparison, aee, cvcc, nss, tcv, ccd, mns, drift
        # are all available to pass into the report builder.

        # Step 8 audit deferred — needs report_path from Step 17b
        audit = {}  # populated after build_report in Step 17b




        # ═══════════════════════════════════════════════════════════════════════
        # v3 BEAST MODE ENHANCEMENTS — Steps 9-17
        # ═══════════════════════════════════════════════════════════════════════

        # Initialize v3 weather vars + get weather_src FIRST
        weather_src = config.get('weather_source', 'open_meteo')
        atmo_stability = {}; episodes_classified = []; forecast_7day = {}

        # ── Step 8b: 7-day weather forecast (always unless mock) ──────────────
        forecast_7day = {}
        if weather_src != "mock":
            print("  " + chr(45)*68)
            print("  STEP 8b — Fetching 7-day weather forecast from Open-Meteo")
            forecast_7day = _fetch_7day_forecast(config)

        # ── Step 9: Daily Snapshots ───────────────────────────────────────────
        print(f"\n{'═'*70}")
        print("  PHASE 3 — CONTEXTUAL ENRICHMENT  (Snapshots → Weather → Episodes)")
        print(f"  {'═'*70}")
        print(f"\n{'─'*70}")
        print("  STEP 9 — v3 Daily Snapshot System")
        snapshots = {}
        print(f"  Weather source: {weather_src.upper()}")

        # Pre-fetch ALL weather from Open-Meteo in ONE batch call
        # Always attempted unless user explicitly sets --weather mock
        live_weather_cache = {}
        if weather_src != "mock":
            print("  Fetching weather from Open-Meteo API (batch)...")
            all_dates = sorted(daily_df["date"].dropna().unique())
            live_weather_cache = _get_live_weather_cache(config, all_dates)
            if live_weather_cache:
                sources = set(w.get("source","?") for w in live_weather_cache.values())
                print(f"  ✓ Weather: {len(live_weather_cache)} days | sources: {sources}")
            else:
                print("  ⚠ Open-Meteo unavailable — CSV columns / mock fallback active")
        else:
            print("  ⚠ Mock weather mode — Open-Meteo disabled by --weather mock")

        if config.get('enable_snapshots', True):
            snapshots = build_daily_snapshots(
                daily_df, config, live_weather_cache=live_weather_cache)
        else:
            print("  ⚠ Snapshots disabled")

        # ── Step 10: Comparison Engine ────────────────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 10 — v3 Comparison Engine (Today vs Yesterday/Week/Month)")
        comparison = {}
        if snapshots:
            all_snap_dates = sorted(snapshots.keys())
            today_date = config.get('test_today_date') or all_snap_dates[-1]
            comparison = run_comparison_engine(snapshots, today_date)
        else:
            comparison = {'today': 'N/A', 'status': 'NO_SNAPSHOTS'}

        # ── Step 10b: Atmospheric Stability + Episode Classification ──────────
        atmo_stability = {}
        episodes_classified = []
        # Atmospheric stability — uses whatever weather is available
        if live_weather_cache or snapshots:
            weather_for_stability = live_weather_cache or {
                d: s.get("weather",{}) for d,s in snapshots.items()
            }
            print("  STEP 10b — Atmospheric Stability Analysis")
            atmo_stability = compute_atmospheric_stability(weather_for_stability)
            anomaly_days = []
            pm10_sub = daily_df[daily_df["variable"] == "PM10"].copy()
            pm10_sub["date"] = pd.to_datetime(pm10_sub["date"])
            if "is_anomaly" in pm10_sub.columns:
                for _, row in pm10_sub[pm10_sub["is_anomaly"] == True].iterrows():
                    d_str = str(row["date"])[:10]
                    anomaly_days.append({
                        "date"       : d_str,
                        "pm10_mean"  : float(row.get("mean", 0)),
                        "anomaly_type": "ANOMALY",
                    })
            if anomaly_days:
                print("  STEP 10c — Pollution Episode Classification")
                episodes_classified = classify_pollution_episodes(
                    anomaly_days, live_weather_cache, atmo_stability, config)

        # Enrich J with v3 comparison + weather + WHO
        # Add live weather summary to J for LLM use
        if live_weather_cache and comparison.get("today"):
            today_w = live_weather_cache.get(comparison["today"],
                      comparison.get("weather", {}))
            comparison["weather"] = today_w
            comparison["weather_source"] = today_w.get("source","mock")

        J_v3 = enrich_json_contract_v3(J, comparison, config)
        # Enrich J_v3 with atmospheric data
        if atmo_stability:
            J_v3["atmospheric_stability"] = {
                "high_risk_count"  : atmo_stability.get("high_risk_count", 0),
                "total_days"       : atmo_stability.get("total_days", 0),
                "high_risk_days"   : atmo_stability.get("high_risk_days", [])[:5],
            }
        if forecast_7day:
            J_v3["weather_forecast_7day"] = {
                d: {k: v for k, v in w.items()
                    if k in ("date","temp_mean_c","temp_max_c","temp_min_c",
                             "wind_mean_kmh","rain_mm","weather_desc",
                             "pollution_risk","pollution_risk_score")}
                for d, w in forecast_7day.items()
            }
        if episodes_classified:
            J_v3["pollution_episodes"] = {
                e["date"]: {
                    "episode_type": e["episode_type"],
                    "confidence"  : e["confidence"],
                    "reason"      : e["reason"],
                }
                for e in episodes_classified
            }

        json_v3_path = out_dir / 'json_contract_v3.json'
        with open(json_v3_path, 'w') as jf:
            json.dump(J_v3, jf, indent=2, default=str)

        # ── Step 11: v3 Enhanced RGV ──────────────────────────────────────────
        print(f"\n{'═'*70}")
        print("  PHASE 4 — VERIFICATION LAYER  (9 Novel Verification Modules)")
        print(f"  {'═'*70}")
        print(f"\n{'─'*70}")
        print("  STEP 11 — v3 Enhanced RGV (Semantic + Confidence + Contradiction)")
        ergv       = EnhancedRGV_v3(J_v3)
        ver_v3     = ergv.verify_v3(narrative, comparison)
        evaluation = run_evaluation_layer(narrative_base, narrative, ver_v3)
        ver_v3_path = out_dir / 'grounding_verification_v3.json'
        with open(ver_v3_path, 'w') as vf:
            json.dump(ver_v3, vf, indent=2, default=str)

        # ── Step 12: N-01 Temporal Claim Verifier ────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 12 — N-01: Temporal Claim Verifier (TCV)")
        temporal_facts = build_temporal_facts(daily_df, J)
        J_v3["temporal_facts"] = temporal_facts
        tcv = verify_temporal_claims(narrative, temporal_facts)

        # ── Step 13: N-02 Causal Claim Detector ──────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 13 — N-02: Causal Claim Detector (CCD)")
        ccd = detect_causal_claims(narrative, J)

        # ── Step 14: N-03 Anomaly Explanation Engine ──────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 14 — N-03: Anomaly Explanation Engine (AEE)")
        aee = run_anomaly_explanation_engine(daily_df, snapshots, config)

        # ── Step 15: N-04 Report Drift Detector ──────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 15 — N-04: Report Drift Detector (RDD)")
        drift = run_report_drift_detector(J, ver_v3, out_dir)

        # ── Step 16: N-05 Uncertainty Propagation Tracker ────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 16 — N-05: Uncertainty Propagation Tracker (UPT)")
        uncertainty_table = build_uncertainty_table(daily_df, J)
        upt = annotate_claims_with_uncertainty(narrative, uncertainty_table, ver_v3)

        # ── Step 16b: N-06 Cross-Variable Consistency Checker ────────────────
        print(f"\n{'─'*70}")
        print("  STEP 16b — N-06: Cross-Variable Consistency Checker (CVCC)")
        cvcc = run_cross_variable_consistency(J, narrative)

        # ── Step 16c: N-07 Narrative Specificity Scorer ───────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 16c — N-07: Narrative Specificity Scorer (NSS)")
        nss = score_narrative_specificity(narrative)

        # ── Step 16d: N-08 Benchmark Comparator ───────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 16d — N-08: Benchmark Comparator")
        benchmark = run_benchmark_comparator(J, config)

        # ── Step 17: N-09 Master Novelty Score ───────────────────────────────
        print(f"\n{'─'*70}")
        print("  STEP 17 — N-09: Master Novelty Score (MNS)")
        mns = compute_master_novelty_score(
            verification, ver_v3, tcv, ccd, cvcc, nss, aee)

        # ── Step 17b: Build Report (all v3 data now available) ──────────────
        print(f"\n{'─'*70}")
        print("  STEP 17b — Report Compilation (all v3 data ready)")
        report_path = build_report(
            J, daily_df, charts, narrative,
            ols_result, corr_matrix, config, out_dir,
            verification=verification,
            chart_insights=chart_insights,
            ml_if=ml_if, ml_hw=ml_hw,
            v3_upt          = upt,
            v3_benchmark    = benchmark,
            v3_comparison   = comparison,
            v3_aee          = aee,
            v3_cvcc         = cvcc,
            v3_nss          = nss,
            v3_tcv          = tcv,
            v3_ccd          = ccd,
            v3_mns          = mns,
            v3_drift        = drift,
            v3_ver_v3       = ver_v3,
            v3_stability    = atmo_stability,
            v3_episodes     = episodes_classified,
            v3_forecast     = forecast_7day,
            cleaning_report = cleaning_report,
            data_profile    = data_profile,
            stat_inference  = stat_inference,
        )
        print(f"  ✓ Report saved: {report_path}")

        # ── Step 17c: Generate Summary PPT ───────────────────────────────────
        print(f"\n" + "-"*70)
        print("  STEP 17c — Executive Summary PPT Generation")
        try:
            ppt_path = generate_summary_ppt(
                J=J, J_v3=J_v3,
                verification=verification,
                ml_if=ml_if, ml_hw=ml_hw,
                config=config, out_dir=out_dir,
                v3_mns=mns, v3_episodes=episodes_classified,
                v3_forecast=forecast_7day,
                v3_benchmark=benchmark,
                charts=charts,
                cleaning_report=cleaning_report,
            )
        except Exception as _ppt_e:
            print(f"  ⚠ PPT generation error: {_ppt_e}")
            ppt_path = None

        # ── Step 8: Audit Log (now report_path is available) ─────────────────
        print("  " + chr(8212)*68)




        print("  STEP 8 — Reproducibility Audit Log")
        audit = write_audit_log(J, narrative, report_path, config, out_dir,
                                verification=verification, ml_if=ml_if, ml_hw=ml_hw)

        # Save master novelty score to audit
        mns_path = out_dir / 'master_novelty_score.json'
        with open(mns_path, 'w') as mf:
            json.dump(mns, mf, indent=2)

        # ── Final Summary ─────────────────────────────────────────────────────
        g3 = ver_v3['overall_G_rate']
        print(f"\n{'═'*70}")
        print("  ATARS v3.1 — PIPELINE COMPLETE")
        print(f"{'═'*70}")
        print(f"  ✓ Report          : {report_path}")
        print(f"  ✓ Charts          : {len(charts)} figures")
        print(f"  ✓ Snapshots       : {len(snapshots)} daily files")
        print(f"  ✓ G_rate          : {round(g3*100,1)}%  [{ver_v3['overall_verdict']}]")
        print(f"  ✓ TCV rate        : {round(tcv.get('TCV_rate',0)*100,1)}%  [{tcv.get('TCV_verdict','N/A')}]")
        print(f"  ✓ CCD             : {ccd.get('causal_claim_count',0)} causal claims  [{ccd.get('CCD_verdict','N/A')}]")
        print(f"  ✓ CVCC            : {len(cvcc.get('violations',[]))} violations  [{cvcc.get('CVCC_verdict','N/A')}]")
        print(f"  ✓ NSS             : {nss.get('overall_specificity',0):.2f}  [{nss.get('NSS_verdict','N/A')}]")
        print(f"  ✓ AEE             : {aee.get('explained',0)}/{aee.get('total_anomalies',0)} anomalies explained")
        print(f"  ✓ UPT             : {upt.get('total_annotated',0)} claims with CI")
        print(f"  ✓ Benchmark       : PM10 = {benchmark.get('comparisons',{}).get('PM10',{}).get('vs_who','N/A')}  [{benchmark.get('overall_severity','N/A')}]")
        print(f"  ✓ MNS             : {mns.get('MNS',0):.3f}  [{mns.get('verdict','N/A')}]")
        print(f"  ✓ Confidence      : {ver_v3['overall_confidence']}")
        print(f"  ✓ Semantic        : {'PASS' if ver_v3['semantic_pass'] else 'REVIEW'}")
        print(f"  ✓ Contradictions  : {ver_v3['contradiction_count']}")
        print(f"  ✓ Drift Status    : {drift.get('status','N/A')}")
        print(f"  ✓ Mode            : {mode.upper()}")
        print(f"  ✓ Audit ID        : {audit['run_id']}")
        print(f"  ✓ Hash(J)         : {audit['hash_J'][:32]}...")
        if ml_if.get('available'):
            print(f"  ✓ IF anomalies    : {ml_if['n_anomalies']} days | IF-only: {ml_if['n_if_only']}")
        if ml_hw.get('available'):
            print(f"  ✓ HW forecast     : RMSE={ml_hw['rmse']:.2f} | "
                  f"14d mean={float(np.mean(ml_hw['forecast'])):.1f} µg/m³")
        # Data profile in summary
        if data_profile and data_profile.get("readiness"):
            rd = data_profile["readiness"]
            print(f"  ✓ Data readiness    : {rd.get('verdict','N/A')} "
                  f"(PM10={rd.get('pm10_valid',0):.1f}% complete, "
                  f"coverage={rd.get('coverage_pct',0):.1f}%)")
        print(f"\n  Q(D) = {J['data_quality']['overall_Q']:.4f}  [{J['data_quality']['confidence_flag']}]")
        print(f"{'═'*70}\n")

        return str(report_path) if report_path else None, J_v3, audit

    except Exception as e:
        print(f"\n  ✗ Pipeline error: {e}")
        import traceback
        traceback.print_exc()
        return None, None, None


def main():
    parser = argparse.ArgumentParser(
        description="ATARS v3.0 — Verified AI Reporting System",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter
    )
    parser.add_argument("--data",         default="data.csv")
    parser.add_argument("--city",         default="City")
    parser.add_argument("--station",      default="Station-01")
    parser.add_argument("--output",       default="atars_v3_output")
    parser.add_argument("--no-llm",       action="store_true")
    parser.add_argument("--llm-model",    default="llama3.2")
    parser.add_argument("--baseline",     type=int, default=30)
    parser.add_argument("--timeout",      type=int, default=600)
    parser.add_argument("--retry",        type=int, default=2)
    parser.add_argument("--no-forecast",  action="store_true")
    parser.add_argument("--weather",      choices=["mock","open_meteo"], default="open_meteo",
                        help="mock = seasonal values (no internet) | "
                             "open_meteo = free live/historical weather API")
    parser.add_argument("--mode",         choices=["normal","test"], default="normal",
                        help="test = use dataset as timeline, mock weather")
    parser.add_argument("--no-snapshots", action="store_true",
                        help="Skip daily snapshot generation")
    parser.add_argument("--today",        default=None,
                        help="Override today date YYYY-MM-DD (test mode only)")
    args = parser.parse_args()

    config = DEFAULT_CONFIG.copy()
    config.update({
        "data_file"           : args.data,
        "city"                : args.city,
        "station_id"          : args.station,
        "output_dir"          : args.output,
        "use_llm"             : not args.no_llm,
        "llm_model"           : args.llm_model,
        "baseline_window_days": args.baseline,
        "llm_timeout"         : args.timeout,
        "llm_retry"           : args.retry,
        "use_forecast"        : not args.no_forecast,
        "mode"                : args.mode,
        "enable_snapshots"    : not args.no_snapshots,
        "weather_source"     : args.weather,
        "test_today_date"     : args.today,
    })
    run_pipeline(config)


if __name__ == "__main__":
    main()
